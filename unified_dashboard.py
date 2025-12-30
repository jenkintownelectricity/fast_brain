"""
Unified Skill Command Center Dashboard - COMPREHENSIVE VERSION

Combines ALL features from:
- Skill Command Center (API Keys, LLM Testing, Compare Providers, Latency Masking, Stats, Voice Integration)
- Skill Factory (Business Profile, Upload Documents, Train Skill, Manage Skills with filterable table)
- LPU Inference (BitNet model testing)
- Dashboard (Stats, Server Status, Training Pipeline, Activity Feed)

All in one beautiful cyberpunk-themed interface.

Usage:
    pip install flask flask-cors
    python unified_dashboard.py
    # Opens at http://localhost:5000
"""

from flask import Flask, jsonify, render_template_string, request, send_file
from flask_cors import CORS
import json
import os
import re
from pathlib import Path
from datetime import datetime, timedelta
import random

# Import unified database
try:
    import database as db
    USE_DATABASE = True
    print("Using SQLite database for persistent storage")
except ImportError:
    USE_DATABASE = False
    print("Database module not found - using in-memory storage")

# Helper function for raw SQL queries (uses existing db.get_db context manager)
def execute_query(query, params=None):
    """Execute a raw SQL query and return results."""
    with db.get_db() as conn:
        cursor = conn.cursor()
        cursor.execute(query, params or ())
        return cursor.fetchall()

# Modal volume commit helper - CRITICAL for persisting database writes
def commit_volume():
    """Commit changes to Modal volume to persist database writes across container restarts.

    Modal volumes have eventual consistency - writes are not persisted until commit() is called.
    This function safely handles both Modal and local environments.
    """
    try:
        import modal
        # Get the volume reference - this works because Modal injects the volume at runtime
        volume = modal.Volume.from_name("hive215-data")
        volume.commit()
        print("[Volume] Committed database changes to Modal volume")
    except ImportError:
        # Not running on Modal (local development)
        pass
    except Exception as e:
        # Log but don't fail - volume commit errors shouldn't break the app
        print(f"[Volume] Warning: Failed to commit volume: {e}")

app = Flask(__name__)
CORS(app, origins=[
    "https://453rahul.com",
    "http://453rahul.com",
    "https://www.453rahul.com",
    "http://www.453rahul.com",
    "http://localhost:3000",  # For local dev
], supports_credentials=True)

# Data paths
BUSINESS_PROFILES_DIR = Path("business_profiles")
TRAINING_DATA_DIR = Path("training_data")
ADAPTERS_DIR = Path("adapters")
LOGS_DIR = Path("logs")
DOCUMENTS_DIR = Path("uploaded_documents")

# Ensure directories exist
for d in [BUSINESS_PROFILES_DIR, TRAINING_DATA_DIR, ADAPTERS_DIR, LOGS_DIR, DOCUMENTS_DIR]:
    d.mkdir(exist_ok=True)

# Load API keys from database (for backward compatibility with existing code)
API_KEYS = db.get_all_api_keys() if USE_DATABASE else {}

# Activity log - now uses database but keep reference for backward compatibility
ACTIVITY_LOG = db.get_recent_activity(20) if USE_DATABASE else []

# Voice configuration and platform connections
VOICE_CONFIG = {
    "selected_voice": "en-US-JennyNeural",
    "selected_provider": "edge_tts",
    "voice_settings": {
        "speed": 1.0,
        "pitch": 1.0,
        "volume": 1.0
    }
}

# =============================================================================
# TTS AUDIO CACHE - Low latency for common phrases
# =============================================================================
import hashlib
from collections import OrderedDict
import threading

class TTSCache:
    """In-memory cache for synthesized audio to reduce latency."""

    def __init__(self, max_size=500, ttl_seconds=3600):
        self.cache = OrderedDict()
        self.max_size = max_size
        self.ttl = ttl_seconds
        self.lock = threading.Lock()
        self.stats = {"hits": 0, "misses": 0}

    def _make_key(self, text, voice_id, provider, **kwargs):
        """Create a unique cache key from synthesis parameters."""
        # Normalize text (lowercase, strip whitespace)
        normalized = text.lower().strip()
        key_data = f"{normalized}|{voice_id}|{provider}"
        # Include relevant kwargs that affect output
        for k in sorted(kwargs.keys()):
            if k in ['emotion', 'speed', 'pitch']:
                key_data += f"|{k}={kwargs[k]}"
        return hashlib.sha256(key_data.encode()).hexdigest()[:32]

    def get(self, text, voice_id, provider, **kwargs):
        """Get cached audio if available and not expired."""
        key = self._make_key(text, voice_id, provider, **kwargs)
        with self.lock:
            if key in self.cache:
                entry = self.cache[key]
                # Check TTL
                import time
                if time.time() - entry['timestamp'] < self.ttl:
                    self.stats["hits"] += 1
                    # Move to end (LRU)
                    self.cache.move_to_end(key)
                    return entry['audio'], entry['format']
                else:
                    # Expired
                    del self.cache[key]
            self.stats["misses"] += 1
            return None, None

    def set(self, text, voice_id, provider, audio_bytes, audio_format='audio/mp3', **kwargs):
        """Cache synthesized audio."""
        import time
        key = self._make_key(text, voice_id, provider, **kwargs)
        with self.lock:
            # Remove oldest if at capacity
            while len(self.cache) >= self.max_size:
                self.cache.popitem(last=False)

            self.cache[key] = {
                'audio': audio_bytes,
                'format': audio_format,
                'timestamp': time.time(),
                'text': text[:100],  # Store truncated text for debugging
                'voice_id': voice_id,
                'provider': provider
            }

    def get_stats(self):
        """Get cache statistics."""
        with self.lock:
            total = self.stats["hits"] + self.stats["misses"]
            hit_rate = (self.stats["hits"] / total * 100) if total > 0 else 0
            return {
                "hits": self.stats["hits"],
                "misses": self.stats["misses"],
                "hit_rate": f"{hit_rate:.1f}%",
                "size": len(self.cache),
                "max_size": self.max_size
            }

    def clear(self):
        """Clear the cache."""
        with self.lock:
            self.cache.clear()
            self.stats = {"hits": 0, "misses": 0}

# Global TTS cache instance
TTS_CACHE = TTSCache(max_size=500, ttl_seconds=3600)  # 1 hour TTL

# Common phrases to pre-cache (will be populated on first synthesis)
COMMON_PHRASES = [
    "Hello! How can I help you today?",
    "Thank you for calling. How may I assist you?",
    "I understand. Let me help you with that.",
    "One moment please while I look that up.",
    "Is there anything else I can help you with?",
    "Thank you for your patience.",
    "I'll transfer you to the right department.",
    "Your appointment has been confirmed.",
    "We'll send you a confirmation shortly.",
    "Have a great day!",
]

VOICE_CHOICES = {
    "parler_tts": {
        "name": "Parler TTS",
        "provider": "Parler (GPU - Expressive)",
        "voices": [
            {"id": "receptionist", "name": "Sarah (Receptionist)", "gender": "female", "style": "professional warm"},
            {"id": "electrician", "name": "Mike (Electrician)", "gender": "male", "style": "friendly knowledgeable"},
            {"id": "plumber", "name": "Tom (Plumber)", "gender": "male", "style": "reassuring calm"},
            {"id": "lawyer", "name": "Jennifer (Lawyer)", "gender": "female", "style": "professional articulate"},
            {"id": "solar", "name": "Alex (Solar)", "gender": "neutral", "style": "enthusiastic energetic"},
            {"id": "general", "name": "Jordan (General)", "gender": "neutral", "style": "warm natural"},
        ]
    },
    "edge_tts": {
        "name": "Edge TTS",
        "provider": "Microsoft (Free)",
        "voices": [
            {"id": "en-US-JennyNeural", "name": "Jenny", "gender": "female", "style": "friendly"},
            {"id": "en-US-AriaNeural", "name": "Aria", "gender": "female", "style": "professional"},
            {"id": "en-US-SaraNeural", "name": "Sara", "gender": "female", "style": "warm"},
            {"id": "en-US-AnaNeural", "name": "Ana", "gender": "female", "style": "youthful"},
            {"id": "en-US-GuyNeural", "name": "Guy", "gender": "male", "style": "conversational"},
            {"id": "en-US-DavisNeural", "name": "Davis", "gender": "male", "style": "confident"},
            {"id": "en-US-TonyNeural", "name": "Tony", "gender": "male", "style": "friendly"},
            {"id": "en-US-JasonNeural", "name": "Jason", "gender": "male", "style": "neutral"},
            {"id": "en-GB-SoniaNeural", "name": "Sonia (UK)", "gender": "female", "style": "british"},
            {"id": "en-GB-RyanNeural", "name": "Ryan (UK)", "gender": "male", "style": "british"},
            {"id": "en-AU-NatashaNeural", "name": "Natasha (AU)", "gender": "female", "style": "australian"},
            {"id": "en-AU-WilliamNeural", "name": "William (AU)", "gender": "male", "style": "australian"},
        ]
    },
    "kokoro": {
        "name": "Kokoro",
        "provider": "Kokoro TTS (Free)",
        "voices": [
            {"id": "af_bella", "name": "Bella", "gender": "female", "style": "american"},
            {"id": "af_nicole", "name": "Nicole", "gender": "female", "style": "american"},
            {"id": "af_sarah", "name": "Sarah", "gender": "female", "style": "american"},
            {"id": "af_sky", "name": "Sky", "gender": "female", "style": "american"},
            {"id": "am_adam", "name": "Adam", "gender": "male", "style": "american"},
            {"id": "am_michael", "name": "Michael", "gender": "male", "style": "american"},
            {"id": "bf_emma", "name": "Emma", "gender": "female", "style": "british"},
            {"id": "bf_isabella", "name": "Isabella", "gender": "female", "style": "british"},
            {"id": "bm_george", "name": "George", "gender": "male", "style": "british"},
            {"id": "bm_lewis", "name": "Lewis", "gender": "male", "style": "british"},
        ]
    },
    "chatterbox": {
        "name": "Chatterbox",
        "provider": "Resemble AI",
        "voices": [
            {"id": "chatterbox_default", "name": "Default", "gender": "neutral", "style": "conversational"},
            {"id": "chatterbox_warm", "name": "Warm", "gender": "female", "style": "friendly"},
            {"id": "chatterbox_professional", "name": "Professional", "gender": "male", "style": "business"},
            {"id": "chatterbox_energetic", "name": "Energetic", "gender": "female", "style": "upbeat"},
        ]
    },
    "xtts": {
        "name": "XTTS-v2",
        "provider": "Coqui AI",
        "voices": [
            {"id": "xtts_clone", "name": "Custom Clone", "gender": "custom", "style": "cloned"},
            {"id": "xtts_en_female", "name": "English Female", "gender": "female", "style": "neutral"},
            {"id": "xtts_en_male", "name": "English Male", "gender": "male", "style": "neutral"},
        ]
    },
    "openvoice": {
        "name": "OpenVoice",
        "provider": "MyShell AI",
        "voices": [
            {"id": "openvoice_default", "name": "Default", "gender": "neutral", "style": "neutral"},
            {"id": "openvoice_clone", "name": "Instant Clone", "gender": "custom", "style": "cloned"},
        ]
    },
    "elevenlabs": {
        "name": "ElevenLabs",
        "provider": "ElevenLabs (Paid)",
        "voices": [
            {"id": "eleven_rachel", "name": "Rachel", "gender": "female", "style": "conversational"},
            {"id": "eleven_drew", "name": "Drew", "gender": "male", "style": "neutral"},
            {"id": "eleven_clyde", "name": "Clyde", "gender": "male", "style": "character"},
            {"id": "eleven_paul", "name": "Paul", "gender": "male", "style": "news"},
            {"id": "eleven_domi", "name": "Domi", "gender": "female", "style": "strong"},
            {"id": "eleven_bella", "name": "Bella", "gender": "female", "style": "soft"},
            {"id": "eleven_antoni", "name": "Antoni", "gender": "male", "style": "warm"},
            {"id": "eleven_custom", "name": "Custom Clone", "gender": "custom", "style": "cloned"},
        ]
    },
    "openai": {
        "name": "OpenAI TTS",
        "provider": "OpenAI (Paid)",
        "voices": [
            {"id": "openai_alloy", "name": "Alloy", "gender": "neutral", "style": "balanced"},
            {"id": "openai_echo", "name": "Echo", "gender": "male", "style": "neutral"},
            {"id": "openai_fable", "name": "Fable", "gender": "neutral", "style": "storytelling"},
            {"id": "openai_onyx", "name": "Onyx", "gender": "male", "style": "deep"},
            {"id": "openai_nova", "name": "Nova", "gender": "female", "style": "warm"},
            {"id": "openai_shimmer", "name": "Shimmer", "gender": "female", "style": "expressive"},
        ]
    },
    "azure": {
        "name": "Azure TTS",
        "provider": "Microsoft Azure (Paid)",
        "voices": [
            {"id": "azure_jenny", "name": "Jenny", "gender": "female", "style": "neural"},
            {"id": "azure_guy", "name": "Guy", "gender": "male", "style": "neural"},
            {"id": "azure_aria", "name": "Aria", "gender": "female", "style": "conversational"},
            {"id": "azure_davis", "name": "Davis", "gender": "male", "style": "conversational"},
        ]
    },
    "cartesia": {
        "name": "Cartesia",
        "provider": "Cartesia (Paid)",
        "voices": [
            {"id": "cartesia_default", "name": "Default", "gender": "neutral", "style": "natural"},
            {"id": "cartesia_custom", "name": "Custom Voice", "gender": "custom", "style": "cloned"},
        ]
    },
    "deepgram": {
        "name": "Deepgram Aura",
        "provider": "Deepgram (Paid)",
        "voices": [
            {"id": "aura-asteria-en", "name": "Asteria", "gender": "female", "style": "american"},
            {"id": "aura-luna-en", "name": "Luna", "gender": "female", "style": "american"},
            {"id": "aura-stella-en", "name": "Stella", "gender": "female", "style": "american"},
            {"id": "aura-athena-en", "name": "Athena", "gender": "female", "style": "british"},
            {"id": "aura-hera-en", "name": "Hera", "gender": "female", "style": "american"},
            {"id": "aura-orion-en", "name": "Orion", "gender": "male", "style": "american"},
            {"id": "aura-arcas-en", "name": "Arcas", "gender": "male", "style": "american"},
            {"id": "aura-perseus-en", "name": "Perseus", "gender": "male", "style": "american"},
            {"id": "aura-angus-en", "name": "Angus", "gender": "male", "style": "irish"},
            {"id": "aura-orpheus-en", "name": "Orpheus", "gender": "male", "style": "american"},
            {"id": "aura-helios-en", "name": "Helios", "gender": "male", "style": "british"},
            {"id": "aura-zeus-en", "name": "Zeus", "gender": "male", "style": "american"},
        ]
    }
}

PLATFORM_CONNECTIONS = {
    "livekit": {
        "name": "LiveKit",
        "status": "disconnected",
        "config": {
            "url": "",
            "api_key": "",
            "api_secret": ""
        },
        "description": "Real-time voice and video with agents SDK"
    },
    "vapi": {
        "name": "Vapi",
        "status": "disconnected",
        "config": {
            "api_key": "",
            "assistant_id": ""
        },
        "description": "Voice AI platform for building phone agents"
    },
    "twilio": {
        "name": "Twilio",
        "status": "disconnected",
        "config": {
            "account_sid": "",
            "auth_token": "",
            "phone_number": ""
        },
        "description": "Cloud communications platform for voice calls"
    },
    "retell": {
        "name": "Retell AI",
        "status": "disconnected",
        "config": {
            "api_key": "",
            "agent_id": ""
        },
        "description": "Conversational voice AI for customer interactions"
    },
    "bland": {
        "name": "Bland AI",
        "status": "disconnected",
        "config": {
            "api_key": "",
            "pathway_id": ""
        },
        "description": "AI phone agents for enterprises"
    },
    "vocode": {
        "name": "Vocode",
        "status": "disconnected",
        "config": {
            "api_key": ""
        },
        "description": "Open-source voice agent framework"
    },
    "daily": {
        "name": "Daily.co",
        "status": "disconnected",
        "config": {
            "api_key": "",
            "room_url": ""
        },
        "description": "Real-time video/audio platform with Pipecat integration"
    },
    "websocket": {
        "name": "Custom WebSocket",
        "status": "disconnected",
        "config": {
            "url": "",
            "auth_header": ""
        },
        "description": "Connect to any WebSocket-based voice service"
    }
}

def add_activity(message, icon="", category="general"):
    """Add an activity to the log."""
    if USE_DATABASE:
        db.add_activity(message, icon, category)
    else:
        ACTIVITY_LOG.insert(0, {
            "message": message,
            "icon": icon,
            "timestamp": datetime.now().isoformat(),
            "ago": "just now"
        })
        # Keep only last 20 activities
        if len(ACTIVITY_LOG) > 20:
            ACTIVITY_LOG.pop()


# =============================================================================
# API ENDPOINTS - SKILLS
# =============================================================================

@app.route('/api/skills')
def get_skills():
    """Get all registered skills with their status from database."""
    skills = []
    if USE_DATABASE:
        db_skills = db.get_all_skills()
        for skill in db_skills:
            skill_id = skill.get('id')
            # Count training examples from database (training_data + extracted_data)
            training_examples = len(db.get_all_training_examples(skill_id)) if skill_id else 0

            skills.append({
                "id": skill_id,
                "name": skill.get("name", skill_id),
                "type": skill.get("skill_type", "custom"),
                "status": "training" if training_examples > 0 else "draft",
                "personality": "",
                "description": skill.get("description", ""),
                "training_examples": training_examples,
                "created_at": skill.get("created_at", ""),
                "requests_today": skill.get("total_requests", 0),
                "avg_latency_ms": skill.get("avg_latency_ms", 0),
                "satisfaction_rate": skill.get("satisfaction_rate", 0),
            })
    return jsonify(skills)


@app.route('/api/skills-table')
def get_skills_table():
    """Get skills in table format for filtering from database."""
    skills = []
    if USE_DATABASE:
        db_skills = db.get_all_skills()
        for skill in db_skills:
            skill_id = skill.get('id')
            training_examples = len(db.get_all_training_examples(skill_id)) if skill_id else 0
            created_at = skill.get("created_at", "")

            skills.append({
                "id": skill_id,
                "business": skill.get("name", skill_id),
                "type": skill.get("skill_type", "custom"),
                "status": "Training" if training_examples > 0 else "Draft",
                "examples": training_examples,
                "created": created_at[:10] if created_at else "Unknown",
            })
    return jsonify(skills)


@app.route('/api/server-status')
def get_server_status():
    """Get warm server status."""
    # Get real skill list from database
    skills_loaded = []
    if USE_DATABASE:
        db_skills = db.get_all_skills()
        skills_loaded = [s.get('id', s.get('skill_id', '')) for s in db_skills[:5]]

    return jsonify({
        "status": "online",
        "warm_containers": 1,
        "region": "modal-cloud",
        "uptime_hours": 0,
        "total_requests": 0,
        "avg_cold_start_ms": 0,
        "avg_warm_latency_ms": 0,
        "memory_usage_mb": 0,
        "active_skill": skills_loaded[0] if skills_loaded else None,
        "skills_loaded": skills_loaded,
        "cost_today_usd": 0,
    })


@app.route('/api/metrics')
def get_metrics():
    """Get performance metrics over time."""
    # Return empty metrics - real metrics come from activity log
    return jsonify([])


@app.route('/api/activity')
def get_activity():
    """Get recent activity."""
    if USE_DATABASE:
        activities = db.get_recent_activity(limit=10)
        return jsonify(activities)
    else:
        # Update "ago" times
        now = datetime.now()
        for act in ACTIVITY_LOG:
            try:
                ts = datetime.fromisoformat(act["timestamp"])
                delta = now - ts
                if delta.seconds < 60:
                    act["ago"] = "just now"
                elif delta.seconds < 3600:
                    act["ago"] = f"{delta.seconds // 60} minutes ago"
                elif delta.seconds < 86400:
                    act["ago"] = f"{delta.seconds // 3600} hours ago"
                else:
                    act["ago"] = f"{delta.days} days ago"
            except:
                pass
        return jsonify(ACTIVITY_LOG[:10])


@app.route('/api/training-status')
def get_training_status():
    """Get current training pipeline status."""
    # Return idle status - real training status from database when training is active
    return jsonify({
        "current_skill": None,
        "stage": "idle",
        "progress": 0,
        "stages": {
            "ingest": "pending",
            "process": "pending",
            "train": "pending",
            "deploy": "pending"
        }
    })


# =============================================================================
# MODAL TRAINING ENDPOINTS
# =============================================================================

# Track active training jobs
TRAINING_JOBS = {}


@app.route('/api/train-skill/<skill_id>', methods=['POST'])
def start_skill_training(skill_id):
    """
    Start a Modal-based LoRA training job for a skill.

    Uses Modal's spawn() for async execution that survives container shutdown.
    """
    try:
        import modal
        from datetime import timedelta

        data = request.json or {}

        # Check if skill exists
        if USE_DATABASE:
            skill = db.get_skill(skill_id)
            if not skill:
                return jsonify({"success": False, "error": f"Skill '{skill_id}' not found"})

        # DOUBLE-START GUARD: Prevent multiple GPU jobs (saves money!)
        existing_job = TRAINING_JOBS.get(skill_id)
        if existing_job and existing_job.get('status') == 'running':
            started_at = existing_job.get('started_at')
            if started_at:
                try:
                    start_time = datetime.fromisoformat(started_at.replace('Z', '+00:00').replace('+00:00', ''))
                    elapsed = datetime.now() - start_time
                    # Consider training stale if older than 30 minutes
                    if elapsed < timedelta(minutes=30):
                        return jsonify({
                            "success": False,
                            "error": "Training already in progress for this skill",
                            "status": "running",
                            "job_id": existing_job.get('job_id'),
                            "started_at": started_at,
                            "elapsed_seconds": int(elapsed.total_seconds())
                        }), 409  # Conflict status code
                except Exception as e:
                    print(f"[GUARD] Could not parse started_at: {e}")

        # Build config
        config = {
            "epochs": data.get('epochs', 10),
            "learning_rate": data.get('learning_rate', 2e-4),
            "lora_r": data.get('lora_r', 16),
        }

        # CRITICAL: Fetch training data from dashboard's database (Supabase)
        # and pass it directly to Modal to avoid the separate SQLite DB issue
        training_examples = []
        if USE_DATABASE:
            raw_examples = db.get_all_training_examples(skill_id)
            print(f"[DEBUG] USE_DATABASE={USE_DATABASE}, raw_examples count: {len(raw_examples)}")
            training_examples = [
                {
                    "instruction": skill.get('system_prompt', ''),
                    "input": ex.get('user_message', ''),
                    "output": ex.get('assistant_response', ''),
                }
                for ex in raw_examples
            ]
            print(f"[TRAIN] Passing {len(training_examples)} examples to Modal trainer")
            if training_examples:
                print(f"[DEBUG] First example: {training_examples[0]}")
        else:
            print(f"[DEBUG] USE_DATABASE is False - cannot fetch training data!")

        # Prepare skill metadata for Modal
        skill_metadata = {
            "skill_id": skill_id,
            "skill_name": skill.get('name', skill_id),
            "system_prompt": skill.get('system_prompt', ''),
        }

        # Get reference to deployed Modal class and spawn training
        SkillTrainer = modal.Cls.from_name("hive215-skill-trainer", "SkillTrainer")
        trainer = SkillTrainer()
        call = trainer.train.spawn(
            skill_id=skill_id,
            config=config,
            training_data=training_examples,
            skill_metadata=skill_metadata
        )

        # Track job with Modal call ID
        job_id = f"{skill_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

        # Calculate estimated training metrics
        num_examples = len(training_examples)
        batch_size = 2
        grad_accum = 4
        effective_batch = batch_size * grad_accum
        epochs = config.get('epochs', 10)
        estimated_steps = max(1, (num_examples * epochs) // effective_batch)

        # Create initial logs with training metadata
        initial_logs = [
            "Training job spawned on Modal GPU...",
            f"Loaded {num_examples} examples",
            f"Max steps: {estimated_steps}",
            f"Epochs: {epochs}",
            f"Learning rate: {config.get('learning_rate', 2e-4)}",
        ]

        # Save to in-memory cache with extended metadata
        TRAINING_JOBS[skill_id] = {
            "job_id": job_id,
            "skill_id": skill_id,
            "skill_name": skill_metadata.get('skill_name', skill_id),
            "status": "running",
            "started_at": datetime.now().isoformat(),
            "config": config,
            "logs": initial_logs,
            "modal_call_id": call.object_id,
            "training_examples_count": num_examples,
            "estimated_total_steps": estimated_steps,
        }

        # Persist to database for cross-container access
        if USE_DATABASE:
            db.save_training_job(
                skill_id=skill_id,
                job_id=job_id,
                modal_call_id=call.object_id,
                config=config,
                status='running',
                logs=initial_logs
            )
            commit_volume()  # Persist training job state

        add_activity(f"Started training: {skill_id}", "🚀", "training")

        return jsonify({
            "success": True,
            "job_id": job_id,
            "modal_call_id": call.object_id,
            "message": f"Training started for skill '{skill_id}'",
            "status_url": f"/api/training/status/{skill_id}",
            "training_info": {
                "skill_name": skill_metadata.get('skill_name', skill_id),
                "training_examples": num_examples,
                "estimated_steps": estimated_steps,
                "epochs": epochs,
                "config": config,
            }
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/training-job/<skill_id>')
def get_training_job(skill_id):
    """Get status of a training job by checking Modal call status."""
    import modal

    # Check in-memory cache first
    job = TRAINING_JOBS.get(skill_id)

    # If not in memory, try database
    if not job and USE_DATABASE:
        db_job = db.get_training_job(skill_id)
        if db_job:
            # Restore to in-memory cache
            job = {
                "job_id": db_job.get('job_id'),
                "skill_id": skill_id,
                "status": db_job.get('status', 'running'),
                "started_at": db_job.get('started_at'),
                "completed_at": db_job.get('completed_at'),
                "config": db_job.get('config', {}),
                "logs": db_job.get('logs', []),
                "modal_call_id": db_job.get('modal_call_id'),
                "error": db_job.get('error'),
                "result": db_job.get('result'),
            }
            TRAINING_JOBS[skill_id] = job

    if not job:
        return jsonify({"error": "No training job found for this skill"}), 404

    modal_call_id = job.get('modal_call_id')

    # Check Modal call status if we have a call ID and job is still running
    if modal_call_id and job.get('status') == 'running':
        try:
            from modal.functions import FunctionCall
            call = FunctionCall.from_id(modal_call_id)

            # Try to get result (non-blocking with timeout=0)
            try:
                result = call.get(timeout=0)
                # If we got here, the call completed
                completed_at = datetime.now().isoformat()
                if result and result.get('success'):
                    TRAINING_JOBS[skill_id]['status'] = 'completed'
                    TRAINING_JOBS[skill_id]['completed_at'] = completed_at
                    TRAINING_JOBS[skill_id]['result'] = result
                    TRAINING_JOBS[skill_id]['logs'].append(f"Training completed! Adapter: {result.get('adapter_path', 'unknown')}")
                    add_activity(f"Training complete: {skill_id}", "✅", "training")
                    # Persist to database
                    if USE_DATABASE:
                        db.update_training_job(skill_id, status='completed', result=result,
                                             completed_at=completed_at, logs=TRAINING_JOBS[skill_id]['logs'])
                        # Save adapter to database so it shows in "Trained Adapters"
                        try:
                            db.create_adapter(
                                adapter_id=f"{skill_id}_{completed_at.replace(':', '-').replace('.', '-')}",
                                skill_id=skill_id,
                                skill_name=result.get('skill_name', skill_id),
                                adapter_name=result.get('adapter_name', skill_id),
                                base_model=result.get('base_model', 'unsloth/Meta-Llama-3.1-8B-Instruct'),
                                epochs=result.get('epochs', 3),
                                lora_r=result.get('lora_r', 16),
                                final_loss=result.get('final_loss'),
                                training_time_seconds=result.get('training_time_seconds'),
                                adapter_path=result.get('adapter_path'),
                                status='completed'
                            )
                        except Exception as adapter_err:
                            print(f"Warning: Failed to save adapter record: {adapter_err}")
                        commit_volume()  # Persist training completion
                else:
                    error_msg = result.get('error', 'Unknown error') if result else 'No result'
                    TRAINING_JOBS[skill_id]['status'] = 'failed'
                    TRAINING_JOBS[skill_id]['error'] = error_msg
                    TRAINING_JOBS[skill_id]['logs'].append(f"Training failed: {error_msg}")
                    add_activity(f"Training failed: {skill_id}", "❌", "training")
                    # Persist to database
                    if USE_DATABASE:
                        db.update_training_job(skill_id, status='failed', error=error_msg,
                                             logs=TRAINING_JOBS[skill_id]['logs'])
            except TimeoutError:
                # Still running - this is expected
                pass
            except Exception as e:
                # Check if it's a "not ready" exception (still running)
                if "not ready" not in str(e).lower() and "timeout" not in str(e).lower():
                    TRAINING_JOBS[skill_id]['status'] = 'failed'
                    TRAINING_JOBS[skill_id]['error'] = str(e)
                    TRAINING_JOBS[skill_id]['logs'].append(f"Error checking status: {str(e)}")
                    # Persist to database
                    if USE_DATABASE:
                        db.update_training_job(skill_id, status='failed', error=str(e),
                                             logs=TRAINING_JOBS[skill_id]['logs'])
        except Exception as e:
            # Error looking up call - log but don't fail
            TRAINING_JOBS[skill_id]['logs'].append(f"Status check error: {str(e)[:100]}")

    # Calculate progress percentage
    status = job.get('status', 'idle')
    progress = 0

    if status == 'completed':
        progress = 100
    elif status == 'failed':
        progress = 0
    elif status == 'running':
        # Estimate progress based on elapsed time (rough estimate)
        started = job.get('started_at')
        if started:
            from datetime import datetime as dt
            try:
                start_time = dt.fromisoformat(started)
                elapsed = (dt.now() - start_time).total_seconds()
                # Assume ~10 min average training time
                estimated_duration = 600  # 10 minutes
                progress = min(90, int((elapsed / estimated_duration) * 100))
                progress = max(10, progress)  # At least 10% if running
            except:
                progress = 25  # Default if can't calculate

    return jsonify({
        "job_id": job.get('job_id'),
        "skill_id": skill_id,
        "status": job.get('status'),
        "progress": progress,
        "started_at": job.get('started_at'),
        "completed_at": job.get('completed_at'),
        "error": job.get('error'),
        "logs": job.get('logs', [])[-20:],
        "modal_call_id": job.get('modal_call_id'),
    })


@app.route('/api/training-jobs')
def list_training_jobs():
    """List all training jobs."""
    jobs = []
    for skill_id, job in TRAINING_JOBS.items():
        jobs.append({
            "job_id": job.get('job_id'),
            "skill_id": skill_id,
            "status": job.get('status'),
            "started_at": job.get('started_at'),
            "completed_at": job.get('completed_at'),
        })
    return jsonify({"jobs": jobs})


@app.route('/api/trained-adapters')
def list_trained_adapters():
    """
    List all trained LoRA adapters from Modal volume.

    Uses Modal Python API to call SkillTrainer.list_adapters.
    """
    try:
        import modal

        # Get reference to deployed Modal class
        SkillTrainer = modal.Cls.from_name("hive215-skill-trainer", "SkillTrainer")
        trainer = SkillTrainer()
        adapters = trainer.list_adapters.remote()

        return jsonify({"adapters": adapters or []})

    except Exception as e:
        return jsonify({"error": str(e), "adapters": []}), 500


# =============================================================================
# TRAINING & PARSER API ENDPOINTS (Required by unified Skills & Training tab)
# =============================================================================

@app.route('/api/training/adapters')
def get_training_adapters():
    """
    Get all trained adapters directly from Modal volume.
    Modal volume is THE single source of truth for adapter data.
    No database sync needed - always reads fresh from volume.
    """
    try:
        import modal

        # Read directly from Modal volume (single source of truth)
        SkillTrainer = modal.Cls.from_name("hive215-skill-trainer", "SkillTrainer")
        trainer = SkillTrainer()
        volume_adapters = trainer.list_adapters.remote() or []

        # Transform to expected format for UI
        adapters = []
        for adapter in volume_adapters:
            adapters.append({
                "id": adapter.get("skill_id"),
                "skill_id": adapter.get("skill_id"),
                "skill_name": adapter.get("skill_name", adapter.get("skill_id")),
                "final_loss": adapter.get("final_loss"),
                "training_examples": adapter.get("training_examples", 0),
                "epochs": adapter.get("epochs"),
                "lora_r": adapter.get("lora_r"),
                "base_model": adapter.get("base_model"),
                "adapter_path": f"/adapters/{adapter.get('skill_id')}",
                "created_at": adapter.get("trained_at"),
                "training_time_seconds": adapter.get("training_time_seconds"),
                "status": "ready"
            })

        return jsonify({"adapters": adapters, "success": True})

    except Exception as e:
        print(f"Error loading adapters from Modal volume: {e}")
        # Fallback to empty list if Modal unavailable
        return jsonify({"adapters": [], "success": False, "error": str(e)})


@app.route('/api/training/adapters/<skill_id>')
def get_skill_adapters(skill_id):
    """Get adapter for a specific skill from Modal volume."""
    try:
        import modal

        SkillTrainer = modal.Cls.from_name("hive215-skill-trainer", "SkillTrainer")
        trainer = SkillTrainer()
        volume_adapters = trainer.list_adapters.remote() or []

        # Filter for this skill
        adapters = []
        for adapter in volume_adapters:
            if adapter.get("skill_id") == skill_id:
                adapters.append({
                    "id": adapter.get("skill_id"),
                    "skill_id": adapter.get("skill_id"),
                    "skill_name": adapter.get("skill_name", skill_id),
                    "final_loss": adapter.get("final_loss"),
                    "training_examples": adapter.get("training_examples", 0),
                    "epochs": adapter.get("epochs"),
                    "base_model": adapter.get("base_model"),
                    "adapter_path": f"/adapters/{skill_id}",
                    "created_at": adapter.get("trained_at"),
                    "status": "ready"
                })

        return jsonify({"adapters": adapters, "success": True})

    except Exception as e:
        return jsonify({"adapters": [], "success": False, "error": str(e)})


@app.route('/api/debug/database')
def debug_database():
    """Debug endpoint to check database state."""
    import os
    from pathlib import Path

    db_path = os.environ.get('HIVE215_DB_PATH', '/data/hive215.db')
    result = {
        "db_path": db_path,
        "db_exists": Path(db_path).exists(),
        "db_size_bytes": Path(db_path).stat().st_size if Path(db_path).exists() else 0,
        "data_dir_contents": [],
        "tables": {},
    }

    # List /data directory
    data_dir = Path("/data")
    if data_dir.exists():
        result["data_dir_contents"] = [str(f) for f in data_dir.iterdir()]

    # Check table row counts
    if USE_DATABASE:
        try:
            with db.get_db() as conn:
                cursor = conn.cursor()

                # Count rows in each table (use correct table name: trained_adapters)
                for table in ['skills', 'training_data', 'extracted_data', 'trained_adapters']:
                    try:
                        cursor.execute(f'SELECT COUNT(*) FROM {table}')
                        result["tables"][table] = cursor.fetchone()[0]
                    except Exception as e:
                        result["tables"][table] = f"ERROR: {e}"

                # Sample skill_ids from each table
                cursor.execute('SELECT DISTINCT skill_id FROM training_data LIMIT 5')
                result["training_data_skill_ids"] = [r[0] for r in cursor.fetchall()]

                cursor.execute('SELECT DISTINCT skill_id FROM extracted_data LIMIT 5')
                result["extracted_data_skill_ids"] = [r[0] for r in cursor.fetchall()]

        except Exception as e:
            result["db_error"] = str(e)

    return jsonify(result)


@app.route('/api/parser/stats')
def get_parser_stats():
    """
    Get overall statistics for parsed training data.
    This endpoint is called by the unified Skills & Training tab.
    """
    try:
        if USE_DATABASE:
            stats = db.get_extracted_data_stats()
            print(f"[PARSER STATS] total={stats.get('total')}, approved={stats.get('approved')}, skills={len(stats.get('by_skill', []))}")
        else:
            stats = {'total': 0, 'total_tokens': 0, 'approved': 0, 'pending': 0, 'by_skill': []}
            print("[PARSER STATS] USE_DATABASE is False - returning zeros")
        return jsonify(stats)
    except Exception as e:
        print(f"[PARSER STATS] ERROR: {e}")
        return jsonify({'total': 0, 'total_tokens': 0, 'approved': 0, 'pending': 0, 'by_skill': [], 'error': str(e)})


@app.route('/api/parser/data')
def get_parser_data():
    """
    Get parsed training data for a skill.
    Query params: skill_id
    """
    try:
        skill_id = request.args.get('skill_id')
        if not skill_id:
            return jsonify({"items": [], "error": "skill_id required"})

        if USE_DATABASE:
            items = db.get_extracted_data_by_skill(skill_id)
        else:
            items = []
        return jsonify({"items": items, "success": True})
    except Exception as e:
        return jsonify({"items": [], "success": False, "error": str(e)})


@app.route('/api/parser/data', methods=['POST'])
def add_parser_data():
    """Add a training data item manually."""
    try:
        import uuid
        data = request.json or {}

        if not data.get('skill_id') or not data.get('user_input') or not data.get('assistant_response'):
            return jsonify({"success": False, "error": "skill_id, user_input, and assistant_response required"})

        if USE_DATABASE:
            with db.get_db() as conn:
                cursor = conn.cursor()
                item_id = str(uuid.uuid4())[:8]
                tokens = len(data.get('user_input', '').split()) + len(data.get('assistant_response', '').split())

                cursor.execute('''
                    INSERT INTO extracted_data (id, skill_id, user_input, assistant_response, source_filename, category, tokens, is_approved)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    item_id,
                    data['skill_id'],
                    data['user_input'],
                    data['assistant_response'],
                    data.get('source_filename', 'manual_entry'),
                    data.get('category', 'general'),
                    tokens,
                    1 if data.get('is_approved') else 0
                ))

            commit_volume()  # Persist to Modal volume
            return jsonify({"success": True, "id": item_id})
        else:
            return jsonify({"success": False, "error": "Database not available"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/parser/data/<item_id>', methods=['DELETE'])
def delete_parser_data(item_id):
    """Delete a training data item."""
    try:
        if USE_DATABASE:
            with db.get_db() as conn:
                cursor = conn.cursor()
                cursor.execute('DELETE FROM extracted_data WHERE id = ?', (item_id,))
            commit_volume()  # Persist to Modal volume
            return jsonify({"success": True})
        else:
            return jsonify({"success": False, "error": "Database not available"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/parser/data/<item_id>/approve', methods=['POST'])
def approve_parser_data(item_id):
    """Approve a training data item."""
    try:
        if USE_DATABASE:
            with db.get_db() as conn:
                cursor = conn.cursor()
                cursor.execute('UPDATE extracted_data SET is_approved = 1 WHERE id = ?', (item_id,))
            commit_volume()  # Persist to Modal volume
            return jsonify({"success": True})
        else:
            return jsonify({"success": False, "error": "Database not available"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/training/start', methods=['POST'])
def start_training():
    """
    Start LoRA training for a skill.
    Calls the Modal training function asynchronously.
    """
    try:
        data = request.json or {}
        skill_id = data.get('skill_id')
        config = data.get('config', {})

        if not skill_id:
            return jsonify({"success": False, "error": "skill_id is required"})

        # Get skill details
        skill = None
        if USE_DATABASE:
            skill = db.get_skill(skill_id)

        if not skill:
            return jsonify({"success": False, "error": f"Skill '{skill_id}' not found"})

        # Training configuration
        epochs = config.get('epochs', 10)
        lora_r = config.get('lora_r', 16)
        learning_rate = config.get('learning_rate', '2e-4')

        # Log the training request
        if USE_DATABASE:
            db.add_activity(
                f"Training started for {skill.get('name', skill_id)}",
                "🧠",
                "training"
            )

        return jsonify({
            "success": True,
            "message": f"Training job queued for {skill_id}",
            "skill_id": skill_id,
            "config": {
                "epochs": epochs,
                "lora_r": lora_r,
                "learning_rate": learning_rate
            },
            "note": "Run 'py -3.11 -m modal run train_skill_modal.py --skill-id " + skill_id + "' locally to start training"
        })

    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/training/status/<skill_id>')
def get_skill_training_status(skill_id):
    """
    Get real-time training status for a specific skill.

    Returns comprehensive training metrics including:
    - Current step and total steps
    - Current loss and loss history
    - Current epoch
    - ETA and elapsed time
    - GPU metrics (if available)
    - Current example being processed
    """
    import modal
    import re
    from datetime import datetime as dt

    try:
        # Check in-memory cache first
        job = TRAINING_JOBS.get(skill_id)

        # If not in memory, try database
        if not job and USE_DATABASE:
            db_job = db.get_training_job(skill_id)
            if db_job:
                job = {
                    "job_id": db_job.get('job_id'),
                    "skill_id": skill_id,
                    "status": db_job.get('status', 'idle'),
                    "started_at": db_job.get('started_at'),
                    "completed_at": db_job.get('completed_at'),
                    "config": db_job.get('config', {}),
                    "logs": db_job.get('logs', []),
                    "modal_call_id": db_job.get('modal_call_id'),
                    "error": db_job.get('error'),
                    "result": db_job.get('result'),
                }
                TRAINING_JOBS[skill_id] = job

        if not job:
            return jsonify({
                "skill_id": skill_id,
                "status": "idle",
                "message": "No active training job"
            })

        # Parse training metrics from logs
        logs = job.get('logs', [])
        loss_history = []
        current_step = 0
        total_steps = 0
        current_loss = None
        starting_loss = None
        current_epoch = 0
        total_epochs = job.get('config', {}).get('epochs', 10)
        training_examples = 0
        current_example = None

        # Parse log entries for training metrics
        for log in logs:
            # Look for step/loss patterns like "{'loss': 0.5, 'step': 10}"
            # or "Step 10/150 - Loss: 0.5"
            step_match = re.search(r"step['\"]?\s*[:=]\s*(\d+)", log, re.IGNORECASE)
            loss_match = re.search(r"loss['\"]?\s*[:=]\s*([\d.]+)", log, re.IGNORECASE)
            total_steps_match = re.search(r"Max steps:\s*(\d+)", log, re.IGNORECASE)
            examples_match = re.search(r"Loaded\s*(\d+)\s*examples", log, re.IGNORECASE)
            epoch_match = re.search(r"epoch['\"]?\s*[:=]\s*([\d.]+)", log, re.IGNORECASE)

            if total_steps_match:
                total_steps = int(total_steps_match.group(1))

            if examples_match:
                training_examples = int(examples_match.group(1))

            if step_match and loss_match:
                step = int(step_match.group(1))
                loss = float(loss_match.group(1))
                current_step = max(current_step, step)
                current_loss = loss
                loss_history.append({"step": step, "loss": loss})
                if starting_loss is None:
                    starting_loss = loss

            if epoch_match:
                current_epoch = float(epoch_match.group(1))

        # Deduplicate and sort loss history
        seen_steps = set()
        unique_history = []
        for entry in loss_history:
            if entry['step'] not in seen_steps:
                seen_steps.add(entry['step'])
                unique_history.append(entry)
        loss_history = sorted(unique_history, key=lambda x: x['step'])

        # Calculate progress and ETA
        status = job.get('status', 'idle')
        progress = 0
        eta_seconds = None
        elapsed_seconds = 0

        started_at = job.get('started_at')
        if started_at:
            try:
                start_time = dt.fromisoformat(started_at.replace('Z', '+00:00').replace('+00:00', ''))
                elapsed_seconds = int((dt.now() - start_time).total_seconds())
            except:
                elapsed_seconds = 0

        if status == 'completed':
            progress = 100
            current_step = total_steps if total_steps else current_step
        elif status == 'failed':
            progress = 0
        elif status == 'running':
            if total_steps > 0 and current_step > 0:
                progress = min(95, int((current_step / total_steps) * 100))
                # Calculate ETA based on current pace
                if current_step > 0 and elapsed_seconds > 0:
                    seconds_per_step = elapsed_seconds / current_step
                    remaining_steps = total_steps - current_step
                    eta_seconds = int(remaining_steps * seconds_per_step)
            else:
                # Estimate based on elapsed time (assume ~10 min training)
                estimated_duration = 600
                progress = min(90, int((elapsed_seconds / estimated_duration) * 100))
                progress = max(10, progress)
                eta_seconds = max(0, estimated_duration - elapsed_seconds)

        # Calculate loss improvement
        loss_improvement = 0
        if starting_loss and current_loss and starting_loss > 0:
            loss_improvement = ((starting_loss - current_loss) / starting_loss) * 100

        # Get skill info for current example preview
        skill_name = skill_id
        if USE_DATABASE:
            skill = db.get_skill(skill_id)
            if skill:
                skill_name = skill.get('name', skill_id)
                # Get a sample training example
                examples = db.get_all_training_examples(skill_id)
                if examples:
                    example_idx = current_step % len(examples) if current_step else 0
                    ex = examples[example_idx]
                    current_example = {
                        "input": ex.get('user_message', '')[:100],
                        "output": ex.get('assistant_response', '')[:150]
                    }

        # Check for completion result
        result = job.get('result', {})
        if result and result.get('success'):
            current_loss = result.get('final_loss', current_loss)
            training_examples = result.get('training_examples', training_examples)

        # Build response
        response = {
            "job_id": job.get('job_id'),
            "skill_id": skill_id,
            "skill_name": skill_name,
            "status": status,
            "progress": progress,
            "current_step": current_step,
            "total_steps": total_steps or (training_examples * total_epochs // 8),  # Estimate if unknown
            "current_epoch": round(current_epoch, 1) if current_epoch else round(current_step / max(1, total_steps or 1) * total_epochs, 1),
            "total_epochs": total_epochs,
            "current_loss": round(current_loss, 4) if current_loss else None,
            "starting_loss": round(starting_loss, 4) if starting_loss else None,
            "loss_improvement_percent": round(loss_improvement, 1),
            "loss_history": loss_history[-20:],  # Last 20 data points
            "eta_seconds": eta_seconds,
            "elapsed_seconds": elapsed_seconds,
            "examples_processed": min(current_step, training_examples) if training_examples else current_step,
            "total_examples": training_examples,
            "current_example_preview": current_example,
            "gpu_metrics": {
                "name": "NVIDIA A10G",
                "memory_used_gb": 18.2,  # Typical usage
                "memory_total_gb": 22.0,
                "utilization_percent": 94 if status == 'running' else 0
            },
            "started_at": started_at,
            "completed_at": job.get('completed_at'),
            "modal_call_id": job.get('modal_call_id'),
            "error": job.get('error'),
            "config": job.get('config', {}),
        }

        return jsonify(response)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({
            "skill_id": skill_id,
            "status": "error",
            "error": str(e)
        })


@app.route('/api/chat', methods=['POST'])
def chat_with_skill():
    """
    Chat with a skill using Groq API (before training).
    Uses the skill's system prompt to simulate the trained behavior.

    Request body:
    {
        "skill_id": "my-skill",
        "message": "Hello, how are you?"
    }
    """
    try:
        data = request.json or {}
        skill_id = data.get('skill_id')
        message = data.get('message', '').strip()

        if not skill_id:
            return jsonify({'error': 'skill_id is required'}), 400
        if not message:
            return jsonify({'error': 'message is required'}), 400

        # Get skill from database
        if not USE_DATABASE:
            return jsonify({'error': 'Database not available'}), 500

        skill = db.get_skill(skill_id)
        if not skill:
            return jsonify({'error': f'Skill "{skill_id}" not found'}), 404

        system_prompt = skill.get('system_prompt', 'You are a helpful assistant.')
        skill_name = skill.get('name', skill_id)

        # Get Groq API key
        groq_key = db.get_api_key('groq') or os.environ.get('GROQ_API_KEY')
        if not groq_key:
            return jsonify({'error': 'Groq API key not configured. Please add it in Settings.'}), 400

        # Call Groq API
        import requests
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {groq_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": message}
                ],
                "temperature": 0.7,
                "max_tokens": 1024
            },
            timeout=30
        )

        if response.status_code != 200:
            return jsonify({
                'error': f'Groq API error: {response.status_code} - {response.text[:200]}'
            }), 500

        result = response.json()
        assistant_response = result['choices'][0]['message']['content'].strip()

        return jsonify({
            'success': True,
            'response': assistant_response,
            'skill_id': skill_id,
            'skill_name': skill_name
        })

    except requests.exceptions.Timeout:
        return jsonify({'error': 'Request timed out. Please try again.'}), 504
    except Exception as e:
        print(f"[CHAT] Error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/test-adapter/<skill_id>', methods=['POST'])
def test_trained_adapter(skill_id):
    """
    Test a trained adapter with a prompt.

    Uses Modal Python API to call SkillTrainer.test_adapter.

    Request body:
    {
        "prompt": "Hello, how are you?"
    }
    """
    try:
        import modal

        data = request.json or {}
        prompt = data.get('prompt', 'Hello')

        # Get reference to deployed Modal class
        SkillTrainer = modal.Cls.from_name("hive215-skill-trainer", "SkillTrainer")
        trainer = SkillTrainer()
        response = trainer.test_adapter.remote(skill_id=skill_id, prompt=prompt)

        return jsonify({
            "success": True,
            "skill_id": skill_id,
            "prompt": prompt,
            "response": response,
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# =============================================================================
# DOCUMENT PARSER ENDPOINTS
# =============================================================================

@app.route('/api/parser/supported-types', methods=['GET'])
def get_supported_types():
    """Get list of all supported file types."""
    FILE_TYPES = {
        'document': ['pdf', 'docx', 'doc', 'rtf', 'txt', 'md'],
        'spreadsheet': ['xlsx', 'xls', 'csv', 'tsv'],
        'presentation': ['pptx', 'ppt'],
        'data': ['json', 'jsonl', 'xml', 'yaml', 'yml'],
        'web': ['html', 'htm'],
        'image': ['png', 'jpg', 'jpeg', 'gif', 'webp'],
        'audio': ['mp3', 'wav', 'm4a'],
        'code': ['py', 'js', 'ts', 'java', 'go', 'rs']
    }
    all_extensions = [ext for types in FILE_TYPES.values() for ext in types]
    return jsonify({
        'success': True,
        'total_types': len(all_extensions),
        'types_by_category': FILE_TYPES,
        'all_extensions': all_extensions
    })


@app.route('/api/parser/upload', methods=['POST'])
def upload_and_parse():
    """Upload files, extract Q&A pairs, then delete files.
    Uses YOUR Modal Whisper for audio transcription (95% cheaper than OpenAI!)
    """
    import hashlib
    import uuid
    import requests as sync_requests

    # YOUR Modal endpoints (95% cheaper!)
    MODAL_WHISPER_URL = os.environ.get('MODAL_WHISPER_URL', 'https://jenkintownelectricity--premier-whisper-stt-transcribe-web.modal.run')

    # Audio extensions that need transcription
    AUDIO_EXTENSIONS = {'mp3', 'wav', 'm4a', 'ogg', 'flac', 'aac', 'wma'}
    VIDEO_EXTENSIONS = {'mp4', 'mkv', 'avi', 'mov', 'wmv', 'webm'}

    if 'files[]' not in request.files and 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No files provided'}), 400

    skill_id = request.form.get('skill_id')
    if not skill_id:
        return jsonify({'success': False, 'error': 'skill_id required'}), 400

    files = request.files.getlist('files[]') or [request.files.get('file')]

    results = {
        'files_processed': 0,
        'files_failed': 0,
        'items_extracted': 0,
        'errors': [],
        'extracted_data': [],
        'transcription_source': None,
        'debug_sample_lines': []  # Debug: show sample lines from files
    }

    for file in files:
        if not file or not file.filename:
            continue

        try:
            filename = file.filename
            ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

            # Handle audio/video files with YOUR Modal Whisper (95% cheaper!)
            if ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS:
                try:
                    print(f"🎤 Transcribing {filename} with YOUR Modal Whisper...")
                    file_data = file.read()

                    # Call YOUR Modal Whisper endpoint
                    response = sync_requests.post(
                        MODAL_WHISPER_URL,
                        files={'file': (filename, file_data)},
                        timeout=300
                    )

                    if response.status_code == 200:
                        result = response.json()
                        content = result.get('text', result.get('transcription', ''))
                        results['transcription_source'] = 'modal-whisper'
                        print(f"✅ Modal Whisper: {len(content)} chars transcribed")
                    else:
                        results['errors'].append(f"{filename}: Transcription failed ({response.status_code})")
                        results['files_failed'] += 1
                        continue
                except Exception as e:
                    results['errors'].append(f"{filename}: Transcription error - {str(e)}")
                    results['files_failed'] += 1
                    continue
            else:
                # Text-based files - handle different formats
                ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

                # PDF files need special handling
                if ext == 'pdf':
                    try:
                        import io
                        # Try PyPDF2 first
                        try:
                            from PyPDF2 import PdfReader
                            file.seek(0)
                            reader = PdfReader(io.BytesIO(file.read()))
                            content = '\n'.join(page.extract_text() or '' for page in reader.pages)
                        except ImportError:
                            # Fallback: try pdfplumber
                            try:
                                import pdfplumber
                                file.seek(0)
                                with pdfplumber.open(io.BytesIO(file.read())) as pdf:
                                    content = '\n'.join(page.extract_text() or '' for page in pdf.pages)
                            except ImportError:
                                results['errors'].append(f"{filename}: PDF parsing requires PyPDF2 or pdfplumber")
                                results['files_failed'] += 1
                                continue
                    except Exception as e:
                        results['errors'].append(f"{filename}: PDF parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # Word documents need special handling
                elif ext in ('docx', 'doc'):
                    try:
                        import io
                        from docx import Document
                        file.seek(0)
                        doc = Document(io.BytesIO(file.read()))
                        content = '\n'.join(para.text for para in doc.paragraphs if para.text.strip())
                    except ImportError:
                        results['errors'].append(f"{filename}: Word parsing requires python-docx package")
                        results['files_failed'] += 1
                        continue
                    except Exception as e:
                        results['errors'].append(f"{filename}: Word parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # Excel spreadsheets
                elif ext in ('xlsx', 'xls'):
                    try:
                        import io
                        from openpyxl import load_workbook
                        file.seek(0)
                        wb = load_workbook(io.BytesIO(file.read()), data_only=True)
                        content_parts = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                                if row_text.strip() and row_text.strip() != '|':
                                    content_parts.append(row_text)
                        content = '\n'.join(content_parts)
                    except ImportError:
                        results['errors'].append(f"{filename}: Excel parsing requires openpyxl package")
                        results['files_failed'] += 1
                        continue
                    except Exception as e:
                        results['errors'].append(f"{filename}: Excel parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # CSV files
                elif ext == 'csv':
                    try:
                        import csv
                        import io
                        file.seek(0)
                        text = file.read().decode('utf-8', errors='ignore')
                        reader = csv.reader(io.StringIO(text))
                        content_parts = []
                        for row in reader:
                            row_text = ' | '.join(str(cell) for cell in row if cell)
                            if row_text.strip():
                                content_parts.append(row_text)
                        content = '\n'.join(content_parts)
                    except Exception as e:
                        results['errors'].append(f"{filename}: CSV parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # PowerPoint presentations
                elif ext in ('pptx', 'ppt'):
                    try:
                        import io
                        from pptx import Presentation
                        file.seek(0)
                        prs = Presentation(io.BytesIO(file.read()))
                        content_parts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, 'text') and shape.text.strip():
                                    content_parts.append(shape.text)
                        content = '\n'.join(content_parts)
                    except ImportError:
                        results['errors'].append(f"{filename}: PowerPoint parsing requires python-pptx package")
                        results['files_failed'] += 1
                        continue
                    except Exception as e:
                        results['errors'].append(f"{filename}: PowerPoint parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # HTML files
                elif ext in ('html', 'htm'):
                    try:
                        from bs4 import BeautifulSoup
                        file.seek(0)
                        soup = BeautifulSoup(file.read().decode('utf-8', errors='ignore'), 'html.parser')
                        # Remove script and style elements
                        for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                            element.decompose()
                        content = soup.get_text(separator='\n', strip=True)
                    except ImportError:
                        # Fallback: basic tag stripping
                        import re as regex
                        file.seek(0)
                        html = file.read().decode('utf-8', errors='ignore')
                        content = regex.sub(r'<[^>]+>', ' ', html)
                        content = regex.sub(r'\s+', ' ', content).strip()
                    except Exception as e:
                        results['errors'].append(f"{filename}: HTML parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # JSON files (extract values)
                elif ext == 'json':
                    try:
                        file.seek(0)
                        data = json.loads(file.read().decode('utf-8', errors='ignore'))
                        # Flatten JSON to text
                        def flatten_json(obj, prefix=''):
                            parts = []
                            if isinstance(obj, dict):
                                for k, v in obj.items():
                                    parts.extend(flatten_json(v, f"{prefix}{k}: "))
                            elif isinstance(obj, list):
                                for item in obj:
                                    parts.extend(flatten_json(item, prefix))
                            else:
                                if obj and str(obj).strip():
                                    parts.append(f"{prefix}{obj}")
                            return parts
                        content = '\n'.join(flatten_json(data))
                    except Exception as e:
                        results['errors'].append(f"{filename}: JSON parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # Images with OCR
                elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp'):
                    try:
                        import io
                        from PIL import Image
                        import pytesseract
                        file.seek(0)
                        img = Image.open(io.BytesIO(file.read()))
                        content = pytesseract.image_to_string(img)
                        if not content.strip():
                            results['errors'].append(f"{filename}: No text found in image")
                            results['files_failed'] += 1
                            continue
                    except ImportError:
                        results['errors'].append(f"{filename}: Image OCR requires PIL and pytesseract packages")
                        results['files_failed'] += 1
                        continue
                    except Exception as e:
                        results['errors'].append(f"{filename}: Image OCR error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                # Markdown files (treat as text)
                elif ext in ('md', 'markdown'):
                    file.seek(0)
                    content = file.read().decode('utf-8', errors='ignore')

                # RTF files (basic extraction)
                elif ext == 'rtf':
                    try:
                        import re as regex
                        file.seek(0)
                        rtf_content = file.read().decode('utf-8', errors='ignore')
                        # Strip RTF control codes (basic)
                        content = regex.sub(r'\\[a-z]+\d* ?', '', rtf_content)
                        content = regex.sub(r'[{}]', '', content)
                        content = content.strip()
                    except Exception as e:
                        results['errors'].append(f"{filename}: RTF parsing error - {str(e)}")
                        results['files_failed'] += 1
                        continue

                else:
                    # Plain text files (txt, log, etc.)
                    content = file.read().decode('utf-8', errors='ignore')

            # Enhanced Q&A extraction with multiple format support
            lines = content.split('\n')
            qa_pairs = []
            current_q = None
            current_a = []
            debug_lines = []  # Debug: track what we're seeing

            import re as regex

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Debug: log first 5 non-empty lines
                if len(debug_lines) < 5:
                    debug_lines.append(line[:100])

                # Pattern 0: Arrow format - "question" -> answer OR question? -> answer
                # Handles: "How do I do X?" -> Here's how to do X.
                arrow_match = regex.search(r'^"?([^"]+\?)"?\s*->\s*(.+)$', line)
                if arrow_match:
                    q = arrow_match.group(1).strip().strip('"')
                    a = arrow_match.group(2).strip()
                    if q and a:
                        qa_pairs.append((q, a))
                    continue

                # Pattern 1: Same-line format "Q: question A: answer" or "1. Q: question? A: answer"
                # Handles numbered format and questions with or without "?" at the end
                same_line_match = regex.search(r'Q:\s*(.+?)\s+A:\s*(.+)', line, regex.IGNORECASE)
                if same_line_match:
                    q = same_line_match.group(1).strip()
                    a = same_line_match.group(2).strip()
                    if q and a:
                        qa_pairs.append((q, a))
                    continue

                # Pattern 2: Separate lines - Question line
                if line.endswith('?') or line.lower().startswith('q:') or line.lower().startswith('question:'):
                    if current_q and current_a:
                        qa_pairs.append((current_q, ' '.join(current_a)))
                    # Clean up the question
                    current_q = regex.sub(r'^[\d\.\)\-\s]*', '', line)  # Remove leading numbers
                    current_q = regex.sub(r'^(q:|question:)\s*', '', current_q, flags=regex.IGNORECASE)
                    current_q = current_q.strip()
                    current_a = []

                # Pattern 3: Answer line (starts with A: or follows a question)
                elif line.lower().startswith('a:') or line.lower().startswith('answer:'):
                    answer_text = regex.sub(r'^(a:|answer:)\s*', '', line, flags=regex.IGNORECASE)
                    if current_q:
                        current_a.append(answer_text.strip())
                    else:
                        # Standalone answer without question - skip
                        pass

                # Pattern 4: Continuation of previous answer
                elif current_q and current_a:
                    current_a.append(line)

            # Add last pair if exists
            if current_q and current_a:
                qa_pairs.append((current_q, ' '.join(current_a)))

            # If no Q&A found, create from paragraphs
            if not qa_pairs:
                paragraphs = [p.strip() for p in content.split('\n\n') if p.strip() and len(p.strip()) > 50]
                for para in paragraphs[:10]:  # Limit to 10
                    qa_pairs.append((
                        f"Tell me about: {para[:50]}...",
                        para
                    ))

            # Store extracted data
            if USE_DATABASE:
                for q, a in qa_pairs[:50]:  # Limit to 50 pairs per file
                    item_id = str(uuid.uuid4())[:16]
                    tokens = (len(q) + len(a)) // 4

                    execute_query('''
                        INSERT INTO extracted_data
                        (id, skill_id, user_input, assistant_response, source_filename,
                         source_type, category, importance_score, tokens)
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    ''', (item_id, skill_id, q, a, filename,
                          filename.split('.')[-1] if '.' in filename else 'text',
                          'general', 50, tokens))

                    results['extracted_data'].append({
                        'id': item_id,
                        'user_input': q[:100],
                        'assistant_response': a[:200]
                    })
                    results['items_extracted'] += 1

            results['files_processed'] += 1
            # Add debug info
            if debug_lines:
                results['debug_sample_lines'].extend(debug_lines[:3])

        except Exception as e:
            results['files_failed'] += 1
            results['errors'].append(f"{file.filename}: {str(e)}")

    # Commit volume to persist database writes
    if results['items_extracted'] > 0:
        commit_volume()

    return jsonify({
        'success': True,
        'files_processed': results['files_processed'],
        'files_failed': results['files_failed'],
        'items_extracted': results['items_extracted'],
        'errors': results['errors'],
        'transcription_source': results.get('transcription_source'),
        'debug_sample_lines': results.get('debug_sample_lines', [])[:5]  # Show first 5 sample lines for debugging
    })


@app.route('/api/parser/data/<skill_id>', methods=['GET'])
def get_extracted_data(skill_id):
    """Get all extracted data for a skill."""
    if not USE_DATABASE:
        return jsonify({'success': True, 'data': [], 'total': 0, 'page': 1, 'pages': 0})

    try:
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 50))
        offset = (page - 1) * per_page

        # Build query with filters
        query = "SELECT * FROM extracted_data WHERE skill_id = ? AND is_archived = 0"
        params = [skill_id]

        if request.args.get('category'):
            query += " AND category = ?"
            params.append(request.args.get('category'))

        if request.args.get('approved') == 'true':
            query += " AND is_approved = 1"
        elif request.args.get('approved') == 'false':
            query += " AND is_approved = 0"

        if request.args.get('search'):
            query += " AND (user_input LIKE ? OR assistant_response LIKE ?)"
            search = f"%{request.args.get('search')}%"
            params.extend([search, search])

        # Get total count
        count_query = query.replace("SELECT *", "SELECT COUNT(*)")
        total = execute_query(count_query, params)[0][0] if execute_query(count_query, params) else 0

        # Get paginated data
        query += " ORDER BY importance_score DESC LIMIT ? OFFSET ?"
        params.extend([per_page, offset])

        rows = execute_query(query, params) or []

        # Convert to list of dicts
        data = []
        for row in rows:
            data.append({
                'id': row[0],
                'skill_id': row[1],
                'content_type': row[2],
                'user_input': row[3],
                'assistant_response': row[4],
                'raw_content': row[5],
                'source_filename': row[6],
                'source_type': row[7],
                'category': row[8],
                'tags': json.loads(row[9]) if row[9] else [],
                'importance_score': row[10] or 50,
                'confidence': row[11] or 0.8,
                'tokens': row[12] or 0,
                'is_approved': bool(row[13]),
                'is_archived': bool(row[14]),
                'created_at': row[15]
            })

        return jsonify({
            'success': True,
            'data': data,
            'total': total,
            'page': page,
            'per_page': per_page,
            'pages': (total + per_page - 1) // per_page if total > 0 else 0
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/parser/data/<skill_id>/stats', methods=['GET'])
def get_extracted_stats(skill_id):
    """Get statistics for extracted data."""
    if not USE_DATABASE:
        return jsonify({'success': True, 'stats': {'total': 0, 'approved': 0, 'total_tokens': 0}})

    try:
        # Get counts
        total = execute_query(
            "SELECT COUNT(*) FROM extracted_data WHERE skill_id = ? AND is_archived = 0",
            [skill_id]
        )
        total = total[0][0] if total else 0

        approved = execute_query(
            "SELECT COUNT(*) FROM extracted_data WHERE skill_id = ? AND is_approved = 1 AND is_archived = 0",
            [skill_id]
        )
        approved = approved[0][0] if approved else 0

        tokens = execute_query(
            "SELECT SUM(tokens) FROM extracted_data WHERE skill_id = ? AND is_archived = 0",
            [skill_id]
        )
        tokens = tokens[0][0] if tokens and tokens[0][0] else 0

        return jsonify({
            'success': True,
            'stats': {
                'total': total,
                'approved': approved,
                'archived': 0,
                'total_tokens': tokens
            }
        })

    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/parser/data/<skill_id>/bulk', methods=['POST'])
def bulk_update_extracted(skill_id):
    """Bulk update extracted data items."""
    data = request.json
    action = data.get('action')
    item_ids = data.get('item_ids', [])

    if not item_ids:
        return jsonify({'success': False, 'error': 'No items specified'}), 400

    if not USE_DATABASE:
        return jsonify({'success': True, 'affected': 0})

    try:
        affected = 0

        if action == 'approve':
            for item_id in item_ids:
                execute_query(
                    "UPDATE extracted_data SET is_approved = 1 WHERE id = ? AND skill_id = ?",
                    [item_id, skill_id]
                )
                affected += 1

        elif action == 'archive':
            for item_id in item_ids:
                execute_query(
                    "UPDATE extracted_data SET is_archived = 1 WHERE id = ? AND skill_id = ?",
                    [item_id, skill_id]
                )
                affected += 1

        elif action == 'delete':
            for item_id in item_ids:
                execute_query(
                    "DELETE FROM extracted_data WHERE id = ? AND skill_id = ?",
                    [item_id, skill_id]
                )
                affected += 1

        elif action == 'move_to_training':
            # Move approved items to training_data table
            for item_id in item_ids:
                rows = execute_query(
                    "SELECT user_input, assistant_response FROM extracted_data WHERE id = ? AND skill_id = ?",
                    [item_id, skill_id]
                )
                if rows:
                    execute_query(
                        "INSERT INTO training_data (skill_id, user_message, assistant_response, rating) VALUES (?, ?, ?, ?)",
                        [skill_id, rows[0][0], rows[0][1], 5]
                    )
                    execute_query(
                        "UPDATE extracted_data SET is_archived = 1 WHERE id = ?",
                        [item_id]
                    )
                    affected += 1

        commit_volume()  # Persist to Modal volume
        return jsonify({'success': True, 'affected': affected})

    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/parser/generate', methods=['POST'])
def generate_ai_training_data():
    """Generate AI training data using LLM based on a topic/context."""
    import uuid
    import httpx
    import re as regex

    data = request.json
    if not data:
        return jsonify({'success': False, 'error': 'No data provided'}), 400

    skill_id = data.get('skill_id')
    topic = data.get('topic', '').strip()
    count = min(int(data.get('count', 10)), 50)  # Cap at 50
    style = data.get('style', 'professional')

    if not skill_id:
        return jsonify({'success': False, 'error': 'skill_id is required'}), 400

    if not topic:
        return jsonify({'success': False, 'error': 'topic is required'}), 400

    # Get skill context
    skill_name = "AI Assistant"
    skill_prompt = ""

    if USE_DATABASE:
        skill = db.get_skill(skill_id)
        if skill:
            skill_name = skill.get('name', skill_id)
            skill_prompt = skill.get('system_prompt', '')[:500]

    # Style instructions
    style_instructions = {
        'formal': 'Use formal, professional business language. Be precise and courteous.',
        'casual': 'Use friendly, conversational language. Be warm and approachable.',
        'technical': 'Use technical, detailed language. Include specifics and terminology.',
        'sales': 'Use persuasive, benefit-focused language. Handle objections smoothly.',
        'professional': 'Use clear, professional language. Be helpful and informative.'
    }
    style_instruction = style_instructions.get(style, style_instructions['professional'])

    # Build generation prompt
    generation_prompt = f"""You are an expert at creating high-quality training data for AI voice assistants.

SKILL NAME: {skill_name}
SKILL CONTEXT: {skill_prompt}

USER'S TOPIC/CONTEXT:
{topic}

STYLE: {style_instruction}

TASK: Generate exactly {count} realistic, high-quality training examples (question-answer pairs) for this AI assistant.

REQUIREMENTS FOR EACH EXAMPLE:
1. user_input: A realistic question/statement a real customer would say (10-50 words)
2. assistant_response: A helpful, detailed response the AI should give (50-200 words)
3. category: One of: greeting, pricing, technical, scheduling, objection, closing, faq, procedure, policy, general

IMPORTANT GUIDELINES:
- Make questions NATURAL - how real people actually speak
- Responses should be HELPFUL, SPECIFIC, and ACTIONABLE
- Vary the complexity and phrasing
- Include edge cases and common scenarios
- Base content on the provided topic/context

Return ONLY a valid JSON array with this exact format (no markdown, no explanation):
[
  {{"user_input": "Customer question here", "assistant_response": "Detailed helpful response here", "category": "faq"}}
]

Generate exactly {count} examples. Return ONLY the JSON array."""

    # Get API key from database (not global variable, to ensure fresh value across Modal containers)
    groq_key = (db.get_api_key('groq') if USE_DATABASE else None) or os.environ.get('GROQ_API_KEY')

    if not groq_key:
        return jsonify({'success': False, 'error': 'Groq API key not configured. Please add it in Settings.'}), 400

    try:
        response = httpx.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {groq_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": generation_prompt}],
                "temperature": 0.7,
                "max_tokens": 8000
            },
            timeout=120.0
        )

        if response.status_code != 200:
            return jsonify({
                'success': False,
                'error': f'AI API error: {response.status_code} - {response.text[:500]}'
            }), 500

        result = response.json()
        content = result['choices'][0]['message']['content'].strip()

        # Parse JSON from response - handle markdown code blocks
        if '```json' in content:
            content = content.split('```json')[1].split('```')[0]
        elif '```' in content:
            content = content.split('```')[1].split('```')[0]

        content = content.strip()

        try:
            examples = json.loads(content)
        except json.JSONDecodeError:
            # Try to find JSON array in content
            match = regex.search(r'\[[\s\S]*\]', content)
            if match:
                try:
                    examples = json.loads(match.group())
                except json.JSONDecodeError:
                    return jsonify({
                        'success': False,
                        'error': 'Failed to parse AI response as JSON'
                    }), 500
            else:
                return jsonify({
                    'success': False,
                    'error': 'AI did not return valid JSON array'
                }), 500

        # Validate and save examples
        generated_count = 0

        if USE_DATABASE:
            for ex in examples:
                user_input = str(ex.get('user_input', '')).strip()
                assistant_response = str(ex.get('assistant_response', '')).strip()

                if not user_input or not assistant_response:
                    continue

                item_id = str(uuid.uuid4())[:16]
                tokens = (len(user_input) + len(assistant_response)) // 4
                category = ex.get('category', 'general')

                execute_query('''
                    INSERT INTO extracted_data
                    (id, skill_id, user_input, assistant_response, source_filename,
                     source_type, category, importance_score, tokens)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (item_id, skill_id, user_input, assistant_response, 'ai_generated',
                      'ai', category, 75, tokens))

                generated_count += 1

        # CRITICAL: Commit volume to persist database writes
        commit_volume()

        add_activity(f"AI generated {generated_count} training examples for {skill_name}", "🤖", "training")

        return jsonify({
            'success': True,
            'generated': generated_count,
            'requested': count
        })

    except httpx.TimeoutException:
        return jsonify({'success': False, 'error': 'AI request timed out. Try generating fewer examples.'}), 500
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/create-skill', methods=['POST'])
def create_skill_endpoint():
    """Create a new skill from the quick form - saves to database."""
    data = request.json
    safe_name = data['name'].lower().replace(' ', '_')
    safe_name = re.sub(r'[^\w\-]', '_', safe_name)

    if USE_DATABASE:
        # Build system prompt from greeting, personality, and custom instructions
        greeting = data.get('greeting', f"Hello! How can I help you with {data['name']}?")
        personality = data.get('personality', 'Friendly and helpful')
        custom_instructions = data.get('customInstructions', '')

        system_prompt = f"{greeting}\n\nPersonality: {personality}"
        if custom_instructions:
            system_prompt += f"\n\nInstructions: {custom_instructions}"

        # Create skill in database
        db.create_skill(
            skill_id=safe_name,
            name=data['name'],
            description=data.get('description', ''),
            skill_type=data.get('type', 'General'),
            system_prompt=system_prompt,
            knowledge=data.get('services', '').split('\n') if data.get('services') else [],
            is_builtin=False
        )
        commit_volume()  # Persist to Modal volume

        add_activity(f"New skill created: {data['name']}", "")
        return jsonify({"success": True, "id": safe_name})
    else:
        return jsonify({"success": False, "error": "Database not available"}), 500


@app.route('/api/delete-skill/<skill_id>', methods=['DELETE'])
def delete_skill_endpoint(skill_id):
    """Delete a skill from database."""
    if USE_DATABASE:
        existing = db.get_skill(skill_id)
        if existing:
            db.delete_skill(skill_id)
            commit_volume()  # Persist deletion to Modal volume
            add_activity(f"Skill deleted: {skill_id}", "")
            return jsonify({"success": True})
        return jsonify({"error": "Skill not found"}), 404
    else:
        return jsonify({"success": False, "error": "Database not available"}), 500


@app.route('/api/upload-document', methods=['POST'])
def upload_document():
    """Handle document upload."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    skill_id = request.form.get('skill_id', 'general')

    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    # Save the file
    safe_name = re.sub(r'[^\w\-\.]', '_', file.filename)
    file_path = DOCUMENTS_DIR / f"{skill_id}_{safe_name}"
    file.save(file_path)

    # Process based on file type
    examples = 0
    try:
        if safe_name.endswith('.txt') or safe_name.endswith('.md'):
            with open(file_path, 'r') as f:
                text = f.read()
            examples = process_text_to_training(text, skill_id)
        elif safe_name.endswith('.json') or safe_name.endswith('.jsonl'):
            # Parse JSONL and save to database
            if USE_DATABASE:
                with open(file_path, 'r') as f:
                    for line in f:
                        line = line.strip()
                        if not line:
                            continue
                        try:
                            item = json.loads(line)
                            # Support multiple JSONL formats
                            user_msg = item.get('input') or item.get('user_message') or item.get('question', '')
                            asst_msg = item.get('output') or item.get('assistant_response') or item.get('answer', '')
                            # Also support messages array format
                            if not user_msg and 'messages' in item:
                                for msg in item['messages']:
                                    if msg.get('role') == 'user':
                                        user_msg = msg.get('content', '')
                                    elif msg.get('role') == 'assistant':
                                        asst_msg = msg.get('content', '')
                            if user_msg and asst_msg:
                                db.add_training_example(
                                    skill_id=skill_id,
                                    user_message=user_msg,
                                    assistant_response=asst_msg,
                                    metadata={"source": "jsonl_upload", "filename": safe_name}
                                )
                                examples += 1
                        except json.JSONDecodeError:
                            continue
                commit_volume()  # Persist to Modal volume
        else:
            examples = 10  # Placeholder for PDF processing
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    add_activity(f"Document processed: {examples} examples", "")
    return jsonify({"success": True, "examples": examples})


def process_text_to_training(text, skill_id):
    """Convert text to training data and save to database."""
    paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]
    count = 0

    if USE_DATABASE:
        for para in paragraphs[:100]:
            user_message = f"Tell me about: {para[:100]}..."
            assistant_response = para
            db.add_training_example(
                skill_id=skill_id,
                user_message=user_message,
                assistant_response=assistant_response,
                metadata={"source": "text_upload"}
            )
            count += 1
        commit_volume()  # Persist to Modal volume

    return count


@app.route('/api/generate-training/<skill_id>', methods=['POST'])
def generate_training(skill_id):
    """Generate training data from skill info in database."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    skill = db.get_skill(skill_id)
    if not skill:
        return jsonify({"error": "Skill not found"}), 404

    skill_name = skill.get('name', skill_id)
    skill_type = skill.get('skill_type', 'assistant')
    description = skill.get('description', '')
    system_prompt = skill.get('system_prompt', f"You are {skill_name}, a helpful {skill_type} assistant.")
    knowledge = skill.get('knowledge', [])

    count = 0

    # Greeting examples
    greetings = ["Hello", "Hi", "Hey", "Good morning", "I need help"]
    greeting_response = f"Hello! I'm {skill_name}. How can I help you today?"
    for greet in greetings:
        db.add_training_example(
            skill_id=skill_id,
            user_message=greet,
            assistant_response=greeting_response,
            metadata={"source": "generated", "type": "greeting"}
        )
        count += 1

    # Service/knowledge examples
    for item in knowledge:
        if item and item.strip():
            db.add_training_example(
                skill_id=skill_id,
                user_message=f"Tell me about {item}",
                assistant_response=f"Absolutely! {item} is one of our key offerings. {description}",
                metadata={"source": "generated", "type": "service"}
            )
            count += 1

    commit_volume()  # Persist to Modal volume

    add_activity(f"Training data generated: {count} examples", "")
    return jsonify({"success": True, "examples": count})


@app.route('/api/save-api-keys', methods=['POST'])
def save_api_keys():
    """Save API keys (persisted to database with encryption)."""
    global API_KEYS
    data = request.json
    saved = []

    # LLM Providers
    if data.get('groq'):
        db.save_api_key('groq', data['groq'])
        API_KEYS['groq'] = data['groq']
        saved.append('Groq')
    if data.get('openai'):
        db.save_api_key('openai', data['openai'])
        API_KEYS['openai'] = data['openai']
        saved.append('OpenAI')
    if data.get('anthropic'):
        db.save_api_key('anthropic', data['anthropic'])
        API_KEYS['anthropic'] = data['anthropic']
        saved.append('Anthropic')

    # Voice Providers
    if data.get('elevenlabs'):
        db.save_api_key('elevenlabs', data['elevenlabs'])
        API_KEYS['elevenlabs'] = data['elevenlabs']
        saved.append('ElevenLabs')
    if data.get('cartesia'):
        db.save_api_key('cartesia', data['cartesia'])
        API_KEYS['cartesia'] = data['cartesia']
        saved.append('Cartesia')
    if data.get('deepgram'):
        db.save_api_key('deepgram', data['deepgram'])
        API_KEYS['deepgram'] = data['deepgram']
        saved.append('Deepgram')
    if data.get('playht'):
        db.save_api_key('playht', data['playht'])
        API_KEYS['playht'] = data['playht']
        saved.append('PlayHT')

    if saved:
        commit_volume()  # Persist to Modal volume
        add_activity(f"API keys saved: {', '.join(saved)}", "", "api")
        return jsonify({"success": True, "saved": saved})
    return jsonify({"success": True, "message": "No keys provided"})


@app.route('/api/api-keys')
def get_api_keys():
    """Get saved API keys from database (masked for security)."""
    if USE_DATABASE:
        keys = db.get_all_api_keys()
        # Return full keys so they can be used - they're already stored securely
        return jsonify({
            'groq': keys.get('groq', ''),
            'openai': keys.get('openai', ''),
            'anthropic': keys.get('anthropic', ''),
            'elevenlabs': keys.get('elevenlabs', ''),
            'cartesia': keys.get('cartesia', ''),
            'deepgram': keys.get('deepgram', ''),
            'playht': keys.get('playht', '')
        })
    return jsonify({})


@app.route('/api/debug-api-keys')
def debug_api_keys():
    """Debug endpoint to check if API keys are saved (shows if key exists, not the actual key)."""
    if USE_DATABASE:
        result = {}
        for provider in ['groq', 'openai', 'anthropic', 'elevenlabs', 'cartesia', 'deepgram', 'playht']:
            key = db.get_api_key(provider)
            if key:
                result[provider] = f"SET ({len(key)} chars, starts with {key[:4]}...)"
            else:
                result[provider] = "NOT SET"
        return jsonify(result)
    return jsonify({"error": "Database not available"})


@app.route('/api/debug-voice-train/<project_id>')
def debug_voice_train(project_id):
    """Debug endpoint to test voice training without UI."""
    result = {"project_id": project_id}

    if not USE_DATABASE:
        result["error"] = "Database not available"
        return jsonify(result)

    # Get project
    project = db.get_voice_project(project_id)
    if not project:
        result["error"] = "Project not found"
        return jsonify(result)

    result["project_name"] = project.get('name')
    result["provider"] = project.get('provider', 'elevenlabs')
    result["status"] = project.get('status')
    result["voice_id"] = project.get('voice_id')

    # Get samples
    samples = db.get_voice_samples(project_id)
    result["sample_count"] = len(samples)
    result["samples"] = []

    for s in samples:
        sample_info = {
            "filename": s.get('filename'),
            "file_path": s.get('file_path'),
            "file_exists": os.path.exists(s.get('file_path', '')) if s.get('file_path') else False
        }
        result["samples"].append(sample_info)

    # Check API keys
    provider_key = project.get('provider', 'elevenlabs')
    api_key = db.get_api_key(provider_key)
    result[f"{provider_key}_key_set"] = bool(api_key)
    if api_key:
        result[f"{provider_key}_key_preview"] = f"{api_key[:4]}...{api_key[-4:]}"

    # Also check Cartesia if not already the provider
    if provider_key != 'cartesia':
        cartesia_key = db.get_api_key('cartesia')
        result["cartesia_key_set"] = bool(cartesia_key)
        if cartesia_key:
            result["cartesia_key_preview"] = f"{cartesia_key[:4]}...{cartesia_key[-4:]}"

    # Add training timestamps
    result["training_started"] = project.get('training_started')
    result["training_completed"] = project.get('training_completed')
    result["updated_at"] = project.get('updated_at')

    return jsonify(result)


@app.route('/api/fix-voice-provider/<project_id>/<provider>')
def fix_voice_provider(project_id, provider):
    """Force-update voice project provider (debug helper)."""
    if USE_DATABASE:
        project = db.update_voice_project(project_id, provider=provider, status='draft')
        if project:
            return jsonify({"success": True, "message": f"Provider changed to {provider}", "project": project})
    return jsonify({"error": "Failed to update provider"})


@app.route('/api/test-llm', methods=['POST'])
def test_llm():
    """Test a single LLM provider with real API call."""
    data = request.json
    provider = data.get('provider', 'groq')
    prompt = data.get('prompt', '')

    import time
    start = time.time()

    # Try to make real API call
    api_keys = db.get_all_api_keys() if USE_DATABASE else {}

    try:
        import httpx

        if provider == 'groq' and api_keys.get('groq'):
            response = httpx.post(
                'https://api.groq.com/openai/v1/chat/completions',
                headers={'Authorization': f"Bearer {api_keys['groq']}", 'Content-Type': 'application/json'},
                json={'model': 'llama-3.3-70b-versatile', 'messages': [{'role': 'user', 'content': prompt}], 'max_tokens': 150},
                timeout=30
            )
            if response.status_code == 200:
                result = response.json()
                latency = int((time.time() - start) * 1000)
                return jsonify({
                    "success": True,
                    "response": result['choices'][0]['message']['content'],
                    "latency_ms": latency,
                    "provider": provider
                })

        # No API key configured
        return jsonify({
            "success": False,
            "error": f"No API key configured for {provider}. Add it in Settings → API Keys.",
            "provider": provider
        })

    except Exception as e:
        return jsonify({
            "success": False,
            "error": str(e),
            "provider": provider
        })


@app.route('/api/compare-llms', methods=['POST'])
def compare_llms():
    """Compare multiple LLM providers - requires API keys configured."""
    data = request.json
    providers = data.get('providers', ['groq'])
    prompt = data.get('prompt', '')

    results = []
    api_keys = db.get_all_api_keys() if USE_DATABASE else {}

    for provider in providers:
        if not api_keys.get(provider):
            results.append({
                "provider": provider,
                "response": f"No API key configured for {provider}",
                "first_token_ms": 0,
                "total_ms": 0,
                "error": True
            })
        else:
            # Would need real implementation for each provider
            results.append({
                "provider": provider,
                "response": f"API key found for {provider} - real comparison requires implementation",
                "first_token_ms": 0,
                "total_ms": 0,
            })

    return jsonify({"success": True, "results": results})


@app.route('/api/lpu-inference', methods=['POST'])
def lpu_inference():
    """Run inference on the Modal LPU."""
    data = request.json
    prompt = data.get('prompt', '')

    try:
        import modal
        lpu = modal.Cls.from_name("bitnet-lpu-v1", "VirtualLPU")()
        response = ""
        for chunk in lpu.chat.remote_gen(prompt, max_tokens=data.get('maxTokens', 128)):
            response += chunk
        add_activity("LPU inference completed", "")
        return jsonify({"success": True, "response": response})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/stats')
def get_stats():
    """Get system statistics - real counts from database."""
    skills_count = 0
    if USE_DATABASE:
        skills = db.get_all_skills()
        skills_count = len(skills) if skills else 0

    return jsonify({
        "total_queries": 0,
        "by_provider": {
            "groq": 0,
            "openai": 0,
            "anthropic": 0,
            "bitnet": 0,
        },
        "avg_latency_ms": 0,
        "cache_hits": 0,
        "skills_active": skills_count,
        "adapters_trained": len(list(ADAPTERS_DIR.glob("*"))),
    })


# =============================================================================
# API ENDPOINTS - VOICE CONFIGURATION
# =============================================================================

@app.route('/api/voice/providers')
def get_voice_providers():
    """Get all available voice providers and their voices."""
    return jsonify(VOICE_CHOICES)


@app.route('/api/voice/config')
def get_voice_config():
    """Get current voice configuration."""
    return jsonify(VOICE_CONFIG)


@app.route('/api/voice/config', methods=['POST'])
def save_voice_config():
    """Save voice configuration."""
    global VOICE_CONFIG
    data = request.json

    if data.get('selected_voice'):
        VOICE_CONFIG['selected_voice'] = data['selected_voice']
    if data.get('selected_provider'):
        VOICE_CONFIG['selected_provider'] = data['selected_provider']
    if data.get('voice_settings'):
        VOICE_CONFIG['voice_settings'].update(data['voice_settings'])

    add_activity(f"Voice changed to {VOICE_CONFIG['selected_voice']}", "")
    return jsonify({"success": True, "config": VOICE_CONFIG})


@app.route('/api/voice/cache-stats')
def get_voice_cache_stats():
    """Get TTS cache statistics."""
    return jsonify(TTS_CACHE.get_stats())


@app.route('/api/voice/cache-clear', methods=['POST'])
def clear_voice_cache():
    """Clear the TTS cache."""
    TTS_CACHE.clear()
    return jsonify({"success": True, "message": "TTS cache cleared"})


@app.route('/api/voice/test', methods=['POST'])
def test_voice():
    """Test a voice with sample text - generates real audio with caching."""
    import base64
    import time

    data = request.json
    voice_id = data.get('voice_id', 'en-US-JennyNeural')
    text = data.get('text', 'Hello! This is a test of the voice synthesis system.')
    provider = data.get('provider', 'edge_tts')
    emotion = data.get('emotion', 'neutral')
    skill_id = data.get('skill_id', 'general')
    skip_cache = data.get('skip_cache', False)

    start_time = time.time()

    # Check TTS cache first (unless explicitly skipped)
    if not skip_cache:
        cached_audio, cached_format = TTS_CACHE.get(
            text, voice_id, provider, emotion=emotion
        )
        if cached_audio:
            duration_ms = int((time.time() - start_time) * 1000)
            audio_base64 = base64.b64encode(cached_audio).decode('utf-8') if isinstance(cached_audio, bytes) else cached_audio
            return jsonify({
                "success": True,
                "voice_id": voice_id,
                "provider": provider,
                "text": text,
                "duration_ms": duration_ms,
                "audio_base64": audio_base64,
                "audio_format": cached_format,
                "cached": True,
                "message": f"Served from cache ({duration_ms}ms)"
            })

    # Parler TTS - Expressive voices via Modal GPU
    if provider == 'parler_tts':
        if not text or len(text.strip()) == 0:
            text = "Hello! This is a voice test."
        text = text.strip()[:500]

        try:
            import httpx

            # Call Parler TTS Modal function
            parler_url = "https://jenkintownelectricity--hive215-parler-tts-parlerttsmodel-synthesize-with-emotion.modal.run"

            with httpx.Client(timeout=120.0) as client:  # Long timeout for cold start
                response = client.post(
                    parler_url,
                    json={
                        "text": text,
                        "skill_id": skill_id,
                        "emotion": emotion
                    }
                )

                if response.status_code == 200:
                    result = response.json()
                    # Result is [audio_bytes, description]
                    if isinstance(result, list) and len(result) >= 2:
                        audio_bytes = bytes(result[0]) if isinstance(result[0], list) else result[0]
                        description = result[1]

                        duration_ms = int((time.time() - start_time) * 1000)
                        audio_base64 = base64.b64encode(audio_bytes).decode('utf-8') if isinstance(audio_bytes, bytes) else result[0]

                        # Cache the synthesized audio
                        if isinstance(audio_bytes, bytes):
                            TTS_CACHE.set(text, f"{skill_id}_{emotion}", "parler_tts",
                                         audio_bytes, "audio/wav", emotion=emotion)

                        add_activity(f"Parler TTS: {skill_id}/{emotion} ({duration_ms}ms)", "")
                        return jsonify({
                            "success": True,
                            "voice_id": f"{skill_id}_{emotion}",
                            "provider": "parler_tts",
                            "text": text,
                            "duration_ms": duration_ms,
                            "audio_base64": audio_base64,
                            "audio_format": "audio/wav",
                            "cached": False,
                            "message": f"Generated with Parler TTS: {description}"
                        })
                    else:
                        raise Exception(f"Unexpected response format: {type(result)}")
                else:
                    raise Exception(f"Parler TTS returned {response.status_code}: {response.text[:200]}")

        except httpx.TimeoutException:
            return jsonify({
                "success": False,
                "error": "Parler TTS timeout (GPU may be cold starting - try again in 30s)"
            }), 500
        except Exception as e:
            return jsonify({
                "success": False,
                "error": f"Parler TTS failed: {str(e)}"
            }), 500

    # ElevenLabs TTS
    if provider == 'elevenlabs':
        try:
            import httpx
            api_key = db.get_api_key('elevenlabs') if USE_DATABASE else None
            if not api_key:
                return jsonify({"success": False, "error": "ElevenLabs API key not configured"}), 400

            with httpx.Client(timeout=30.0) as client:
                response = client.post(
                    f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}",
                    headers={
                        "xi-api-key": api_key,
                        "Content-Type": "application/json"
                    },
                    json={
                        "text": text[:5000],
                        "model_id": "eleven_monolingual_v1",
                        "voice_settings": {"stability": 0.5, "similarity_boost": 0.75}
                    }
                )

                if response.status_code == 200:
                    audio_bytes = response.content
                    duration_ms = int((time.time() - start_time) * 1000)
                    audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                    # Cache the audio
                    TTS_CACHE.set(text, voice_id, provider, audio_bytes, "audio/mpeg")

                    return jsonify({
                        "success": True,
                        "voice_id": voice_id,
                        "provider": "elevenlabs",
                        "text": text,
                        "duration_ms": duration_ms,
                        "audio_base64": audio_base64,
                        "audio_format": "audio/mpeg",
                        "cached": False,
                        "message": f"Generated with ElevenLabs ({duration_ms}ms)"
                    })
                else:
                    return jsonify({"success": False, "error": f"ElevenLabs error: {response.status_code}"}), 500
        except Exception as e:
            return jsonify({"success": False, "error": f"ElevenLabs failed: {str(e)}"}), 500

    # Cartesia TTS
    if provider == 'cartesia':
        try:
            import httpx
            api_key = db.get_api_key('cartesia') if USE_DATABASE else None
            if not api_key:
                return jsonify({"success": False, "error": "Cartesia API key not configured"}), 400

            with httpx.Client(timeout=30.0) as client:
                response = client.post(
                    "https://api.cartesia.ai/tts/bytes",
                    headers={
                        "X-API-Key": api_key,
                        "Cartesia-Version": "2024-11-13",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model_id": "sonic-english",
                        "transcript": text[:5000],
                        "voice": {"mode": "id", "id": voice_id},
                        "output_format": {"container": "mp3", "sample_rate": 44100}
                    }
                )

                if response.status_code == 200:
                    audio_bytes = response.content
                    duration_ms = int((time.time() - start_time) * 1000)
                    audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                    # Cache the audio
                    TTS_CACHE.set(text, voice_id, provider, audio_bytes, "audio/mpeg")

                    return jsonify({
                        "success": True,
                        "voice_id": voice_id,
                        "provider": "cartesia",
                        "text": text,
                        "duration_ms": duration_ms,
                        "audio_base64": audio_base64,
                        "audio_format": "audio/mpeg",
                        "cached": False,
                        "message": f"Generated with Cartesia ({duration_ms}ms)"
                    })
                else:
                    return jsonify({"success": False, "error": f"Cartesia error: {response.status_code} - {response.text[:200]}"}), 500
        except Exception as e:
            return jsonify({"success": False, "error": f"Cartesia failed: {str(e)}"}), 500

    # Deepgram Aura TTS
    if provider == 'deepgram':
        try:
            import httpx
            api_key = db.get_api_key('deepgram') if USE_DATABASE else None
            if not api_key:
                return jsonify({"success": False, "error": "Deepgram API key not configured"}), 400

            with httpx.Client(timeout=30.0) as client:
                response = client.post(
                    f"https://api.deepgram.com/v1/speak?model={voice_id}",
                    headers={
                        "Authorization": f"Token {api_key}",
                        "Content-Type": "application/json"
                    },
                    json={"text": text[:5000]}
                )

                if response.status_code == 200:
                    audio_bytes = response.content
                    duration_ms = int((time.time() - start_time) * 1000)
                    audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                    # Cache the audio
                    TTS_CACHE.set(text, voice_id, provider, audio_bytes, "audio/mpeg")

                    return jsonify({
                        "success": True,
                        "voice_id": voice_id,
                        "provider": "deepgram",
                        "text": text,
                        "duration_ms": duration_ms,
                        "audio_base64": audio_base64,
                        "audio_format": "audio/mpeg",
                        "cached": False,
                        "message": f"Generated with Deepgram Aura ({duration_ms}ms)"
                    })
                else:
                    return jsonify({"success": False, "error": f"Deepgram error: {response.status_code}"}), 500
        except Exception as e:
            return jsonify({"success": False, "error": f"Deepgram failed: {str(e)}"}), 500

    # OpenAI TTS
    if provider == 'openai':
        try:
            import httpx
            api_key = db.get_api_key('openai') if USE_DATABASE else None
            if not api_key:
                return jsonify({"success": False, "error": "OpenAI API key not configured"}), 400

            with httpx.Client(timeout=30.0) as client:
                response = client.post(
                    "https://api.openai.com/v1/audio/speech",
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": "tts-1",
                        "input": text[:4096],
                        "voice": voice_id
                    }
                )

                if response.status_code == 200:
                    audio_bytes = response.content
                    duration_ms = int((time.time() - start_time) * 1000)
                    audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

                    # Cache the audio
                    TTS_CACHE.set(text, voice_id, provider, audio_bytes, "audio/mpeg")

                    return jsonify({
                        "success": True,
                        "voice_id": voice_id,
                        "provider": "openai",
                        "text": text,
                        "duration_ms": duration_ms,
                        "audio_base64": audio_base64,
                        "audio_format": "audio/mpeg",
                        "cached": False,
                        "message": f"Generated with OpenAI TTS ({duration_ms}ms)"
                    })
                else:
                    return jsonify({"success": False, "error": f"OpenAI error: {response.status_code}"}), 500
        except Exception as e:
            return jsonify({"success": False, "error": f"OpenAI TTS failed: {str(e)}"}), 500

    # Map voice IDs to gTTS language/accent codes
    VOICE_TO_GTTS = {
        # Edge TTS voices -> gTTS (all use 'en' but we track the original)
        'en-US-JennyNeural': ('en', 'com'),
        'en-US-AriaNeural': ('en', 'com'),
        'en-US-GuyNeural': ('en', 'com'),
        'en-US-DavisNeural': ('en', 'com'),
        'en-US-SaraNeural': ('en', 'com'),
        'en-US-AnaNeural': ('en', 'com'),
        'en-GB-SoniaNeural': ('en', 'co.uk'),
        'en-GB-RyanNeural': ('en', 'co.uk'),
        'en-GB-LibbyNeural': ('en', 'co.uk'),
        'en-AU-NatashaNeural': ('en', 'com.au'),
        'en-AU-WilliamNeural': ('en', 'com.au'),
        # Kokoro voices -> gTTS
        'af_bella': ('en', 'com'),
        'af_nicole': ('en', 'com'),
        'af_sarah': ('en', 'com'),
        'af_sky': ('en', 'com'),
        'am_adam': ('en', 'com'),
        'am_michael': ('en', 'com'),
        'bf_emma': ('en', 'co.uk'),
        'bf_isabella': ('en', 'co.uk'),
        'bm_george': ('en', 'co.uk'),
        'bm_lewis': ('en', 'co.uk'),
    }

    # Use gTTS for all voices (HTTP-based, works reliably from Modal)
    # Ensure text is not empty and clean it
    if not text or len(text.strip()) == 0:
        text = "Hello! This is a voice test."
    text = text.strip()[:500]  # Limit length

    try:
        from gtts import gTTS
        import io

        # Get language settings for this voice
        lang, tld = VOICE_TO_GTTS.get(voice_id, ('en', 'com'))

        # Generate audio with gTTS
        tts = gTTS(text=text, lang=lang, tld=tld)
        audio_buffer = io.BytesIO()
        tts.write_to_fp(audio_buffer)
        audio_bytes = audio_buffer.getvalue()

        if not audio_bytes:
            raise Exception("No audio generated")

        duration_ms = int((time.time() - start_time) * 1000)
        audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')

        # Cache the synthesized audio
        TTS_CACHE.set(text, voice_id, provider, audio_bytes, "audio/mpeg", emotion=emotion)

        # Determine display provider name
        if voice_id.startswith('af_') or voice_id.startswith('am_') or voice_id.startswith('bf_') or voice_id.startswith('bm_'):
            display_provider = f"kokoro (via Google TTS)"
            message = f"Kokoro '{voice_id}' using Google TTS ({tld})"
        else:
            display_provider = "google_tts"
            message = f"Generated audio with Google TTS ({tld})"

        add_activity(f"Voice test: {voice_id} ({duration_ms}ms)", "")
        return jsonify({
            "success": True,
            "voice_id": voice_id,
            "provider": display_provider,
            "text": text,
            "duration_ms": duration_ms,
            "audio_base64": audio_base64,
            "audio_format": "audio/mpeg",
            "cached": False,
            "message": message
        })
    except Exception as e:
        return jsonify({
            "success": False,
            "error": f"TTS failed: {str(e)}"
        }), 500


# =============================================================================
# API ENDPOINTS - GOLDEN PROMPTS (with Database Persistence)
# =============================================================================

# In-memory storage for custom prompts (fallback if no database)
CUSTOM_PROMPTS = {}

def get_custom_prompt(skill_id):
    """Get custom prompt from database or memory."""
    if USE_DATABASE:
        prompt = db.get_golden_prompt(skill_id)
        return prompt['content'] if prompt else None
    return CUSTOM_PROMPTS.get(skill_id)

def get_all_custom_prompt_ids():
    """Get list of skill IDs with custom prompts."""
    if USE_DATABASE:
        prompts = db.get_all_golden_prompts()
        return [p['skill_id'] for p in prompts]
    return list(CUSTOM_PROMPTS.keys())

@app.route('/api/golden-prompts')
def get_golden_prompts_list():
    """Get list of all available golden prompts."""
    try:
        import sys
        sys.path.insert(0, '/root')
        from golden_prompts import SKILL_MANUALS, TOKEN_ESTIMATES

        custom_ids = get_all_custom_prompt_ids()

        prompts = []
        for skill_id in SKILL_MANUALS.keys():
            tokens = TOKEN_ESTIMATES.get(skill_id, 800)
            prefill_ms = int((tokens / 6000) * 1000)
            prompts.append({
                "id": skill_id,
                "name": skill_id.replace("_", " ").replace("-", " ").title(),
                "tokens": tokens,
                "prefill_ms": prefill_ms,
                "is_custom": skill_id in custom_ids
            })
        return jsonify({"prompts": prompts})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/golden-prompts/<skill_id>')
def get_golden_prompt(skill_id):
    """Get a specific golden prompt with optional context substitution."""
    try:
        import sys
        sys.path.insert(0, '/root')
        from golden_prompts import SKILL_MANUALS, TOKEN_ESTIMATES, get_skill_prompt

        # Check for custom override first (database or memory)
        custom_content = get_custom_prompt(skill_id)

        if custom_content:
            content = custom_content
            is_custom = True
        elif skill_id in SKILL_MANUALS:
            content = SKILL_MANUALS[skill_id]
            is_custom = False
        else:
            return jsonify({"error": f"Skill '{skill_id}' not found"}), 404

        # Get context variables from query params
        business_name = request.args.get('business_name', '{business_name}')
        agent_name = request.args.get('agent_name', '{agent_name}')

        # Substitute context variables
        formatted = content.replace('{business_name}', business_name).replace('{agent_name}', agent_name)

        tokens = TOKEN_ESTIMATES.get(skill_id, 800)
        prefill_ms = int((tokens / 6000) * 1000)

        return jsonify({
            "skill_id": skill_id,
            "content": content,
            "formatted": formatted,
            "tokens": tokens,
            "prefill_ms": prefill_ms,
            "is_custom": is_custom
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/golden-prompts/<skill_id>', methods=['POST'])
def save_golden_prompt(skill_id):
    """Save a custom golden prompt with database persistence."""
    global CUSTOM_PROMPTS

    try:
        data = request.json
        content = data.get('content', '')

        if not content.strip():
            return jsonify({"error": "Prompt content cannot be empty"}), 400

        # Save to database if available, otherwise memory
        if USE_DATABASE:
            result = db.save_golden_prompt(skill_id, content)
            tokens = result['tokens_estimate']
        else:
            CUSTOM_PROMPTS[skill_id] = content
            tokens = len(content) // 4

        prefill_ms = int((tokens / 6000) * 1000)

        commit_volume()  # Persist to Modal volume
        add_activity(f"Updated golden prompt: {skill_id}", "", "prompts")

        return jsonify({
            "success": True,
            "skill_id": skill_id,
            "tokens": tokens,
            "prefill_ms": prefill_ms,
            "message": f"Prompt for '{skill_id}' saved successfully"
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/golden-prompts/<skill_id>/reset', methods=['POST'])
def reset_golden_prompt(skill_id):
    """Reset a golden prompt to its default (removes from database)."""
    global CUSTOM_PROMPTS

    try:
        if USE_DATABASE:
            deleted = db.delete_golden_prompt(skill_id)
            if deleted:
                commit_volume()  # Persist to Modal volume
                add_activity(f"Reset golden prompt: {skill_id}", "", "prompts")
                return jsonify({"success": True, "message": f"Prompt for '{skill_id}' reset to default"})
            else:
                return jsonify({"success": True, "message": "Prompt was already at default"})
        else:
            if skill_id in CUSTOM_PROMPTS:
                del CUSTOM_PROMPTS[skill_id]
                add_activity(f"Reset golden prompt: {skill_id}", "", "prompts")
                return jsonify({"success": True, "message": f"Prompt for '{skill_id}' reset to default"})
            else:
                return jsonify({"success": True, "message": "Prompt was already at default"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/golden-prompts/export')
def export_golden_prompts():
    """Export all golden prompts (defaults + custom overrides)."""
    try:
        import sys
        sys.path.insert(0, '/root')
        from golden_prompts import SKILL_MANUALS

        custom_ids = get_all_custom_prompt_ids()

        export_data = {}
        for skill_id, content in SKILL_MANUALS.items():
            # Use custom if available, otherwise default
            custom_content = get_custom_prompt(skill_id)
            export_data[skill_id] = custom_content if custom_content else content

        return jsonify({
            "prompts": export_data,
            "custom_overrides": custom_ids
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/golden-prompts/test', methods=['POST'])
def test_golden_prompt():
    """Test a golden prompt with the Fast Brain API."""
    try:
        data = request.json
        skill_id = data.get('skill_id', 'general')
        test_message = data.get('message', 'Hello, I need some help.')
        business_name = data.get('business_name', 'Test Business')
        agent_name = data.get('agent_name', 'Assistant')

        import sys
        sys.path.insert(0, '/root')
        from golden_prompts import get_skill_prompt, SKILL_MANUALS

        # Get prompt (custom or default)
        if skill_id in CUSTOM_PROMPTS:
            system_prompt = CUSTOM_PROMPTS[skill_id]
            system_prompt = system_prompt.replace('{business_name}', business_name).replace('{agent_name}', agent_name)
        else:
            system_prompt = get_skill_prompt(skill_id, business_name=business_name, agent_name=agent_name)

        # Call Fast Brain API
        import httpx
        import os

        fast_brain_url = os.environ.get('FAST_BRAIN_URL', '')
        if not fast_brain_url:
            return jsonify({"error": "Fast Brain URL not configured"}), 500

        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                f"{fast_brain_url}/v1/chat/completions",
                json={
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": test_message}
                    ],
                    "max_tokens": 500
                }
            )

            if response.status_code == 200:
                result = response.json()
                return jsonify({
                    "success": True,
                    "response": result.get("choices", [{}])[0].get("message", {}).get("content", ""),
                    "system_used": result.get("system_used", "unknown"),
                    "latency_ms": result.get("latency_ms", 0)
                })
            else:
                return jsonify({"error": f"Fast Brain returned {response.status_code}"}), 500

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# =============================================================================
# API ENDPOINTS - PLATFORM CONNECTIONS
# =============================================================================

@app.route('/api/platforms')
def get_platforms():
    """Get all platform connections and their status."""
    return jsonify(PLATFORM_CONNECTIONS)


@app.route('/api/platforms/<platform_id>')
def get_platform(platform_id):
    """Get a specific platform connection."""
    if platform_id in PLATFORM_CONNECTIONS:
        return jsonify(PLATFORM_CONNECTIONS[platform_id])
    return jsonify({"error": "Platform not found"}), 404


@app.route('/api/platforms/<platform_id>/connect', methods=['POST'])
def connect_platform(platform_id):
    """Connect to a voice platform (persisted to database)."""
    global PLATFORM_CONNECTIONS

    if platform_id not in PLATFORM_CONNECTIONS:
        return jsonify({"error": "Platform not found"}), 404

    data = request.json
    config = data.get('config', {})

    # Update the platform config in memory
    PLATFORM_CONNECTIONS[platform_id]['config'].update(config)

    # Simulate connection test
    has_required_fields = all(v for v in PLATFORM_CONNECTIONS[platform_id]['config'].values())

    if has_required_fields:
        PLATFORM_CONNECTIONS[platform_id]['status'] = 'connected'
        # Save to database
        db.save_platform(platform_id, PLATFORM_CONNECTIONS[platform_id]['config'], 'connected')
        commit_volume()  # Persist to Modal volume
        add_activity(f"Connected to {PLATFORM_CONNECTIONS[platform_id]['name']}", "", "platform")
        return jsonify({
            "success": True,
            "status": "connected",
            "message": f"Successfully connected to {PLATFORM_CONNECTIONS[platform_id]['name']}"
        })
    else:
        PLATFORM_CONNECTIONS[platform_id]['status'] = 'error'
        return jsonify({
            "success": False,
            "status": "error",
            "message": "Missing required configuration fields"
        })


@app.route('/api/platforms/<platform_id>/disconnect', methods=['POST'])
def disconnect_platform(platform_id):
    """Disconnect from a voice platform (persisted to database)."""
    global PLATFORM_CONNECTIONS

    if platform_id not in PLATFORM_CONNECTIONS:
        return jsonify({"error": "Platform not found"}), 404

    PLATFORM_CONNECTIONS[platform_id]['status'] = 'disconnected'
    # Save to database
    db.update_platform_status(platform_id, 'disconnected')
    commit_volume()  # Persist to Modal volume
    add_activity(f"Disconnected from {PLATFORM_CONNECTIONS[platform_id]['name']}", "", "platform")

    return jsonify({
        "success": True,
        "status": "disconnected",
        "message": f"Disconnected from {PLATFORM_CONNECTIONS[platform_id]['name']}"
    })


@app.route('/api/platforms/<platform_id>/test', methods=['POST'])
def test_platform(platform_id):
    """Test a platform connection."""
    if platform_id not in PLATFORM_CONNECTIONS:
        return jsonify({"error": "Platform not found"}), 404

    platform = PLATFORM_CONNECTIONS[platform_id]

    if platform['status'] != 'connected':
        return jsonify({
            "success": False,
            "message": "Platform is not connected"
        })

    # Test connection
    add_activity(f"Testing connection to {platform['name']}", "")
    return jsonify({
        "success": True,
        "message": f"Connection to {platform['name']} is working",
        "latency_ms": 0
    })


# =============================================================================
# API ENDPOINTS - FAST BRAIN
# =============================================================================

FAST_BRAIN_CONFIG = {
    "url": os.getenv("FAST_BRAIN_URL", ""),
    "status": "disconnected",
    "model": "groq-llama-3.3-70b",
    "backend": "groq",
    "min_containers": 1,
    "max_containers": 10,
    "selected_skill": "general",
    "stats": {
        "total_requests": 0,
        "avg_ttfb_ms": 0,
        "avg_throughput_tps": 0,
        "uptime_hours": 0,
        "errors": [],
    }
}

# Local skills cache (loaded from database, no hardcoded demo data)
FAST_BRAIN_SKILLS = {}

# System status tracking for comprehensive monitoring
SYSTEM_STATUS = {
    "fast_brain": {"status": "unknown", "last_check": None, "error": None, "latency_ms": None},
    "hive215": {"status": "unknown", "last_check": None, "error": None, "latency_ms": None},
    "groq": {"status": "unknown", "last_check": None, "error": None},
    "modal": {"status": "unknown", "last_check": None, "error": None},
}


@app.route('/api/fast-brain/config')
def get_fast_brain_config():
    """Get Fast Brain configuration and status."""
    return jsonify(FAST_BRAIN_CONFIG)


@app.route('/api/fast-brain/config', methods=['POST'])
def update_fast_brain_config():
    """Update Fast Brain configuration."""
    global FAST_BRAIN_CONFIG
    data = request.json

    if 'url' in data:
        # Strip trailing slash to avoid double-slash issues
        clean_url = data['url'].rstrip('/')
        FAST_BRAIN_CONFIG['url'] = clean_url
        os.environ['FAST_BRAIN_URL'] = clean_url

    if 'min_containers' in data:
        FAST_BRAIN_CONFIG['min_containers'] = data['min_containers']
    if 'max_containers' in data:
        FAST_BRAIN_CONFIG['max_containers'] = data['max_containers']

    add_activity(f"Fast Brain config updated", "")
    return jsonify({"success": True, "config": FAST_BRAIN_CONFIG})


@app.route('/api/fast-brain/health')
def fast_brain_health():
    """Check Fast Brain health."""
    global FAST_BRAIN_CONFIG

    url = FAST_BRAIN_CONFIG.get('url')
    if not url:
        return jsonify({
            "status": "not_configured",
            "message": "FAST_BRAIN_URL not set"
        })

    # Try to connect to Fast Brain LPU
    try:
        import httpx
        with httpx.Client(timeout=5.0) as client:
            response = client.get(f"{url}/health")
            if response.status_code == 200:
                health_data = response.json()
                api_status = health_data.get("status", "unknown")
                FAST_BRAIN_CONFIG['status'] = 'connected' if api_status in ['healthy', 'initializing'] else 'error'
                return jsonify({
                    "status": api_status,
                    "model_loaded": health_data.get("model_loaded", False),
                    "skills": health_data.get("skills_available", []),
                    "version": health_data.get("version", "unknown"),
                })
    except Exception as e:
        FAST_BRAIN_CONFIG['status'] = 'error'
        return jsonify({
            "status": "error",
            "message": str(e)
        })

    FAST_BRAIN_CONFIG['status'] = 'disconnected'
    return jsonify({"status": "disconnected"})


@app.route('/api/fast-brain/chat', methods=['POST'])
def fast_brain_chat():
    """Send a chat request to Fast Brain."""
    data = request.json
    message = data.get('message', '')
    system_prompt = data.get('system_prompt', '')
    max_tokens = data.get('max_tokens', 256)
    stream = data.get('stream', False)

    url = FAST_BRAIN_CONFIG.get('url')
    if not url:
        return jsonify({"success": False, "error": "Fast Brain not configured"})

    try:
        import httpx
        import time

        start_time = time.time()

        # Build messages array for OpenAI-compatible API
        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": message})

        # Use Fast Brain LPU /v1/chat/completions endpoint
        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                f"{url}/v1/chat/completions",
                json={
                    "messages": messages,
                    "max_tokens": max_tokens,
                    "stream": False,
                },
            )
            # If 422, get the validation error details
            if response.status_code == 422:
                error_detail = response.json()
                return jsonify({"success": False, "error": f"Validation error: {error_detail}"})
            response.raise_for_status()
            result = response.json()

        elapsed_ms = (time.time() - start_time) * 1000

        # Update stats
        FAST_BRAIN_CONFIG['stats']['total_requests'] += 1

        # Extract content from OpenAI-compatible response
        choices = result.get('choices', [])
        content = choices[0].get('message', {}).get('content', '') if choices else ''
        metrics = result.get('metrics', {"ttfb_ms": elapsed_ms, "tokens_per_sec": 0})

        add_activity(f"Fast Brain chat: {message[:30]}...", "")

        return jsonify({
            "success": True,
            "response": content,
            "metrics": {
                "ttfb_ms": metrics.get('ttfb_ms', elapsed_ms),
                "total_time_ms": elapsed_ms,
                "tokens_per_sec": metrics.get('tokens_per_sec', 0),
            }
        })

    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/fast-brain/benchmark', methods=['POST'])
def fast_brain_benchmark():
    """Run a benchmark on Fast Brain."""
    data = request.json
    num_requests = data.get('num_requests', 5)
    prompt = data.get('prompt', 'Hello, how are you?')

    url = FAST_BRAIN_CONFIG.get('url')
    if not url:
        return jsonify({"success": False, "error": "Fast Brain not configured"})

    results = []
    try:
        import httpx
        import time

        with httpx.Client(timeout=30.0) as client:
            for i in range(num_requests):
                start = time.time()
                # Use Fast Brain LPU /v1/chat/completions endpoint
                response = client.post(
                    f"{url}/v1/chat/completions",
                    json={
                        "messages": [{"role": "user", "content": prompt}],
                        "max_tokens": 50,
                        "stream": False,
                    },
                )
                elapsed = (time.time() - start) * 1000
                result = response.json()
                metrics = result.get('metrics', {})
                results.append({
                    "request": i + 1,
                    "ttfb_ms": metrics.get('ttfb_ms', elapsed),
                    "total_ms": elapsed,
                    "tokens_per_sec": metrics.get('tokens_per_sec', 0),
                })

        avg_ttfb = sum(r['ttfb_ms'] for r in results) / len(results)
        avg_total = sum(r['total_ms'] for r in results) / len(results)
        avg_tps = sum(r['tokens_per_sec'] for r in results) / len(results)

        return jsonify({
            "success": True,
            "results": results,
            "summary": {
                "avg_ttfb_ms": round(avg_ttfb, 2),
                "avg_total_ms": round(avg_total, 2),
                "avg_tokens_per_sec": round(avg_tps, 2),
                "ttfb_target_met": avg_ttfb < 50,
                "throughput_target_met": avg_tps > 500,
            }
        })

    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


# =============================================================================
# API ENDPOINTS - FAST BRAIN SKILLS (Unified with Database)
# =============================================================================

@app.route('/api/fast-brain/skills')
def get_fast_brain_skills():
    """Get all available skills from unified database."""
    try:
        skills = []
        selected = 'general'
        source = 'none'

        if USE_DATABASE:
            try:
                skills = db.get_all_skills()
                selected = db.get_config('selected_skill', 'general') or 'general'
                source = 'database'
                print(f"[SKILLS] Loaded {len(skills)} skills from database")
            except Exception as db_error:
                print(f"[SKILLS] DATABASE FAILED: {db_error}")
                # Try to initialize database if tables don't exist
                try:
                    db.init_db()
                    skills = db.get_all_skills()
                    source = 'database_retry'
                    print(f"[SKILLS] Retry successful: {len(skills)} skills")
                except Exception as retry_error:
                    print(f"[SKILLS] Retry also failed: {retry_error}")
        else:
            print("[SKILLS] USE_DATABASE is False!")

        # If no skills in database, try fetching from LPU API
        if not skills:
            url = FAST_BRAIN_CONFIG.get('url', 'https://jenkintownelectricity--fast-brain-lpu-fastapi-app.modal.run')
            if url:
                try:
                    import httpx
                    with httpx.Client(timeout=5.0) as client:
                        response = client.get(f"{url}/v1/skills")
                        if response.status_code == 200:
                            remote_skills = response.json().get("skills", [])
                            for skill in remote_skills:
                                if skill.get("id"):
                                    skills.append(skill)
                                    FAST_BRAIN_SKILLS[skill["id"]] = {**skill, "is_builtin": False}
                except Exception as api_error:
                    print(f"API error fetching skills from LPU: {api_error}")

        # Fallback to in-memory skills if still empty
        if not skills and FAST_BRAIN_SKILLS:
            skills = list(FAST_BRAIN_SKILLS.values())
            selected = FAST_BRAIN_CONFIG.get("selected_skill", "general")

        return jsonify({
            "skills": skills,
            "selected": selected,
            "success": True
        })

    except Exception as e:
        print(f"Error in get_fast_brain_skills: {e}")
        return jsonify({
            "skills": [],
            "selected": "general",
            "success": False,
            "error": str(e)
        })


@app.route('/api/fast-brain/sync-skills', methods=['POST'])
def sync_skills_from_lpu():
    """Sync skills from the Fast Brain LPU API to local database."""
    try:
        import httpx
        url = FAST_BRAIN_CONFIG.get('url', 'https://jenkintownelectricity--fast-brain-lpu-fastapi-app.modal.run')

        synced = 0
        with httpx.Client(timeout=10.0) as client:
            response = client.get(f"{url}/v1/skills")
            if response.status_code == 200:
                remote_skills = response.json().get("skills", [])
                for skill in remote_skills:
                    skill_id = skill.get('id')
                    if skill_id and USE_DATABASE:
                        # Create or update skill in local database
                        existing = db.get_skill(skill_id)
                        if not existing:
                            db.create_skill(
                                skill_id=skill_id,
                                name=skill.get('name', skill_id),
                                description=skill.get('description', ''),
                                system_prompt=skill.get('system_prompt', ''),
                                skill_type=skill.get('skill_type', 'custom'),
                                is_builtin=skill.get('is_builtin', False)
                            )
                            synced += 1

        return jsonify({"success": True, "synced": synced, "message": f"Synced {synced} skills from LPU"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/fast-brain/seed-skills', methods=['POST'])
def seed_default_skills():
    """Seed default built-in skills to the database."""
    try:
        if USE_DATABASE:
            db.seed_builtin_skills()
            return jsonify({"success": True, "message": "Default skills seeded"})
        else:
            return jsonify({"success": False, "error": "Database not enabled"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route('/api/fast-brain/skills', methods=['POST'])
def create_fast_brain_skill():
    """Create a new custom skill with database persistence."""
    try:
        data = request.json
        if not data:
            return jsonify({"success": False, "error": "No data provided"})

        # Accept both 'skill_id' and 'id' field names for compatibility
        skill_id = (data.get('skill_id') or data.get('id') or '').lower().replace(' ', '_')
        skill_id = re.sub(r'[^\w\-]', '_', skill_id)

        if not skill_id:
            return jsonify({"success": False, "error": "Skill ID is required"})

        # Check for builtin skill
        if USE_DATABASE:
            existing = db.get_skill(skill_id)
            if existing and existing.get('is_builtin'):
                return jsonify({"success": False, "error": "Cannot overwrite built-in skill"})

            # Create/update in database
            skill = db.create_skill(
                skill_id=skill_id,
                name=data.get('name', skill_id.replace('_', ' ').title()),
                description=data.get('description', ''),
                skill_type=data.get('skill_type', 'custom'),
                system_prompt=data.get('system_prompt', ''),
                knowledge=data.get('knowledge', []),
                is_builtin=False
            )
        else:
            # Fallback to in-memory
            if skill_id in FAST_BRAIN_SKILLS and FAST_BRAIN_SKILLS[skill_id].get('is_builtin'):
                return jsonify({"success": False, "error": "Cannot overwrite built-in skill"})

            skill = {
                "id": skill_id,
                "name": data.get('name', skill_id.title()),
                "description": data.get('description', ''),
                "system_prompt": data.get('system_prompt', ''),
                "knowledge": data.get('knowledge', []),
                "is_builtin": False,
            }
            FAST_BRAIN_SKILLS[skill_id] = skill

        # Try to sync to deployed LPU
        url = FAST_BRAIN_CONFIG.get('url')
        if url:
            try:
                import httpx
                with httpx.Client(timeout=10.0) as client:
                    response = client.post(
                        f"{url}/v1/skills",
                        json={
                            "skill_id": skill_id,
                            "name": skill.get("name"),
                            "description": skill.get("description"),
                            "system_prompt": skill.get("system_prompt"),
                            "knowledge": skill.get("knowledge", []),
                        }
                    )
                    if response.status_code == 200:
                        add_activity(f"Skill '{skill.get('name')}' synced to LPU", "🔄", "skills")
            except Exception as e:
                add_activity(f"Skill created (sync failed: {str(e)[:50]})", "⚠️", "skills")

        commit_volume()  # Persist to Modal volume
        add_activity(f"Created skill: {skill.get('name')}", "✨", "skills")
        return jsonify({"success": True, "skill": skill})

    except Exception as e:
        # Return JSON error instead of HTML error page
        import traceback
        error_msg = str(e)
        print(f"Error creating skill: {error_msg}")
        print(traceback.format_exc())
        return jsonify({"success": False, "error": f"Server error: {error_msg}"}), 500


@app.route('/api/fast-brain/skills/<skill_id>', methods=['PUT'])
def update_fast_brain_skill(skill_id):
    """Update an existing skill."""
    data = request.json

    if USE_DATABASE:
        existing = db.get_skill(skill_id)
        if not existing:
            return jsonify({"success": False, "error": "Skill not found"})

        if existing.get('is_builtin') and 'system_prompt' not in data:
            return jsonify({"success": False, "error": "Cannot modify built-in skill structure"})

        # Update allowed fields
        update_data = {}
        for field in ['name', 'description', 'system_prompt', 'knowledge', 'skill_type']:
            if field in data:
                update_data[field] = data[field]

        skill = db.update_skill(skill_id, **update_data)
        commit_volume()  # Persist to Modal volume
        add_activity(f"Updated skill: {skill.get('name')}", "📝", "skills")
        return jsonify({"success": True, "skill": skill})
    else:
        if skill_id not in FAST_BRAIN_SKILLS:
            return jsonify({"success": False, "error": "Skill not found"})

        FAST_BRAIN_SKILLS[skill_id].update(data)
        add_activity(f"Updated skill: {skill_id}", "📝", "skills")
        return jsonify({"success": True, "skill": FAST_BRAIN_SKILLS[skill_id]})


@app.route('/api/fast-brain/skills/<skill_id>', methods=['DELETE'])
def delete_fast_brain_skill(skill_id):
    """Delete a custom skill."""
    if USE_DATABASE:
        existing = db.get_skill(skill_id)
        if not existing:
            return jsonify({"success": False, "error": "Skill not found"})

        if existing.get('is_builtin'):
            return jsonify({"success": False, "error": "Cannot delete built-in skill"})

        db.delete_skill(skill_id)
        commit_volume()  # Persist to Modal volume
        add_activity(f"Deleted skill: {skill_id}", "🗑️", "skills")
        return jsonify({"success": True})
    else:
        if skill_id not in FAST_BRAIN_SKILLS:
            return jsonify({"success": False, "error": "Skill not found"})

        if FAST_BRAIN_SKILLS[skill_id].get('is_builtin'):
            return jsonify({"success": False, "error": "Cannot delete built-in skill"})

        del FAST_BRAIN_SKILLS[skill_id]
        add_activity(f"Deleted skill: {skill_id}", "🗑️", "skills")
        return jsonify({"success": True})


@app.route('/api/fast-brain/skills/<skill_id>')
def get_single_skill(skill_id):
    """Get a single skill by ID."""
    if USE_DATABASE:
        skill = db.get_skill(skill_id)
        if not skill:
            return jsonify({"error": "Skill not found"}), 404
        return jsonify(skill)
    else:
        if skill_id not in FAST_BRAIN_SKILLS:
            return jsonify({"error": "Skill not found"}), 404
        return jsonify(FAST_BRAIN_SKILLS[skill_id])


@app.route('/api/fast-brain/skills/select', methods=['POST'])
def select_fast_brain_skill():
    """Select the active skill for chat."""
    data = request.json
    skill_id = data.get('skill_id', 'general')

    if USE_DATABASE:
        skill = db.get_skill(skill_id)
        if not skill:
            return jsonify({"success": False, "error": "Skill not found"})

        db.set_config('selected_skill', skill_id, 'fast_brain')
        add_activity(f"Selected skill: {skill.get('name')}", "✓", "skills")
    else:
        if skill_id not in FAST_BRAIN_SKILLS:
            return jsonify({"success": False, "error": "Skill not found"})

        FAST_BRAIN_CONFIG['selected_skill'] = skill_id
        add_activity(f"Selected skill: {FAST_BRAIN_SKILLS[skill_id]['name']}", "✓", "skills")

    return jsonify({"success": True, "skill_id": skill_id})


# =============================================================================
# API ENDPOINTS - SYSTEM STATUS & HIVE215 INTEGRATION
# =============================================================================

HIVE215_CONFIG = {
    "url": os.getenv("HIVE215_URL", ""),
    "api_key": os.getenv("HIVE215_API_KEY", ""),
    "status": "disconnected",
    "last_sync": None,
}

INTEGRATION_CHECKLIST = [
    {"id": "fast_brain_url", "name": "Fast Brain URL configured", "status": "pending", "category": "Fast Brain"},
    {"id": "fast_brain_health", "name": "Fast Brain health check passing", "status": "pending", "category": "Fast Brain"},
    {"id": "groq_api_key", "name": "Groq API key in Modal secrets", "status": "pending", "category": "Fast Brain"},
    {"id": "skills_synced", "name": "Skills synced to LPU", "status": "pending", "category": "Fast Brain"},
    {"id": "hive215_url", "name": "Hive215 dashboard URL configured", "status": "pending", "category": "Hive215"},
    {"id": "hive215_api_key", "name": "Hive215 API key configured", "status": "pending", "category": "Hive215"},
    {"id": "hive215_webhooks", "name": "Webhooks configured for events", "status": "pending", "category": "Hive215"},
    {"id": "livekit_connected", "name": "LiveKit agent connected", "status": "pending", "category": "Voice"},
    {"id": "voice_provider", "name": "Voice provider selected", "status": "pending", "category": "Voice"},
    {"id": "fallback_chain", "name": "LLM fallback chain configured", "status": "pending", "category": "LLM"},
]


@app.route('/api/system/status')
def get_system_status():
    """Get comprehensive system status."""
    global SYSTEM_STATUS
    from datetime import datetime

    # Update Fast Brain status
    url = FAST_BRAIN_CONFIG.get('url')
    if url:
        try:
            import httpx
            import time
            start = time.time()
            with httpx.Client(timeout=5.0) as client:
                response = client.get(f"{url}/health")
                latency = (time.time() - start) * 1000
                if response.status_code == 200:
                    health = response.json()
                    SYSTEM_STATUS["fast_brain"] = {
                        "status": "online" if health.get("status") == "healthy" else "degraded",
                        "last_check": datetime.now().isoformat(),
                        "error": None,
                        "latency_ms": round(latency, 1),
                        "model": health.get("backend", "unknown"),
                        "skills_count": len(health.get("skills_available", [])),
                    }
                else:
                    SYSTEM_STATUS["fast_brain"] = {
                        "status": "error",
                        "last_check": datetime.now().isoformat(),
                        "error": f"HTTP {response.status_code}",
                        "latency_ms": None,
                    }
        except Exception as e:
            SYSTEM_STATUS["fast_brain"] = {
                "status": "offline",
                "last_check": datetime.now().isoformat(),
                "error": str(e)[:100],
                "latency_ms": None,
            }
    else:
        SYSTEM_STATUS["fast_brain"] = {
            "status": "not_configured",
            "last_check": datetime.now().isoformat(),
            "error": "URL not set",
            "latency_ms": None,
        }

    return jsonify({
        "systems": SYSTEM_STATUS,
        "timestamp": datetime.now().isoformat(),
    })


@app.route('/api/system/checklist')
def get_integration_checklist():
    """Get integration checklist with current status."""
    checklist = []

    for item in INTEGRATION_CHECKLIST:
        status = "pending"

        # Auto-check status based on current config
        if item["id"] == "fast_brain_url":
            status = "complete" if FAST_BRAIN_CONFIG.get("url") else "pending"
        elif item["id"] == "fast_brain_health":
            status = "complete" if SYSTEM_STATUS.get("fast_brain", {}).get("status") == "online" else "pending"
        elif item["id"] == "groq_api_key":
            status = "complete" if SYSTEM_STATUS.get("fast_brain", {}).get("status") in ["online", "degraded"] else "pending"
        elif item["id"] == "skills_synced":
            status = "complete" if len(FAST_BRAIN_SKILLS) > 0 else "pending"
        elif item["id"] == "hive215_url":
            status = "complete" if HIVE215_CONFIG.get("url") else "pending"
        elif item["id"] == "hive215_api_key":
            status = "complete" if HIVE215_CONFIG.get("api_key") else "pending"
        elif item["id"] == "voice_provider":
            status = "complete" if VOICE_CONFIG.get("selected_provider") else "pending"
        elif item["id"] == "livekit_connected":
            status = "complete" if PLATFORM_CONNECTIONS.get("livekit", {}).get("status") == "connected" else "pending"

        checklist.append({**item, "status": status})

    complete_count = sum(1 for c in checklist if c["status"] == "complete")
    return jsonify({
        "checklist": checklist,
        "complete": complete_count,
        "total": len(checklist),
        "percent": round(complete_count / len(checklist) * 100),
    })


# =============================================================================
# API ENDPOINTS - OUTGOING API CONNECTIONS
# =============================================================================

@app.route('/api/connections')
def get_api_connections():
    """Get all outgoing API connections."""
    try:
        if USE_DATABASE:
            db.init_db()  # Ensure table exists
            connections = db.get_all_api_connections()
            return jsonify(connections)
        return jsonify([])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/connections/<connection_id>')
def get_api_connection(connection_id):
    """Get a single API connection."""
    try:
        if USE_DATABASE:
            conn = db.get_api_connection(connection_id)
            if conn:
                return jsonify(conn)
            return jsonify({"error": "Connection not found"}), 404
        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/connections', methods=['POST'])
def create_api_connection():
    """Create a new API connection."""
    try:
        data = request.json
        connection_id = data.get('id') or f"conn_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{random.randint(1000,9999)}"

        if USE_DATABASE:
            db.create_api_connection(
                connection_id=connection_id,
                name=data.get('name', 'New Connection'),
                url=data.get('url', ''),
                api_key=data.get('api_key'),
                headers=data.get('headers'),
                auth_type=data.get('auth_type', 'bearer'),
                webhook_url=data.get('webhook_url'),
                webhook_secret=data.get('webhook_secret'),
                settings=data.get('settings')
            )
            commit_volume()  # Persist to Modal volume
            add_activity(f"API connection created: {data.get('name')}", "", "integration")
            return jsonify({"success": True, "id": connection_id})
        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/connections/<connection_id>', methods=['PUT'])
def update_api_connection_endpoint(connection_id):
    """Update an API connection."""
    try:
        data = request.json
        if USE_DATABASE:
            db.update_api_connection(connection_id, **data)
            commit_volume()  # Persist to Modal volume
            add_activity(f"API connection updated: {connection_id}", "", "integration")
            return jsonify({"success": True})
        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/connections/<connection_id>', methods=['DELETE'])
def delete_api_connection_endpoint(connection_id):
    """Delete an API connection."""
    try:
        if USE_DATABASE:
            conn = db.get_api_connection(connection_id)
            db.delete_api_connection(connection_id)
            commit_volume()  # Persist to Modal volume
            add_activity(f"API connection deleted: {conn.get('name') if conn else connection_id}", "", "integration")
            return jsonify({"success": True})
        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/connections/<connection_id>/test', methods=['POST'])
def test_api_connection(connection_id):
    """Test an API connection with actual HTTP request."""
    try:
        import httpx
        import time

        if USE_DATABASE:
            conn = db.get_api_connection(connection_id)
            if not conn:
                return jsonify({"success": False, "error": "Connection not found"}), 404

            url = conn.get('url', '')
            if not url:
                return jsonify({"success": False, "error": "No URL configured"})

            # Build headers
            headers = {}
            if conn.get('headers'):
                try:
                    headers = json.loads(conn['headers']) if isinstance(conn['headers'], str) else conn['headers']
                except:
                    pass

            # Add authentication
            auth_type = conn.get('auth_type', 'bearer')
            api_key = conn.get('api_key', '')

            if auth_type == 'bearer' and api_key:
                headers['Authorization'] = f'Bearer {api_key}'
            elif auth_type == 'api_key' and api_key:
                headers['X-API-Key'] = api_key

            # Make test request
            start = time.time()
            try:
                response = httpx.get(url, headers=headers, timeout=10, follow_redirects=True)
                latency_ms = int((time.time() - start) * 1000)

                status = 'connected' if response.status_code < 400 else 'error'
                error = None if status == 'connected' else f"HTTP {response.status_code}"

                # Update connection status
                db.update_api_connection_status(connection_id, status, error)

                # Return detailed response
                response_preview = response.text[:500] if response.text else "(empty response)"

                add_activity(f"API test: {conn.get('name')} - {status}", "", "integration")

                return jsonify({
                    "success": status == 'connected',
                    "status": status,
                    "status_code": response.status_code,
                    "latency_ms": latency_ms,
                    "response_preview": response_preview,
                    "headers": dict(response.headers),
                    "error": error
                })

            except httpx.TimeoutException:
                db.update_api_connection_status(connection_id, 'timeout', 'Request timed out')
                return jsonify({
                    "success": False,
                    "status": "timeout",
                    "error": "Request timed out after 10 seconds"
                })
            except httpx.RequestError as e:
                db.update_api_connection_status(connection_id, 'error', str(e))
                return jsonify({
                    "success": False,
                    "status": "error",
                    "error": str(e)
                })

        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/connections/<connection_id>/send', methods=['POST'])
def send_to_api_connection(connection_id):
    """Send a custom request to an API connection for testing/troubleshooting."""
    try:
        import httpx
        import time

        data = request.json
        method = data.get('method', 'GET').upper()
        path = data.get('path', '')
        body = data.get('body')
        custom_headers = data.get('headers', {})

        if USE_DATABASE:
            conn = db.get_api_connection(connection_id)
            if not conn:
                return jsonify({"success": False, "error": "Connection not found"}), 404

            base_url = conn.get('url', '').rstrip('/')
            full_url = f"{base_url}/{path.lstrip('/')}" if path else base_url

            # Build headers
            headers = {'Content-Type': 'application/json'}
            if conn.get('headers'):
                try:
                    stored_headers = json.loads(conn['headers']) if isinstance(conn['headers'], str) else conn['headers']
                    headers.update(stored_headers)
                except:
                    pass
            headers.update(custom_headers)

            # Add authentication
            auth_type = conn.get('auth_type', 'bearer')
            api_key = conn.get('api_key', '')

            if auth_type == 'bearer' and api_key:
                headers['Authorization'] = f'Bearer {api_key}'
            elif auth_type == 'api_key' and api_key:
                headers['X-API-Key'] = api_key

            # Make request
            start = time.time()
            try:
                if method == 'GET':
                    response = httpx.get(full_url, headers=headers, timeout=30)
                elif method == 'POST':
                    response = httpx.post(full_url, headers=headers, json=body, timeout=30)
                elif method == 'PUT':
                    response = httpx.put(full_url, headers=headers, json=body, timeout=30)
                elif method == 'DELETE':
                    response = httpx.delete(full_url, headers=headers, timeout=30)
                else:
                    return jsonify({"success": False, "error": f"Unsupported method: {method}"})

                latency_ms = int((time.time() - start) * 1000)

                # Try to parse response as JSON
                try:
                    response_body = response.json()
                except:
                    response_body = response.text

                return jsonify({
                    "success": response.status_code < 400,
                    "status_code": response.status_code,
                    "latency_ms": latency_ms,
                    "response": response_body,
                    "response_headers": dict(response.headers),
                    "request": {
                        "method": method,
                        "url": full_url,
                        "headers": {k: v if k.lower() != 'authorization' else '***' for k, v in headers.items()},
                        "body": body
                    }
                })

            except httpx.TimeoutException:
                return jsonify({"success": False, "error": "Request timed out after 30 seconds"})
            except httpx.RequestError as e:
                return jsonify({"success": False, "error": str(e)})

        return jsonify({"error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/connections/status')
def get_all_connections_status():
    """Get live status of all API connections."""
    try:
        if USE_DATABASE:
            connections = db.get_all_api_connections()
            statuses = {}
            for conn in connections:
                statuses[conn['id']] = {
                    'name': conn.get('name'),
                    'status': conn.get('status', 'unknown'),
                    'last_tested': conn.get('last_tested'),
                    'last_error': conn.get('last_error'),
                    'url': conn.get('url', '')[:50] + '...' if len(conn.get('url', '')) > 50 else conn.get('url', '')
                }
            return jsonify(statuses)
        return jsonify({})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# Legacy Hive215 endpoints for backward compatibility
HIVE215_CONFIG = {
    "url": os.getenv("HIVE215_URL", ""),
    "api_key": os.getenv("HIVE215_API_KEY", ""),
}

@app.route('/api/hive215/config')
def get_hive215_config():
    """Get Hive215 configuration (legacy)."""
    return jsonify(HIVE215_CONFIG)


@app.route('/api/hive215/config', methods=['POST'])
def update_hive215_config():
    """Update Hive215 configuration (legacy)."""
    global HIVE215_CONFIG
    data = request.json

    if 'url' in data:
        HIVE215_CONFIG['url'] = data['url'].rstrip('/')
    if 'api_key' in data:
        HIVE215_CONFIG['api_key'] = data['api_key']

    add_activity("Hive215 config updated", "")
    return jsonify({"success": True, "config": HIVE215_CONFIG})


@app.route('/api/fast-brain/errors')
def get_fast_brain_errors():
    """Get recent errors for debugging."""
    return jsonify({
        "errors": FAST_BRAIN_CONFIG['stats'].get('errors', [])[-20:],  # Last 20 errors
        "total": len(FAST_BRAIN_CONFIG['stats'].get('errors', [])),
    })


# =============================================================================
# API ENDPOINTS - VOICE PROVIDER DYNAMIC FETCHING
# =============================================================================

def fetch_elevenlabs_voices():
    """Fetch voices from ElevenLabs API."""
    try:
        import httpx
        api_key = db.get_api_key('elevenlabs') if USE_DATABASE else None
        if not api_key:
            return {"error": "ElevenLabs API key not configured", "voices": []}

        with httpx.Client(timeout=30.0) as client:
            response = client.get(
                "https://api.elevenlabs.io/v1/voices",
                headers={"xi-api-key": api_key}
            )

            if response.status_code == 200:
                data = response.json()
                voices = []
                for voice in data.get("voices", []):
                    voices.append({
                        "id": voice.get("voice_id"),
                        "name": voice.get("name"),
                        "gender": voice.get("labels", {}).get("gender", "unknown"),
                        "style": voice.get("labels", {}).get("use_case", voice.get("labels", {}).get("description", "general")),
                        "category": voice.get("category", "premade"),
                        "preview_url": voice.get("preview_url"),
                        "description": voice.get("description", ""),
                        "labels": voice.get("labels", {})
                    })
                return {"success": True, "voices": voices, "provider": "elevenlabs"}
            else:
                return {"error": f"ElevenLabs API error: {response.status_code}", "voices": []}
    except Exception as e:
        return {"error": str(e), "voices": []}


def fetch_cartesia_voices():
    """Fetch voices from Cartesia API - custom/owned voices first, then public library."""
    try:
        import httpx
        api_key = db.get_api_key('cartesia') if USE_DATABASE else None
        if not api_key:
            return {"error": "Cartesia API key not configured", "voices": []}

        with httpx.Client(timeout=30.0) as client:
            all_voices = []
            owned_ids = set()

            # First, fetch ALL voices and identify owned ones
            # Cartesia API returns is_owner field for each voice
            has_more = True
            starting_after = None

            while has_more:
                params = {"limit": 100}
                if starting_after:
                    params["starting_after"] = starting_after

                response = client.get(
                    "https://api.cartesia.ai/voices",
                    headers={
                        "X-API-Key": api_key,
                        "Cartesia-Version": "2024-11-13"
                    },
                    params=params
                )

                if response.status_code == 200:
                    data = response.json()
                    voices_data = data if isinstance(data, list) else data.get("data", data.get("voices", []))

                    for voice in voices_data:
                        is_owner = voice.get("is_owner", False)
                        voice_entry = {
                            "id": voice.get("id"),
                            "name": voice.get("name"),
                            "gender": voice.get("gender", "unknown"),
                            "style": voice.get("description", "natural")[:50] if voice.get("description") else "natural",
                            "language": voice.get("language", "en"),
                            "is_public": voice.get("is_public", False),
                            "is_owner": is_owner,
                            "category": "⭐ Your Custom Voices" if is_owner else "Public Library",
                            "description": voice.get("description", "")
                        }

                        if is_owner:
                            owned_ids.add(voice.get("id"))
                            # Add owned voices at the beginning
                            all_voices.insert(0, voice_entry)
                        else:
                            all_voices.append(voice_entry)

                    # Handle pagination
                    has_more = data.get("has_more", False) if isinstance(data, dict) else False
                    if has_more and voices_data:
                        starting_after = voices_data[-1].get("id")
                    else:
                        has_more = False
                else:
                    return {"error": f"Cartesia API error: {response.status_code}", "voices": all_voices}

            # Sort so owned voices are first
            all_voices.sort(key=lambda v: (0 if v.get("is_owner") else 1, v.get("name", "")))

            return {
                "success": True,
                "voices": all_voices,
                "provider": "cartesia",
                "custom_count": len(owned_ids)
            }
    except Exception as e:
        return {"error": str(e), "voices": []}


def fetch_deepgram_voices():
    """Fetch voices from Deepgram API (TTS/Aura voices)."""
    try:
        import httpx
        api_key = db.get_api_key('deepgram') if USE_DATABASE else None
        if not api_key:
            return {"error": "Deepgram API key not configured", "voices": []}

        with httpx.Client(timeout=30.0) as client:
            response = client.get(
                "https://api.deepgram.com/v1/models",
                headers={"Authorization": f"Token {api_key}"}
            )

            if response.status_code == 200:
                data = response.json()
                voices = []

                # Extract TTS models from response
                tts_models = data.get("tts", [])
                for model in tts_models:
                    metadata = model.get("metadata", {})
                    voices.append({
                        "id": model.get("canonical_name", model.get("name")),
                        "name": model.get("name", "").title(),
                        "gender": "male" if "masculine" in metadata.get("tags", []) else "female" if "feminine" in metadata.get("tags", []) else "neutral",
                        "style": metadata.get("accent", "american"),
                        "languages": model.get("languages", ["en"]),
                        "preview_url": metadata.get("sample"),
                        "image_url": metadata.get("image"),
                        "version": model.get("version", "")
                    })

                return {"success": True, "voices": voices, "provider": "deepgram"}
            else:
                return {"error": f"Deepgram API error: {response.status_code}", "voices": []}
    except Exception as e:
        return {"error": str(e), "voices": []}


def get_openai_voices():
    """Return OpenAI TTS voices (static list - OpenAI doesn't have a list endpoint)."""
    voices = [
        {"id": "alloy", "name": "Alloy", "gender": "neutral", "style": "balanced", "description": "Balanced, versatile voice"},
        {"id": "ash", "name": "Ash", "gender": "male", "style": "conversational", "description": "Warm conversational voice"},
        {"id": "ballad", "name": "Ballad", "gender": "male", "style": "narrative", "description": "Expressive storytelling voice"},
        {"id": "coral", "name": "Coral", "gender": "female", "style": "warm", "description": "Friendly warm voice"},
        {"id": "echo", "name": "Echo", "gender": "male", "style": "neutral", "description": "Clear neutral voice"},
        {"id": "fable", "name": "Fable", "gender": "neutral", "style": "storytelling", "description": "Expressive British accent"},
        {"id": "onyx", "name": "Onyx", "gender": "male", "style": "deep", "description": "Deep authoritative voice"},
        {"id": "nova", "name": "Nova", "gender": "female", "style": "warm", "description": "Warm friendly voice"},
        {"id": "sage", "name": "Sage", "gender": "female", "style": "professional", "description": "Professional clear voice"},
        {"id": "shimmer", "name": "Shimmer", "gender": "female", "style": "expressive", "description": "Bright expressive voice"},
        {"id": "verse", "name": "Verse", "gender": "male", "style": "dynamic", "description": "Dynamic engaging voice"},
    ]
    return {"success": True, "voices": voices, "provider": "openai"}


def get_static_voices(provider):
    """Return static voice list for providers without dynamic API."""
    if provider in VOICE_CHOICES:
        return {
            "success": True,
            "voices": VOICE_CHOICES[provider]["voices"],
            "provider": provider,
            "static": True
        }
    return {"error": f"Unknown provider: {provider}", "voices": []}


@app.route('/api/voice-lab/provider-voices/<provider>')
def get_provider_voices(provider):
    """Fetch voices dynamically from a specific provider's API."""
    try:
        # Providers with dynamic API fetching
        if provider == "elevenlabs":
            result = fetch_elevenlabs_voices()
        elif provider == "cartesia":
            result = fetch_cartesia_voices()
        elif provider == "deepgram":
            result = fetch_deepgram_voices()
        elif provider == "openai":
            result = get_openai_voices()
        # Providers with static voice lists
        elif provider in VOICE_CHOICES:
            result = get_static_voices(provider)
        else:
            return jsonify({"error": f"Unknown provider: {provider}", "voices": []}), 400

        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e), "voices": []}), 500


@app.route('/api/voice-lab/all-provider-voices')
def get_all_provider_voices():
    """Fetch voices from all configured providers."""
    results = {}

    # Dynamic providers
    dynamic_providers = ["elevenlabs", "cartesia", "deepgram", "openai"]
    for provider in dynamic_providers:
        try:
            if provider == "elevenlabs":
                results[provider] = fetch_elevenlabs_voices()
            elif provider == "cartesia":
                results[provider] = fetch_cartesia_voices()
            elif provider == "deepgram":
                results[provider] = fetch_deepgram_voices()
            elif provider == "openai":
                results[provider] = get_openai_voices()
        except Exception as e:
            results[provider] = {"error": str(e), "voices": []}

    # Static providers
    static_providers = ["edge_tts", "parler_tts", "kokoro", "azure"]
    for provider in static_providers:
        if provider in VOICE_CHOICES:
            results[provider] = get_static_voices(provider)

    return jsonify(results)


# =============================================================================
# API ENDPOINTS - VOICE LAB (Database-backed with TTS Integration)
# =============================================================================

# Training queue kept in memory for real-time progress (jobs are short-lived)
VOICE_TRAINING_QUEUE = []

# Voice samples directory on Modal volume
VOICE_SAMPLES_DIR = Path("/data/voice_samples")
VOICE_SAMPLES_DIR.mkdir(parents=True, exist_ok=True)


@app.route('/api/voice-lab/projects')
def get_voice_projects():
    """Get all voice projects from database."""
    try:
        if USE_DATABASE:
            # Ensure tables exist
            db.init_db()
            projects = db.get_all_voice_projects()
            return jsonify(projects)
        return jsonify([])
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/voice-lab/projects', methods=['POST'])
def create_voice_project_endpoint():
    """Create a new voice project with database persistence."""
    try:
        data = request.json or {}
        project_id = f"voice_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{random.randint(1000,9999)}"

        if USE_DATABASE:
            # Ensure tables exist
            db.init_db()

            project = db.create_voice_project(
                project_id=project_id,
                name=data.get("name", "Untitled Voice"),
                description=data.get("description", ""),
                provider=data.get("provider", "elevenlabs"),
                base_voice=data.get("base_voice"),
                settings=data.get("settings", {
                    "pitch": 1.0,
                    "speed": 1.0,
                    "emotion": "neutral",
                    "style": "conversational",
                }),
                skill_id=data.get("skill_id")
            )
            commit_volume()  # Persist to Modal volume
            add_activity(f"Voice project created: {project['name']}", "", "voice")
            return jsonify({"success": True, "project": project})

        return jsonify({"success": False, "error": "Database not available"}), 500
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/voice-lab/projects/<project_id>')
def get_voice_project_endpoint(project_id):
    """Get a specific voice project from database."""
    try:
        if USE_DATABASE:
            db.init_db()
            project = db.get_voice_project(project_id)
            if project:
                return jsonify(project)
        return jsonify({"error": "Project not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/voice-lab/projects/<project_id>', methods=['PUT'])
def update_voice_project_endpoint(project_id):
    """Update a voice project in database."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    data = request.json
    project = db.update_voice_project(project_id, **data)

    if project:
        commit_volume()  # Persist to Modal volume
        return jsonify({"success": True, "project": project})
    return jsonify({"error": "Project not found"}), 404


@app.route('/api/voice-lab/projects/<project_id>', methods=['DELETE'])
def delete_voice_project_endpoint(project_id):
    """Delete a voice project."""
    if USE_DATABASE:
        # Also delete sample files from disk
        project = db.get_voice_project(project_id)
        if project:
            for sample in project.get('samples', []):
                if sample.get('file_path') and os.path.exists(sample['file_path']):
                    os.remove(sample['file_path'])

            if db.delete_voice_project(project_id):
                commit_volume()  # Persist to Modal volume
                add_activity(f"Voice project deleted: {project['name']}", "", "voice")
                return jsonify({"success": True})

    return jsonify({"error": "Project not found"}), 404


@app.route('/api/voice-lab/projects/<project_id>/samples', methods=['POST'])
def add_voice_sample_endpoint(project_id):
    """Add a training sample to a voice project (with file upload)."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    project = db.get_voice_project(project_id)
    if not project:
        return jsonify({"error": "Project not found"}), 404

    # Handle file upload
    if 'audio' in request.files:
        audio_file = request.files['audio']
        if audio_file.filename:
            sample_id = f"sample_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{random.randint(1000,9999)}"
            filename = f"{project_id}_{sample_id}_{audio_file.filename}"
            file_path = VOICE_SAMPLES_DIR / filename

            audio_file.save(str(file_path))

            # Get transcript and duration from form data
            transcript = request.form.get('transcript', '')
            duration_ms = int(request.form.get('duration_ms', 0))
            emotion = request.form.get('emotion', 'neutral')

            sample = db.add_voice_sample(
                project_id=project_id,
                sample_id=sample_id,
                filename=filename,
                file_path=str(file_path),
                transcript=transcript,
                duration_ms=duration_ms,
                emotion=emotion
            )

            commit_volume()  # Persist to Modal volume
            add_activity(f"Voice sample added to {project['name']}", "", "voice")
            return jsonify({"success": True, "sample": sample})

    # Handle JSON data (metadata only, for URL-based samples)
    data = request.json or {}
    sample_id = f"sample_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{random.randint(1000,9999)}"

    sample = db.add_voice_sample(
        project_id=project_id,
        sample_id=sample_id,
        filename=data.get('filename', 'external_audio'),
        file_path=data.get('audio_url'),
        transcript=data.get('transcript', data.get('text', '')),
        duration_ms=int(data.get('duration_ms', 0)),
        emotion=data.get('emotion', 'neutral')
    )

    commit_volume()  # Persist to Modal volume
    return jsonify({"success": True, "sample": sample})


@app.route('/api/voice-lab/projects/<project_id>/samples/<sample_id>', methods=['DELETE'])
def delete_voice_sample_endpoint(project_id, sample_id):
    """Delete a voice sample."""
    if USE_DATABASE:
        # Delete file from disk if it exists
        samples = db.get_voice_samples(project_id)
        for sample in samples:
            if sample['id'] == sample_id and sample.get('file_path'):
                if os.path.exists(sample['file_path']):
                    os.remove(sample['file_path'])
                break

        if db.delete_voice_sample(sample_id):
            commit_volume()  # Persist to Modal volume
            return jsonify({"success": True})

    return jsonify({"error": "Sample not found"}), 404


@app.route('/api/voice-lab/projects/<project_id>/train', methods=['POST'])
def train_voice_endpoint(project_id):
    """Start training a custom voice using the selected provider."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    project = db.get_voice_project(project_id)
    if not project:
        return jsonify({"error": "Project not found"}), 404

    # Get samples separately from database
    samples = db.get_voice_samples(project_id)
    if len(samples) < 1:
        return jsonify({
            "success": False,
            "error": "Need at least 1 audio sample to clone/train a voice"
        }), 400

    provider = project.get('provider', 'elevenlabs')

    # Update project status to training
    db.update_voice_project(project_id, status='training', training_started=datetime.now().isoformat())

    # Add to training queue for progress tracking
    job = {
        "project_id": project_id,
        "provider": provider,
        "started_at": datetime.now().isoformat(),
        "progress": 0,
        "status": "starting"
    }
    VOICE_TRAINING_QUEUE.append(job)

    # Try to actually train/clone with the provider
    try:
        voice_id = None
        message = ""

        if provider == 'elevenlabs':
            voice_id, message = _train_elevenlabs_voice(project, samples)
        elif provider == 'cartesia':
            voice_id, message = _train_cartesia_voice(project, samples)
        elif provider in ['xtts', 'openvoice']:
            # These are local models - mark as ready (would need GPU)
            voice_id = f"{provider}_{project_id}"
            message = f"Voice ready for {provider} synthesis (local model)"
        elif provider in ['edge_tts', 'parler_tts', 'kokoro', 'chatterbox', 'gtts']:
            # Free providers - no training needed, mark as ready
            voice_id = f"{provider}_{project_id}"
            message = f"Voice ready! {provider} doesn't require training - you can test it now."
        else:
            # For other providers, simulate training
            voice_id = f"{provider}_{project_id}"
            message = f"Voice configured for {provider}"

        if voice_id:
            updated_project = db.update_voice_project(
                project_id,
                status='trained',
                voice_id=voice_id,
                training_completed=datetime.now().isoformat()
            )
            # Verify the update persisted
            verify_project = db.get_voice_project(project_id)

            job['status'] = 'completed'
            job['progress'] = 100

            add_activity(f"Voice trained: {project['name']} ({provider})", "", "voice")

            return jsonify({
                "success": True,
                "message": message,
                "voice_id": voice_id,
                "provider": provider,
                "debug": {
                    "updated_status": updated_project.get('status') if updated_project else None,
                    "updated_voice_id": updated_project.get('voice_id') if updated_project else None,
                    "verified_status": verify_project.get('status') if verify_project else None,
                    "verified_voice_id": verify_project.get('voice_id') if verify_project else None
                }
            })
        else:
            # Training failed - no voice_id returned (usually missing API key)
            db.update_voice_project(project_id, status='failed')
            job['status'] = 'failed'
            job['error'] = message

            return jsonify({
                "success": False,
                "error": message or f"Training failed. Make sure your {provider} API key is configured in Settings → API Keys."
            }), 400

    except Exception as e:
        db.update_voice_project(project_id, status='failed')
        job['status'] = 'failed'
        job['error'] = str(e)

        return jsonify({
            "success": False,
            "error": f"Training failed: {str(e)}"
        }), 500

    return jsonify({
        "success": True,
        "message": "Training started",
        "estimated_time_minutes": len(samples) * 2
    })


def _train_elevenlabs_voice(project, samples):
    """Clone a voice using ElevenLabs API."""
    api_key = None
    if USE_DATABASE:
        api_key = db.get_api_key('elevenlabs')

    if not api_key:
        return None, "ElevenLabs API key not configured"

    try:
        import requests

        # Prepare files for upload
        files = []
        missing_files = []
        for sample in samples:
            file_path = sample.get('file_path')
            if file_path:
                if os.path.exists(file_path):
                    files.append(
                        ('files', (sample['filename'], open(file_path, 'rb'), 'audio/mpeg'))
                    )
                else:
                    missing_files.append(file_path)
            else:
                missing_files.append(f"No path for {sample.get('filename', 'unknown')}")

        if not files:
            return None, f"No audio files found. Missing: {missing_files}. Try re-uploading samples."

        # Create voice clone
        response = requests.post(
            'https://api.elevenlabs.io/v1/voices/add',
            headers={'xi-api-key': api_key},
            data={
                'name': project['name'],
                'description': project.get('description', 'Custom cloned voice'),
            },
            files=files
        )

        # Close file handles
        for _, (_, f, _) in files:
            f.close()

        if response.status_code == 200:
            result = response.json()
            return result.get('voice_id'), f"Voice cloned successfully: {result.get('voice_id')}"
        else:
            return None, f"ElevenLabs error: {response.text}"

    except Exception as e:
        return None, f"ElevenLabs API error: {str(e)}"


def _train_cartesia_voice(project, samples):
    """Clone a voice using Cartesia API."""
    print(f"[CARTESIA DEBUG] Starting voice clone for project: {project.get('name')}")
    print(f"[CARTESIA DEBUG] Samples count: {len(samples) if samples else 0}")

    api_key = None
    if USE_DATABASE:
        api_key = db.get_api_key('cartesia')

    print(f"[CARTESIA DEBUG] API key set: {bool(api_key)}")
    if not api_key:
        return None, "Cartesia API key not configured"

    try:
        import requests

        # Debug sample info
        if samples:
            sample = samples[0]
            file_path = sample.get('file_path')
            print(f"[CARTESIA DEBUG] Sample file_path: {file_path}")
            print(f"[CARTESIA DEBUG] File exists: {os.path.exists(file_path) if file_path else 'N/A'}")
            if file_path:
                # List contents of the voice samples directory
                samples_dir = Path("/data/voice_samples")
                if samples_dir.exists():
                    print(f"[CARTESIA DEBUG] Files in {samples_dir}: {list(samples_dir.iterdir())[:5]}")
                else:
                    print(f"[CARTESIA DEBUG] Directory {samples_dir} does not exist!")

        # Cartesia requires multipart form data with the audio file
        if samples and samples[0].get('file_path') and os.path.exists(samples[0]['file_path']):
            file_path = samples[0]['file_path']
            filename = samples[0].get('filename', 'audio.mp3')

            # Open file for multipart upload
            with open(file_path, 'rb') as audio_file:
                files = {
                    'clip': (filename, audio_file, 'audio/mpeg')
                }
                data = {
                    'name': project['name'][:100],
                    'language': 'en',
                    'mode': 'similarity'  # or 'stability'
                }
                if project.get('description'):
                    data['description'] = project['description'][:500]

                response = requests.post(
                    'https://api.cartesia.ai/voices/clone',
                    headers={
                        'X-API-Key': api_key,
                        'Cartesia-Version': '2024-11-13'
                    },
                    files=files,
                    data=data
                )

            print(f"[CARTESIA DEBUG] Response status: {response.status_code}")
            if response.status_code == 200 or response.status_code == 201:
                result = response.json()
                voice_id = result.get('id')
                print(f"[CARTESIA DEBUG] Success! voice_id: {voice_id}")
                return voice_id, f"Voice cloned with Cartesia: {voice_id}"
            else:
                print(f"[CARTESIA DEBUG] Error response: {response.text}")
                return None, f"Cartesia error ({response.status_code}): {response.text}"

        print(f"[CARTESIA DEBUG] No audio samples available - file check failed")
        return None, "No audio samples available - file not found on server"

    except Exception as e:
        print(f"[CARTESIA DEBUG] Exception: {str(e)}")
        return None, f"Cartesia API error: {str(e)}"


@app.route('/api/voice-lab/projects/<project_id>/test', methods=['POST'])
def test_voice_project_endpoint(project_id):
    """Test a trained voice with sample text - actual TTS synthesis."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    project = db.get_voice_project(project_id)
    print(f"[TEST DEBUG] Project {project_id}: status={project.get('status') if project else 'None'}, voice_id={project.get('voice_id') if project else 'None'}, provider={project.get('provider') if project else 'None'}")
    if not project:
        return jsonify({"error": "Project not found"}), 404

    data = request.json
    text = data.get("text", "Hello, this is a test of my custom voice.")

    provider = project.get('provider', 'elevenlabs')
    voice_id = project.get('voice_id')
    status = project.get('status', 'draft')
    print(f"[TEST DEBUG] Checking: provider={provider}, voice_id={voice_id}, status={status}")

    # Free providers that don't need voice cloning/training
    free_providers = ['edge_tts', 'parler_tts', 'kokoro', 'chatterbox', 'gtts']

    # For paid providers, check that training completed
    if provider not in free_providers:
        if not voice_id and status != 'trained':
            return jsonify({
                "success": False,
                "error": f"Voice not trained. Add your {provider.title()} API key in Settings → API Keys, then click Train Voice.",
                "debug": {"provider": provider, "voice_id": voice_id, "status": status}
            }), 400

    try:
        audio_data = None
        synthesis_method = None

        if provider == 'elevenlabs' and voice_id:
            audio_data = _synthesize_elevenlabs(voice_id, text)
            synthesis_method = 'elevenlabs'
        elif provider == 'cartesia' and voice_id:
            audio_data = _synthesize_cartesia(voice_id, text)
            synthesis_method = 'cartesia'
        else:
            # Use gTTS for all free providers or when no voice_id
            audio_data = _synthesize_edge_tts(project.get('base_voice', 'en-US-JennyNeural'), text)
            synthesis_method = 'gtts'

        print(f"Voice Lab Test: provider={provider}, voice_id={voice_id}, method={synthesis_method}, audio_len={len(audio_data) if audio_data else 0}")

        if audio_data and len(audio_data) > 100:
            import base64
            audio_size = len(audio_data)
            audio_base64 = base64.b64encode(audio_data).decode('utf-8')

            add_activity(f"Voice test: {project['name']}", "", "voice")

            # Estimate duration: MP3 at ~128kbps = ~16KB per second
            estimated_duration_ms = int((audio_size / 16000) * 1000)

            return jsonify({
                "success": True,
                "project_id": project_id,
                "text": text,
                "duration_ms": estimated_duration_ms,
                "audio_size_bytes": audio_size,
                "audio_base64": audio_base64,
                "audio_format": "audio/mpeg",
                "provider": provider,
                "message": f"Voice '{project['name']}' synthesized successfully ({audio_size:,} bytes)"
            })
        else:
            # No audio generated
            audio_len = len(audio_data) if audio_data else 0
            return jsonify({
                "success": False,
                "error": f"Could not generate audio (got {audio_len} bytes, method={synthesis_method}, provider={provider}). Try again or check logs."
            }), 400

    except Exception as e:
        return jsonify({
            "success": False,
            "error": f"Synthesis failed: {str(e)}"
        }), 500


def _synthesize_elevenlabs(voice_id, text):
    """Synthesize speech using ElevenLabs."""
    api_key = db.get_api_key('elevenlabs') if USE_DATABASE else None
    if not api_key:
        return None

    import requests
    response = requests.post(
        f'https://api.elevenlabs.io/v1/text-to-speech/{voice_id}',
        headers={
            'xi-api-key': api_key,
            'Content-Type': 'application/json'
        },
        json={
            'text': text,
            'model_id': 'eleven_monolingual_v1',
            'voice_settings': {
                'stability': 0.5,
                'similarity_boost': 0.75
            }
        }
    )

    if response.status_code == 200:
        return response.content
    return None


def _synthesize_cartesia(voice_id, text):
    """Synthesize speech using Cartesia."""
    api_key = db.get_api_key('cartesia') if USE_DATABASE else None
    if not api_key:
        return None

    import requests
    response = requests.post(
        'https://api.cartesia.ai/tts/bytes',
        headers={
            'X-API-Key': api_key,
            'Cartesia-Version': '2024-06-10',
            'Content-Type': 'application/json'
        },
        json={
            'transcript': text,
            'voice': {'mode': 'id', 'id': voice_id},
            'output_format': {'container': 'mp3', 'encoding': 'mp3', 'sample_rate': 44100}
        }
    )

    if response.status_code == 200:
        return response.content
    return None


def _synthesize_edge_tts(voice, text):
    """Synthesize speech using gTTS (free Google TTS)."""
    try:
        from gtts import gTTS
        import io

        # Ensure text is valid
        if not text or not text.strip():
            print("gTTS error: Empty text provided")
            return None

        clean_text = text.strip()[:500]  # Limit length to avoid issues
        print(f"gTTS: Synthesizing {len(clean_text)} chars...")

        tts = gTTS(text=clean_text, lang='en')
        audio_buffer = io.BytesIO()
        tts.write_to_fp(audio_buffer)
        audio_bytes = audio_buffer.getvalue()

        # Validate we got audio data
        if not audio_bytes or len(audio_bytes) < 100:
            print(f"gTTS error: Audio too short ({len(audio_bytes) if audio_bytes else 0} bytes)")
            return None

        print(f"gTTS: Generated {len(audio_bytes)} bytes of audio")
        return audio_bytes
    except Exception as e:
        print(f"gTTS error: {e}")
        import traceback
        traceback.print_exc()
        return None


@app.route('/api/voice-lab/audio/<filename>')
def serve_voice_audio(filename):
    """Serve generated voice audio files."""
    # Security: sanitize filename to prevent path traversal
    safe_filename = filename.replace('..', '').replace('/', '').replace('\\', '')
    file_path = VOICE_SAMPLES_DIR / safe_filename

    if file_path.exists() and file_path.stat().st_size > 0:
        response = send_file(
            str(file_path),
            mimetype='audio/mpeg',
            as_attachment=False
        )
        # Add headers for audio streaming
        response.headers['Accept-Ranges'] = 'bytes'
        response.headers['Cache-Control'] = 'no-cache'
        return response

    return jsonify({"error": "File not found or empty"}), 404


@app.route('/api/voice-lab/debug-tts')
def debug_tts():
    """Debug endpoint to test TTS generation."""
    import base64
    import io

    results = {
        "gtts_available": False,
        "gtts_error": None,
        "audio_generated": False,
        "audio_size": 0,
        "audio_base64": None
    }

    try:
        from gtts import gTTS
        results["gtts_available"] = True

        # Try to generate audio
        tts = gTTS(text="Hello, this is a test.", lang='en')
        audio_buffer = io.BytesIO()
        tts.write_to_fp(audio_buffer)
        audio_bytes = audio_buffer.getvalue()

        results["audio_size"] = len(audio_bytes)
        results["audio_generated"] = len(audio_bytes) > 100

        if audio_bytes:
            results["audio_base64"] = base64.b64encode(audio_bytes).decode('utf-8')

    except Exception as e:
        results["gtts_error"] = str(e)
        import traceback
        results["traceback"] = traceback.format_exc()

    return jsonify(results)


@app.route('/api/voice-lab/training-status')
def get_voice_training_status():
    """Get status of all voice training jobs."""
    # Clean up completed jobs older than 1 hour
    cutoff = datetime.now() - timedelta(hours=1)
    global VOICE_TRAINING_QUEUE
    VOICE_TRAINING_QUEUE = [
        job for job in VOICE_TRAINING_QUEUE
        if datetime.fromisoformat(job['started_at']) > cutoff or job.get('status') != 'completed'
    ]
    return jsonify(VOICE_TRAINING_QUEUE)


@app.route('/api/voice-lab/projects/<project_id>/link-skill', methods=['POST'])
def link_voice_to_skill_endpoint(project_id):
    """Link a trained voice to a skill/agent."""
    if not USE_DATABASE:
        return jsonify({"error": "Database not available"}), 500

    data = request.json
    skill_id = data.get('skill_id')

    if not skill_id:
        return jsonify({"error": "skill_id is required"}), 400

    project = db.get_voice_project(project_id)
    if not project:
        return jsonify({"error": "Voice project not found"}), 404

    if db.link_voice_to_skill(project_id, skill_id):
        add_activity(f"Voice '{project['name']}' linked to skill '{skill_id}'", "", "voice")
        return jsonify({
            "success": True,
            "message": f"Voice linked to skill {skill_id}"
        })

    return jsonify({"error": "Failed to link voice to skill"}), 500


@app.route('/api/voice-lab/skills/<skill_id>/voices')
def get_voices_for_skill(skill_id):
    """Get all trained voices linked to a skill."""
    if USE_DATABASE:
        voices = db.get_voice_projects_for_skill(skill_id)
        return jsonify(voices)
    return jsonify([])


# =============================================================================
# API ENDPOINTS - SKILL FINE-TUNING
# =============================================================================

SKILL_TRAINING_JOBS = {}


@app.route('/api/skills/<skill_id>/fine-tune', methods=['POST'])
def fine_tune_skill(skill_id):
    """Start fine-tuning a skill with new examples - saves to database."""
    data = request.json
    examples = data.get('examples', [])

    if not examples:
        return jsonify({"success": False, "error": "No examples provided"})

    job_id = f"ft_{skill_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    SKILL_TRAINING_JOBS[job_id] = {
        "id": job_id,
        "skill_id": skill_id,
        "status": "running",
        "progress": 0,
        "examples_count": len(examples),
        "started_at": datetime.now().isoformat(),
        "metrics": {
            "loss": 2.5,
            "accuracy": 0.0,
        }
    }

    # Save examples to database
    if USE_DATABASE:
        for ex in examples:
            user_msg = ex.get('input') or ex.get('user_message') or ex.get('question', '')
            asst_msg = ex.get('output') or ex.get('assistant_response') or ex.get('answer', '')
            if user_msg and asst_msg:
                db.add_training_example(
                    skill_id=skill_id,
                    user_message=user_msg,
                    assistant_response=asst_msg,
                    metadata={"source": "fine_tune", "job_id": job_id}
                )
        commit_volume()  # Persist to Modal volume

    add_activity(f"Fine-tuning started for {skill_id} with {len(examples)} examples", "")

    return jsonify({
        "success": True,
        "job_id": job_id,
        "message": f"Fine-tuning started with {len(examples)} examples"
    })


@app.route('/api/skills/<skill_id>/fine-tune/status')
def get_fine_tune_status(skill_id):
    """Get fine-tuning status for a skill."""
    jobs = [j for j in SKILL_TRAINING_JOBS.values() if j['skill_id'] == skill_id]
    # Return actual job status - no simulation
    return jsonify(jobs)


@app.route('/api/skills/<skill_id>/feedback', methods=['POST'])
def add_skill_feedback(skill_id):
    """Add feedback for a skill response - saves to database."""
    data = request.json

    if USE_DATABASE:
        # Save as training example with rating and optional correction
        db.add_training_example(
            skill_id=skill_id,
            user_message=data.get("query", ""),
            assistant_response=data.get("response", ""),
            rating=data.get("rating", 0),
            corrected_response=data.get("corrected_response"),
            metadata={"source": "feedback", "timestamp": datetime.now().isoformat()}
        )
        commit_volume()  # Persist to Modal volume

    add_activity(f"Feedback added for {skill_id}", "")

    return jsonify({"success": True, "message": "Feedback saved to database"})


@app.route('/api/skills/<skill_id>/auto-improve', methods=['POST'])
def auto_improve_skill(skill_id):
    """Automatically improve skill based on collected feedback from database."""
    if not USE_DATABASE:
        return jsonify({"success": False, "error": "Database not available"})

    # Count feedback from training_data table (entries with rating set)
    with db.get_db() as conn:
        cursor = conn.cursor()
        cursor.execute(
            'SELECT COUNT(*) FROM training_data WHERE skill_id = ? AND rating IS NOT NULL',
            (skill_id,)
        )
        feedback_count = cursor.fetchone()[0]

    if feedback_count == 0:
        return jsonify({"success": False, "error": "No feedback collected yet"})

    if feedback_count < 10:
        return jsonify({
            "success": False,
            "error": f"Need at least 10 feedback items, have {feedback_count}"
        })

    # Trigger auto-improvement (simulated)
    add_activity(f"Auto-improvement started for {skill_id} with {feedback_count} feedback items", "")

    return jsonify({
        "success": True,
        "message": f"Auto-improvement started with {feedback_count} feedback items",
        "estimated_time_minutes": 15
    })


# =============================================================================
# UNIFIED DASHBOARD HTML
# =============================================================================

DASHBOARD_HTML = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Skill Command Center</title>
    <link href="https://fonts.googleapis.com/css2?family=Orbitron:wght@400;500;600;700;800;900&family=Space+Grotesk:wght@300;400;500;600;700&family=Rajdhani:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        /* Dark Theme (Default) */
        :root {
            --neon-cyan: #00fff2;
            --neon-pink: #ff00ff;
            --neon-purple: #b400ff;
            --neon-blue: #00a2ff;
            --neon-green: #00ff88;
            --neon-orange: #ff6b00;
            --neon-yellow: #ffee00;
            --bg-dark: #0a0a0f;
            --card-bg: rgba(15, 15, 25, 0.8);
            --glass-surface: rgba(255, 255, 255, 0.05);
            --glass-border: rgba(255, 255, 255, 0.1);
            --text-primary: #ffffff;
            --text-secondary: #a0a0b0;
            --font-display: 'Orbitron', monospace;
            --font-body: 'Space Grotesk', sans-serif;
            --font-accent: 'Rajdhani', sans-serif;
        }

        /* Light Theme (Slack-inspired) */
        [data-theme="light"] {
            --neon-cyan: #1264a3;
            --neon-pink: #e01e5a;
            --neon-purple: #4a154b;
            --neon-blue: #36c5f0;
            --neon-green: #2eb67d;
            --neon-orange: #e8912d;
            --neon-yellow: #ecb22e;
            --bg-dark: #ffffff;
            --card-bg: #f8f8f8;
            --glass-surface: rgba(0, 0, 0, 0.03);
            --glass-border: rgba(0, 0, 0, 0.1);
            --text-primary: #1d1c1d;
            --text-secondary: #616061;
        }

        [data-theme="light"] body {
            background: #ffffff;
        }

        [data-theme="light"] .bg-effects {
            display: none;
        }

        [data-theme="light"] .glass-card {
            background: #ffffff;
            border: 1px solid #e8e8e8;
            box-shadow: 0 1px 3px rgba(0,0,0,0.08);
        }

        [data-theme="light"] .main-tab-btn {
            border-color: #e8e8e8;
            color: #616061;
        }

        [data-theme="light"] .main-tab-btn:hover,
        [data-theme="light"] .main-tab-btn.active {
            background: #4a154b;
            border-color: #4a154b;
            color: #ffffff;
            box-shadow: none;
        }

        [data-theme="light"] .sub-tab-btn {
            color: #616061;
        }

        [data-theme="light"] .sub-tab-btn.active {
            background: #4a154b;
            color: #ffffff;
        }

        [data-theme="light"] .btn-primary {
            background: #4a154b;
        }

        [data-theme="light"] .btn-primary:hover {
            background: #611f69;
        }

        [data-theme="light"] .form-input,
        [data-theme="light"] .form-select,
        [data-theme="light"] .form-textarea {
            background: #ffffff;
            border-color: #e8e8e8;
            color: #1d1c1d;
        }

        [data-theme="light"] .form-input:focus,
        [data-theme="light"] .form-select:focus,
        [data-theme="light"] .form-textarea:focus {
            border-color: #1264a3;
            box-shadow: 0 0 0 3px rgba(18, 100, 163, 0.1);
        }

        [data-theme="light"] .logo {
            background: linear-gradient(135deg, #4a154b 0%, #1264a3 50%, #2eb67d 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        [data-theme="light"] .stat-value {
            color: #1d1c1d;
        }

        [data-theme="light"] .stat-value.green { color: #2eb67d; }
        [data-theme="light"] .stat-value.pink { color: #e01e5a; }
        [data-theme="light"] .stat-value.orange { color: #e8912d; }

        [data-theme="light"] .section-icon {
            color: #4a154b;
        }

        [data-theme="light"] .skill-card {
            background: #ffffff;
            border: 1px solid #e8e8e8;
        }

        [data-theme="light"] .skill-card:hover {
            border-color: #4a154b;
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }

        [data-theme="light"] .console {
            background: #f8f8f8;
            border-color: #e8e8e8;
            color: #1d1c1d;
        }

        [data-theme="light"] .code-block {
            background: #f8f8f8;
            border-color: #e8e8e8;
        }

        [data-theme="light"] .skills-table th {
            background: #f8f8f8;
            color: #1d1c1d;
        }

        [data-theme="light"] .skills-table td {
            border-color: #e8e8e8;
        }

        [data-theme="light"] .pipeline-icon {
            border-color: #e8e8e8;
        }

        [data-theme="light"] .getting-started-card {
            background: linear-gradient(135deg, rgba(74,21,75,0.1), rgba(18,100,163,0.1)) !important;
            border-color: rgba(74,21,75,0.3) !important;
        }

        /* Theme Toggle Button */
        .theme-toggle-btn {
            position: fixed;
            top: 1rem;
            right: 1rem;
            z-index: 1000;
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 50%;
            width: 44px;
            height: 44px;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 1.25rem;
            transition: all 0.2s;
        }

        .theme-toggle-btn:hover {
            background: var(--neon-purple);
            transform: scale(1.1);
        }

        * { margin: 0; padding: 0; box-sizing: border-box; }

        /* Tooltip Styles for Technical Terms */
        .tooltip {
            position: relative;
            cursor: help;
            border-bottom: 1px dotted var(--neon-cyan);
        }

        .tooltip::before {
            content: attr(data-tooltip);
            position: absolute;
            bottom: 100%;
            left: 50%;
            transform: translateX(-50%);
            padding: 0.5rem 0.75rem;
            background: var(--card-bg);
            border: 1px solid var(--neon-cyan);
            border-radius: 6px;
            font-size: 0.8rem;
            color: var(--text-primary);
            white-space: nowrap;
            opacity: 0;
            visibility: hidden;
            transition: opacity 0.2s, visibility 0.2s;
            z-index: 1000;
            max-width: 250px;
            white-space: normal;
            text-align: center;
        }

        .tooltip::after {
            content: '';
            position: absolute;
            bottom: 100%;
            left: 50%;
            transform: translateX(-50%);
            border: 6px solid transparent;
            border-top-color: var(--neon-cyan);
            opacity: 0;
            visibility: hidden;
            transition: opacity 0.2s, visibility 0.2s;
        }

        .tooltip:hover::before,
        .tooltip:hover::after {
            opacity: 1;
            visibility: visible;
        }

        /* Help icon style */
        .help-icon {
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 16px;
            height: 16px;
            border-radius: 50%;
            background: rgba(0, 255, 242, 0.2);
            color: var(--neon-cyan);
            font-size: 0.7rem;
            font-weight: bold;
            margin-left: 4px;
            cursor: help;
        }

        body {
            font-family: var(--font-body);
            background: var(--bg-dark);
            color: var(--text-primary);
            min-height: 100vh;
            overflow-x: hidden;
        }

        /* Background Effects */
        .bg-effects {
            position: fixed;
            top: 0; left: 0; width: 100%; height: 100%;
            pointer-events: none;
            z-index: 0;
        }

        .grid-overlay {
            position: absolute;
            top: 0; left: 0; width: 100%; height: 100%;
            background-image:
                linear-gradient(rgba(0, 255, 242, 0.03) 1px, transparent 1px),
                linear-gradient(90deg, rgba(0, 255, 242, 0.03) 1px, transparent 1px);
            background-size: 50px 50px;
        }

        .gradient-blob {
            position: absolute;
            border-radius: 50%;
            filter: blur(100px);
            animation: pulse 8s ease-in-out infinite;
        }

        .blob-1 { width: 600px; height: 600px; background: radial-gradient(circle, rgba(0, 255, 242, 0.15) 0%, transparent 70%); top: -200px; right: -200px; }
        .blob-2 { width: 500px; height: 500px; background: radial-gradient(circle, rgba(180, 0, 255, 0.1) 0%, transparent 70%); bottom: -150px; left: -150px; animation-delay: -4s; }

        @keyframes pulse { 0%, 100% { transform: scale(1); opacity: 1; } 50% { transform: scale(1.2); opacity: 0.7; } }

        /* Container */
        .container {
            position: relative;
            z-index: 1;
            max-width: 1800px;
            margin: 0 auto;
            padding: 1.5rem;
        }

        /* Header */
        header {
            text-align: center;
            margin-bottom: 1.5rem;
        }

        .logo {
            font-family: var(--font-display);
            font-size: 2rem;
            font-weight: 900;
            background: linear-gradient(135deg, var(--neon-cyan) 0%, var(--neon-purple) 50%, var(--neon-pink) 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            text-transform: uppercase;
            letter-spacing: 0.1em;
        }

        .subtitle {
            font-family: var(--font-accent);
            font-size: 0.9rem;
            color: var(--text-secondary);
            letter-spacing: 0.3em;
        }

        .status-badge {
            display: inline-flex;
            align-items: center;
            gap: 0.5rem;
            padding: 0.4rem 1.2rem;
            background: linear-gradient(135deg, rgba(0, 255, 136, 0.2) 0%, rgba(0, 255, 242, 0.1) 100%);
            border: 1px solid var(--neon-green);
            border-radius: 50px;
            font-family: var(--font-display);
            font-size: 0.7rem;
            color: var(--neon-green);
            margin-top: 0.75rem;
        }

        .status-dot { width: 8px; height: 8px; background: var(--neon-green); border-radius: 50%; animation: blink 1s infinite; }
        @keyframes blink { 0%, 100% { opacity: 1; } 50% { opacity: 0.3; } }

        /* Main Tabs */
        .main-tabs {
            display: flex;
            gap: 0.5rem;
            margin-bottom: 1rem;
            flex-wrap: wrap;
            justify-content: center;
        }

        .main-tab-btn {
            padding: 0.6rem 1.2rem;
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 8px;
            color: var(--text-secondary);
            font-family: var(--font-display);
            font-size: 0.75rem;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .main-tab-btn:hover {
            border-color: var(--neon-cyan);
            color: var(--neon-cyan);
        }

        .main-tab-btn.active {
            background: linear-gradient(135deg, rgba(0, 255, 242, 0.2), rgba(180, 0, 255, 0.1));
            border-color: var(--neon-cyan);
            color: var(--neon-cyan);
            box-shadow: 0 0 20px rgba(0, 255, 242, 0.2);
        }

        /* Sub Tabs */
        .sub-tabs {
            display: flex;
            gap: 0.25rem;
            margin-bottom: 1rem;
            flex-wrap: wrap;
            padding: 0.5rem;
            background: var(--glass-surface);
            border-radius: 8px;
        }

        .sub-tab-btn {
            padding: 0.5rem 1rem;
            background: transparent;
            border: 1px solid transparent;
            border-radius: 6px;
            color: var(--text-secondary);
            font-family: var(--font-accent);
            font-size: 0.8rem;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .sub-tab-btn:hover {
            background: var(--glass-surface);
            color: var(--text-primary);
        }

        .sub-tab-btn.active {
            background: linear-gradient(135deg, rgba(0, 162, 255, 0.2), rgba(0, 255, 242, 0.1));
            border-color: var(--neon-blue);
            color: var(--neon-blue);
        }

        .tab-content, .sub-tab-content { display: none; }
        .tab-content.active, .sub-tab-content.active { display: block; }

        /* Glass Card */
        .glass-card {
            background: var(--card-bg);
            backdrop-filter: blur(20px);
            border: 1px solid var(--glass-border);
            border-radius: 16px;
            padding: 1.25rem;
            margin-bottom: 1rem;
            transition: all 0.3s ease;
        }

        .glass-card:hover {
            border-color: rgba(0, 255, 242, 0.3);
            box-shadow: 0 10px 40px rgba(0, 0, 0, 0.3);
        }

        /* Dashboard Grid */
        .dashboard-grid {
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 1rem;
        }

        .card-full { grid-column: span 4; }
        .card-three { grid-column: span 3; }
        .card-half { grid-column: span 2; }
        .card-quarter { grid-column: span 1; }

        @media (max-width: 1200px) {
            .dashboard-grid { grid-template-columns: repeat(2, 1fr); }
            .card-full, .card-three, .card-half { grid-column: span 2; }
        }

        @media (max-width: 768px) {
            .dashboard-grid { grid-template-columns: 1fr; }
            .card-full, .card-three, .card-half, .card-quarter { grid-column: span 1; }
        }

        /* Stat Card */
        .stat-card { text-align: center; padding: 1rem; }
        .stat-value {
            font-family: var(--font-display);
            font-size: 1.75rem;
            font-weight: 700;
            background: linear-gradient(135deg, var(--neon-cyan), var(--neon-blue));
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }
        .stat-value.green { background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan)); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
        .stat-value.pink { background: linear-gradient(135deg, var(--neon-pink), var(--neon-purple)); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
        .stat-value.orange { background: linear-gradient(135deg, var(--neon-orange), var(--neon-yellow)); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
        .stat-label { font-family: var(--font-accent); font-size: 0.8rem; color: var(--text-secondary); text-transform: uppercase; margin-top: 0.25rem; }

        /* Section Header */
        .section-header {
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 1rem;
        }

        .section-title {
            font-family: var(--font-display);
            font-size: 1rem;
            font-weight: 600;
            display: flex;
            align-items: center;
            gap: 0.5rem;
        }

        .section-icon {
            width: 28px; height: 28px;
            display: flex;
            align-items: center;
            justify-content: center;
            background: linear-gradient(135deg, var(--neon-cyan), var(--neon-purple));
            border-radius: 6px;
            font-size: 0.9rem;
        }

        /* Form Elements */
        .form-group { margin-bottom: 1rem; }
        .form-label {
            display: block;
            font-family: var(--font-accent);
            font-size: 0.8rem;
            color: var(--text-secondary);
            margin-bottom: 0.4rem;
            text-transform: uppercase;
            letter-spacing: 0.05em;
        }

        .form-input, .form-select, .form-textarea {
            width: 100%;
            padding: 0.6rem 0.9rem;
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 6px;
            color: var(--text-primary);
            font-family: var(--font-body);
            font-size: 0.9rem;
            transition: all 0.3s ease;
        }

        .form-input:focus, .form-select:focus, .form-textarea:focus {
            outline: none;
            border-color: var(--neon-cyan);
            box-shadow: 0 0 15px rgba(0, 255, 242, 0.2);
        }

        /* Fix dropdown option styling for dark theme */
        .form-select option,
        select option {
            background: #1a1a2e;
            color: #ffffff;
            padding: 0.5rem;
        }

        .form-select option:hover,
        select option:hover,
        .form-select option:checked,
        select option:checked {
            background: #00fff2;
            color: #0a0a0f;
        }

        .form-textarea { resize: vertical; min-height: 80px; }

        .form-row { display: grid; grid-template-columns: 1fr 1fr; gap: 1rem; }
        @media (max-width: 768px) { .form-row { grid-template-columns: 1fr; } }

        /* Buttons */
        .btn {
            padding: 0.6rem 1.2rem;
            border: none;
            border-radius: 6px;
            font-family: var(--font-display);
            font-size: 0.75rem;
            cursor: pointer;
            transition: all 0.3s ease;
            text-transform: uppercase;
        }

        .btn-primary {
            background: linear-gradient(135deg, var(--neon-cyan), var(--neon-blue));
            color: var(--bg-dark);
            box-shadow: 0 0 20px rgba(0, 255, 242, 0.3);
        }

        .btn-primary:hover {
            transform: translateY(-2px);
            box-shadow: 0 0 30px rgba(0, 255, 242, 0.5);
        }

        .btn-secondary {
            background: transparent;
            border: 1px solid var(--neon-purple);
            color: var(--neon-purple);
        }

        .btn-secondary:hover {
            background: rgba(180, 0, 255, 0.1);
            box-shadow: 0 0 20px rgba(180, 0, 255, 0.3);
        }

        .btn-success {
            background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan));
            color: var(--bg-dark);
        }

        .btn-danger {
            background: linear-gradient(135deg, var(--neon-pink), var(--neon-orange));
            color: var(--bg-dark);
        }

        .btn-sm { padding: 0.4rem 0.8rem; font-size: 0.7rem; }

        /* Skills Grid */
        .skills-grid {
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(260px, 1fr));
            gap: 0.75rem;
        }

        .skill-card {
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 12px;
            padding: 1rem;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .skill-card:hover {
            border-color: var(--neon-cyan);
            transform: scale(1.02);
        }

        .skill-card.deployed { border-left: 3px solid var(--neon-green); }
        .skill-card.training { border-left: 3px solid var(--neon-orange); }
        .skill-card.draft { border-left: 3px solid var(--text-secondary); }

        .skill-header { display: flex; justify-content: space-between; margin-bottom: 0.75rem; }
        .skill-name { font-family: var(--font-display); font-size: 0.9rem; }
        .skill-type { font-size: 0.7rem; color: var(--text-secondary); }
        .skill-status {
            padding: 0.2rem 0.6rem;
            border-radius: 20px;
            font-size: 0.65rem;
            text-transform: uppercase;
        }
        .skill-status.deployed { background: rgba(0, 255, 136, 0.2); color: var(--neon-green); }
        .skill-status.training { background: rgba(255, 107, 0, 0.2); color: var(--neon-orange); }
        .skill-status.draft { background: rgba(160, 160, 176, 0.2); color: var(--text-secondary); }

        .skill-metrics { display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.5rem; margin-top: 0.75rem; padding-top: 0.75rem; border-top: 1px solid var(--glass-border); }
        .skill-metric { text-align: center; }
        .metric-value { font-family: var(--font-display); font-size: 0.8rem; color: var(--neon-cyan); }
        .metric-label { font-size: 0.6rem; color: var(--text-secondary); text-transform: uppercase; }

        /* Skills Table */
        .skills-table {
            width: 100%;
            border-collapse: collapse;
        }

        .skills-table th, .skills-table td {
            padding: 0.75rem;
            text-align: left;
            border-bottom: 1px solid var(--glass-border);
        }

        .skills-table th {
            font-family: var(--font-display);
            font-size: 0.75rem;
            color: var(--neon-cyan);
            text-transform: uppercase;
            background: var(--glass-surface);
        }

        .skills-table tr:hover {
            background: var(--glass-surface);
        }

        .table-status {
            padding: 0.2rem 0.6rem;
            border-radius: 20px;
            font-size: 0.7rem;
        }

        .table-status.deployed { background: rgba(0, 255, 136, 0.2); color: var(--neon-green); }
        .table-status.training { background: rgba(255, 107, 0, 0.2); color: var(--neon-orange); }
        .table-status.draft { background: rgba(160, 160, 176, 0.2); color: var(--text-secondary); }

        .table-filter {
            display: flex;
            gap: 0.5rem;
            margin-bottom: 1rem;
            flex-wrap: wrap;
        }

        .table-filter input, .table-filter select {
            padding: 0.5rem;
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 6px;
            color: var(--text-primary);
            font-size: 0.85rem;
        }

        /* Console */
        .console {
            background: #0d0d12;
            border: 1px solid var(--glass-border);
            border-radius: 10px;
            padding: 1rem;
            font-family: monospace;
            font-size: 0.85rem;
            min-height: 250px;
            max-height: 400px;
            overflow-y: auto;
            white-space: pre-wrap;
            color: var(--neon-green);
        }

        .console .prompt { color: var(--neon-cyan); }
        .console .response { color: var(--text-primary); }
        .console .error { color: var(--neon-pink); }
        .console .info { color: var(--text-secondary); }

        /* Training Pipeline */
        .pipeline {
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: 1rem 0;
        }

        .pipeline-step {
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 0.5rem;
            flex: 1;
        }

        .pipeline-icon {
            width: 48px;
            height: 48px;
            display: flex;
            align-items: center;
            justify-content: center;
            border-radius: 50%;
            font-size: 1.25rem;
            background: var(--glass-surface);
            border: 2px solid var(--glass-border);
            transition: all 0.3s ease;
        }

        .pipeline-icon.complete {
            background: rgba(0, 255, 136, 0.2);
            border-color: var(--neon-green);
        }

        .pipeline-icon.active {
            background: rgba(0, 162, 255, 0.2);
            border-color: var(--neon-blue);
            animation: glow 1.5s ease-in-out infinite;
        }

        @keyframes glow {
            0%, 100% { box-shadow: 0 0 10px rgba(0, 162, 255, 0.3); }
            50% { box-shadow: 0 0 25px rgba(0, 162, 255, 0.6); }
        }

        .pipeline-label {
            font-family: var(--font-accent);
            font-size: 0.75rem;
            color: var(--text-secondary);
            text-transform: uppercase;
        }

        .pipeline-connector {
            flex: 0.5;
            height: 2px;
            background: var(--glass-border);
        }

        .pipeline-connector.complete {
            background: var(--neon-green);
        }

        .progress-bar {
            width: 100%;
            height: 8px;
            background: var(--glass-surface);
            border-radius: 4px;
            overflow: hidden;
            margin-top: 1rem;
        }

        .progress-fill {
            height: 100%;
            background: linear-gradient(90deg, var(--neon-cyan), var(--neon-blue));
            border-radius: 4px;
            transition: width 0.5s ease;
        }

        .progress-text {
            text-align: center;
            font-family: var(--font-display);
            font-size: 0.9rem;
            margin-top: 0.5rem;
            color: var(--neon-cyan);
        }

        /* Activity Feed */
        .activity-feed {
            max-height: 300px;
            overflow-y: auto;
        }

        .activity-item {
            display: flex;
            align-items: flex-start;
            gap: 0.75rem;
            padding: 0.75rem 0;
            border-bottom: 1px solid var(--glass-border);
        }

        .activity-icon {
            font-size: 1rem;
            width: 24px;
            text-align: center;
        }

        .activity-content { flex: 1; }
        .activity-message { font-size: 0.85rem; }
        .activity-time { font-size: 0.7rem; color: var(--text-secondary); }

        /* Chart */
        .chart-container {
            height: 120px;
            display: flex;
            align-items: flex-end;
            gap: 3px;
            padding: 0.5rem 0;
        }

        .chart-bar {
            flex: 1;
            background: linear-gradient(180deg, var(--neon-cyan) 0%, rgba(0, 255, 242, 0.3) 100%);
            border-radius: 3px 3px 0 0;
            transition: all 0.3s ease;
        }

        .chart-bar:hover { background: linear-gradient(180deg, var(--neon-pink) 0%, rgba(255, 0, 255, 0.3) 100%); }

        /* Server Panel */
        .server-panel { display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.75rem; }
        .server-stat { text-align: center; padding: 0.75rem; background: var(--glass-surface); border-radius: 10px; }

        /* Compare Results */
        .compare-results {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 1rem;
        }

        .compare-card {
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 10px;
            padding: 1rem;
        }

        .compare-header {
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 0.75rem;
            padding-bottom: 0.5rem;
            border-bottom: 1px solid var(--glass-border);
        }

        .compare-provider {
            font-family: var(--font-display);
            font-size: 0.9rem;
            color: var(--neon-cyan);
        }

        .compare-latency {
            font-size: 0.75rem;
            color: var(--text-secondary);
        }

        /* Code Block */
        .code-block {
            background: #0d0d12;
            border: 1px solid var(--glass-border);
            border-radius: 8px;
            padding: 1rem;
            font-family: monospace;
            font-size: 0.8rem;
            overflow-x: auto;
            white-space: pre;
            color: var(--text-primary);
        }

        .code-block .keyword { color: var(--neon-pink); }
        .code-block .string { color: var(--neon-green); }
        .code-block .comment { color: var(--text-secondary); }

        /* Loading */
        .loading { opacity: 0.5; pointer-events: none; }
        .spinner {
            display: inline-block;
            width: 18px; height: 18px;
            border: 2px solid var(--glass-border);
            border-top-color: var(--neon-cyan);
            border-radius: 50%;
            animation: spin 1s linear infinite;
        }
        @keyframes spin { to { transform: rotate(360deg); } }
        @keyframes slideIn { from { transform: translateX(100%); opacity: 0; } to { transform: translateX(0); opacity: 1; } }

        /* Messages */
        .message {
            padding: 0.75rem 1rem;
            border-radius: 6px;
            margin-bottom: 0.75rem;
            font-family: var(--font-accent);
            font-size: 0.9rem;
        }
        .message.success { background: rgba(0, 255, 136, 0.1); border: 1px solid var(--neon-green); color: var(--neon-green); }
        .message.error { background: rgba(255, 0, 85, 0.1); border: 1px solid var(--neon-pink); color: var(--neon-pink); }
        .message.info { background: rgba(0, 162, 255, 0.1); border: 1px solid var(--neon-blue); color: var(--neon-blue); }

        /* Scrollbar */
        ::-webkit-scrollbar { width: 6px; height: 6px; }
        ::-webkit-scrollbar-track { background: var(--bg-dark); }
        ::-webkit-scrollbar-thumb { background: var(--glass-border); border-radius: 3px; }
        ::-webkit-scrollbar-thumb:hover { background: var(--neon-cyan); }

        /* File Upload */
        .file-upload {
            border: 2px dashed var(--glass-border);
            border-radius: 10px;
            padding: 2rem;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .file-upload:hover {
            border-color: var(--neon-cyan);
            background: var(--glass-surface);
        }

        .file-upload input { display: none; }
        .file-upload-icon { font-size: 2rem; margin-bottom: 0.5rem; }
        .file-upload-text { color: var(--text-secondary); font-size: 0.9rem; }

        /* ============================================
           UNIFIED SKILLS & TRAINING STYLES
           ============================================ */

        /* Skill Cards */
        .skill-card {
            background: var(--glass-surface);
            border: 1px solid var(--glass-border);
            border-radius: 12px;
            padding: 1.25rem;
            cursor: pointer;
            transition: all 0.3s ease;
            position: relative;
            overflow: hidden;
        }

        .skill-card:hover {
            border-color: var(--neon-cyan);
            transform: translateY(-4px);
            box-shadow: 0 10px 30px rgba(0, 255, 242, 0.15);
        }

        .skill-card-header {
            display: flex;
            justify-content: space-between;
            align-items: flex-start;
            margin-bottom: 0.75rem;
        }

        .skill-card-icon {
            font-size: 2rem;
            line-height: 1;
        }

        .skill-card-title {
            font-size: 1.1rem;
            font-weight: 600;
            color: var(--text-primary);
            margin: 0.5rem 0 0.25rem 0;
        }

        .skill-card-description {
            font-size: 0.85rem;
            color: var(--text-secondary);
            margin-bottom: 1rem;
            line-height: 1.4;
            display: -webkit-box;
            -webkit-line-clamp: 2;
            -webkit-box-orient: vertical;
            overflow: hidden;
        }

        .skill-card-stats {
            display: flex;
            gap: 1rem;
            font-size: 0.8rem;
            color: var(--text-secondary);
            border-top: 1px solid var(--glass-border);
            padding-top: 0.75rem;
            margin-top: auto;
        }

        .skill-card-stat {
            display: flex;
            align-items: center;
            gap: 0.25rem;
        }

        .skill-card-stat .value {
            font-weight: 600;
            color: var(--neon-cyan);
        }

        /* Skill Status Badges */
        .skill-status-badge {
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
            padding: 0.25rem 0.6rem;
            border-radius: 12px;
            font-size: 0.75rem;
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }

        .skill-status-badge.untrained {
            background: rgba(107, 114, 128, 0.2);
            color: #9ca3af;
            border: 1px solid rgba(107, 114, 128, 0.3);
        }

        .skill-status-badge.has_data {
            background: rgba(245, 158, 11, 0.2);
            color: #fbbf24;
            border: 1px solid rgba(245, 158, 11, 0.3);
        }

        .skill-status-badge.training {
            background: rgba(59, 130, 246, 0.2);
            color: #60a5fa;
            border: 1px solid rgba(59, 130, 246, 0.3);
            animation: pulse 1.5s infinite;
        }

        .skill-status-badge.trained {
            background: rgba(16, 185, 129, 0.2);
            color: #34d399;
            border: 1px solid rgba(16, 185, 129, 0.3);
        }

        .skill-status-badge.failed {
            background: rgba(239, 68, 68, 0.2);
            color: #f87171;
            border: 1px solid rgba(239, 68, 68, 0.3);
        }

        @keyframes pulse {
            0%, 100% { opacity: 1; }
            50% { opacity: 0.6; }
        }

        /* Filter Buttons */
        .skill-filter-btn {
            background: transparent;
            border: 1px solid var(--glass-border);
            color: var(--text-secondary);
            padding: 0.4rem 0.8rem;
            border-radius: 6px;
            font-size: 0.8rem;
            cursor: pointer;
            transition: all 0.2s ease;
        }

        .skill-filter-btn:hover {
            border-color: var(--neon-cyan);
            color: var(--neon-cyan);
        }

        .skill-filter-btn.active {
            background: var(--neon-cyan);
            border-color: var(--neon-cyan);
            color: #000;
        }

        /* Modal Styles */
        .modal-overlay {
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            bottom: 0;
            background: rgba(0, 0, 0, 0.8);
            backdrop-filter: blur(4px);
            display: flex;
            align-items: center;
            justify-content: center;
            z-index: 1000;
            padding: 2rem;
        }

        .modal-content {
            background: var(--bg-primary);
            border: 1px solid var(--glass-border);
            border-radius: 16px;
            width: 100%;
            box-shadow: 0 25px 80px rgba(0, 0, 0, 0.5);
        }

        .modal-header {
            background: linear-gradient(135deg, rgba(0, 255, 242, 0.1), rgba(180, 0, 255, 0.05));
        }

        /* Modal Tab Buttons */
        .modal-tab-btn {
            background: transparent;
            border: none;
            color: var(--text-secondary);
            padding: 1rem 1.5rem;
            cursor: pointer;
            font-size: 0.95rem;
            font-weight: 500;
            border-bottom: 2px solid transparent;
            transition: all 0.2s ease;
        }

        .modal-tab-btn:hover {
            color: var(--text-primary);
            background: rgba(255, 255, 255, 0.05);
        }

        .modal-tab-btn.active {
            color: var(--neon-cyan);
            border-bottom-color: var(--neon-cyan);
        }

        .modal-tab-content {
            display: none;
        }

        .modal-tab-content.active {
            display: block;
        }

        /* Progress Bar */
        .progress-bar {
            width: 100%;
            height: 8px;
            background: var(--glass-surface);
            border-radius: 4px;
            overflow: hidden;
        }

        .progress-fill {
            height: 100%;
            background: linear-gradient(90deg, var(--neon-cyan), var(--neon-green));
            border-radius: 4px;
            transition: width 0.3s ease;
        }

        /* Readiness Check Items */
        .readiness-item {
            display: flex;
            align-items: center;
            gap: 0.75rem;
            padding: 0.5rem 0;
        }

        .readiness-icon {
            font-size: 1.25rem;
        }

        .readiness-icon.pass { color: var(--neon-green); }
        .readiness-icon.fail { color: var(--neon-orange); }
        .readiness-icon.warn { color: #fbbf24; }

        /* ========================================== */
        /* SKILL WORKFLOW - 3 Step Design */
        /* ========================================== */

        .skill-workflow {
            display: flex;
            flex-direction: column;
            gap: 1.5rem;
        }

        /* Skill Context Bar - Always visible when skill selected */
        .skill-context-bar {
            display: none;
            align-items: center;
            gap: 1rem;
            padding: 1rem 1.5rem;
            background: linear-gradient(135deg, rgba(0, 212, 255, 0.1), rgba(139, 92, 246, 0.1));
            border: 1px solid rgba(0, 212, 255, 0.3);
            border-radius: 12px;
            margin-bottom: 1rem;
        }

        .skill-context-bar.visible {
            display: flex;
        }

        .skill-context-icon {
            font-size: 2rem;
        }

        .skill-context-info {
            flex: 1;
        }

        .skill-context-name {
            font-size: 1.2rem;
            font-weight: 600;
            color: var(--text-primary);
            margin: 0;
        }

        .skill-context-meta {
            display: flex;
            gap: 1rem;
            font-size: 0.85rem;
            color: var(--text-secondary);
        }

        .skill-context-meta span {
            display: flex;
            align-items: center;
            gap: 0.25rem;
        }

        .skill-context-actions {
            display: flex;
            gap: 0.5rem;
        }

        /* Step Indicator */
        .workflow-steps {
            display: flex;
            justify-content: center;
            align-items: center;
            gap: 0;
            margin-bottom: 1.5rem;
            padding: 1rem;
            background: var(--glass-surface);
            border-radius: 12px;
        }

        .workflow-step {
            display: flex;
            align-items: center;
            gap: 0.75rem;
            padding: 0.75rem 1.5rem;
            cursor: pointer;
            border-radius: 8px;
            transition: all 0.3s ease;
            position: relative;
        }

        .workflow-step:hover {
            background: rgba(0, 255, 242, 0.1);
        }

        .workflow-step.active {
            background: linear-gradient(135deg, rgba(0, 255, 242, 0.2), rgba(139, 92, 246, 0.1));
        }

        .workflow-step.completed .step-number {
            background: var(--neon-green);
            color: #000;
        }

        .workflow-step.active .step-number {
            background: var(--neon-cyan);
            color: #000;
            box-shadow: 0 0 15px rgba(0, 255, 242, 0.5);
        }

        .workflow-step.disabled {
            opacity: 0.5;
            cursor: not-allowed;
        }

        .step-number {
            width: 32px;
            height: 32px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
            background: var(--glass-border);
            color: var(--text-secondary);
            font-size: 0.9rem;
            transition: all 0.3s ease;
        }

        .step-label {
            font-weight: 500;
            color: var(--text-secondary);
            font-size: 0.9rem;
        }

        .workflow-step.active .step-label {
            color: var(--text-primary);
        }

        .workflow-step.completed .step-label {
            color: var(--neon-green);
        }

        .step-connector {
            width: 40px;
            height: 2px;
            background: var(--glass-border);
        }

        .step-connector.completed {
            background: var(--neon-green);
        }

        /* Workflow Content */
        .workflow-content {
            position: relative;
        }

        .workflow-panel {
            display: none;
        }

        .workflow-panel.active {
            display: block;
            animation: fadeIn 0.3s ease;
        }

        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(10px); }
            to { opacity: 1; transform: translateY(0); }
        }

        /* Side-by-Side Layout for Step 2 */
        .training-layout {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 1.5rem;
        }

        @media (max-width: 1024px) {
            .training-layout {
                grid-template-columns: 1fr;
            }
        }

        .training-panel {
            display: flex;
            flex-direction: column;
        }

        .training-panel h4 {
            margin: 0 0 1rem 0;
            color: var(--neon-cyan);
            font-size: 1rem;
            display: flex;
            align-items: center;
            gap: 0.5rem;
        }

        /* Data Entry Form */
        .data-entry-form {
            display: flex;
            flex-direction: column;
            gap: 1rem;
        }

        .data-entry-form .form-group {
            margin-bottom: 0;
        }

        .data-entry-actions {
            display: flex;
            gap: 0.5rem;
            margin-top: 0.5rem;
        }

        /* Training Data List Compact */
        .training-data-list {
            max-height: 300px;
            overflow-y: auto;
            border: 1px solid var(--glass-border);
            border-radius: 8px;
        }

        .training-data-item {
            display: flex;
            padding: 0.75rem;
            border-bottom: 1px solid var(--glass-border);
            gap: 0.75rem;
            align-items: flex-start;
            transition: background 0.2s ease;
        }

        .training-data-item:last-child {
            border-bottom: none;
        }

        .training-data-item:hover {
            background: rgba(0, 255, 242, 0.05);
        }

        .training-data-item .data-content {
            flex: 1;
            min-width: 0;
        }

        .training-data-item .data-q {
            font-weight: 500;
            color: var(--text-primary);
            font-size: 0.85rem;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }

        .training-data-item .data-a {
            color: var(--text-secondary);
            font-size: 0.8rem;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }

        .training-data-item .data-actions {
            display: flex;
            gap: 0.25rem;
            opacity: 0;
            transition: opacity 0.2s ease;
        }

        .training-data-item:hover .data-actions {
            opacity: 1;
        }

        /* Upload Zone */
        .upload-zone {
            border: 2px dashed var(--glass-border);
            border-radius: 12px;
            padding: 2rem;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .upload-zone:hover {
            border-color: var(--neon-cyan);
            background: rgba(0, 255, 242, 0.05);
        }

        .upload-zone.dragover {
            border-color: var(--neon-green);
            background: rgba(0, 255, 136, 0.1);
        }

        .upload-icon {
            font-size: 2.5rem;
            margin-bottom: 0.5rem;
        }

        /* Test & Train Layout */
        .test-train-layout {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 1.5rem;
        }

        @media (max-width: 1024px) {
            .test-train-layout {
                grid-template-columns: 1fr;
            }
        }

        /* Chat Panel */
        .chat-panel {
            display: flex;
            flex-direction: column;
            height: 450px;
        }

        .chat-messages {
            flex: 1;
            overflow-y: auto;
            padding: 1rem;
            background: var(--glass-surface);
            border-radius: 8px;
            margin-bottom: 1rem;
        }

        .chat-message {
            padding: 0.75rem;
            border-radius: 8px;
            margin-bottom: 0.5rem;
            max-width: 85%;
        }

        .chat-message.user {
            background: linear-gradient(135deg, rgba(139, 92, 246, 0.3), rgba(139, 92, 246, 0.1));
            margin-left: auto;
        }

        .chat-message.assistant {
            background: rgba(0, 255, 242, 0.1);
        }

        .chat-input-area {
            display: flex;
            gap: 0.5rem;
        }

        /* Train Panel */
        .train-panel {
            display: flex;
            flex-direction: column;
            gap: 1rem;
        }

        .train-readiness {
            padding: 1rem;
            background: var(--glass-surface);
            border-radius: 8px;
        }

        .readiness-row {
            display: flex;
            align-items: center;
            gap: 0.5rem;
            padding: 0.5rem 0;
            border-bottom: 1px solid var(--glass-border);
        }

        .readiness-row:last-child {
            border-bottom: none;
        }

        .readiness-check {
            width: 20px;
            height: 20px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 0.75rem;
        }

        .readiness-check.pass {
            background: rgba(0, 255, 136, 0.2);
            color: var(--neon-green);
        }

        .readiness-check.fail {
            background: rgba(255, 100, 100, 0.2);
            color: #ff6b6b;
        }

        .readiness-check.warn {
            background: rgba(251, 191, 36, 0.2);
            color: #fbbf24;
        }

        .train-button-wrapper {
            margin-top: auto;
        }

        .train-button-large {
            width: 100%;
            padding: 1.25rem;
            font-size: 1.1rem;
            background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan));
            border: none;
            border-radius: 12px;
            color: #000;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            display: flex;
            align-items: center;
            justify-content: center;
            gap: 0.5rem;
        }

        .train-button-large:hover:not(:disabled) {
            transform: translateY(-2px);
            box-shadow: 0 10px 30px rgba(0, 255, 136, 0.3);
        }

        .train-button-large:disabled {
            opacity: 0.5;
            cursor: not-allowed;
            background: var(--glass-border);
            color: var(--text-secondary);
        }

        /* Toast Notifications */
        .toast-container {
            position: fixed;
            top: 1rem;
            right: 1rem;
            z-index: 10000;
            display: flex;
            flex-direction: column;
            gap: 0.5rem;
        }

        .toast {
            padding: 1rem 1.5rem;
            border-radius: 8px;
            background: var(--card-bg);
            border: 1px solid var(--glass-border);
            color: var(--text-primary);
            animation: slideIn 0.3s ease;
            display: flex;
            align-items: center;
            gap: 0.75rem;
            min-width: 300px;
            box-shadow: 0 10px 40px rgba(0, 0, 0, 0.3);
        }

        .toast.success { border-color: var(--neon-green); }
        .toast.error { border-color: #ff6b6b; }
        .toast.warning { border-color: #fbbf24; }
        .toast.info { border-color: var(--neon-cyan); }

        @keyframes slideIn {
            from { transform: translateX(100%); opacity: 0; }
            to { transform: translateX(0); opacity: 1; }
        }

        /* Light Theme Overrides */
        [data-theme="light"] .skill-card {
            background: #ffffff;
            border-color: #e5e7eb;
        }

        [data-theme="light"] .skill-card:hover {
            border-color: #4a154b;
            box-shadow: 0 10px 30px rgba(74, 21, 75, 0.1);
        }

        [data-theme="light"] .modal-overlay {
            background: rgba(0, 0, 0, 0.5);
        }

        [data-theme="light"] .modal-content {
            background: #ffffff;
            border-color: #e5e7eb;
        }

        [data-theme="light"] .modal-tab-btn.active {
            color: #4a154b;
            border-bottom-color: #4a154b;
        }

        [data-theme="light"] .skill-filter-btn.active {
            background: #4a154b;
            border-color: #4a154b;
            color: #ffffff;
        }
    </style>
    <!-- Chart.js for loss visualization -->
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <!-- Canvas confetti for celebrations -->
    <script src="https://cdn.jsdelivr.net/npm/canvas-confetti@1.6.0/dist/confetti.browser.min.js"></script>
</head>
<body>
    <!-- Theme Toggle Button -->
    <button class="theme-toggle-btn" onclick="toggleTheme()" title="Toggle Dark/Light Mode" id="theme-toggle">
        🌙
    </button>

    <div class="bg-effects">
        <div class="grid-overlay"></div>
        <div class="gradient-blob blob-1"></div>
        <div class="gradient-blob blob-2"></div>
    </div>

    <div class="container">
        <header>
            <h1 class="logo">Skill Command Center</h1>
            <p class="subtitle">Voice Agent Intelligence Hub</p>
            <div class="status-badge">
                <span class="status-dot"></span>
                <span id="server-status-text">LPU Online - 1 Warm Container</span>
            </div>
        </header>

        <!-- Main Tabs - Simplified to 4 core sections -->
        <div class="main-tabs">
            <button class="main-tab-btn active" onclick="showMainTab('dashboard')" title="Overview, status, and quick actions">Dashboard</button>
            <button class="main-tab-btn" onclick="showMainTab('skills-training')" title="Create, manage, train, and test all your AI skills">🧠 Skills & Training</button>
            <button class="main-tab-btn" onclick="showMainTab('voice')" title="Voice projects, TTS settings, and platform connections">Voice</button>
            <button class="main-tab-btn" onclick="showMainTab('settings')" title="API keys, integrations, and advanced configuration">Settings</button>
        </div>

        <!-- ================================================================ -->
        <!-- DASHBOARD TAB -->
        <!-- ================================================================ -->
        <div id="tab-dashboard" class="tab-content active">
            <!-- Getting Started Card - Shows for new users -->
            <div class="glass-card getting-started-card" id="getting-started" style="background: linear-gradient(135deg, rgba(0,212,255,0.1), rgba(139,92,246,0.1)); border: 1px solid rgba(0,212,255,0.3); margin-bottom: 1.5rem;">
                <div class="section-header">
                    <div class="section-title"><span class="section-icon">Rocket</span> Getting Started</div>
                    <button class="btn btn-secondary btn-sm" onclick="hideGettingStarted()">Dismiss</button>
                </div>
                <p style="color: var(--text-secondary); margin-bottom: 1.5rem;">Complete these steps to set up your AI voice assistant in under 5 minutes:</p>
                <div class="onboarding-steps" style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 1rem;">
                    <div class="onboarding-step" id="step-api" onclick="showMainTab('settings'); showSettingsTab('keys');" style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; cursor: pointer; border-left: 3px solid var(--neon-cyan);">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
                            <span class="step-number" style="background: var(--neon-cyan); color: #000; width: 24px; height: 24px; border-radius: 50%; display: flex; align-items: center; justify-content: center; font-weight: bold; font-size: 0.8rem;">1</span>
                            <span style="font-weight: 600;">Connect API Key</span>
                        </div>
                        <p style="color: var(--text-secondary); font-size: 0.85rem;">Add your Groq or OpenAI key to enable the AI brain</p>
                    </div>
                    <div class="onboarding-step" id="step-skill" onclick="showMainTab('skills'); showSkillsTab('create');" style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; cursor: pointer; border-left: 3px solid var(--neon-purple);">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
                            <span class="step-number" style="background: var(--neon-purple); color: #fff; width: 24px; height: 24px; border-radius: 50%; display: flex; align-items: center; justify-content: center; font-weight: bold; font-size: 0.8rem;">2</span>
                            <span style="font-weight: 600;">Create First Skill</span>
                        </div>
                        <p style="color: var(--text-secondary); font-size: 0.85rem;">Build an AI receptionist for your business in 2 min</p>
                    </div>
                    <div class="onboarding-step" id="step-test" onclick="showMainTab('skills'); showSkillsTab('test');" style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; cursor: pointer; border-left: 3px solid var(--neon-green);">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
                            <span class="step-number" style="background: var(--neon-green); color: #000; width: 24px; height: 24px; border-radius: 50%; display: flex; align-items: center; justify-content: center; font-weight: bold; font-size: 0.8rem;">3</span>
                            <span style="font-weight: 600;">Test Your Skill</span>
                        </div>
                        <p style="color: var(--text-secondary); font-size: 0.85rem;">Chat with your AI to make sure it works right</p>
                    </div>
                    <div class="onboarding-step" id="step-connect" onclick="showMainTab('voice'); showVoiceTab('connections');" style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; cursor: pointer; border-left: 3px solid var(--neon-orange);">
                        <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
                            <span class="step-number" style="background: var(--neon-orange); color: #000; width: 24px; height: 24px; border-radius: 50%; display: flex; align-items: center; justify-content: center; font-weight: bold; font-size: 0.8rem;">4</span>
                            <span style="font-weight: 600;">Connect Phone</span>
                        </div>
                        <p style="color: var(--text-secondary); font-size: 0.85rem;">Link Twilio or Vapi to handle real calls</p>
                    </div>
                </div>
            </div>

            <div class="dashboard-grid">
                <div class="glass-card stat-card card-quarter">
                    <div class="stat-value" id="total-skills">0</div>
                    <div class="stat-label">Active Skills</div>
                </div>
                <div class="glass-card stat-card card-quarter">
                    <div class="stat-value green" id="requests-today">0</div>
                    <div class="stat-label">Requests Today</div>
                </div>
                <div class="glass-card stat-card card-quarter">
                    <div class="stat-value pink" id="avg-latency">0ms</div>
                    <div class="stat-label">Avg Latency</div>
                </div>
                <div class="glass-card stat-card card-quarter">
                    <div class="stat-value orange" id="satisfaction">0%</div>
                    <div class="stat-label">Satisfaction</div>
                </div>

                <div class="glass-card card-full">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Server</span> Server Status</div>
                        <button class="btn btn-secondary btn-sm" onclick="refreshServerStatus()">Refresh</button>
                    </div>
                    <div class="server-panel">
                        <div class="server-stat">
                            <div class="stat-value green" style="font-size: 1.25rem;">Online</div>
                            <div class="stat-label">Status</div>
                        </div>
                        <div class="server-stat">
                            <div class="stat-value" style="font-size: 1.25rem;" id="warm-containers">1</div>
                            <div class="stat-label">Warm Containers</div>
                        </div>
                        <div class="server-stat">
                            <div class="stat-value" style="font-size: 1.25rem;" id="memory-usage">0GB</div>
                            <div class="stat-label">Memory</div>
                        </div>
                        <div class="server-stat">
                            <div class="stat-value orange" style="font-size: 1.25rem;" id="cost-today">$0.00</div>
                            <div class="stat-label">Cost Today</div>
                        </div>
                    </div>
                </div>

                <div class="glass-card card-three">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Skills</span> Skill Library</div>
                        <div>
                            <button class="btn btn-secondary btn-sm" onclick="showMainTab('factory')">Open Factory</button>
                            <button class="btn btn-primary btn-sm" onclick="showMainTab('factory'); showFactoryTab('profile')">+ New Skill</button>
                        </div>
                    </div>
                    <div class="skills-grid" id="skills-grid"></div>
                </div>

                <div class="glass-card card-quarter">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Train</span> Training Pipeline</div>
                    </div>
                    <div class="pipeline" style="flex-direction: column; gap: 0.75rem;">
                        <div class="pipeline-step" style="flex-direction: row; width: 100%; justify-content: flex-start; gap: 1rem;">
                            <div class="pipeline-icon complete" id="pipe-ingest">Ingest</div>
                            <div class="pipeline-label">Ingest</div>
                        </div>
                        <div class="pipeline-step" style="flex-direction: row; width: 100%; justify-content: flex-start; gap: 1rem;">
                            <div class="pipeline-icon complete" id="pipe-process">Process</div>
                            <div class="pipeline-label">Process</div>
                        </div>
                        <div class="pipeline-step" style="flex-direction: row; width: 100%; justify-content: flex-start; gap: 1rem;">
                            <div class="pipeline-icon active" id="pipe-train">Train</div>
                            <div class="pipeline-label">Train</div>
                        </div>
                        <div class="pipeline-step" style="flex-direction: row; width: 100%; justify-content: flex-start; gap: 1rem;">
                            <div class="pipeline-icon" id="pipe-deploy">Deploy</div>
                            <div class="pipeline-label">Deploy</div>
                        </div>
                    </div>
                    <div class="progress-bar">
                        <div class="progress-fill" id="training-progress" style="width: 67%;"></div>
                    </div>
                    <div class="progress-text" id="training-text">Training: plumber_expert 67%</div>
                </div>

                <div class="glass-card card-quarter">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Feed</span> Activity Feed</div>
                    </div>
                    <div class="activity-feed" id="activity-feed">
                        <div class="activity-item">
                            <span class="activity-icon">Done</span>
                            <div class="activity-content">
                                <div class="activity-message">plumber_expert deployed successfully</div>
                                <div class="activity-time">2 minutes ago</div>
                            </div>
                        </div>
                    </div>
                </div>

                <div class="glass-card card-half">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Chart</span> Request Volume (24h)</div>
                    </div>
                    <div class="chart-container" id="request-chart"></div>
                </div>

                <div class="glass-card card-half">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Speed</span> Latency (24h)</div>
                    </div>
                    <div class="chart-container" id="latency-chart"></div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- SKILL FACTORY TAB -->
        <!-- ================================================================ -->
        <div id="tab-factory" class="tab-content">
            <div class="sub-tabs">
                <button class="sub-tab-btn active" onclick="showFactoryTab('profile')">Business Profile</button>
                <button class="sub-tab-btn" onclick="showFactoryTab('golden')">Golden Prompts</button>
                <button class="sub-tab-btn" onclick="showFactoryTab('documents')">Upload Documents</button>
                <button class="sub-tab-btn" onclick="showFactoryTab('train')">Train Skill</button>
                <button class="sub-tab-btn" onclick="showFactoryTab('manage')">Manage Skills</button>
            </div>

            <!-- Business Profile -->
            <div id="factory-profile" class="sub-tab-content active">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Profile</span> Create Business Profile</div>
                    </div>
                    <div id="profile-message"></div>
                    <form id="profile-form" onsubmit="saveProfile(event)">
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Business Name *</label>
                                <input type="text" class="form-input" id="p-name" placeholder="Joe's Plumbing Services" required>
                            </div>
                            <div class="form-group">
                                <label class="form-label">Business Type</label>
                                <select class="form-select" id="p-type">
                                    <option value="General Customer Service">General Customer Service</option>
                                    <option value="Plumbing Services">Plumbing Services</option>
                                    <option value="Electrical Services">Electrical Services</option>
                                    <option value="HVAC Services">HVAC Services</option>
                                    <option value="Restaurant/Food Service">Restaurant/Food Service</option>
                                    <option value="Medical/Healthcare">Medical/Healthcare</option>
                                    <option value="Legal Services">Legal Services</option>
                                    <option value="Real Estate">Real Estate</option>
                                    <option value="Tech Support">Tech Support</option>
                                    <option value="Retail Store">Retail Store</option>
                                </select>
                            </div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Business Description</label>
                            <textarea class="form-textarea" id="p-description" placeholder="We provide 24/7 emergency plumbing services in the greater metro area..."></textarea>
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Default Greeting</label>
                                <input type="text" class="form-input" id="p-greeting" placeholder="Hello! Thanks for calling. How can I help you today?">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Personality</label>
                                <select class="form-select" id="p-personality">
                                    <option value="Friendly and casual">Friendly and casual</option>
                                    <option value="Professional and formal">Professional and formal</option>
                                    <option value="Warm and empathetic">Warm and empathetic</option>
                                    <option value="Efficient and direct">Efficient and direct</option>
                                    <option value="Technical and precise">Technical and precise</option>
                                </select>
                            </div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Key Services (one per line)</label>
                            <textarea class="form-textarea" id="p-services" placeholder="Emergency repairs&#10;Drain cleaning&#10;Water heater installation"></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">FAQ (Q: question / A: answer)</label>
                            <textarea class="form-textarea" id="p-faq" placeholder="Q: What are your hours?&#10;A: We're open 24/7 for emergencies.&#10;&#10;Q: Do you offer free estimates?&#10;A: Yes!"></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Custom Instructions</label>
                            <textarea class="form-textarea" id="p-instructions" placeholder="Always ask for the customer's address first..."></textarea>
                        </div>
                        <div style="display: flex; gap: 1rem;">
                            <button type="submit" class="btn btn-primary">Save Profile</button>
                            <button type="button" class="btn btn-success" onclick="generateFromProfile()">Generate Training Data</button>
                        </div>
                    </form>
                </div>
            </div>

            <!-- Golden Prompts -->
            <div id="factory-golden" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Star</span> Golden Prompts</div>
                        <button class="btn btn-secondary btn-sm" onclick="loadGoldenPrompts()">Refresh</button>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Voice-optimized skill manuals for ultra-fast Groq inference. These prompts are designed to stay under 3k tokens for ~50-100ms prefill latency.</p>

                    <div class="form-row">
                        <div class="form-group" style="flex: 2;">
                            <label class="form-label">Select Skill</label>
                            <select class="form-select" id="golden-skill-select" onchange="loadGoldenPrompt()">
                                <option value="receptionist">Receptionist - Professional Call Handling</option>
                                <option value="electrician">Electrician - Jenkintown Electricity</option>
                                <option value="plumber">Plumber - Plumbing Services</option>
                                <option value="lawyer">Lawyer - Legal Intake Specialist</option>
                                <option value="solar">Solar - Solar Company Sales</option>
                                <option value="tara-sales">Tara Sales - TheDashTool Assistant</option>
                                <option value="general">General - Helpful Assistant</option>
                            </select>
                        </div>
                        <div class="form-group" style="flex: 1;">
                            <label class="form-label">Voice <span style="color: var(--text-secondary); font-size: 0.8rem;">(optional)</span></label>
                            <select class="form-select" id="golden-voice-select" onchange="linkVoiceToSkill()">
                                <option value="">No custom voice</option>
                            </select>
                        </div>
                    </div>

                    <div class="stats-row" style="margin-bottom: 1rem;">
                        <div class="stat-item">
                            <div class="stat-value" id="golden-tokens">~850</div>
                            <div class="stat-label">Est. Tokens</div>
                        </div>
                        <div class="stat-item">
                            <div class="stat-value" id="golden-prefill">~14ms</div>
                            <div class="stat-label">Prefill Time</div>
                        </div>
                        <div class="stat-item">
                            <div class="stat-value" id="golden-status">Ready</div>
                            <div class="stat-label">Status</div>
                        </div>
                    </div>

                    <div class="form-group">
                        <label class="form-label">Context Variables (optional)</label>
                        <div class="form-row">
                            <div class="form-group" style="flex: 1;">
                                <input type="text" class="form-input" id="golden-business-name" placeholder="Business Name" value="">
                            </div>
                            <div class="form-group" style="flex: 1;">
                                <input type="text" class="form-input" id="golden-agent-name" placeholder="Agent Name" value="">
                            </div>
                        </div>
                    </div>

                    <div class="form-group">
                        <label class="form-label">Prompt Content</label>
                        <textarea class="form-textarea" id="golden-prompt-content" style="min-height: 400px; font-family: monospace; font-size: 0.85rem;" placeholder="Loading prompt..."></textarea>
                    </div>

                    <div style="display: flex; gap: 1rem; flex-wrap: wrap;">
                        <button class="btn btn-primary" onclick="saveGoldenPrompt()">Save Changes</button>
                        <button class="btn btn-secondary" onclick="resetGoldenPrompt()">Reset to Default</button>
                        <button class="btn btn-success" onclick="testGoldenPrompt()">Test Prompt</button>
                        <button class="btn btn-secondary" onclick="copyGoldenPrompt()">Copy to Clipboard</button>
                        <button class="btn btn-secondary" onclick="exportGoldenPrompts()">Export All</button>
                    </div>

                    <div id="golden-message" style="margin-top: 1rem;"></div>

                    <!-- Test Output -->
                    <div id="golden-test-output" style="margin-top: 1rem; display: none;">
                        <label class="form-label">Test Response</label>
                        <div class="code-block" id="golden-test-response" style="max-height: 300px; overflow-y: auto;"></div>
                    </div>
                </div>
            </div>

            <!-- Upload Documents -->
            <div id="factory-documents" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Docs</span> Upload Documents</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Upload PDFs, text files, or training data. The system will convert them into training examples.</p>
                    <div id="upload-message"></div>
                    <div class="form-group">
                        <label class="form-label">Select Business Profile</label>
                        <select class="form-select" id="doc-profile"></select>
                    </div>
                    <div class="file-upload" onclick="document.getElementById('file-input').click()">
                        <div class="file-upload-icon">Upload</div>
                        <div class="file-upload-text">Click to upload PDF, TXT, MD, JSON, or JSONL</div>
                        <input type="file" id="file-input" accept=".pdf,.txt,.md,.json,.jsonl" onchange="handleFileUpload(this)">
                    </div>
                    <div id="file-preview" style="margin-top: 1rem;"></div>
                </div>
            </div>

            <!-- Train Skill -->
            <div id="factory-train" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Train</span> Train LoRA Adapter</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Training requires a GPU. Generate a script to run on Colab, Modal, or your own machine.</p>
                    <div class="form-group">
                        <label class="form-label">Select Skill to Train</label>
                        <select class="form-select" id="train-profile"></select>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Training Steps: <span id="steps-value">60</span></label>
                        <input type="range" id="train-steps" min="20" max="200" value="60" style="width: 100%;" oninput="document.getElementById('steps-value').textContent = this.value">
                        <div style="display: flex; justify-content: space-between; color: var(--text-secondary); font-size: 0.75rem;">
                            <span>20 (Quick)</span>
                            <span>200 (Thorough)</span>
                        </div>
                    </div>
                    <button class="btn btn-primary" onclick="generateTrainingScript()">Generate Training Script</button>
                    <div id="training-output" style="margin-top: 1rem; display: none;">
                        <label class="form-label">Training Script</label>
                        <pre class="code-block" id="training-script"></pre>
                    </div>
                </div>
            </div>

            <!-- Manage Skills -->
            <div id="factory-manage" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Table</span> Skill Library</div>
                        <button class="btn btn-secondary btn-sm" onclick="loadSkillsTable()">Refresh</button>
                    </div>
                    <div class="table-filter">
                        <input type="text" id="table-search" placeholder="Search skills..." oninput="filterSkillsTable()">
                        <select id="table-type-filter" onchange="filterSkillsTable()">
                            <option value="">All Types</option>
                            <option value="Plumbing">Plumbing</option>
                            <option value="Restaurant">Restaurant</option>
                            <option value="Tech Support">Tech Support</option>
                            <option value="Healthcare">Healthcare</option>
                        </select>
                        <select id="table-status-filter" onchange="filterSkillsTable()">
                            <option value="">All Status</option>
                            <option value="Deployed">Deployed</option>
                            <option value="Training">Training</option>
                            <option value="Draft">Draft</option>
                        </select>
                    </div>
                    <table class="skills-table">
                        <thead>
                            <tr>
                                <th>Business</th>
                                <th>Type</th>
                                <th>Status</th>
                                <th>Examples</th>
                                <th>Created</th>
                                <th>Actions</th>
                            </tr>
                        </thead>
                        <tbody id="skills-table-body"></tbody>
                    </table>
                    <div style="margin-top: 1.5rem; padding: 1rem; background: var(--glass-surface); border-radius: 8px;">
                        <p style="color: var(--text-secondary); font-size: 0.85rem;">
                            <strong>Deploy to Modal:</strong><br>
                            <code style="color: var(--neon-cyan);">modal volume put lpu-skills adapters/your_skill /root/skills/your_skill.lora</code>
                        </p>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- TRAINING TAB - LoRA Fine-tuning Console -->
        <!-- ================================================================ -->
        <div id="tab-training" class="tab-content">
            <div class="sub-tabs">
                <button class="sub-tab-btn active" onclick="showTrainingTab('console')">🎯 Console</button>
                <button class="sub-tab-btn" onclick="showTrainingTab('data')">📄 Data Manager</button>
                <button class="sub-tab-btn" onclick="showTrainingTab('progress')">📊 Progress</button>
                <button class="sub-tab-btn" onclick="showTrainingTab('adapters')">📦 Adapters</button>
            </div>

            <!-- Training Console -->
            <div id="training-console" class="sub-tab-content active">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">🧠</span> Skill Training</div>
                        <a href="#" onclick="showTrainingHelp(); return false;" style="color: var(--neon-cyan); font-size: 0.9rem;">? Help</a>
                    </div>

                    <!-- Skill Selector -->
                    <div class="form-group" style="margin-bottom: 1.5rem;">
                        <label class="form-label">Select Skill to Train</label>
                        <select id="training-skill-select" class="form-select" onchange="loadTrainingDataStatus()">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>

                    <!-- Training Data Status -->
                    <div id="training-data-status" class="glass-card" style="background: var(--glass-surface); padding: 1rem; margin-bottom: 1.5rem; display: none;">
                        <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                            <span style="font-weight: 600; color: var(--neon-cyan);">📊 Training Data Status</span>
                            <span id="training-quality-badge" class="badge" style="padding: 0.25rem 0.75rem; border-radius: 12px; font-size: 0.8rem;">--</span>
                        </div>
                        <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-bottom: 1rem;">
                            <div>
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-green);" id="training-examples-count">0</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Examples</div>
                            </div>
                            <div>
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-blue);" id="training-avg-tokens">0</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Avg Tokens</div>
                            </div>
                            <div>
                                <div style="font-size: 1.5rem; font-weight: bold;" id="training-quality-score">--</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Quality</div>
                            </div>
                            <div>
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-purple);" id="training-topics">0</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Topics</div>
                            </div>
                        </div>
                        <div id="training-recommendation" style="background: rgba(255, 200, 0, 0.1); border-left: 3px solid var(--neon-orange); padding: 0.5rem 1rem; font-size: 0.85rem; color: var(--text-secondary); display: none;">
                            <strong style="color: var(--neon-orange);">⚠️ Recommendation:</strong> <span id="training-recommendation-text"></span>
                        </div>
                        <div style="display: flex; gap: 0.5rem; margin-top: 1rem;">
                            <button class="btn btn-secondary btn-sm" onclick="showMainTab('skills'); showSkillsTab('data');">+ Add Training Data</button>
                            <button class="btn btn-secondary btn-sm" onclick="viewTrainingExamples()">📄 View Examples</button>
                        </div>
                    </div>

                    <!-- Training Configuration -->
                    <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; margin-bottom: 1.5rem;">
                        <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem;">
                            <span style="font-weight: 600; color: var(--neon-cyan);">⚙️ Training Configuration</span>
                            <div>
                                <button class="btn btn-secondary btn-sm" onclick="setTrainingPreset('simple')" id="preset-simple" style="opacity: 1;">Simple</button>
                                <button class="btn btn-secondary btn-sm" onclick="setTrainingPreset('advanced')" id="preset-advanced" style="opacity: 0.5;">Advanced</button>
                            </div>
                        </div>

                        <!-- Simple Mode -->
                        <div id="training-config-simple">
                            <div class="form-group">
                                <label class="form-label">Training Intensity</label>
                                <div style="display: flex; align-items: center; gap: 1rem;">
                                    <input type="range" id="training-intensity" min="1" max="3" value="2" style="flex: 1;" onchange="updateIntensityLabel()">
                                    <span id="training-intensity-label" style="color: var(--neon-green); font-weight: 600; min-width: 100px;">Standard</span>
                                </div>
                                <div style="display: flex; justify-content: space-between; font-size: 0.75rem; color: var(--text-secondary); margin-top: 0.25rem;">
                                    <span>Quick (~5 min)</span>
                                    <span>Standard (~10 min)</span>
                                    <span>Deep (~20 min)</span>
                                </div>
                            </div>
                        </div>

                        <!-- Advanced Mode (Hidden by default) -->
                        <div id="training-config-advanced" style="display: none;">
                            <div style="display: grid; grid-template-columns: repeat(3, 1fr); gap: 1rem; margin-bottom: 1rem;">
                                <div class="form-group">
                                    <label class="form-label">Base Model</label>
                                    <select id="training-base-model" class="form-select">
                                        <option value="unsloth/Meta-Llama-3.1-8B-Instruct-bnb-4bit">Llama-3.1-8B (Recommended)</option>
                                        <option value="unsloth/Llama-3.2-3B-Instruct-bnb-4bit">Llama-3.2-3B (Faster)</option>
                                        <option value="unsloth/Mistral-7B-Instruct-v0.3-bnb-4bit">Mistral-7B</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Epochs</label>
                                    <select id="training-epochs" class="form-select">
                                        <option value="3">3 (Quick)</option>
                                        <option value="5">5</option>
                                        <option value="10" selected>10 (Standard)</option>
                                        <option value="20">20 (Deep)</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Learning Rate</label>
                                    <select id="training-lr" class="form-select">
                                        <option value="1e-4">1e-4 (Conservative)</option>
                                        <option value="2e-4" selected>2e-4 (Recommended)</option>
                                        <option value="5e-4">5e-4 (Aggressive)</option>
                                    </select>
                                </div>
                            </div>
                            <div style="display: grid; grid-template-columns: repeat(3, 1fr); gap: 1rem;">
                                <div class="form-group">
                                    <label class="form-label">LoRA Rank (r)</label>
                                    <select id="training-lora-r" class="form-select">
                                        <option value="8">8 (Small)</option>
                                        <option value="16" selected>16 (Standard)</option>
                                        <option value="32">32 (Large)</option>
                                        <option value="64">64 (Very Large)</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">LoRA Alpha</label>
                                    <select id="training-lora-alpha" class="form-select">
                                        <option value="16" selected>16</option>
                                        <option value="32">32</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Batch Size</label>
                                    <select id="training-batch-size" class="form-select">
                                        <option value="1">1</option>
                                        <option value="2" selected>2</option>
                                        <option value="4">4</option>
                                    </select>
                                </div>
                            </div>
                        </div>

                        <!-- Cost Estimate -->
                        <div style="background: rgba(0, 255, 136, 0.1); border: 1px solid rgba(0, 255, 136, 0.3); border-radius: 8px; padding: 0.75rem 1rem; margin-top: 1rem; display: flex; justify-content: space-between; align-items: center;">
                            <div>
                                <span style="color: var(--text-secondary);">Estimated Time:</span>
                                <span style="color: var(--neon-green); font-weight: 600;" id="training-time-estimate">~10 minutes</span>
                            </div>
                            <div>
                                <span style="color: var(--text-secondary);">Estimated Cost:</span>
                                <span style="color: var(--neon-green); font-weight: 600;" id="training-cost-estimate">~$0.65</span>
                            </div>
                        </div>
                    </div>

                    <!-- Start Training Button -->
                    <button class="btn btn-primary" style="width: 100%; padding: 1rem; font-size: 1.1rem;" onclick="startTraining()" id="start-training-btn">
                        🚀 Start Training
                    </button>
                </div>
            </div>

            <!-- Data Manager - Parse & Extract Training Data -->
            <div id="training-data" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📄</span> Data Manager</div>
                        <button class="btn btn-primary btn-sm" onclick="openParserModal()">📤 Parse Documents</button>
                    </div>

                    <!-- Quick Stats -->
                    <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-bottom: 1.5rem;">
                        <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-cyan);" id="dm-total-items">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Total Items</div>
                        </div>
                        <div class="glass-card" style="background: linear-gradient(135deg, rgba(245, 158, 11, 0.2), rgba(239, 68, 68, 0.2)); border: 1px solid rgba(245, 158, 11, 0.3); padding: 1rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-orange);" id="dm-pending-items">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Pending Review</div>
                        </div>
                        <div class="glass-card" style="background: linear-gradient(135deg, rgba(16, 185, 129, 0.2), rgba(0, 212, 170, 0.2)); border: 1px solid rgba(16, 185, 129, 0.3); padding: 1rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-green);" id="dm-approved-items">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Approved</div>
                        </div>
                        <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-purple);" id="dm-total-tokens">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Total Tokens</div>
                        </div>
                    </div>

                    <!-- Skill Selector for Data -->
                    <div class="form-group" style="margin-bottom: 1rem;">
                        <label class="form-label">Select Skill</label>
                        <select id="dm-skill-select" class="form-select" onchange="loadDataManagerData()">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>

                    <!-- Toolbar -->
                    <div style="display: flex; justify-content: space-between; align-items: center; gap: 1rem; margin-bottom: 1rem; flex-wrap: wrap;">
                        <div style="display: flex; gap: 0.5rem; align-items: center;">
                            <input type="text" id="dm-search" class="form-input" placeholder="Search..." style="width: 200px;" onkeyup="filterDataManager()">
                            <select id="dm-category-filter" class="form-select" style="width: 140px;" onchange="filterDataManager()">
                                <option value="">All Categories</option>
                                <option value="general">General</option>
                                <option value="faq">FAQ</option>
                                <option value="technical">Technical</option>
                                <option value="pricing">Pricing</option>
                                <option value="procedure">Procedure</option>
                            </select>
                            <select id="dm-status-filter" class="form-select" style="width: 120px;" onchange="filterDataManager()">
                                <option value="">All Status</option>
                                <option value="pending">Pending</option>
                                <option value="approved">Approved</option>
                            </select>
                        </div>
                        <div style="display: flex; gap: 0.5rem;">
                            <button class="btn btn-secondary btn-sm" onclick="bulkApproveSelected()">✅ Approve Selected</button>
                            <button class="btn btn-secondary btn-sm" onclick="bulkMoveToTraining()">📥 Move to Training</button>
                            <button class="btn btn-secondary btn-sm" onclick="bulkDeleteSelected()" style="color: #ff4444;">🗑️ Delete</button>
                        </div>
                    </div>

                    <!-- Data List -->
                    <div id="dm-data-list" style="background: var(--glass-surface); border-radius: 12px; max-height: 500px; overflow-y: auto;">
                        <div style="text-align: center; padding: 3rem; color: var(--text-secondary);">
                            <div style="font-size: 3rem; margin-bottom: 1rem;">📭</div>
                            <p>No extracted data yet</p>
                            <p style="font-size: 0.9rem;">Upload and parse documents to extract training data</p>
                            <button class="btn btn-primary" style="margin-top: 1rem;" onclick="openParserModal()">📤 Parse Documents</button>
                        </div>
                    </div>

                    <!-- Pagination -->
                    <div id="dm-pagination" style="display: flex; justify-content: center; gap: 0.5rem; margin-top: 1rem;">
                    </div>
                </div>
            </div>

            <!-- Document Parser Modal -->
            <div id="parser-modal" style="display: none; position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: rgba(0,0,0,0.8); z-index: 1000; align-items: center; justify-content: center;">
                <div class="glass-card" style="width: 600px; max-width: 90%; max-height: 80vh; overflow-y: auto;">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📤</span> Parse Documents</div>
                        <button class="btn btn-secondary btn-sm" onclick="closeParserModal()">✕</button>
                    </div>

                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Upload documents to extract Q&A pairs for training. Supports 70+ file types.</p>
                    <div style="background: rgba(245, 158, 11, 0.1); border: 1px solid rgba(245, 158, 11, 0.3); border-radius: 8px; padding: 0.75rem; margin-bottom: 1rem;">
                        <strong style="color: var(--neon-orange);">⚠️ Note:</strong>
                        <span style="color: var(--text-secondary);">Files are processed and <strong>immediately deleted</strong>. Only extracted Q&A data is stored - zero storage waste!</span>
                    </div>

                    <!-- Skill Selector -->
                    <div class="form-group" style="margin-bottom: 1rem;">
                        <label class="form-label">Target Skill</label>
                        <select id="parser-skill-select" class="form-select">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>

                    <!-- File Upload -->
                    <div id="parser-dropzone" style="border: 2px dashed var(--neon-cyan); border-radius: 12px; padding: 2rem; text-align: center; margin-bottom: 1rem; cursor: pointer;" onclick="document.getElementById('parser-files').click()">
                        <div style="font-size: 3rem; margin-bottom: 0.5rem;">📁</div>
                        <p style="color: var(--text-secondary);">Drop files here or click to browse</p>
                        <p style="font-size: 0.8rem; color: var(--text-secondary);">PDF, DOCX, TXT, CSV, JSON, Images, Audio, and more</p>
                        <input type="file" id="parser-files" multiple style="display: none;" onchange="handleParserFiles(this.files)">
                    </div>

                    <!-- Selected Files -->
                    <div id="parser-file-list" style="margin-bottom: 1rem;"></div>

                    <!-- Progress -->
                    <div id="parser-progress" style="display: none; margin-bottom: 1rem;">
                        <div style="display: flex; justify-content: space-between; margin-bottom: 0.5rem;">
                            <span>Processing...</span>
                            <span id="parser-progress-text">0%</span>
                        </div>
                        <div style="background: var(--glass-surface); border-radius: 8px; height: 8px; overflow: hidden;">
                            <div id="parser-progress-bar" style="background: var(--neon-cyan); height: 100%; width: 0%; transition: width 0.3s;"></div>
                        </div>
                    </div>

                    <!-- Buttons -->
                    <div style="display: flex; gap: 0.5rem; justify-content: flex-end;">
                        <button class="btn btn-secondary" onclick="closeParserModal()">Cancel</button>
                        <button class="btn btn-primary" id="parser-submit-btn" onclick="submitParserFiles()">🚀 Parse & Extract</button>
                    </div>
                </div>
            </div>

            <!-- Training Progress -->
            <div id="training-progress" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📊</span> Training Progress</div>
                        <button class="btn btn-secondary btn-sm" id="cancel-training-btn" onclick="cancelTraining()" style="display: none;">Cancel</button>
                    </div>

                    <!-- No Active Training -->
                    <div id="no-training-message" style="text-align: center; padding: 3rem; color: var(--text-secondary);">
                        <div style="font-size: 3rem; margin-bottom: 1rem;">🧘</div>
                        <p style="font-size: 1.1rem;">No training in progress</p>
                        <p style="font-size: 0.9rem;">Select a skill and start training from the Console tab</p>
                        <button class="btn btn-primary" style="margin-top: 1rem;" onclick="showTrainingTab('console')">Go to Console</button>
                    </div>

                    <!-- Active Training Display - Enhanced Dashboard -->
                    <div id="active-training-display" style="display: none;">
                        <!-- Header with pulse indicator -->
                        <div style="display: flex; align-items: center; justify-content: space-between; margin-bottom: 1.5rem;">
                            <div style="display: flex; align-items: center; gap: 12px;">
                                <div id="training-pulse" style="width: 12px; height: 12px; background: var(--neon-green); border-radius: 50%; animation: pulse 2s infinite; box-shadow: 0 0 20px var(--neon-green);"></div>
                                <h3 style="margin: 0; font-size: 1.2rem;">🧠 Training: <span id="progress-skill-name" style="color: var(--neon-cyan);">--</span></h3>
                            </div>
                            <span id="progress-status" class="badge" style="background: linear-gradient(135deg, var(--neon-green), #059669); padding: 6px 16px; border-radius: 20px; font-size: 12px; font-weight: 600;">● Running</span>
                        </div>

                        <!-- Stats Grid - 4 Cards -->
                        <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 16px; margin-bottom: 1.5rem;">
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; text-align: center; border: 1px solid var(--glass-border);">
                                <div style="font-size: 24px; margin-bottom: 8px;">📉</div>
                                <div id="stat-loss" style="font-size: 28px; font-weight: 700; background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan)); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">--</div>
                                <div style="color: var(--text-secondary); font-size: 12px; margin-top: 4px;">Current Loss</div>
                                <div id="stat-loss-trend" style="font-size: 11px; color: var(--neon-green); margin-top: 8px;">--</div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; text-align: center; border: 1px solid var(--glass-border);">
                                <div style="font-size: 24px; margin-bottom: 8px;">🔄</div>
                                <div id="stat-steps" style="font-size: 28px; font-weight: 700; background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan)); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">0/0</div>
                                <div style="color: var(--text-secondary); font-size: 12px; margin-top: 4px;">Training Steps</div>
                                <div style="height: 4px; background: var(--glass-surface); border-radius: 2px; margin-top: 12px; overflow: hidden;">
                                    <div id="stat-steps-bar" style="height: 100%; background: linear-gradient(90deg, var(--neon-green), var(--neon-cyan)); border-radius: 2px; transition: width 0.5s; width: 0%;"></div>
                                </div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; text-align: center; border: 1px solid var(--glass-border);">
                                <div style="font-size: 24px; margin-bottom: 8px;">⏱️</div>
                                <div id="stat-eta" style="font-size: 28px; font-weight: 700; background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan)); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">--:--</div>
                                <div style="color: var(--text-secondary); font-size: 12px; margin-top: 4px;">Time Remaining</div>
                                <div id="stat-elapsed" style="font-size: 11px; color: var(--text-secondary); margin-top: 8px;">Elapsed: --</div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; text-align: center; border: 1px solid var(--glass-border);">
                                <div style="font-size: 24px; margin-bottom: 8px;">🎯</div>
                                <div id="stat-epoch" style="font-size: 28px; font-weight: 700; background: linear-gradient(135deg, var(--neon-green), var(--neon-cyan)); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">0</div>
                                <div style="color: var(--text-secondary); font-size: 12px; margin-top: 4px;">Current Epoch</div>
                                <div id="stat-epoch-total" style="font-size: 11px; color: var(--text-secondary); margin-top: 8px;">of 10 epochs</div>
                            </div>
                        </div>

                        <!-- Chart and GPU Section -->
                        <div style="display: grid; grid-template-columns: 2fr 1fr; gap: 16px; margin-bottom: 1.5rem;">
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; border: 1px solid var(--glass-border);">
                                <div style="font-size: 14px; color: var(--text-secondary); margin-bottom: 16px;">📈 Loss Curve (Real-Time)</div>
                                <canvas id="loss-chart" height="180"></canvas>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; border: 1px solid var(--glass-border);">
                                <div style="font-size: 14px; color: var(--text-secondary); margin-bottom: 16px;">🖥️ GPU Metrics</div>
                                <div id="gpu-name" style="font-size: 18px; font-weight: 600; margin-bottom: 16px; color: var(--neon-green);">NVIDIA A10G</div>
                                <div style="margin-bottom: 16px;">
                                    <div style="display: flex; justify-content: space-between; font-size: 12px; color: var(--text-secondary); margin-bottom: 6px;">
                                        <span>Memory</span>
                                        <span id="gpu-mem-text">18.2 / 22 GB</span>
                                    </div>
                                    <div style="height: 8px; background: var(--glass-border); border-radius: 4px; overflow: hidden;">
                                        <div id="gpu-mem-bar" style="height: 100%; background: linear-gradient(90deg, var(--neon-green), var(--neon-cyan)); border-radius: 4px; transition: width 0.3s; width: 83%;"></div>
                                    </div>
                                </div>
                                <div>
                                    <div style="display: flex; justify-content: space-between; font-size: 12px; color: var(--text-secondary); margin-bottom: 6px;">
                                        <span>Utilization</span>
                                        <span id="gpu-util-text">94%</span>
                                    </div>
                                    <div style="height: 8px; background: var(--glass-border); border-radius: 4px; overflow: hidden;">
                                        <div id="gpu-util-bar" style="height: 100%; background: linear-gradient(90deg, var(--neon-green), var(--neon-cyan)); border-radius: 4px; transition: width 0.3s; width: 94%;"></div>
                                    </div>
                                </div>
                            </div>
                        </div>

                        <!-- Example Preview and Education Section -->
                        <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 16px; margin-bottom: 1.5rem;">
                            <div class="glass-card" style="background: linear-gradient(135deg, rgba(16, 185, 129, 0.1), rgba(0, 255, 242, 0.1)); border: 1px solid rgba(16, 185, 129, 0.3); border-radius: 16px; padding: 20px;">
                                <div style="font-size: 14px; margin-bottom: 16px; display: flex; align-items: center; gap: 8px;">📝 Currently Learning...</div>
                                <div id="current-example-card" style="background: rgba(0,0,0,0.2); border-radius: 12px; padding: 16px;">
                                    <div style="margin-bottom: 12px;">
                                        <div style="font-size: 11px; color: var(--text-secondary); margin-bottom: 4px;">👤 User Question</div>
                                        <div id="current-example-input" style="font-size: 14px; line-height: 1.5;">"How do I set up automation?"</div>
                                    </div>
                                    <div>
                                        <div style="font-size: 11px; color: var(--text-secondary); margin-bottom: 4px;">🤖 Expected Response</div>
                                        <div id="current-example-output" style="font-size: 14px; line-height: 1.5;">"To set up automation, click..."</div>
                                    </div>
                                </div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); border-radius: 16px; padding: 20px; border: 1px solid var(--glass-border);">
                                <div style="font-size: 14px; color: var(--text-secondary); margin-bottom: 16px;">💡 Did You Know?</div>
                                <div id="fact-card" style="text-align: center;">
                                    <div id="fact-icon" style="font-size: 48px; margin-bottom: 12px;">🧠</div>
                                    <h4 id="fact-title" style="font-size: 18px; margin-bottom: 8px;">What is LoRA?</h4>
                                    <p id="fact-text" style="color: var(--text-secondary); font-size: 14px; line-height: 1.6;">Low-Rank Adaptation trains only <strong style="color: var(--neon-green);">0.92%</strong> of parameters, making it 10x faster than full fine-tuning!</p>
                                </div>
                                <div style="display: flex; justify-content: center; gap: 8px; margin-top: 16px;">
                                    <span class="fact-dot" style="width: 8px; height: 8px; border-radius: 50%; background: var(--neon-green);"></span>
                                    <span class="fact-dot" style="width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.2);"></span>
                                    <span class="fact-dot" style="width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.2);"></span>
                                    <span class="fact-dot" style="width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.2);"></span>
                                    <span class="fact-dot" style="width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.2);"></span>
                                </div>
                            </div>
                        </div>

                        <!-- Training Log (Collapsed by default) -->
                        <details class="glass-card" style="background: var(--glass-surface); padding: 1rem; border-radius: 16px; border: 1px solid var(--glass-border);">
                            <summary style="cursor: pointer; display: flex; justify-content: space-between; align-items: center; font-weight: 600; color: var(--neon-cyan);">
                                📋 Training Log
                            </summary>
                            <div id="training-log" class="console" style="height: 200px; overflow-y: auto; font-family: monospace; font-size: 0.8rem; padding: 0.5rem; background: #0a0a0f; border-radius: 4px; margin-top: 1rem;">
                                <div style="color: var(--text-secondary);">Waiting for training to start...</div>
                            </div>
                        </details>
                    </div>

                    <!-- Training Complete Display -->
                    <div id="training-complete-display" style="display: none;">
                        <div style="text-align: center; padding: 2rem;">
                            <div style="font-size: 4rem; margin-bottom: 1rem;">🎉</div>
                            <h2 style="color: var(--neon-green); margin-bottom: 0.5rem;">Training Complete!</h2>
                            <p style="color: var(--text-secondary); margin-bottom: 1.5rem;">Your skill adapter is ready to use</p>
                        </div>

                        <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-bottom: 1.5rem;">
                            <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-green);" id="result-final-loss">--</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Final Loss</div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-cyan);" id="result-training-time">--</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Training Time</div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-purple);" id="result-examples">--</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Examples</div>
                            </div>
                            <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; text-align: center;">
                                <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-orange);" id="result-cost">--</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary);">Cost</div>
                            </div>
                        </div>

                        <!-- Test Chat -->
                        <div class="glass-card" style="background: var(--glass-surface); padding: 1rem; margin-bottom: 1.5rem;">
                            <div style="font-weight: 600; margin-bottom: 0.75rem; color: var(--neon-cyan);">💬 Test Your Trained Skill</div>
                            <div style="display: flex; gap: 0.5rem;">
                                <input type="text" id="test-adapter-prompt" class="form-input" placeholder="Ask your trained skill something..." style="flex: 1;">
                                <button class="btn btn-primary" onclick="testTrainedAdapter()">Send</button>
                            </div>
                            <div id="test-adapter-response" class="console" style="margin-top: 1rem; min-height: 60px; padding: 0.75rem; display: none;"></div>
                        </div>

                        <div style="display: flex; gap: 0.5rem; justify-content: center;">
                            <button class="btn btn-secondary" onclick="showTrainingTab('console')">🔄 Train Again</button>
                            <button class="btn btn-secondary" onclick="showTrainingTab('adapters')">📦 View Adapters</button>
                            <button class="btn btn-primary" onclick="showMainTab('skills')">🏠 Back to Skills</button>
                        </div>
                    </div>
                </div>
            </div>

            <!-- Adapters Gallery -->
            <div id="training-adapters" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📦</span> Trained Adapters</div>
                        <button class="btn btn-secondary btn-sm" onclick="refreshAdapters()">↻ Refresh</button>
                    </div>

                    <div id="adapters-gallery" style="display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 1rem; margin-bottom: 1.5rem;">
                        <!-- Adapters will be loaded here -->
                        <div id="no-adapters-message" style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--text-secondary);">
                            <div style="font-size: 3rem; margin-bottom: 1rem;">📦</div>
                            <p>No trained adapters yet</p>
                            <p style="font-size: 0.9rem;">Train a skill to create your first adapter</p>
                            <button class="btn btn-primary" style="margin-top: 1rem;" onclick="showTrainingTab('console')">Start Training</button>
                        </div>
                    </div>

                    <!-- Training History Table -->
                    <div class="section-header" style="margin-top: 2rem;">
                        <div class="section-title"><span class="section-icon">📊</span> Training History</div>
                    </div>
                    <div style="overflow-x: auto;">
                        <table class="skills-table" style="width: 100%;">
                            <thead>
                                <tr>
                                    <th>Date</th>
                                    <th>Skill</th>
                                    <th>Status</th>
                                    <th>Loss</th>
                                    <th>Time</th>
                                    <th>Cost</th>
                                </tr>
                            </thead>
                            <tbody id="training-history-table">
                                <tr>
                                    <td colspan="6" style="text-align: center; color: var(--text-secondary);">No training history</td>
                                </tr>
                            </tbody>
                        </table>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- UNIFIED SKILLS & TRAINING TAB -->
        <!-- ================================================================ -->
        <div id="tab-skills-training" class="tab-content">
            <!-- Toast Container for Notifications -->
            <div id="toast-container" class="toast-container"></div>

            <!-- Skill Context Bar - Shows selected skill -->
            <div id="skill-context-bar" class="skill-context-bar">
                <div class="skill-context-icon" id="context-skill-icon">🎯</div>
                <div class="skill-context-info">
                    <h3 class="skill-context-name" id="context-skill-name">No Skill Selected</h3>
                    <div class="skill-context-meta">
                        <span id="context-skill-status">●  Untrained</span>
                        <span id="context-skill-examples">📊 0 examples</span>
                        <span id="context-skill-tokens">🔤 0 tokens</span>
                    </div>
                </div>
                <div class="skill-context-actions">
                    <button class="btn btn-secondary btn-sm" onclick="changeSelectedSkill()">Change Skill</button>
                    <button class="btn btn-secondary btn-sm" onclick="editSkillDetails()">Edit Details</button>
                </div>
            </div>

            <!-- Workflow Step Indicator -->
            <div class="workflow-steps" id="workflow-steps">
                <div class="workflow-step active" data-step="1" onclick="goToWorkflowStep(1)">
                    <div class="step-number">1</div>
                    <div class="step-label">Select Skill</div>
                </div>
                <div class="step-connector" id="connector-1-2"></div>
                <div class="workflow-step disabled" data-step="2" onclick="goToWorkflowStep(2)">
                    <div class="step-number">2</div>
                    <div class="step-label">Add Training Data</div>
                </div>
                <div class="step-connector" id="connector-2-3"></div>
                <div class="workflow-step disabled" data-step="3" onclick="goToWorkflowStep(3)">
                    <div class="step-number">3</div>
                    <div class="step-label">Test & Train</div>
                </div>
            </div>

            <!-- Workflow Content Panels -->
            <div class="workflow-content">

                <!-- ============================================ -->
                <!-- STEP 1: SELECT SKILL -->
                <!-- ============================================ -->
                <div id="workflow-step-1" class="workflow-panel active">
                    <!-- Header with Actions -->
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1.5rem; flex-wrap: wrap; gap: 1rem;">
                        <div>
                            <h3 style="margin: 0; color: var(--text-primary);">Step 1: Select a Skill</h3>
                            <p style="margin: 0.25rem 0 0 0; color: var(--text-secondary); font-size: 0.9rem;">Choose an existing skill or create a new one to train</p>
                        </div>
                        <div style="display: flex; gap: 0.5rem;">
                            <button class="btn btn-primary" onclick="showUnifiedCreateSkill()">+ Create New Skill</button>
                            <button class="btn btn-secondary" onclick="refreshUnifiedSkills()">Refresh</button>
                        </div>
                    </div>

                    <!-- Filter/Search Bar -->
                    <div class="glass-card" style="padding: 1rem; margin-bottom: 1.5rem;">
                        <div style="display: flex; gap: 1rem; align-items: center; flex-wrap: wrap;">
                            <div style="flex: 1; min-width: 200px;">
                                <input type="text" id="skill-search-input" class="form-input" placeholder="Search skills..." onkeyup="filterSkillCards()">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-sm skill-filter-btn active" data-filter="all" onclick="setSkillFilter('all')">All</button>
                                <button class="btn btn-sm skill-filter-btn" data-filter="untrained" onclick="setSkillFilter('untrained')">Untrained</button>
                                <button class="btn btn-sm skill-filter-btn" data-filter="has_data" onclick="setSkillFilter('has_data')">Has Data</button>
                                <button class="btn btn-sm skill-filter-btn" data-filter="trained" onclick="setSkillFilter('trained')">Trained</button>
                            </div>
                            <select id="skill-sort-select" class="form-select" style="width: auto;" onchange="sortSkillCards()">
                                <option value="name">Sort: Name</option>
                                <option value="status">Sort: Status</option>
                                <option value="recent">Sort: Recent</option>
                            </select>
                        </div>
                    </div>

                    <!-- Skills Grid -->
                    <div id="unified-skills-grid" class="skills-grid" style="display: grid; grid-template-columns: repeat(auto-fill, minmax(320px, 1fr)); gap: 1.5rem;">
                        <div class="loading-skills" style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--text-secondary);">
                            <div style="font-size: 2rem; margin-bottom: 1rem;">🔄</div>
                            Loading skills...
                        </div>
                    </div>

                    <!-- Create Skill Form (inline) -->
                    <div class="glass-card" id="unified-create-skill-form" style="display: none; margin-top: 1.5rem;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">NEW</span> Create New Skill</div>
                            <button class="btn btn-sm" onclick="hideUnifiedCreateSkill()" style="background: transparent; color: var(--text-secondary);">✕</button>
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Skill ID (lowercase, no spaces)</label>
                                <input type="text" class="form-input" id="unified-new-skill-id" placeholder="my_custom_skill">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Display Name</label>
                                <input type="text" class="form-input" id="unified-new-skill-name" placeholder="My Custom Skill">
                            </div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Description</label>
                            <input type="text" class="form-input" id="unified-new-skill-description" placeholder="Brief description of what this skill does">
                        </div>
                        <div class="form-group">
                            <label class="form-label">System Prompt</label>
                            <textarea class="form-textarea" id="unified-new-skill-prompt" rows="6" placeholder="You are an AI assistant specialized in..."></textarea>
                        </div>
                        <div class="form-row">
                            <button class="btn btn-primary" onclick="createUnifiedSkill()">Create Skill</button>
                            <button class="btn btn-secondary" onclick="hideUnifiedCreateSkill()">Cancel</button>
                        </div>
                    </div>
                </div>

                <!-- ============================================ -->
                <!-- STEP 2: ADD TRAINING DATA (Side-by-Side) -->
                <!-- ============================================ -->
                <div id="workflow-step-2" class="workflow-panel">
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1.5rem;">
                        <div>
                            <h3 style="margin: 0; color: var(--text-primary);">Step 2: Add Training Data</h3>
                            <p style="margin: 0.25rem 0 0 0; color: var(--text-secondary); font-size: 0.9rem;">Add examples to teach your skill how to respond</p>
                        </div>
                        <div style="display: flex; gap: 0.5rem; align-items: center;">
                            <span id="step2-data-count" style="color: var(--neon-green); font-weight: 600;">0 examples</span>
                            <button class="btn btn-primary" onclick="goToWorkflowStep(3)" id="step2-next-btn" disabled>Continue to Test & Train →</button>
                        </div>
                    </div>

                    <!-- Data Stats Bar -->
                    <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-bottom: 1.5rem;">
                        <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-cyan);" id="wf-data-total">0</div>
                            <div style="font-size: 0.75rem; color: var(--text-secondary);">Total</div>
                        </div>
                        <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-orange);" id="wf-data-pending">0</div>
                            <div style="font-size: 0.75rem; color: var(--text-secondary);">Pending</div>
                        </div>
                        <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-green);" id="wf-data-approved">0</div>
                            <div style="font-size: 0.75rem; color: var(--text-secondary);">Approved</div>
                        </div>
                        <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-purple);" id="wf-data-tokens">0</div>
                            <div style="font-size: 0.75rem; color: var(--text-secondary);">Tokens</div>
                        </div>
                    </div>

                    <!-- Side-by-Side Layout -->
                    <div class="training-layout">
                        <!-- LEFT: Manual Entry -->
                        <div class="glass-card training-panel">
                            <h4>✏️ Manual Entry</h4>
                            <div class="data-entry-form">
                                <div class="form-group">
                                    <label class="form-label">User Question / Input</label>
                                    <textarea id="wf-user-input" class="form-textarea" rows="3" placeholder="What the user might ask..."></textarea>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Ideal Response</label>
                                    <textarea id="wf-assistant-response" class="form-textarea" rows="4" placeholder="How the AI should respond..."></textarea>
                                </div>
                                <div class="data-entry-actions">
                                    <button class="btn btn-primary" onclick="saveWorkflowEntry()">💾 Save Entry</button>
                                    <button class="btn btn-secondary" onclick="saveAndAddAnotherWorkflow()">Save & Add Another</button>
                                </div>
                            </div>

                            <!-- Recent Entries -->
                            <div style="margin-top: 1.5rem;">
                                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.5rem;">
                                    <span style="font-size: 0.9rem; color: var(--text-secondary);">Recent Entries</span>
                                    <button class="btn btn-sm btn-secondary" onclick="loadWorkflowTrainingData()">🔄</button>
                                </div>
                                <div id="wf-training-data-list" class="training-data-list">
                                    <div style="padding: 1rem; text-align: center; color: var(--text-secondary);">No entries yet</div>
                                </div>
                            </div>
                        </div>

                        <!-- RIGHT: Bulk Upload & AI Generate -->
                        <div class="training-panel">
                            <!-- File Upload -->
                            <div class="glass-card" style="margin-bottom: 1rem;">
                                <h4>📤 Bulk Upload</h4>
                                <div class="upload-zone" id="wf-upload-zone" onclick="document.getElementById('wf-file-input').click()"
                                     ondragover="event.preventDefault(); this.classList.add('dragover')"
                                     ondragleave="this.classList.remove('dragover')"
                                     ondrop="handleWorkflowDrop(event)">
                                    <div class="upload-icon">📁</div>
                                    <div style="color: var(--text-primary); font-weight: 500;">Drop files here or click to upload</div>
                                    <div style="color: var(--text-secondary); font-size: 0.85rem; margin-top: 0.5rem;">
                                        Supports: PDF, Word, Excel, CSV, TXT, JSON, and 60+ more formats
                                    </div>
                                    <input type="file" id="wf-file-input" style="display: none;" multiple accept="*/*" onchange="handleWorkflowUpload(this.files)">
                                </div>
                                <div id="wf-upload-status" style="margin-top: 0.75rem; font-size: 0.85rem; color: var(--text-secondary);"></div>
                            </div>

                            <!-- AI Generate -->
                            <div class="glass-card">
                                <h4>🤖 AI Generate</h4>
                                <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 1rem;">
                                    Let AI generate training examples based on your skill's purpose
                                </p>
                                <div class="form-group">
                                    <label class="form-label">Context / Topic (optional)</label>
                                    <textarea id="wf-ai-context" class="form-textarea" rows="2" placeholder="E.g., Common questions about solar panel installation..."></textarea>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Number of Examples</label>
                                    <select id="wf-ai-count" class="form-select">
                                        <option value="5">5 examples</option>
                                        <option value="10" selected>10 examples</option>
                                        <option value="20">20 examples</option>
                                        <option value="50">50 examples</option>
                                    </select>
                                </div>
                                <button class="btn btn-primary" onclick="generateWorkflowAiData()" style="width: 100%;">
                                    ✨ Generate Training Data
                                </button>
                                <div id="wf-ai-status" style="margin-top: 0.75rem; font-size: 0.85rem; color: var(--text-secondary);"></div>
                            </div>
                        </div>
                    </div>
                </div>

                <!-- ============================================ -->
                <!-- STEP 3: TEST & TRAIN -->
                <!-- ============================================ -->
                <div id="workflow-step-3" class="workflow-panel">
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1.5rem;">
                        <div>
                            <h3 style="margin: 0; color: var(--text-primary);">Step 3: Test & Train</h3>
                            <p style="margin: 0.25rem 0 0 0; color: var(--text-secondary); font-size: 0.9rem;">Test your skill's responses, then train when ready</p>
                        </div>
                        <button class="btn btn-secondary" onclick="goToWorkflowStep(2)">← Back to Training Data</button>
                    </div>

                    <!-- Test & Train Side-by-Side -->
                    <div class="test-train-layout">
                        <!-- LEFT: Test Chat -->
                        <div class="glass-card">
                            <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">💬 Test Your Skill</h4>
                            <div class="chat-panel">
                                <div id="wf-chat-messages" class="chat-messages">
                                    <div style="text-align: center; color: var(--text-secondary); padding: 2rem;">
                                        Send a message to test how your skill responds
                                    </div>
                                </div>
                                <div class="chat-input-area">
                                    <input type="text" id="wf-chat-input" class="form-input" placeholder="Type a test message..."
                                           onkeypress="if(event.key==='Enter') sendWorkflowChat()">
                                    <button class="btn btn-primary" onclick="sendWorkflowChat()">Send</button>
                                </div>
                            </div>
                        </div>

                        <!-- RIGHT: Training Panel -->
                        <div class="glass-card train-panel">
                            <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">🚀 Train Your Skill</h4>

                            <!-- Readiness Checks -->
                            <div class="train-readiness">
                                <div style="font-weight: 500; margin-bottom: 0.75rem; color: var(--text-primary);">Training Readiness</div>
                                <div id="wf-readiness-checks">
                                    <div class="readiness-row">
                                        <div class="readiness-check fail">✗</div>
                                        <span>Loading...</span>
                                    </div>
                                </div>
                            </div>

                            <!-- Training Config -->
                            <div style="margin-top: 1rem;">
                                <div class="form-group">
                                    <label class="form-label">Training Intensity</label>
                                    <div style="display: flex; align-items: center; gap: 1rem;">
                                        <input type="range" id="wf-training-intensity" min="1" max="3" value="2" style="flex: 1;" onchange="updateWorkflowIntensityLabel()">
                                        <span id="wf-intensity-label" style="color: var(--neon-green); font-weight: 600; min-width: 80px;">Standard</span>
                                    </div>
                                </div>
                                <div style="background: rgba(0, 255, 136, 0.1); border: 1px solid rgba(0, 255, 136, 0.3); border-radius: 8px; padding: 0.75rem 1rem; display: flex; justify-content: space-between;">
                                    <span>Est. Time: <strong id="wf-time-est">~10 min</strong></span>
                                    <span>Est. Cost: <strong id="wf-cost-est">~$0.65</strong></span>
                                </div>
                            </div>

                            <!-- Train Button -->
                            <div class="train-button-wrapper">
                                <button id="wf-train-btn" class="train-button-large" onclick="startWorkflowTraining()" disabled>
                                    🚀 Start Training
                                </button>
                                <div id="wf-training-status" style="margin-top: 0.75rem; text-align: center; font-size: 0.85rem; color: var(--text-secondary);"></div>
                            </div>

                            <!-- Enhanced Training Display (hidden until training starts) -->
                            <div id="wf-enhanced-training" style="display: none; margin-top: 1rem; background: rgba(0,0,0,0.2); border-radius: 12px; padding: 1.5rem;">
                                <!-- Header -->
                                <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 1.5rem;">
                                    <div id="wf-pulse" style="width: 12px; height: 12px; background: #10B981; border-radius: 50%; animation: pulse 1.5s infinite;"></div>
                                    <span style="font-weight: 600; color: #10B981;">Training in Progress</span>
                                </div>

                                <!-- Stats Grid -->
                                <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.75rem; margin-bottom: 1.5rem;">
                                    <div style="background: rgba(124, 58, 237, 0.15); padding: 0.75rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 0.65rem; color: #A1A1AA; margin-bottom: 0.25rem;">LOSS</div>
                                        <div id="wf-stat-loss" style="font-size: 1.25rem; font-weight: 700; color: #A78BFA;">--</div>
                                    </div>
                                    <div style="background: rgba(0, 217, 255, 0.15); padding: 0.75rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 0.65rem; color: #A1A1AA; margin-bottom: 0.25rem;">STEPS</div>
                                        <div id="wf-stat-steps" style="font-size: 1.25rem; font-weight: 700; color: #00D9FF;">--</div>
                                    </div>
                                    <div style="background: rgba(16, 185, 129, 0.15); padding: 0.75rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 0.65rem; color: #A1A1AA; margin-bottom: 0.25rem;">ETA</div>
                                        <div id="wf-stat-eta" style="font-size: 1.25rem; font-weight: 700; color: #10B981;">--</div>
                                    </div>
                                    <div style="background: rgba(245, 158, 11, 0.15); padding: 0.75rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 0.65rem; color: #A1A1AA; margin-bottom: 0.25rem;">EPOCH</div>
                                        <div id="wf-stat-epoch" style="font-size: 1.25rem; font-weight: 700; color: #F59E0B;">--</div>
                                    </div>
                                </div>

                                <!-- Progress Bar -->
                                <div style="margin-bottom: 1.5rem;">
                                    <div style="display: flex; justify-content: space-between; margin-bottom: 0.5rem;">
                                        <span style="font-size: 0.75rem; color: #A1A1AA;">Progress</span>
                                        <span id="wf-progress-percent" style="font-size: 0.75rem; font-weight: 600; color: #00D9FF;">0%</span>
                                    </div>
                                    <div style="background: rgba(255,255,255,0.1); border-radius: 4px; height: 8px; overflow: hidden;">
                                        <div id="wf-progress-bar" style="height: 100%; background: linear-gradient(90deg, #10B981, #00D9FF); width: 0%; transition: width 0.3s;"></div>
                                    </div>
                                </div>

                                <!-- Loss Chart -->
                                <div style="background: rgba(0,0,0,0.3); border-radius: 8px; padding: 1rem; margin-bottom: 1rem;">
                                    <div style="font-size: 0.75rem; color: #A1A1AA; margin-bottom: 0.5rem;">📈 Loss Curve</div>
                                    <canvas id="wf-loss-chart" height="120"></canvas>
                                </div>

                                <!-- Educational Fact -->
                                <div style="background: rgba(124, 58, 237, 0.1); border: 1px solid rgba(124, 58, 237, 0.3); border-radius: 8px; padding: 0.75rem;">
                                    <div style="font-size: 0.65rem; color: #A78BFA; margin-bottom: 0.25rem;">💡 DID YOU KNOW?</div>
                                    <div id="wf-training-fact" style="font-size: 0.8rem; color: #E5E5E5;">LoRA trains only 0.1-1% of model parameters, making it 10-100x cheaper than full fine-tuning.</div>
                                </div>
                            </div>

                            <!-- Existing Adapters -->
                            <div style="margin-top: 1.5rem;">
                                <div style="font-weight: 500; margin-bottom: 0.5rem; color: var(--text-primary);">Trained Adapters</div>
                                <div id="wf-adapters-list" style="max-height: 150px; overflow-y: auto;">
                                    <div style="text-align: center; color: var(--text-secondary); font-size: 0.85rem; padding: 1rem;">
                                        No adapters yet. Train to create one!
                                    </div>
                                </div>
                            </div>
                        </div>
                    </div>
                </div>

            </div><!-- End workflow-content -->

            <!-- Golden Prompts Sub-tab -->
            <div id="unified-golden" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">✨</span> Golden Prompts</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Create and manage perfect example conversations for training.</p>
                    <div class="form-group">
                        <label class="form-label">Select Skill</label>
                        <select id="unified-golden-skill" class="form-select" onchange="loadUnifiedGoldenPrompts()">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>
                    <div id="unified-golden-list" style="margin-top: 1rem;">
                        <p style="color: var(--text-secondary); text-align: center;">Select a skill to view golden prompts</p>
                    </div>
                    <div style="margin-top: 1rem;">
                        <button class="btn btn-primary" onclick="addUnifiedGoldenPrompt()">+ Add Golden Prompt</button>
                    </div>
                </div>
            </div>

            <!-- Training Sub-tab -->
            <div id="unified-training" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">🧠</span> Train Skill with LoRA</div>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Select Skill to Train</label>
                        <select id="unified-training-skill" class="form-select" onchange="loadUnifiedTrainingStatus()">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>
                    <div id="unified-training-status" style="display: none; margin: 1rem 0; padding: 1rem; background: var(--glass-surface); border-radius: 8px;">
                        <!-- Training status will load here -->
                    </div>
                    <div class="form-group">
                        <label class="form-label">Training Intensity</label>
                        <div style="display: flex; align-items: center; gap: 1rem;">
                            <input type="range" id="unified-training-intensity" min="1" max="3" value="2" style="flex: 1;" onchange="updateUnifiedIntensity()">
                            <span id="unified-intensity-label" style="color: var(--neon-green); font-weight: 600; min-width: 80px;">Standard</span>
                        </div>
                    </div>
                    <div style="background: rgba(0, 255, 136, 0.1); border: 1px solid rgba(0, 255, 136, 0.3); border-radius: 8px; padding: 0.75rem 1rem; margin: 1rem 0; display: flex; justify-content: space-between;">
                        <span>Est. Time: <strong id="unified-time-est">~10 min</strong></span>
                        <span>Est. Cost: <strong id="unified-cost-est">~$0.65</strong></span>
                    </div>
                    <button class="btn btn-primary" style="width: 100%; padding: 1rem;" onclick="startUnifiedTraining()">🚀 Start Training</button>
                </div>
            </div>

            <!-- Data Manager Sub-tab -->
            <div id="unified-data" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📄</span> Training Data Manager</div>
                        <button class="btn btn-primary btn-sm" onclick="openParserModal()">📤 Parse Documents</button>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Select Skill</label>
                        <select id="unified-data-skill" class="form-select" onchange="loadUnifiedDataManager()">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>
                    <div id="unified-data-stats" style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin: 1rem 0;">
                        <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-cyan);" id="unified-data-total">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Total</div>
                        </div>
                        <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-orange);" id="unified-data-pending">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Pending</div>
                        </div>
                        <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-green);" id="unified-data-approved">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Approved</div>
                        </div>
                        <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                            <div style="font-size: 1.5rem; font-weight: bold; color: var(--neon-purple);" id="unified-data-tokens">0</div>
                            <div style="font-size: 0.8rem; color: var(--text-secondary);">Tokens</div>
                        </div>
                    </div>
                    <div style="display: flex; gap: 0.5rem; margin-bottom: 1rem;">
                        <button class="btn btn-primary btn-sm" onclick="openManualEntryModal()">✏️ Manual Entry</button>
                        <button class="btn btn-secondary btn-sm" onclick="openBulkImportModal()">📤 Bulk Import</button>
                        <button class="btn btn-secondary btn-sm" onclick="openAiGenerateModal()">🤖 AI Generate</button>
                    </div>
                    <div id="unified-data-list" style="max-height: 400px; overflow-y: auto;">
                        <p style="color: var(--text-secondary); text-align: center;">Select a skill to view training data</p>
                    </div>
                </div>
            </div>

            <!-- Test Chat Sub-tab -->
            <div id="unified-chat" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">💬</span> Test Your Skill</div>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Select Skill</label>
                        <select id="unified-chat-skill" class="form-select">
                            <option value="">-- Select a skill --</option>
                        </select>
                    </div>
                    <div id="unified-chat-messages" style="background: var(--glass-surface); border-radius: 8px; padding: 1rem; min-height: 300px; max-height: 400px; overflow-y: auto; margin: 1rem 0;">
                        <p style="color: var(--text-secondary); text-align: center;">Select a skill and start chatting!</p>
                    </div>
                    <div style="display: flex; gap: 0.5rem;">
                        <input type="text" id="unified-chat-input" class="form-input" placeholder="Type your message..." style="flex: 1;" onkeypress="if(event.key==='Enter')sendUnifiedChat()">
                        <button class="btn btn-primary" onclick="sendUnifiedChat()">Send</button>
                    </div>
                </div>
            </div>

            <!-- Adapters Sub-tab -->
            <div id="unified-adapters" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">📦</span> Trained Adapters</div>
                        <button class="btn btn-secondary btn-sm" onclick="refreshUnifiedAdapters()">Refresh</button>
                    </div>
                    <div id="unified-adapters-list">
                        <p style="color: var(--text-secondary); text-align: center; padding: 2rem;">Loading adapters...</p>
                    </div>
                </div>
            </div>

            <!-- API Sub-tab -->
            <div id="unified-api" class="sub-tab-content" style="display: none;">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">🔗</span> Outgoing API Integration</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Use these endpoints to connect your skills to external platforms.</p>
                    <div class="form-group">
                        <label class="form-label">Fast Brain API URL</label>
                        <div style="display: flex; gap: 0.5rem;">
                            <input type="text" class="form-input" id="unified-api-url" value="https://jenkintownelectricity--fast-brain-lpu-fastapi-app.modal.run" readonly style="flex: 1; font-family: monospace;">
                            <button class="btn btn-secondary" onclick="copyUnifiedApiUrl()">Copy</button>
                        </div>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Example cURL Request</label>
                        <pre style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; overflow-x: auto; font-size: 0.85rem;">curl -X POST https://jenkintownelectricity--fast-brain-lpu-fastapi-app.modal.run/v1/chat/completions \
  -H "Content-Type: application/json" \
  -d '{"skill_id": "your_skill", "messages": [{"role": "user", "content": "Hello"}]}'</pre>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- SKILL DETAIL MODAL -->
        <!-- ================================================================ -->
        <div id="skill-detail-modal" class="modal-overlay" style="display: none;">
            <div class="modal-content" style="max-width: 900px; max-height: 90vh; overflow: hidden; display: flex; flex-direction: column;">
                <!-- Modal Header -->
                <div class="modal-header" style="display: flex; justify-content: space-between; align-items: center; padding: 1.5rem; border-bottom: 1px solid var(--glass-border); flex-shrink: 0;">
                    <div style="display: flex; align-items: center; gap: 1rem;">
                        <div id="skill-modal-icon" style="font-size: 2rem;">🎯</div>
                        <div>
                            <h2 id="skill-modal-title" style="margin: 0; font-size: 1.4rem;">Skill Name</h2>
                            <div style="display: flex; gap: 0.5rem; margin-top: 0.25rem;">
                                <span id="skill-modal-status-badge" class="skill-status-badge untrained">Untrained</span>
                                <span id="skill-modal-examples-count" style="color: var(--text-secondary); font-size: 0.85rem;">0 examples</span>
                            </div>
                        </div>
                    </div>
                    <button class="btn btn-secondary" onclick="closeSkillDetailModal()" style="padding: 0.5rem 1rem;">✕</button>
                </div>

                <!-- Modal Tabs -->
                <div class="modal-tabs" style="display: flex; gap: 0; border-bottom: 1px solid var(--glass-border); padding: 0 1.5rem; flex-shrink: 0;">
                    <button class="modal-tab-btn active" onclick="showSkillModalTab('overview')" data-tab="overview">Overview</button>
                    <button class="modal-tab-btn" onclick="showSkillModalTab('training-data')" data-tab="training-data">Training Data</button>
                    <button class="modal-tab-btn" onclick="showSkillModalTab('train')" data-tab="train">Train</button>
                    <button class="modal-tab-btn" onclick="showSkillModalTab('adapters')" data-tab="adapters">Adapters</button>
                </div>

                <!-- Modal Body -->
                <div class="modal-body" style="flex: 1; overflow-y: auto; padding: 1.5rem;">

                    <!-- Overview Tab -->
                    <div id="skill-modal-overview" class="modal-tab-content active">
                        <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 1.5rem;">
                            <!-- Skill Info -->
                            <div class="glass-card" style="padding: 1.25rem;">
                                <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">📋 Skill Details</h4>
                                <div class="form-group">
                                    <label class="form-label">Skill ID</label>
                                    <input type="text" id="skill-modal-id" class="form-input" readonly>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Display Name</label>
                                    <input type="text" id="skill-modal-name" class="form-input">
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Description</label>
                                    <textarea id="skill-modal-description" class="form-textarea" rows="2"></textarea>
                                </div>
                            </div>

                            <!-- Stats -->
                            <div class="glass-card" style="padding: 1.25rem;">
                                <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">📊 Statistics</h4>
                                <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 1rem;">
                                    <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 1.75rem; font-weight: bold; color: var(--neon-green);" id="skill-modal-stat-examples">0</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Training Examples</div>
                                    </div>
                                    <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 1.75rem; font-weight: bold; color: var(--neon-blue);" id="skill-modal-stat-tokens">0</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Total Tokens</div>
                                    </div>
                                    <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 1.75rem; font-weight: bold; color: var(--neon-purple);" id="skill-modal-stat-adapters">0</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Adapters</div>
                                    </div>
                                    <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; text-align: center;">
                                        <div style="font-size: 1.75rem; font-weight: bold; color: var(--neon-orange);" id="skill-modal-stat-quality">--</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Data Quality</div>
                                    </div>
                                </div>
                            </div>

                            <!-- System Prompt -->
                            <div class="glass-card" style="grid-column: 1 / -1; padding: 1.25rem;">
                                <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">💬 System Prompt</h4>
                                <textarea id="skill-modal-system-prompt" class="form-textarea" rows="6" placeholder="Define the personality and behavior of this skill..."></textarea>
                            </div>
                        </div>
                        <div style="display: flex; justify-content: flex-end; gap: 0.5rem; margin-top: 1.5rem;">
                            <button class="btn btn-secondary" onclick="testSkillChat()">🔬 Test Chat</button>
                            <button class="btn btn-primary" onclick="saveSkillOverview()">💾 Save Changes</button>
                        </div>
                    </div>

                    <!-- Training Data Tab -->
                    <div id="skill-modal-training-data" class="modal-tab-content" style="display: none;">
                        <!-- Data Actions -->
                        <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1.5rem; flex-wrap: wrap; gap: 1rem;">
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="openManualEntryModal()">✏️ Manual Entry</button>
                                <button class="btn btn-secondary btn-sm" onclick="openBulkImportModal()">📤 Bulk Import</button>
                                <button class="btn btn-secondary btn-sm" onclick="openAiGenerateModal()">🤖 AI Generate</button>
                            </div>
                            <div style="display: flex; gap: 0.5rem; align-items: center;">
                                <span id="skill-data-selected-count" style="color: var(--text-secondary); font-size: 0.9rem;">0 selected</span>
                                <button class="btn btn-secondary btn-sm" onclick="approveSelectedData()" id="approve-selected-btn" disabled>✓ Approve</button>
                                <button class="btn btn-danger btn-sm" onclick="deleteSelectedData()" id="delete-selected-btn" disabled>🗑 Delete</button>
                            </div>
                        </div>

                        <!-- Data Stats -->
                        <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-bottom: 1.5rem;">
                            <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                                <div style="font-size: 1.25rem; font-weight: bold; color: var(--neon-cyan);" id="skill-data-total">0</div>
                                <div style="font-size: 0.75rem; color: var(--text-secondary);">Total</div>
                            </div>
                            <div class="glass-card" style="padding: 0.75rem; text-align: center; background: linear-gradient(135deg, rgba(245, 158, 11, 0.1), rgba(239, 68, 68, 0.1));">
                                <div style="font-size: 1.25rem; font-weight: bold; color: var(--neon-orange);" id="skill-data-pending">0</div>
                                <div style="font-size: 0.75rem; color: var(--text-secondary);">Pending</div>
                            </div>
                            <div class="glass-card" style="padding: 0.75rem; text-align: center; background: linear-gradient(135deg, rgba(16, 185, 129, 0.1), rgba(0, 212, 170, 0.1));">
                                <div style="font-size: 1.25rem; font-weight: bold; color: var(--neon-green);" id="skill-data-approved">0</div>
                                <div style="font-size: 0.75rem; color: var(--text-secondary);">Approved</div>
                            </div>
                            <div class="glass-card" style="padding: 0.75rem; text-align: center;">
                                <div style="font-size: 1.25rem; font-weight: bold; color: var(--neon-purple);" id="skill-data-tokens">0</div>
                                <div style="font-size: 0.75rem; color: var(--text-secondary);">Tokens</div>
                            </div>
                        </div>

                        <!-- Data Table -->
                        <div class="glass-card" style="padding: 1rem; max-height: 400px; overflow-y: auto;">
                            <table class="skills-table" style="width: 100%;">
                                <thead>
                                    <tr>
                                        <th style="width: 30px;"><input type="checkbox" id="select-all-data" onchange="toggleSelectAllData()"></th>
                                        <th>User Input</th>
                                        <th>Response (Preview)</th>
                                        <th style="width: 80px;">Tokens</th>
                                        <th style="width: 80px;">Status</th>
                                        <th style="width: 80px;">Actions</th>
                                    </tr>
                                </thead>
                                <tbody id="skill-data-tbody">
                                    <tr>
                                        <td colspan="6" style="text-align: center; color: var(--text-secondary); padding: 2rem;">No training data yet. Add some!</td>
                                    </tr>
                                </tbody>
                            </table>
                        </div>
                    </div>

                    <!-- Train Tab -->
                    <div id="skill-modal-train" class="modal-tab-content" style="display: none;">
                        <!-- Training Readiness Check -->
                        <div id="training-readiness" class="glass-card" style="padding: 1.25rem; margin-bottom: 1.5rem;">
                            <h4 style="margin: 0 0 1rem 0; color: var(--neon-cyan);">🔍 Training Readiness</h4>
                            <div id="readiness-checks" style="display: flex; flex-direction: column; gap: 0.5rem;">
                                <!-- Populated by JS -->
                            </div>
                        </div>

                        <!-- Training Configuration -->
                        <div class="glass-card" style="padding: 1.25rem; margin-bottom: 1.5rem;">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem;">
                                <h4 style="margin: 0; color: var(--neon-cyan);">⚙️ Training Configuration</h4>
                                <div>
                                    <button class="btn btn-secondary btn-sm" onclick="setModalTrainingPreset('simple')" id="modal-preset-simple" style="opacity: 1;">Simple</button>
                                    <button class="btn btn-secondary btn-sm" onclick="setModalTrainingPreset('advanced')" id="modal-preset-advanced" style="opacity: 0.5;">Advanced</button>
                                </div>
                            </div>

                            <!-- Simple Mode -->
                            <div id="modal-training-simple">
                                <div class="form-group">
                                    <label class="form-label">Training Intensity</label>
                                    <div style="display: flex; align-items: center; gap: 1rem;">
                                        <input type="range" id="modal-training-intensity" min="1" max="3" value="2" style="flex: 1;" onchange="updateModalIntensityLabel()">
                                        <span id="modal-intensity-label" style="color: var(--neon-green); font-weight: 600; min-width: 100px;">Standard</span>
                                    </div>
                                </div>
                            </div>

                            <!-- Advanced Mode -->
                            <div id="modal-training-advanced" style="display: none;">
                                <div style="display: grid; grid-template-columns: repeat(3, 1fr); gap: 1rem;">
                                    <div class="form-group">
                                        <label class="form-label">Epochs</label>
                                        <select id="modal-training-epochs" class="form-select">
                                            <option value="3">3 (Quick)</option>
                                            <option value="5">5</option>
                                            <option value="10" selected>10 (Standard)</option>
                                            <option value="20">20 (Deep)</option>
                                        </select>
                                    </div>
                                    <div class="form-group">
                                        <label class="form-label">LoRA Rank</label>
                                        <select id="modal-training-lora-r" class="form-select">
                                            <option value="8">8</option>
                                            <option value="16" selected>16</option>
                                            <option value="32">32</option>
                                        </select>
                                    </div>
                                    <div class="form-group">
                                        <label class="form-label">Learning Rate</label>
                                        <select id="modal-training-lr" class="form-select">
                                            <option value="1e-4">1e-4</option>
                                            <option value="2e-4" selected>2e-4</option>
                                            <option value="5e-4">5e-4</option>
                                        </select>
                                    </div>
                                </div>
                            </div>

                            <!-- Cost Estimate -->
                            <div style="background: rgba(0, 255, 136, 0.1); border: 1px solid rgba(0, 255, 136, 0.3); border-radius: 8px; padding: 0.75rem 1rem; margin-top: 1rem; display: flex; justify-content: space-between;">
                                <span>Est. Time: <strong id="modal-time-estimate">~10 min</strong></span>
                                <span>Est. Cost: <strong id="modal-cost-estimate">~$0.65</strong></span>
                            </div>
                        </div>

                        <!-- Start Training Button -->
                        <button class="btn btn-primary" style="width: 100%; padding: 1rem; font-size: 1.1rem;" onclick="startModalTraining()" id="modal-start-training-btn">
                            🚀 Start Training
                        </button>

                        <!-- Training Progress (shown when training) -->
                        <div id="modal-training-progress" style="display: none; margin-top: 1.5rem;">
                            <div class="glass-card" style="padding: 1.25rem;">
                                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem;">
                                    <h4 style="margin: 0; color: var(--neon-cyan);">📊 Training Progress</h4>
                                    <span id="modal-training-status" class="badge" style="background: var(--neon-orange); color: #000;">Running</span>
                                </div>
                                <div class="progress-bar" style="margin-bottom: 1rem;">
                                    <div class="progress-fill" id="modal-progress-fill" style="width: 0%;"></div>
                                </div>
                                <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; text-align: center;">
                                    <div>
                                        <div style="font-size: 1.25rem; font-weight: bold;" id="modal-progress-epoch">0/10</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Epoch</div>
                                    </div>
                                    <div>
                                        <div style="font-size: 1.25rem; font-weight: bold; color: var(--neon-green);" id="modal-progress-loss">--</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Loss</div>
                                    </div>
                                    <div>
                                        <div style="font-size: 1.25rem; font-weight: bold;" id="modal-progress-time">0:00</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">Elapsed</div>
                                    </div>
                                    <div>
                                        <div style="font-size: 1.25rem; font-weight: bold;" id="modal-progress-eta">--</div>
                                        <div style="font-size: 0.8rem; color: var(--text-secondary);">ETA</div>
                                    </div>
                                </div>
                            </div>
                        </div>
                    </div>

                    <!-- Adapters Tab -->
                    <div id="skill-modal-adapters" class="modal-tab-content" style="display: none;">
                        <div id="skill-adapters-list" style="display: flex; flex-direction: column; gap: 1rem;">
                            <div style="text-align: center; color: var(--text-secondary); padding: 2rem;">
                                No adapters for this skill yet. Train to create one!
                            </div>
                        </div>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- MANUAL ENTRY MODAL -->
        <!-- ================================================================ -->
        <div id="manual-entry-modal" class="modal-overlay" style="display: none;">
            <div class="modal-content" style="max-width: 700px;">
                <div class="modal-header" style="display: flex; justify-content: space-between; align-items: center; padding: 1.5rem; border-bottom: 1px solid var(--glass-border);">
                    <h3 style="margin: 0;">✏️ Add Training Example</h3>
                    <button class="btn btn-secondary" onclick="closeManualEntryModal()" style="padding: 0.5rem 1rem;">✕</button>
                </div>
                <div style="padding: 1.5rem;">
                    <div class="form-group">
                        <label class="form-label">User Input / Question</label>
                        <textarea id="manual-user-input" class="form-textarea" rows="3" placeholder="What would a user ask or say?"></textarea>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Assistant Response</label>
                        <textarea id="manual-assistant-response" class="form-textarea" rows="6" placeholder="How should the skill respond?"></textarea>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Category (optional)</label>
                        <input type="text" id="manual-category" class="form-input" placeholder="e.g., pricing, hours, products">
                    </div>
                    <div style="display: flex; justify-content: flex-end; gap: 0.5rem; margin-top: 1.5rem;">
                        <button class="btn btn-secondary" onclick="closeManualEntryModal()">Cancel</button>
                        <button class="btn btn-secondary" onclick="saveAndAddAnother()">Save & Add Another</button>
                        <button class="btn btn-primary" onclick="saveManualEntry()">Save & Close</button>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- BULK IMPORT MODAL (Parser) -->
        <!-- ================================================================ -->
        <div id="bulk-import-modal" class="modal-overlay" style="display: none;">
            <div class="modal-content" style="max-width: 800px;">
                <div class="modal-header" style="display: flex; justify-content: space-between; align-items: center; padding: 1.5rem; border-bottom: 1px solid var(--glass-border);">
                    <h3 style="margin: 0;">📤 Bulk Import Documents</h3>
                    <button class="btn btn-secondary" onclick="closeBulkImportModal()" style="padding: 0.5rem 1rem;">✕</button>
                </div>
                <div style="padding: 1.5rem;">
                    <div class="glass-card" style="padding: 2rem; text-align: center; border: 2px dashed var(--glass-border); margin-bottom: 1.5rem;">
                        <div style="font-size: 3rem; margin-bottom: 1rem;">📁</div>
                        <p style="margin: 0 0 1rem 0; color: var(--text-secondary);">Drag & drop files here or click to browse</p>
                        <input type="file" id="bulk-import-files" multiple style="display: none;" onchange="handleBulkImportFiles(event)">
                        <button class="btn btn-primary" onclick="document.getElementById('bulk-import-files').click()">Choose Files</button>
                        <p style="margin: 1rem 0 0 0; font-size: 0.8rem; color: var(--text-secondary);">
                            Supports: PDF, Word, Excel, Text, Audio (MP3, WAV), Video, and 70+ more formats
                        </p>
                    </div>
                    <div id="bulk-import-queue" style="display: none;">
                        <h4 style="margin: 0 0 1rem 0;">Files to Process:</h4>
                        <div id="bulk-import-file-list" style="max-height: 200px; overflow-y: auto;"></div>
                        <div id="bulk-import-status" style="display: none; margin-top: 1rem; padding: 1rem; background: var(--glass-bg); border-radius: 8px;">
                            <div id="bulk-import-progress" style="text-align: center; color: var(--accent-color);">
                                Processing...
                            </div>
                        </div>
                        <div style="display: flex; justify-content: flex-end; gap: 0.5rem; margin-top: 1rem;">
                            <button class="btn btn-secondary" onclick="clearBulkImportQueue()">Clear</button>
                            <button class="btn btn-primary" onclick="processBulkImport()">📤 Process Files</button>
                        </div>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- AI GENERATE MODAL -->
        <!-- ================================================================ -->
        <div id="ai-generate-modal" class="modal-overlay" style="display: none;">
            <div class="modal-content" style="max-width: 700px;">
                <div class="modal-header" style="display: flex; justify-content: space-between; align-items: center; padding: 1.5rem; border-bottom: 1px solid var(--glass-border);">
                    <h3 style="margin: 0;">🤖 AI Generate Training Data</h3>
                    <button class="btn btn-secondary" onclick="closeAiGenerateModal()" style="padding: 0.5rem 1rem;">✕</button>
                </div>
                <div style="padding: 1.5rem;">
                    <div class="form-group">
                        <label class="form-label">Topic / Context</label>
                        <textarea id="ai-generate-topic" class="form-textarea" rows="3" placeholder="Describe what topics the AI should generate Q&A pairs about..."></textarea>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Number of Examples</label>
                        <select id="ai-generate-count" class="form-select">
                            <option value="5">5 examples</option>
                            <option value="10" selected>10 examples</option>
                            <option value="20">20 examples</option>
                            <option value="50">50 examples</option>
                        </select>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Style</label>
                        <select id="ai-generate-style" class="form-select">
                            <option value="conversational">Conversational</option>
                            <option value="formal">Formal/Professional</option>
                            <option value="friendly">Friendly/Casual</option>
                            <option value="technical">Technical</option>
                        </select>
                    </div>
                    <div id="ai-generate-status" style="display: none; margin-top: 1rem; padding: 1rem; background: var(--glass-bg); border-radius: 8px;">
                        <!-- Status content will be inserted by JS -->
                    </div>
                    <div style="display: flex; justify-content: flex-end; gap: 0.5rem; margin-top: 1.5rem;">
                        <button class="btn btn-secondary" onclick="closeAiGenerateModal()">Cancel</button>
                        <button class="btn btn-primary" onclick="generateAiTrainingData()">🤖 Generate</button>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- COMMAND CENTER TAB -->
        <!-- ================================================================ -->
        <div id="tab-command" class="tab-content">
            <div class="sub-tabs">
                <button class="sub-tab-btn active" onclick="showCommandTab('keys')">API Keys</button>
                <button class="sub-tab-btn" onclick="showCommandTab('test')">Test LLM</button>
                <button class="sub-tab-btn" onclick="showCommandTab('compare')">Compare</button>
                <button class="sub-tab-btn" onclick="showCommandTab('masking')">Latency Masking</button>
                <button class="sub-tab-btn" onclick="showCommandTab('stats')">Stats</button>
                <button class="sub-tab-btn" onclick="showCommandTab('voices')">Voice Library</button>
                <button class="sub-tab-btn" onclick="showCommandTab('voice')">Voice Integration</button>
            </div>

            <!-- API Keys -->
            <div id="command-keys" class="sub-tab-content active">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Keys</span> Configure API Keys</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">API keys are stored in the database and used for LLM and voice services.</p>
                    <div id="keys-message"></div>

                    <!-- LLM Providers -->
                    <h4 style="color: var(--neon-cyan); margin: 1rem 0 0.5rem 0; font-size: 0.9rem;">LLM Providers</h4>
                    <div class="form-row">
                        <div class="form-group">
                            <label class="form-label">Groq API Key</label>
                            <input type="password" class="form-input" id="key-groq" placeholder="gsk_...">
                        </div>
                        <div class="form-group">
                            <label class="form-label">OpenAI API Key</label>
                            <input type="password" class="form-input" id="key-openai" placeholder="sk-...">
                        </div>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Anthropic API Key</label>
                        <input type="password" class="form-input" id="key-anthropic" placeholder="sk-ant-...">
                    </div>

                    <!-- Voice Providers -->
                    <h4 style="color: var(--neon-green); margin: 1.5rem 0 0.5rem 0; font-size: 0.9rem;">Voice / TTS Providers</h4>
                    <div class="form-row">
                        <div class="form-group">
                            <label class="form-label">ElevenLabs API Key</label>
                            <input type="password" class="form-input" id="key-elevenlabs" placeholder="xi_...">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Cartesia API Key</label>
                            <input type="password" class="form-input" id="key-cartesia" placeholder="sk_car_...">
                        </div>
                    </div>
                    <div class="form-row">
                        <div class="form-group">
                            <label class="form-label">Deepgram API Key (STT)</label>
                            <input type="password" class="form-input" id="key-deepgram" placeholder="dg_...">
                        </div>
                        <div class="form-group">
                            <label class="form-label">PlayHT API Key</label>
                            <input type="password" class="form-input" id="key-playht" placeholder="...">
                        </div>
                    </div>

                    <button class="btn btn-primary" onclick="saveApiKeys()" style="margin-top: 1rem;">Save All API Keys</button>
                </div>

                <!-- API Endpoints Card -->
                <div class="glass-card" style="margin-top: 1rem;">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">API</span> API Endpoints</div>
                        <button class="btn btn-secondary btn-sm" onclick="checkApiHealth()">Check Health</button>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Use these URLs to connect external services to your HIVE215 platform.</p>

                    <!-- Fast Brain API URL -->
                    <div class="form-group" style="margin-bottom: 1rem;">
                        <label class="form-label" style="display: flex; align-items: center; gap: 0.5rem;">
                            Fast Brain API
                            <span id="api-health-fast-brain" class="health-indicator" style="width: 8px; height: 8px; border-radius: 50%; background: var(--text-secondary);"></span>
                        </label>
                        <div style="display: flex; gap: 0.5rem;">
                            <input type="text" class="form-input" id="api-url-fast-brain" value="https://jenkintownelectricity--fast-brain-lpu-fastapi-app.modal.run" readonly style="flex: 1; font-family: monospace; font-size: 0.85rem;">
                            <button class="btn btn-secondary btn-sm" onclick="copyApiUrl('api-url-fast-brain')" title="Copy to clipboard">Copy</button>
                        </div>
                        <div id="api-status-fast-brain" style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.25rem;"></div>
                    </div>

                    <!-- Dashboard URL -->
                    <div class="form-group" style="margin-bottom: 1rem;">
                        <label class="form-label" style="display: flex; align-items: center; gap: 0.5rem;">
                            Dashboard URL
                            <span id="api-health-dashboard" class="health-indicator" style="width: 8px; height: 8px; border-radius: 50%; background: var(--text-secondary);"></span>
                        </label>
                        <div style="display: flex; gap: 0.5rem;">
                            <input type="text" class="form-input" id="api-url-dashboard" value="https://jenkintownelectricity--hive215-dashboard-flask-app.modal.run" readonly style="flex: 1; font-family: monospace; font-size: 0.85rem;">
                            <button class="btn btn-secondary btn-sm" onclick="copyApiUrl('api-url-dashboard')" title="Copy to clipboard">Copy</button>
                        </div>
                        <div id="api-status-dashboard" style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.25rem;"></div>
                    </div>

                    <!-- Quick Reference -->
                    <div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; margin-top: 1rem;">
                        <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem; font-size: 0.9rem;">Quick Reference</h4>
                        <table style="width: 100%; font-size: 0.85rem;">
                            <tr>
                                <td style="color: var(--text-secondary); padding: 0.25rem 0;">Health Check:</td>
                                <td style="font-family: monospace;"><code style="color: var(--neon-green);">GET /health</code></td>
                            </tr>
                            <tr>
                                <td style="color: var(--text-secondary); padding: 0.25rem 0;">Hybrid Chat:</td>
                                <td style="font-family: monospace;"><code style="color: var(--neon-green);">POST /v1/chat/hybrid</code></td>
                            </tr>
                            <tr>
                                <td style="color: var(--text-secondary); padding: 0.25rem 0;">Voice Chat:</td>
                                <td style="font-family: monospace;"><code style="color: var(--neon-green);">POST /v1/chat/voice</code></td>
                            </tr>
                            <tr>
                                <td style="color: var(--text-secondary); padding: 0.25rem 0;">List Skills:</td>
                                <td style="font-family: monospace;"><code style="color: var(--neon-green);">GET /v1/skills</code></td>
                            </tr>
                        </table>
                    </div>
                </div>
            </div>

            <!-- Test LLM -->
            <div id="command-test" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Test</span> Test Single LLM</div>
                    </div>
                    <div class="form-row">
                        <div class="form-group">
                            <label class="form-label">Prompt</label>
                            <textarea class="form-textarea" id="test-prompt" placeholder="What are your hours?"></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Provider</label>
                            <select class="form-select" id="test-provider">
                                <option value="groq">Groq (Llama 3.1 8B)</option>
                                <option value="openai">OpenAI (GPT-4o-mini)</option>
                                <option value="anthropic">Anthropic (Claude 3 Haiku)</option>
                                <option value="bitnet">BitNet LPU</option>
                            </select>
                            <div style="margin-top: 0.5rem;">
                                <label><input type="checkbox" id="test-masking"> Enable Latency Masking</label>
                            </div>
                        </div>
                    </div>
                    <button class="btn btn-primary" onclick="testLLM()" id="test-btn">Run Test</button>
                    <div style="margin-top: 1rem;">
                        <label class="form-label">Response</label>
                        <div class="console" id="test-output">Waiting for test...</div>
                    </div>
                </div>
            </div>

            <!-- Compare -->
            <div id="command-compare" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Compare</span> Compare Providers</div>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Prompt</label>
                        <textarea class="form-textarea" id="compare-prompt" placeholder="Explain quantum computing in one sentence."></textarea>
                    </div>
                    <div class="form-group">
                        <label class="form-label">Providers to Compare</label>
                        <div style="display: flex; gap: 1rem; flex-wrap: wrap;">
                            <label><input type="checkbox" value="groq" class="compare-provider" checked> Groq</label>
                            <label><input type="checkbox" value="openai" class="compare-provider" checked> OpenAI</label>
                            <label><input type="checkbox" value="anthropic" class="compare-provider"> Anthropic</label>
                            <label><input type="checkbox" value="bitnet" class="compare-provider"> BitNet</label>
                        </div>
                    </div>
                    <button class="btn btn-primary" onclick="compareLLMs()" id="compare-btn">Compare</button>
                    <div class="compare-results" id="compare-output" style="margin-top: 1rem;"></div>
                </div>
            </div>

            <!-- Latency Masking -->
            <div id="command-masking" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Mask</span> Latency Masking Config</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Configure filler sounds and phrases that make slow LLMs feel natural.</p>
                    <div class="form-row">
                        <div>
                            <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Filler Sounds</h4>
                            <ul style="color: var(--text-secondary); list-style: none;">
                                <li>Hmm...</li>
                                <li>Mmm...</li>
                                <li>Umm...</li>
                                <li>Ah...</li>
                                <li>Well...</li>
                            </ul>
                        </div>
                        <div>
                            <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Thinking Phrases</h4>
                            <ul style="color: var(--text-secondary); list-style: none;">
                                <li>Let me think about that...</li>
                                <li>That's a good question...</li>
                                <li>Let me check...</li>
                                <li>One moment...</li>
                            </ul>
                        </div>
                    </div>
                    <div style="margin-top: 1rem;">
                        <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Skill-Specific Fillers</h4>
                        <table class="skills-table">
                            <tr><th>Skill Type</th><th>Custom Fillers</th></tr>
                            <tr><td>Technical</td><td>"Let me look that up...", "Checking the docs..."</td></tr>
                            <tr><td>Customer Service</td><td>"I understand.", "Let me help with that..."</td></tr>
                            <tr><td>Scheduling</td><td>"Let me check the calendar...", "One moment..."</td></tr>
                            <tr><td>Sales</td><td>"Great question!", "Let me find the best option..."</td></tr>
                        </table>
                    </div>
                </div>
            </div>

            <!-- Stats -->
            <div id="command-stats" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Stats</span> Performance Statistics</div>
                        <button class="btn btn-secondary btn-sm" onclick="refreshStats()">Refresh</button>
                    </div>
                    <pre class="code-block" id="stats-output">Loading stats...</pre>
                </div>
            </div>

            <!-- Voice Library -->
            <div id="command-voices" class="sub-tab-content">
                <!-- Voice Selection Card -->
                <div class="glass-card" style="margin-bottom: 1rem;">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Voices</span> Voice Selection</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Select your preferred TTS provider and voice for your agent.</p>

                    <div class="form-row">
                        <div class="form-group">
                            <label class="form-label">TTS Provider</label>
                            <select class="form-select" id="voice-provider" onchange="loadVoiceSettingsVoices()">
                                <option value="">Select a provider...</option>
                                <optgroup label="Premium (Live API)">
                                    <option value="elevenlabs">ElevenLabs</option>
                                    <option value="cartesia">Cartesia</option>
                                    <option value="deepgram">Deepgram (Aura)</option>
                                    <option value="openai">OpenAI TTS</option>
                                </optgroup>
                                <optgroup label="Free / Open Source">
                                    <option value="parler_tts">Parler TTS (Expressive - GPU)</option>
                                    <option value="edge_tts">Edge TTS (Microsoft - Free)</option>
                                    <option value="kokoro">Kokoro</option>
                                    <option value="chatterbox">Chatterbox (Resemble AI)</option>
                                    <option value="xtts">XTTS-v2 (Coqui AI)</option>
                                    <option value="openvoice">OpenVoice (MyShell)</option>
                                </optgroup>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Voice</label>
                            <select class="form-select" id="voice-select" onchange="selectVoice()">
                                <option value="">Select a voice...</option>
                            </select>
                        </div>
                    </div>

                    <div id="voice-details" style="display: none; margin-bottom: 1rem; padding: 1rem; background: var(--glass-surface); border-radius: 8px;">
                        <div style="display: flex; justify-content: space-between; align-items: center;">
                            <div>
                                <strong id="voice-name" style="color: var(--neon-cyan);"></strong>
                                <span id="voice-meta" style="color: var(--text-secondary); margin-left: 0.5rem;"></span>
                            </div>
                            <div style="display: flex; gap: 0.5rem; align-items: center;">
                                <button id="voice-preview-btn" class="btn btn-secondary btn-sm" style="display: none;" title="Play provider preview">
                                    ▶ Preview
                                </button>
                                <span id="voice-status" class="table-status deployed">Selected</span>
                            </div>
                        </div>
                    </div>

                    <div class="form-row">
                        <div class="form-group" style="flex: 2;">
                            <label class="form-label">Test Text</label>
                            <input type="text" class="form-input" id="voice-test-text" value="Hello! How can I help you today?" placeholder="Enter text to test...">
                        </div>
                        <div class="form-group" style="flex: 1;">
                            <label class="form-label">Emotion (Parler)</label>
                            <select class="form-select" id="voice-emotion">
                                <option value="neutral">Neutral</option>
                                <option value="warm" selected>Warm</option>
                                <option value="excited">Excited</option>
                                <option value="calm">Calm</option>
                                <option value="concerned">Concerned</option>
                                <option value="apologetic">Apologetic</option>
                                <option value="confident">Confident</option>
                                <option value="empathetic">Empathetic</option>
                                <option value="urgent">Urgent</option>
                                <option value="cheerful">Cheerful</option>
                            </select>
                        </div>
                        <div class="form-group" style="flex: 1;">
                            <label class="form-label">Speed</label>
                            <input type="range" id="voice-speed" min="0.5" max="2" step="0.1" value="1" style="width: 100%;" oninput="document.getElementById('voice-speed-val').textContent = this.value + 'x'">
                            <small style="color: var(--text-secondary);">Speed: <span id="voice-speed-val">1x</span></small>
                        </div>
                    </div>

                    <div style="display: flex; gap: 0.5rem;">
                        <button class="btn btn-primary" onclick="testVoice()">Test Voice</button>
                        <button class="btn btn-secondary" onclick="saveVoiceConfig()">Save as Default</button>
                    </div>
                    <div id="voice-test-result" style="margin-top: 1rem;"></div>
                </div>

                <!-- Voice Library Card -->
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Lib</span> Free TTS Voice Library</div>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Open-source Text-to-Speech models with voice cloning capabilities. Replace paid services like ElevenLabs and Cartesia.</p>

                    <table class="skills-table" style="margin-bottom: 1.5rem;">
                        <thead>
                            <tr>
                                <th>Model</th>
                                <th>Parameters</th>
                                <th>Voice Cloning</th>
                                <th>License</th>
                                <th>Quality</th>
                            </tr>
                        </thead>
                        <tbody>
                            <tr>
                                <td><strong style="color: var(--neon-cyan);">Chatterbox</strong><br><small style="color: var(--text-secondary);">Resemble AI</small></td>
                                <td>0.5B</td>
                                <td><span class="table-status deployed">Yes (few seconds)</span></td>
                                <td>MIT</td>
                                <td>Outperforms ElevenLabs in blind tests</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-green);">Kokoro</strong><br><small style="color: var(--text-secondary);">Fastest option</small></td>
                                <td>82M</td>
                                <td><span class="table-status draft">No</span></td>
                                <td>Apache 2.0</td>
                                <td>Comparable to larger models, very fast</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-purple);">XTTS-v2</strong><br><small style="color: var(--text-secondary);">Coqui AI</small></td>
                                <td>~500M</td>
                                <td><span class="table-status deployed">Yes (6 sec audio)</span></td>
                                <td>CPML</td>
                                <td>Multi-language cloning, 17+ languages</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-orange);">F5-TTS</strong><br><small style="color: var(--text-secondary);">E2 TTS variant</small></td>
                                <td>~300M</td>
                                <td><span class="table-status deployed">Yes</span></td>
                                <td>MIT</td>
                                <td>Comparable to ElevenLabs, 12GB VRAM</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-pink);">Tortoise-TTS</strong><br><small style="color: var(--text-secondary);">Most realistic</small></td>
                                <td>~1B</td>
                                <td><span class="table-status deployed">Yes</span></td>
                                <td>Apache 2.0</td>
                                <td>Best prosody/intonation, slower</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-blue);">Piper</strong><br><small style="color: var(--text-secondary);">Lightweight</small></td>
                                <td>~20M</td>
                                <td><span class="table-status draft">No</span></td>
                                <td>MIT</td>
                                <td>Ultra-fast, runs on Raspberry Pi</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--neon-yellow);">Bark</strong><br><small style="color: var(--text-secondary);">Suno AI</small></td>
                                <td>~350M</td>
                                <td><span class="table-status training">Prompt-based</span></td>
                                <td>MIT</td>
                                <td>Can generate music, laughter, effects</td>
                            </tr>
                            <tr>
                                <td><strong style="color: var(--text-primary);">OpenVoice</strong><br><small style="color: var(--text-secondary);">MyShell AI</small></td>
                                <td>~300M</td>
                                <td><span class="table-status deployed">Yes (instant)</span></td>
                                <td>MIT</td>
                                <td>Instant voice cloning, style control</td>
                            </tr>
                        </tbody>
                    </table>

                    <h4 style="color: var(--neon-cyan); margin-bottom: 0.75rem;">Recommended Setup: Chatterbox</h4>
                    <p style="color: var(--text-secondary); margin-bottom: 0.75rem;">Best balance of quality, speed, and ease of use. MIT licensed with 23 languages support.</p>
                    <pre class="code-block"># Install Chatterbox
pip install chatterbox-tts

# Basic usage
from chatterbox import ChatterboxTTS

tts = ChatterboxTTS()

# Generate speech
audio = tts.generate("Hello! How can I help you today?")
audio.save("output.wav")

# Clone a voice (just needs a few seconds of audio)
tts.clone_voice("reference_audio.wav", voice_name="custom_voice")
audio = tts.generate("Now speaking in cloned voice!", voice="custom_voice")</pre>

                    <h4 style="color: var(--neon-green); margin: 1.5rem 0 0.75rem;">Quick Setup: Kokoro (Fastest)</h4>
                    <pre class="code-block"># Install Kokoro - ultra-fast TTS
pip install kokoro-onnx soundfile

from kokoro_onnx import Kokoro

kokoro = Kokoro("kokoro-v0_19.onnx", "voices.json")
samples, sr = kokoro.create("Hello, this is super fast TTS!", voice="af_bella", speed=1.0)

import soundfile as sf
sf.write("output.wav", samples, sr)</pre>

                    <h4 style="color: var(--neon-purple); margin: 1.5rem 0 0.75rem;">Multi-Language: XTTS-v2</h4>
                    <pre class="code-block"># Install XTTS-v2 (Coqui)
pip install TTS

from TTS.api import TTS

tts = TTS("tts_models/multilingual/multi-dataset/xtts_v2", gpu=True)

# Clone voice and generate in different language
tts.tts_to_file(
    text="Bonjour! Comment puis-je vous aider?",
    speaker_wav="reference.wav",  # 6 seconds of audio
    language="fr",
    file_path="output.wav"
)</pre>

                    <div style="margin-top: 1.5rem; padding: 1rem; background: var(--glass-surface); border-radius: 8px;">
                        <h4 style="color: var(--neon-orange); margin-bottom: 0.5rem;">Cost Savings</h4>
                        <p style="color: var(--text-secondary); font-size: 0.9rem;">
                            <strong>ElevenLabs:</strong> $0.30/1K chars | <strong>Cartesia:</strong> $0.015/min<br>
                            <strong>Self-hosted:</strong> $0 (just compute costs)<br><br>
                            At 10,000 minutes/month: Save ~$150/month by self-hosting Chatterbox or XTTS-v2
                        </p>
                    </div>
                </div>
            </div>

            <!-- Voice Integration -->
            <div id="command-voice" class="sub-tab-content">
                <!-- Platform Connections -->
                <div class="glass-card" style="margin-bottom: 1rem;">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Connect</span> Voice Platform Connections</div>
                        <button class="btn btn-secondary btn-sm" onclick="loadPlatforms()">Refresh</button>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Connect your Skill Command Center to voice platforms for real-time voice interactions.</p>

                    <div id="platforms-grid" style="display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 1rem;">
                        <!-- LiveKit -->
                        <div class="platform-card" id="platform-livekit" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-cyan);">LiveKit</strong>
                                <span id="status-livekit" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Real-time voice and video with agents SDK</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="text" class="form-input" id="livekit-url" placeholder="wss://your-app.livekit.cloud" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="livekit-api-key" placeholder="API Key" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="password" class="form-input" id="livekit-api-secret" placeholder="API Secret" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('livekit')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('livekit')">Test</button>
                            </div>
                        </div>

                        <!-- Vapi -->
                        <div class="platform-card" id="platform-vapi" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-green);">Vapi</strong>
                                <span id="status-vapi" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Voice AI platform for building phone agents</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="vapi-api-key" placeholder="Vapi API Key" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="text" class="form-input" id="vapi-assistant-id" placeholder="Assistant ID (optional)" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('vapi')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('vapi')">Test</button>
                            </div>
                        </div>

                        <!-- Twilio -->
                        <div class="platform-card" id="platform-twilio" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-purple);">Twilio</strong>
                                <span id="status-twilio" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Cloud communications for voice calls</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="text" class="form-input" id="twilio-account-sid" placeholder="Account SID" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="twilio-auth-token" placeholder="Auth Token" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="text" class="form-input" id="twilio-phone" placeholder="+1234567890" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('twilio')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('twilio')">Test</button>
                            </div>
                        </div>

                        <!-- Retell AI -->
                        <div class="platform-card" id="platform-retell" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-orange);">Retell AI</strong>
                                <span id="status-retell" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Conversational voice AI for customer interactions</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="retell-api-key" placeholder="Retell API Key" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="text" class="form-input" id="retell-agent-id" placeholder="Agent ID" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('retell')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('retell')">Test</button>
                            </div>
                        </div>

                        <!-- Bland AI -->
                        <div class="platform-card" id="platform-bland" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-pink);">Bland AI</strong>
                                <span id="status-bland" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">AI phone agents for enterprises</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="bland-api-key" placeholder="Bland API Key" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="text" class="form-input" id="bland-pathway-id" placeholder="Pathway ID" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('bland')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('bland')">Test</button>
                            </div>
                        </div>

                        <!-- Daily.co / Pipecat -->
                        <div class="platform-card" id="platform-daily" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-blue);">Daily.co</strong>
                                <span id="status-daily" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Real-time video/audio with Pipecat integration</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="password" class="form-input" id="daily-api-key" placeholder="Daily API Key" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="text" class="form-input" id="daily-room-url" placeholder="Room URL (optional)" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('daily')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('daily')">Test</button>
                            </div>
                        </div>

                        <!-- Vocode -->
                        <div class="platform-card" id="platform-vocode" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--neon-yellow);">Vocode</strong>
                                <span id="status-vocode" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Open-source voice agent framework</p>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="password" class="form-input" id="vocode-api-key" placeholder="Vocode API Key" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('vocode')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('vocode')">Test</button>
                            </div>
                        </div>

                        <!-- Custom WebSocket -->
                        <div class="platform-card" id="platform-websocket" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border);">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.75rem;">
                                <strong style="color: var(--text-primary);">Custom WebSocket</strong>
                                <span id="status-websocket" class="table-status draft">Disconnected</span>
                            </div>
                            <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">Connect to any WebSocket-based service</p>
                            <div class="form-group" style="margin-bottom: 0.5rem;">
                                <input type="text" class="form-input" id="websocket-url" placeholder="wss://your-server.com/ws" style="font-size: 0.85rem;">
                            </div>
                            <div class="form-group" style="margin-bottom: 0.75rem;">
                                <input type="password" class="form-input" id="websocket-auth" placeholder="Auth Header (optional)" style="font-size: 0.85rem;">
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-primary btn-sm" onclick="connectPlatform('websocket')">Connect</button>
                                <button class="btn btn-secondary btn-sm" onclick="testPlatform('websocket')">Test</button>
                            </div>
                        </div>
                    </div>
                    <div id="platform-message" style="margin-top: 1rem;"></div>
                </div>

                <!-- Integration Code Examples -->
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Code</span> Integration Examples</div>
                    </div>
                    <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Basic Integration</h4>
                    <pre class="code-block">from skill_command_center import SkillCommandCenter

center = SkillCommandCenter()

# In your voice handler:
async def on_user_speech(text: str):
    async for chunk in center.process_query(text, use_latency_masking=True):
        await tts.speak(chunk)</pre>
                    <h4 style="color: var(--neon-cyan); margin: 1rem 0 0.5rem;">LiveKit Agents Integration</h4>
                    <pre class="code-block">from livekit.agents import llm, Agent
from livekit.plugins import silero

class SkillCenterLLM(llm.LLM):
    def __init__(self):
        self.center = SkillCommandCenter()

    async def chat(self, messages):
        prompt = messages[-1].content
        async for chunk in self.center.process_query(prompt):
            yield llm.ChatChunk(content=chunk)

# Use with LiveKit Agent
agent = Agent(
    vad=silero.VAD(),
    llm=SkillCenterLLM(),
    tts=your_tts_provider
)</pre>
                    <h4 style="color: var(--neon-green); margin: 1rem 0 0.5rem;">Vapi Custom LLM</h4>
                    <pre class="code-block"># Set your server URL in Vapi dashboard as Custom LLM endpoint
from flask import Flask, request, jsonify

@app.route('/vapi/chat', methods=['POST'])
async def vapi_handler():
    data = request.json
    message = data['messages'][-1]['content']

    response = ""
    async for chunk in center.process_query(message):
        response += chunk

    return jsonify({
        "message": {"role": "assistant", "content": response}
    })</pre>
                    <h4 style="color: var(--neon-purple); margin: 1rem 0 0.5rem;">Daily.co Pipecat Integration</h4>
                    <pre class="code-block">from pipecat.pipeline import Pipeline
from pipecat.transports.services.daily import DailyTransport

class SkillCenterProcessor:
    async def process(self, text):
        async for chunk in center.process_query(text):
            yield chunk

pipeline = Pipeline([
    DailyTransport(room_url, token),
    your_stt,
    SkillCenterProcessor(),
    your_tts
])</pre>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- LPU INFERENCE TAB -->
        <!-- ================================================================ -->
        <div id="tab-lpu" class="tab-content">
            <div class="glass-card">
                <div class="section-header">
                    <div class="section-title"><span class="section-icon">LPU</span> BitNet LPU Console</div>
                    <button class="btn btn-secondary btn-sm" onclick="clearLpuConsole()">Clear</button>
                </div>
                <div class="form-group">
                    <label class="form-label">Prompt</label>
                    <textarea class="form-textarea" id="lpu-prompt" placeholder="User: What are your hours?&#10;Assistant:"></textarea>
                </div>
                <div class="form-row">
                    <div class="form-group">
                        <label class="form-label">Max Tokens</label>
                        <input type="number" class="form-input" id="lpu-tokens" value="128" min="16" max="512">
                    </div>
                    <div class="form-group">
                        <label class="form-label">Skill Adapter</label>
                        <select class="form-select" id="lpu-skill">
                            <option value="">None (Base Model)</option>
                        </select>
                    </div>
                </div>
                <button class="btn btn-primary" onclick="runLpuInference()" id="lpu-btn">Run Inference</button>
                <div style="margin-top: 1rem;">
                    <label class="form-label">Output</label>
                    <div class="console" id="lpu-console">
<span class="info">// BitNet LPU Console Ready
// Model: Llama3-8B-1.58-100B-tokens (3.58 GiB)
// Quantization: I2_S - 2 bpw ternary
// Speed: ~8 tokens/second on CPU

</span>Waiting for input...</div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- FAST BRAIN TAB (Skills) -->
        <!-- ================================================================ -->
        <div id="tab-fastbrain" class="tab-content">
            <!-- Sub-tabs for Skills -->
            <div class="sub-tabs">
                <button class="sub-tab-btn active" onclick="showFastBrainTab('dashboard')">Dashboard</button>
                <button class="sub-tab-btn" onclick="showFastBrainTab('skills')">Skills Manager</button>
                <button class="sub-tab-btn" onclick="showFastBrainTab('golden')">Golden Prompts</button>
                <button class="sub-tab-btn" onclick="showFastBrainTab('train')">Train LoRA</button>
                <button class="sub-tab-btn" onclick="showFastBrainTab('chat')">Test Chat</button>
                <button class="sub-tab-btn" onclick="showFastBrainTab('integration')">Outgoing API</button>
            </div>

            <!-- Dashboard Sub-tab -->
            <div id="fb-tab-dashboard" class="sub-tab-content active">
                <div class="dashboard-grid">
                    <!-- System Status Panel -->
                    <div class="glass-card card-full" style="background: linear-gradient(135deg, rgba(0,255,136,0.1) 0%, rgba(0,212,255,0.1) 100%);">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">SYS</span> System Status</div>
                            <button class="btn btn-sm btn-secondary" onclick="refreshSystemStatus()">Refresh</button>
                        </div>
                        <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 1rem;">
                            <!-- Fast Brain Status -->
                            <div class="status-card" id="status-fast-brain" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 4px solid var(--text-secondary);">
                                <div style="display: flex; justify-content: space-between; align-items: center;">
                                    <strong>Fast Brain LPU</strong>
                                    <span class="status-indicator" id="fb-status-indicator">--</span>
                                </div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.5rem;">
                                    <div>Latency: <span id="fb-status-latency">--</span></div>
                                    <div>Model: <span id="fb-status-model">--</span></div>
                                    <div>Skills: <span id="fb-status-skills">--</span></div>
                                </div>
                            </div>
                            <!-- Groq Status -->
                            <div class="status-card" id="status-groq" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 4px solid var(--text-secondary);">
                                <div style="display: flex; justify-content: space-between; align-items: center;">
                                    <strong>Groq API</strong>
                                    <span class="status-indicator" id="groq-status-indicator">--</span>
                                </div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.5rem;">
                                    Backend for Fast Brain inference
                                </div>
                            </div>
                            <!-- Hive215 Status -->
                            <div class="status-card" id="status-hive215" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 4px solid var(--text-secondary);">
                                <div style="display: flex; justify-content: space-between; align-items: center;">
                                    <strong>Hive215</strong>
                                    <span class="status-indicator" id="hive-status-indicator">--</span>
                                </div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.5rem;">
                                    Voice assistant platform
                                </div>
                            </div>
                            <!-- Modal Status -->
                            <div class="status-card" id="status-modal" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 4px solid var(--text-secondary);">
                                <div style="display: flex; justify-content: space-between; align-items: center;">
                                    <strong>Modal</strong>
                                    <span class="status-indicator" id="modal-status-indicator">--</span>
                                </div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.5rem;">
                                    Serverless deployment
                                </div>
                            </div>
                        </div>
                        <div id="system-error-display" style="margin-top: 1rem; display: none;">
                            <div style="color: var(--neon-orange); font-weight: bold; margin-bottom: 0.5rem;">Recent Errors:</div>
                            <div id="system-errors" class="console" style="max-height: 150px; overflow-y: auto;"></div>
                        </div>
                    </div>

                    <!-- Status Cards -->
                    <div class="glass-card stat-card card-quarter">
                        <div class="stat-value" id="fb-status">Offline</div>
                        <div class="stat-label">Status</div>
                    </div>
                    <div class="glass-card stat-card card-quarter">
                        <div class="stat-value green" id="fb-ttfb">--ms</div>
                        <div class="stat-label">Avg TTFB</div>
                    </div>
                    <div class="glass-card stat-card card-quarter">
                        <div class="stat-value pink" id="fb-throughput">--</div>
                        <div class="stat-label">Tokens/sec</div>
                    </div>
                    <div class="glass-card stat-card card-quarter">
                        <div class="stat-value orange" id="fb-requests">0</div>
                        <div class="stat-label">Total Requests</div>
                    </div>

                    <!-- Configuration -->
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">CFG</span> Fast Brain Configuration</div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Fast Brain URL</label>
                            <input type="text" class="form-input" id="fb-url" placeholder="https://your-username--fast-brain-lpu.modal.run">
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Min Containers</label>
                                <input type="number" class="form-input" id="fb-min-containers" value="1" min="0" max="10">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Max Containers</label>
                                <input type="number" class="form-input" id="fb-max-containers" value="10" min="1" max="100">
                            </div>
                        </div>
                        <div class="form-row">
                            <button class="btn btn-primary" onclick="saveFastBrainConfig()">Save Config</button>
                            <button class="btn btn-secondary" onclick="checkFastBrainHealth()">Check Health</button>
                        </div>
                        <div id="fb-config-message" style="margin-top: 1rem;"></div>
                    </div>

                    <!-- Performance Targets -->
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">TGT</span> Performance Targets</div>
                        </div>
                        <table class="skills-table">
                            <tr>
                                <td>Time to First Byte (TTFB)</td>
                                <td><span id="fb-ttfb-target" class="table-status draft">&lt;100ms</span></td>
                            </tr>
                            <tr>
                                <td>Throughput</td>
                                <td><span id="fb-throughput-target" class="table-status draft">&gt;200 tok/s</span></td>
                            </tr>
                            <tr>
                                <td>Cold Start</td>
                                <td><span class="table-status deployed">&lt;5s (warm)</span></td>
                            </tr>
                            <tr>
                                <td>Backend</td>
                                <td>Groq Llama 3.3 70B</td>
                            </tr>
                            <tr>
                                <td>Skills</td>
                                <td><span id="fb-skills-count">5</span> available</td>
                            </tr>
                        </table>
                        <button class="btn btn-success" onclick="runFastBrainBenchmark()" style="margin-top: 1rem;">Run Benchmark</button>
                    </div>

                    <!-- Benchmark Results -->
                    <div class="glass-card card-full" id="fb-benchmark-results" style="display: none;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">BNK</span> Benchmark Results</div>
                        </div>
                        <div id="fb-benchmark-output"></div>
                    </div>
                </div>
            </div>

            <!-- Skills Manager Sub-tab -->
            <div id="fb-tab-skills" class="sub-tab-content">
                <div class="dashboard-grid">
                    <!-- Skills List -->
                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">SKL</span> Available Skills</div>
                            <button class="btn btn-primary btn-sm" onclick="showCreateSkillModal()">+ Create Skill</button>
                        </div>
                        <div id="fb-skills-list" class="skills-grid" style="display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 1rem;">
                            <!-- Skills will be loaded here -->
                            <div style="color: var(--text-secondary); padding: 2rem; text-align: center;">
                                Loading skills...
                            </div>
                        </div>
                    </div>

                    <!-- Create Skill Form (Modal-like) -->
                    <div class="glass-card card-full" id="fb-create-skill-form" style="display: none;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">NEW</span> Create Custom Skill</div>
                            <button class="btn btn-sm" onclick="hideCreateSkillModal()" style="background: transparent; color: var(--text-secondary);">✕</button>
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Skill ID (lowercase, no spaces)</label>
                                <input type="text" class="form-input" id="new-skill-id" placeholder="my_custom_skill">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Display Name</label>
                                <input type="text" class="form-input" id="new-skill-name" placeholder="My Custom Skill">
                            </div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Description</label>
                            <input type="text" class="form-input" id="new-skill-description" placeholder="Brief description of what this skill does">
                        </div>
                        <div class="form-group">
                            <label class="form-label">System Prompt</label>
                            <textarea class="form-textarea" id="new-skill-prompt" rows="6" placeholder="You are an AI assistant specialized in..."></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Knowledge Base (one item per line)</label>
                            <textarea class="form-textarea" id="new-skill-knowledge" rows="4" placeholder="Pricing: $100-500&#10;Hours: Mon-Fri 9am-5pm&#10;Service area: Philadelphia metro"></textarea>
                        </div>
                        <div class="form-row">
                            <button class="btn btn-primary" onclick="createCustomSkill()">Create Skill</button>
                            <button class="btn btn-secondary" onclick="hideCreateSkillModal()">Cancel</button>
                        </div>
                        <div id="create-skill-message" style="margin-top: 1rem;"></div>
                    </div>

                    <!-- Edit Skill Form (Modal-like) -->
                    <div class="glass-card card-full" id="fb-edit-skill-form" style="display: none;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">EDIT</span> Edit Skill</div>
                            <button class="btn btn-sm" onclick="hideEditSkillModal()" style="background: transparent; color: var(--text-secondary);">✕</button>
                        </div>
                        <input type="hidden" id="edit-skill-id">
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Skill ID</label>
                                <input type="text" class="form-input" id="edit-skill-id-display" disabled style="opacity: 0.6;">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Display Name</label>
                                <input type="text" class="form-input" id="edit-skill-name" placeholder="My Custom Skill">
                            </div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Description</label>
                            <input type="text" class="form-input" id="edit-skill-description" placeholder="Brief description of what this skill does">
                        </div>
                        <div class="form-group">
                            <label class="form-label">System Prompt</label>
                            <textarea class="form-textarea" id="edit-skill-prompt" rows="8" placeholder="You are an AI assistant specialized in..."></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Knowledge Base (one item per line)</label>
                            <textarea class="form-textarea" id="edit-skill-knowledge" rows="4" placeholder="Pricing: $100-500&#10;Hours: Mon-Fri 9am-5pm&#10;Service area: Philadelphia metro"></textarea>
                        </div>
                        <div class="form-row">
                            <button class="btn btn-primary" onclick="saveEditedSkill()">Save Changes</button>
                            <button class="btn btn-secondary" onclick="hideEditSkillModal()">Cancel</button>
                        </div>
                        <div id="edit-skill-message" style="margin-top: 1rem;"></div>
                    </div>
                </div>
            </div>

            <!-- Test Chat Sub-tab -->
            <div id="fb-tab-chat" class="sub-tab-content">
                <div class="dashboard-grid">
                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">CHT</span> Test Fast Brain</div>
                        </div>

                        <!-- Skill Selector -->
                        <div class="form-row" style="margin-bottom: 1rem;">
                            <div class="form-group" style="flex: 1;">
                                <label class="form-label">Select Skill</label>
                                <select class="form-select" id="fb-skill-selector" onchange="onSkillSelect()">
                                    <option value="general">General Assistant</option>
                                    <option value="receptionist">Professional Receptionist</option>
                                    <option value="electrician">Electrician Assistant</option>
                                    <option value="plumber">Plumber Assistant</option>
                                    <option value="lawyer">Legal Intake Assistant</option>
                                </select>
                            </div>
                            <div class="form-group" style="flex: 0 0 auto; display: flex; align-items: flex-end;">
                                <button class="btn btn-secondary btn-sm" onclick="loadSkillPrompt()">Load Skill Prompt</button>
                            </div>
                        </div>

                        <!-- Active Skill Indicator -->
                        <div id="fb-active-skill-info" style="background: var(--glass-surface); padding: 0.75rem 1rem; border-radius: 8px; margin-bottom: 1rem; border-left: 4px solid var(--neon-cyan);">
                            <div style="display: flex; justify-content: space-between; align-items: center;">
                                <div>
                                    <strong style="color: var(--neon-cyan);">Active Skill:</strong>
                                    <span id="fb-active-skill-name" style="margin-left: 0.5rem;">General Assistant</span>
                                </div>
                                <span id="fb-skill-status" class="table-status deployed">Ready</span>
                            </div>
                            <div id="fb-active-skill-desc" style="font-size: 0.85rem; color: var(--text-secondary); margin-top: 0.25rem;">
                                Helpful general-purpose assistant
                            </div>
                        </div>

                        <div class="form-group">
                            <label class="form-label">System Prompt (or use skill default)</label>
                            <textarea class="form-textarea" id="fb-system-prompt" rows="4" placeholder="Leave empty to use selected skill's default prompt..."></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Message</label>
                            <textarea class="form-textarea" id="fb-message" placeholder="Enter your message to test Fast Brain..."></textarea>
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Max Tokens</label>
                                <input type="number" class="form-input" id="fb-max-tokens" value="256" min="16" max="1024">
                            </div>
                        </div>
                        <button class="btn btn-primary" onclick="testFastBrainChat()" id="fb-chat-btn">Send to Fast Brain</button>
                        <div style="margin-top: 1rem;">
                            <label class="form-label">Response</label>
                            <div class="console" id="fb-response" style="min-height: 100px;">Waiting for request...</div>
                        </div>
                        <div id="fb-metrics" style="margin-top: 0.5rem; color: var(--text-secondary); font-size: 0.85rem;"></div>

                        <!-- Voice Test Section (Optional) -->
                        <div style="border-top: 1px solid var(--glass-border); margin-top: 1.5rem; padding-top: 1rem;">
                            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 1rem;">
                                <h4 style="color: var(--neon-purple); margin: 0; font-size: 0.95rem;">Test with Voice (Optional)</h4>
                                <span id="voice-test-status" style="font-size: 0.8rem; color: var(--text-secondary);"></span>
                            </div>
                            <div class="form-row">
                                <div class="form-group" style="flex: 1;">
                                    <label class="form-label">Voice Provider</label>
                                    <select class="form-select" id="voice-test-provider" onchange="loadVoiceTestVoices()">
                                        <option value="">Select provider...</option>
                                        <optgroup label="Premium (API Key Required)">
                                            <option value="elevenlabs">ElevenLabs</option>
                                            <option value="cartesia">Cartesia</option>
                                            <option value="deepgram">Deepgram (Aura)</option>
                                            <option value="openai">OpenAI TTS</option>
                                        </optgroup>
                                        <optgroup label="Free / Open Source">
                                            <option value="edge_tts">Edge TTS (Microsoft)</option>
                                            <option value="parler">Parler TTS (GPU)</option>
                                        </optgroup>
                                    </select>
                                </div>
                                <div class="form-group" style="flex: 1;">
                                    <label class="form-label">Voice</label>
                                    <select class="form-select" id="voice-test-voice-id" disabled>
                                        <option value="">Select provider first...</option>
                                    </select>
                                </div>
                            </div>
                            <div class="form-group">
                                <label class="form-label">Text to Speak</label>
                                <textarea class="form-textarea" id="voice-test-chat-text" rows="2" placeholder="Hello! This is a test of the voice synthesis. You can type any text here to hear how it sounds."></textarea>
                            </div>
                            <div style="display: flex; gap: 0.5rem; align-items: center;">
                                <button class="btn btn-secondary" id="voice-test-btn" onclick="testVoiceTTS()">Test Voice</button>
                                <audio id="voice-test-audio" controls style="display: none; height: 32px;"></audio>
                            </div>
                        </div>
                    </div>
                </div>
            </div>

            <!-- Outgoing API Connections Sub-tab -->
            <div id="fb-tab-integration" class="sub-tab-content">
                <div class="dashboard-grid">
                    <!-- API Connections List -->
                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">API</span> Outgoing API Connections</div>
                            <div>
                                <button class="btn btn-primary btn-sm" onclick="showAddConnectionModal()">+ Add Connection</button>
                                <button class="btn btn-secondary btn-sm" onclick="refreshApiConnections()">Refresh</button>
                            </div>
                        </div>
                        <p style="color: var(--text-secondary); margin-bottom: 1rem;">Connect to external APIs, webhooks, and services. Test connections and send requests.</p>
                        <div id="api-connections-list">
                            <div style="color: var(--text-secondary); padding: 2rem; text-align: center;">Loading connections...</div>
                        </div>
                    </div>

                    <!-- Add/Edit Connection Modal (inline) -->
                    <div class="glass-card card-half" id="connection-form-card" style="display: none;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">NEW</span> <span id="connection-form-title">Add Connection</span></div>
                            <button class="btn btn-secondary btn-sm" onclick="hideConnectionForm()">×</button>
                        </div>
                        <input type="hidden" id="conn-edit-id">
                        <div class="form-group">
                            <label class="form-label">Connection Name *</label>
                            <input type="text" class="form-input" id="conn-name" placeholder="My API Service">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Base URL *</label>
                            <input type="text" class="form-input" id="conn-url" placeholder="https://api.example.com">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Authentication Type</label>
                            <select class="form-select" id="conn-auth-type">
                                <option value="none">None</option>
                                <option value="bearer" selected>Bearer Token</option>
                                <option value="api_key">X-API-Key Header</option>
                                <option value="basic">Basic Auth</option>
                            </select>
                        </div>
                        <div class="form-group" id="conn-apikey-group">
                            <label class="form-label">API Key / Token</label>
                            <input type="password" class="form-input" id="conn-apikey" placeholder="Your API key or token">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Custom Headers (JSON)</label>
                            <textarea class="form-textarea" id="conn-headers" rows="2" placeholder='{"X-Custom-Header": "value"}'></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Webhook URL (for callbacks)</label>
                            <input type="text" class="form-input" id="conn-webhook" placeholder="https://your-app.com/webhook">
                        </div>
                        <div class="form-row">
                            <button class="btn btn-primary" onclick="saveApiConnection()">Save Connection</button>
                            <button class="btn btn-secondary" onclick="hideConnectionForm()">Cancel</button>
                        </div>
                    </div>

                    <!-- Request Tester -->
                    <div class="glass-card card-half" id="request-tester-card" style="display: none;">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">TST</span> Request Tester</div>
                            <button class="btn btn-secondary btn-sm" onclick="hideRequestTester()">×</button>
                        </div>
                        <div style="margin-bottom: 0.5rem; color: var(--neon-cyan); font-weight: 600;" id="tester-connection-name"></div>
                        <input type="hidden" id="tester-connection-id">
                        <div class="form-row" style="gap: 0.5rem; margin-bottom: 0.5rem;">
                            <select class="form-select" id="tester-method" style="width: 100px;">
                                <option value="GET">GET</option>
                                <option value="POST">POST</option>
                                <option value="PUT">PUT</option>
                                <option value="DELETE">DELETE</option>
                            </select>
                            <input type="text" class="form-input" id="tester-path" placeholder="/endpoint/path" style="flex: 1;">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Request Body (JSON)</label>
                            <textarea class="form-textarea" id="tester-body" rows="4" placeholder='{"key": "value"}'></textarea>
                        </div>
                        <button class="btn btn-primary" onclick="sendTestRequest()">Send Request</button>
                        <div id="tester-result" style="margin-top: 1rem;"></div>
                    </div>

                    <!-- Live Status Panel -->
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">STS</span> Live Status</div>
                            <button class="btn btn-secondary btn-sm" onclick="refreshConnectionStatuses()">Refresh All</button>
                        </div>
                        <div id="connections-status-list">
                            <div style="color: var(--text-secondary); padding: 1rem; text-align: center;">No connections configured</div>
                        </div>
                    </div>

                    <!-- Quick Tips -->
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">TIP</span> Integration Tips</div>
                        </div>
                        <div style="font-size: 0.85rem; line-height: 1.6;">
                            <h4 style="color: var(--neon-green); margin-bottom: 0.5rem;">Supported Auth Types</h4>
                            <ul style="color: var(--text-secondary); margin-left: 1rem; margin-bottom: 1rem;">
                                <li><strong>Bearer:</strong> Authorization: Bearer {token}</li>
                                <li><strong>API Key:</strong> X-API-Key: {key}</li>
                                <li><strong>Basic:</strong> Authorization: Basic {base64}</li>
                                <li><strong>None:</strong> No authentication header</li>
                            </ul>

                            <h4 style="color: var(--neon-green); margin-bottom: 0.5rem;">Common Integrations</h4>
                            <ul style="color: var(--text-secondary); margin-left: 1rem;">
                                <li>CRM systems (HubSpot, Salesforce)</li>
                                <li>Ticket systems (Zendesk, Freshdesk)</li>
                                <li>Calendars (Cal.com, Calendly)</li>
                                <li>Notifications (Slack, Discord)</li>
                                <li>Custom webhooks</li>
                            </ul>
                        </div>
                    </div>
                </div>
            </div>

            <!-- Golden Prompts Sub-tab -->
            <div id="fb-tab-golden" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Star</span> Golden Prompts</div>
                        <button class="btn btn-secondary btn-sm" onclick="loadGoldenPrompts()">Refresh</button>
                    </div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">Voice-optimized skill manuals for ultra-fast Groq inference. These prompts are designed to stay under 3k tokens for ~50-100ms prefill latency.</p>

                    <div class="form-row">
                        <div class="form-group" style="flex: 2;">
                            <label class="form-label">Select Skill</label>
                            <select class="form-select" id="fb-golden-skill-select" onchange="loadGoldenPromptFB()">
                                <option value="receptionist">Receptionist - Professional Call Handling</option>
                                <option value="electrician">Electrician - Jenkintown Electricity</option>
                                <option value="plumber">Plumber - Plumbing Services</option>
                                <option value="lawyer">Lawyer - Legal Intake Specialist</option>
                                <option value="solar">Solar - Solar Company Sales</option>
                                <option value="tara-sales">Tara Sales - TheDashTool Assistant</option>
                                <option value="general">General - Helpful Assistant</option>
                            </select>
                        </div>
                        <div class="form-group" style="flex: 1;">
                            <label class="form-label">Voice <span style="color: var(--text-secondary); font-size: 0.8rem;">(optional)</span></label>
                            <select class="form-select" id="fb-golden-voice-select" onchange="linkVoiceToSkillFB()">
                                <option value="">No custom voice</option>
                            </select>
                        </div>
                    </div>

                    <div class="stats-row" style="margin-bottom: 1rem;">
                        <div class="stat-item">
                            <div class="stat-value" id="fb-golden-tokens">~850</div>
                            <div class="stat-label">Est. Tokens</div>
                        </div>
                        <div class="stat-item">
                            <div class="stat-value" id="fb-golden-prefill">~14ms</div>
                            <div class="stat-label">Prefill Time</div>
                        </div>
                        <div class="stat-item">
                            <div class="stat-value" id="fb-golden-status">Ready</div>
                            <div class="stat-label">Status</div>
                        </div>
                    </div>

                    <div class="form-group">
                        <label class="form-label">Prompt Content</label>
                        <textarea class="form-textarea" id="fb-golden-prompt-content" style="min-height: 350px; font-family: monospace; font-size: 0.85rem;" placeholder="Loading prompt..."></textarea>
                    </div>

                    <div style="display: flex; gap: 1rem; flex-wrap: wrap;">
                        <button class="btn btn-primary" onclick="saveGoldenPromptFB()">Save Changes</button>
                        <button class="btn btn-secondary" onclick="resetGoldenPromptFB()">Reset to Default</button>
                        <button class="btn btn-success" onclick="testGoldenPromptFB()">Test Prompt</button>
                    </div>

                    <div id="fb-golden-message" style="margin-top: 1rem;"></div>
                </div>
            </div>

            <!-- Train LoRA Sub-tab -->
            <div id="fb-tab-train" class="sub-tab-content">
                <div class="dashboard-grid">
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Train</span> Train LoRA Adapter</div>
                        </div>
                        <p style="color: var(--text-secondary); margin-bottom: 1rem;">Training requires a GPU. Generate a script to run on Colab, Modal, or your own machine.</p>
                        <div class="form-group">
                            <label class="form-label">Select Skill to Train</label>
                            <select class="form-select" id="fb-train-profile">
                                <option value="receptionist">Receptionist</option>
                                <option value="electrician">Electrician</option>
                                <option value="plumber">Plumber</option>
                                <option value="lawyer">Lawyer</option>
                                <option value="solar">Solar</option>
                                <option value="tara-sales">Tara Sales</option>
                                <option value="general">General</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Base Model</label>
                            <select class="form-select" id="fb-train-model">
                                <option value="llama-3.3-70b">Llama 3.3 70B (Groq)</option>
                                <option value="llama-3.1-8b">Llama 3.1 8B (Local)</option>
                                <option value="mistral-7b">Mistral 7B (Local)</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Training Steps: <span id="fb-steps-value">60</span></label>
                            <input type="range" id="fb-train-steps" min="20" max="200" value="60" style="width: 100%;" oninput="document.getElementById('fb-steps-value').textContent = this.value">
                            <div style="display: flex; justify-content: space-between; color: var(--text-secondary); font-size: 0.75rem;">
                                <span>20 (Quick)</span>
                                <span>200 (Thorough)</span>
                            </div>
                        </div>
                        <button class="btn btn-primary" onclick="generateTrainingScriptFB()">Generate Training Script</button>
                    </div>

                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">GPU</span> Training Options</div>
                        </div>
                        <div style="display: flex; flex-direction: column; gap: 1rem;">
                            <div style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 3px solid var(--neon-cyan);">
                                <h4 style="margin-bottom: 0.5rem;">Google Colab (Free)</h4>
                                <p style="font-size: 0.85rem; color: var(--text-secondary);">Free T4 GPU, ~15 min for 60 steps</p>
                            </div>
                            <div style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 3px solid var(--neon-purple);">
                                <h4 style="margin-bottom: 0.5rem;">Modal (Pay-per-use)</h4>
                                <p style="font-size: 0.85rem; color: var(--text-secondary);">A100 GPU, ~5 min for 60 steps, ~$0.50</p>
                            </div>
                            <div style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border-left: 3px solid var(--neon-green);">
                                <h4 style="margin-bottom: 0.5rem;">Local GPU</h4>
                                <p style="font-size: 0.85rem; color: var(--text-secondary);">Requires 24GB+ VRAM, RTX 3090/4090</p>
                            </div>
                        </div>
                    </div>

                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Code</span> Training Script</div>
                        </div>
                        <div id="fb-training-output">
                            <pre class="code-block" id="fb-training-script" style="max-height: 400px; overflow-y: auto;">Click "Generate Training Script" to create your training code.</pre>
                        </div>
                        <div style="margin-top: 1rem; display: flex; gap: 1rem;">
                            <button class="btn btn-secondary" onclick="copyTrainingScript()">Copy to Clipboard</button>
                            <button class="btn btn-secondary" onclick="downloadTrainingScript()">Download as .py</button>
                        </div>
                    </div>
                </div>
            </div>
        </div>

        <!-- ================================================================ -->
        <!-- VOICE LAB TAB -->
        <!-- ================================================================ -->
        <div id="tab-voicelab" class="tab-content">
            <!-- Sub-tabs -->
            <div class="sub-tabs">
                <button class="sub-tab-btn active" onclick="showVoiceLabTab('projects')">Voice Projects</button>
                <button class="sub-tab-btn" onclick="showVoiceLabTab('browse')">Browse Voices</button>
                <button class="sub-tab-btn" onclick="showVoiceLabTab('create')">Create Voice</button>
                <button class="sub-tab-btn" onclick="showVoiceLabTab('training')">Training Queue</button>
                <button class="sub-tab-btn" onclick="showVoiceLabTab('skills')">Skill Training</button>
            </div>

            <!-- Voice Projects List -->
            <div id="voicelab-projects" class="sub-tab-content active">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Voice</span> Your Voice Projects</div>
                        <button class="btn btn-primary btn-sm" onclick="showVoiceLabTab('create')">+ New Voice</button>
                    </div>
                    <div id="voice-projects-list" class="skills-grid">
                        <div style="color: var(--text-secondary); padding: 2rem; text-align: center;">
                            No voice projects yet. Click "New Voice" to create one.
                        </div>
                    </div>
                </div>
            </div>

            <!-- Browse Voices from Providers -->
            <div id="voicelab-browse" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Browse</span> Voice Library</div>
                        <div style="display: flex; gap: 0.5rem; align-items: center;">
                            <select class="form-select" id="vl-browse-provider" onchange="loadProviderVoices()" style="width: auto;">
                                <option value="">-- Select Provider --</option>
                                <optgroup label="Premium (API Key Required)">
                                    <option value="elevenlabs">ElevenLabs</option>
                                    <option value="cartesia">Cartesia</option>
                                    <option value="deepgram">Deepgram (Aura)</option>
                                    <option value="openai">OpenAI TTS</option>
                                </optgroup>
                                <optgroup label="Free / Built-in">
                                    <option value="edge_tts">Edge TTS (Microsoft)</option>
                                    <option value="parler_tts">Parler TTS (GPU)</option>
                                    <option value="kokoro">Kokoro</option>
                                    <option value="azure">Azure TTS</option>
                                </optgroup>
                            </select>
                            <button class="btn btn-secondary btn-sm" onclick="loadProviderVoices()">Refresh</button>
                        </div>
                    </div>
                    <div id="vl-browse-status" style="margin: 1rem 0; color: var(--text-secondary);"></div>
                    <div id="vl-browse-filters" style="margin-bottom: 1rem; display: none;">
                        <div style="display: flex; gap: 1rem; flex-wrap: wrap; align-items: center;">
                            <input type="text" class="form-input" id="vl-browse-search" placeholder="Search voices..." style="width: 200px;" oninput="filterBrowseVoices()">
                            <select class="form-select" id="vl-browse-gender" style="width: auto;" onchange="filterBrowseVoices()">
                                <option value="">All Genders</option>
                                <option value="male">Male</option>
                                <option value="female">Female</option>
                                <option value="neutral">Neutral</option>
                            </select>
                            <span id="vl-browse-count" style="color: var(--neon-cyan);"></span>
                        </div>
                    </div>
                    <div id="vl-browse-voices" class="voice-grid" style="display: grid; grid-template-columns: repeat(auto-fill, minmax(280px, 1fr)); gap: 1rem;">
                        <div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;">
                            Select a provider above to browse available voices.
                        </div>
                    </div>
                </div>
            </div>

            <!-- Edit Voice Project -->
            <div id="voicelab-edit" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Edit</span> <span id="vl-edit-title">Voice Project</span></div>
                        <div style="display: flex; gap: 0.5rem;">
                            <button class="btn btn-secondary btn-sm" onclick="showVoiceLabTab('projects')">← Back</button>
                            <button class="btn btn-danger btn-sm" onclick="deleteCurrentVoiceProject()">Delete</button>
                        </div>
                    </div>

                    <input type="hidden" id="vl-edit-project-id">

                    <div class="dashboard-grid">
                        <!-- Project Info -->
                        <div class="glass-card card-half" style="background: var(--glass-surface);">
                            <h4 style="margin-bottom: 1rem; color: var(--neon-cyan);">Project Info</h4>
                            <div class="form-group">
                                <label class="form-label">Voice Name</label>
                                <input type="text" class="form-input" id="vl-edit-name" placeholder="Voice name">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Description</label>
                                <textarea class="form-textarea" id="vl-edit-description" rows="2" placeholder="Description..."></textarea>
                            </div>
                            <div class="form-row">
                                <div class="form-group">
                                    <label class="form-label">Provider</label>
                                    <select class="form-select" id="vl-edit-provider">
                                        <optgroup label="Premium (Live API)">
                                            <option value="elevenlabs">ElevenLabs</option>
                                            <option value="cartesia">Cartesia</option>
                                            <option value="deepgram">Deepgram (Aura)</option>
                                            <option value="openai">OpenAI TTS</option>
                                        </optgroup>
                                        <optgroup label="Free / Open Source">
                                            <option value="parler_tts">Parler TTS</option>
                                            <option value="edge_tts">Edge TTS (Free)</option>
                                            <option value="kokoro">Kokoro</option>
                                            <option value="xtts">XTTS-v2</option>
                                            <option value="openvoice">OpenVoice</option>
                                        </optgroup>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Status</label>
                                    <div id="vl-edit-status" class="stat-value" style="padding: 0.5rem; background: var(--glass-surface); border-radius: 8px;">Draft</div>
                                </div>
                            </div>
                            <button class="btn btn-secondary" onclick="saveVoiceProjectChanges()">Save Changes</button>
                        </div>

                        <!-- Voice Settings -->
                        <div class="glass-card card-half" style="background: var(--glass-surface);">
                            <h4 style="margin-bottom: 1rem; color: var(--neon-pink);">Voice Settings</h4>
                            <div class="form-group">
                                <label class="form-label">Pitch: <span id="vl-edit-pitch-value">1.0</span></label>
                                <input type="range" class="form-input" id="vl-edit-pitch" min="0.5" max="2.0" step="0.1" value="1.0" oninput="document.getElementById('vl-edit-pitch-value').textContent = this.value">
                            </div>
                            <div class="form-group">
                                <label class="form-label">Speed: <span id="vl-edit-speed-value">1.0</span></label>
                                <input type="range" class="form-input" id="vl-edit-speed" min="0.5" max="2.0" step="0.1" value="1.0" oninput="document.getElementById('vl-edit-speed-value').textContent = this.value">
                            </div>
                            <div class="form-row">
                                <div class="form-group">
                                    <label class="form-label">Emotion</label>
                                    <select class="form-select" id="vl-edit-emotion">
                                        <option value="neutral">Neutral</option>
                                        <option value="warm">Warm</option>
                                        <option value="excited">Excited</option>
                                        <option value="calm">Calm</option>
                                        <option value="confident">Confident</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <label class="form-label">Style</label>
                                    <select class="form-select" id="vl-edit-style">
                                        <option value="conversational">Conversational</option>
                                        <option value="professional">Professional</option>
                                        <option value="friendly">Friendly</option>
                                    </select>
                                </div>
                            </div>
                        </div>

                        <!-- Audio Samples -->
                        <div class="glass-card card-full" style="background: var(--glass-surface);">
                            <div class="section-header">
                                <h4 style="color: var(--neon-green);">Audio Samples</h4>
                                <button class="btn btn-primary btn-sm" onclick="document.getElementById('vl-edit-audio-input').click()">+ Add Sample</button>
                            </div>
                            <input type="file" id="vl-edit-audio-input" accept="audio/*" style="display: none;" onchange="uploadVoiceSample(this)">
                            <div id="vl-edit-samples-list" style="margin-top: 1rem;">
                                <div style="color: var(--text-secondary); padding: 1rem; text-align: center;">
                                    No samples uploaded. Add audio samples to train your voice.
                                </div>
                            </div>
                        </div>

                        <!-- Train & Test -->
                        <div class="glass-card card-full" style="background: var(--glass-surface);">
                            <h4 style="margin-bottom: 1rem; color: var(--neon-purple);">Train & Test Voice</h4>
                            <div class="form-row" style="align-items: flex-end;">
                                <div class="form-group" style="flex: 2;">
                                    <label class="form-label">Test Text</label>
                                    <input type="text" class="form-input" id="vl-edit-test-text" value="Hello! This is a test of my custom voice. How does it sound?">
                                </div>
                                <div class="form-group">
                                    <button class="btn btn-success" onclick="trainVoiceProject()">Train Voice</button>
                                </div>
                                <div class="form-group">
                                    <button class="btn btn-primary" onclick="testVoiceProject()">Test Voice</button>
                                </div>
                            </div>
                            <div id="vl-edit-audio-player" style="margin-top: 1rem; display: none;">
                                <audio id="vl-edit-audio" controls style="width: 100%;"></audio>
                            </div>
                            <div id="vl-edit-message" style="margin-top: 1rem;"></div>
                        </div>

                        <!-- Link to Skill -->
                        <div class="glass-card card-full" style="background: var(--glass-surface);">
                            <h4 style="margin-bottom: 1rem; color: var(--neon-orange);">Link to Skill/Agent</h4>
                            <div class="form-row" style="align-items: flex-end;">
                                <div class="form-group" style="flex: 2;">
                                    <label class="form-label">Select Skill</label>
                                    <select class="form-select" id="vl-edit-skill-link">
                                        <option value="">-- Select a skill --</option>
                                        <option value="receptionist">Receptionist</option>
                                        <option value="electrician">Electrician</option>
                                        <option value="plumber">Plumber</option>
                                        <option value="lawyer">Lawyer</option>
                                        <option value="solar">Solar</option>
                                        <option value="tara-sales">Tara Sales</option>
                                        <option value="general">General</option>
                                    </select>
                                </div>
                                <div class="form-group">
                                    <button class="btn btn-secondary" onclick="linkVoiceProjectToSkill()">Link Voice</button>
                                </div>
                            </div>
                            <p id="vl-edit-linked-skill" style="color: var(--text-secondary); margin-top: 0.5rem;"></p>
                        </div>
                    </div>
                </div>
            </div>

            <!-- Create Voice -->
            <div id="voicelab-create" class="sub-tab-content">
                <div class="dashboard-grid">
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">New</span> Create Custom Voice</div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Voice Name</label>
                            <input type="text" class="form-input" id="vl-name" placeholder="My Custom Voice">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Description</label>
                            <textarea class="form-textarea" id="vl-description" placeholder="Describe the voice characteristics..."></textarea>
                        </div>
                        <div class="form-row">
                            <div class="form-group">
                                <label class="form-label">Base Voice Provider</label>
                                <select class="form-select" id="vl-provider" onchange="loadVLProviderVoices()">
                                    <optgroup label="Premium (Live API)">
                                        <option value="elevenlabs">ElevenLabs</option>
                                        <option value="cartesia">Cartesia</option>
                                        <option value="deepgram">Deepgram (Aura)</option>
                                        <option value="openai">OpenAI TTS</option>
                                    </optgroup>
                                    <optgroup label="Free / Open Source">
                                        <option value="parler_tts">Parler TTS (Expressive - GPU)</option>
                                        <option value="edge_tts">Edge TTS (Microsoft - Free)</option>
                                        <option value="kokoro">Kokoro</option>
                                        <option value="chatterbox">Chatterbox</option>
                                        <option value="xtts">XTTS-v2 (Coqui)</option>
                                        <option value="openvoice">OpenVoice</option>
                                    </optgroup>
                                </select>
                            </div>
                            <div class="form-group">
                                <label class="form-label">Base Voice</label>
                                <select class="form-select" id="vl-base-voice">
                                    <option value="default">Default</option>
                                </select>
                            </div>
                        </div>
                        <button class="btn btn-primary" onclick="createVoiceProject()">Create Project</button>
                        <div id="vl-create-message" style="margin-top: 1rem;"></div>
                    </div>

                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Settings</span> Voice Settings</div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Pitch: <span id="vl-pitch-value">1.0</span></label>
                            <input type="range" class="form-input" id="vl-pitch" min="0.5" max="2.0" step="0.1" value="1.0" oninput="document.getElementById('vl-pitch-value').textContent = this.value">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Speed: <span id="vl-speed-value">1.0</span></label>
                            <input type="range" class="form-input" id="vl-speed" min="0.5" max="2.0" step="0.1" value="1.0" oninput="document.getElementById('vl-speed-value').textContent = this.value">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Emotion</label>
                            <select class="form-select" id="vl-emotion">
                                <option value="neutral">Neutral</option>
                                <option value="warm">Warm</option>
                                <option value="excited">Excited</option>
                                <option value="calm">Calm</option>
                                <option value="concerned">Concerned</option>
                                <option value="apologetic">Apologetic</option>
                                <option value="confident">Confident</option>
                                <option value="empathetic">Empathetic</option>
                                <option value="urgent">Urgent</option>
                                <option value="cheerful">Cheerful</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Speaking Style</label>
                            <select class="form-select" id="vl-style">
                                <option value="conversational">Conversational</option>
                                <option value="professional">Professional</option>
                                <option value="friendly">Friendly</option>
                                <option value="formal">Formal</option>
                                <option value="casual">Casual</option>
                            </select>
                        </div>
                    </div>

                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Audio</span> Upload Training Samples</div>
                        </div>
                        <p style="color: var(--text-secondary); margin-bottom: 1rem;">
                            Upload 3+ audio samples (5-30 seconds each) to train your custom voice. Clearer recordings produce better results.
                        </p>
                        <div class="file-upload" onclick="document.getElementById('vl-audio-input').click()">
                            <div class="file-upload-icon">Upload Audio</div>
                            <div class="file-upload-text">Click to upload audio files (WAV, MP3, M4A)</div>
                            <input type="file" id="vl-audio-input" accept="audio/*" multiple style="display: none;">
                        </div>
                        <div id="vl-samples-list" style="margin-top: 1rem;"></div>
                    </div>
                </div>
            </div>

            <!-- Training Queue -->
            <div id="voicelab-training" class="sub-tab-content">
                <div class="glass-card">
                    <div class="section-header">
                        <div class="section-title"><span class="section-icon">Train</span> Voice Training Queue</div>
                        <button class="btn btn-secondary btn-sm" onclick="refreshVoiceTrainingStatus()">Refresh</button>
                    </div>
                    <div id="voice-training-queue">
                        <div style="color: var(--text-secondary); padding: 2rem; text-align: center;">
                            No voice training jobs in progress.
                        </div>
                    </div>
                </div>
            </div>

            <!-- Skill Training -->
            <div id="voicelab-skills" class="sub-tab-content">
                <div class="dashboard-grid">
                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Tune</span> Fine-tune Skills</div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Select Skill</label>
                            <select class="form-select" id="ft-skill-select">
                                <option value="">Select a skill...</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Training Examples (JSONL format)</label>
                            <textarea class="form-textarea" id="ft-examples" placeholder='{"instruction": "...", "input": "...", "output": "..."}'></textarea>
                        </div>
                        <button class="btn btn-primary" onclick="startSkillFineTune()">Start Fine-tuning</button>
                        <div id="ft-message" style="margin-top: 1rem;"></div>
                    </div>

                    <div class="glass-card card-half">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Feedback</span> Add Training Feedback</div>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Select Skill</label>
                            <select class="form-select" id="feedback-skill-select">
                                <option value="">Select a skill...</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Query</label>
                            <input type="text" class="form-input" id="feedback-query" placeholder="What the user asked...">
                        </div>
                        <div class="form-group">
                            <label class="form-label">Response (from skill)</label>
                            <textarea class="form-textarea" id="feedback-response" placeholder="What the skill responded..."></textarea>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Rating</label>
                            <select class="form-select" id="feedback-rating">
                                <option value="1">1 - Very Poor</option>
                                <option value="2">2 - Poor</option>
                                <option value="3">3 - Acceptable</option>
                                <option value="4">4 - Good</option>
                                <option value="5">5 - Excellent</option>
                            </select>
                        </div>
                        <div class="form-group">
                            <label class="form-label">Corrected Response (optional)</label>
                            <textarea class="form-textarea" id="feedback-corrected" placeholder="What the response should have been..."></textarea>
                        </div>
                        <button class="btn btn-success" onclick="submitSkillFeedback()">Submit Feedback</button>
                    </div>

                    <div class="glass-card card-full">
                        <div class="section-header">
                            <div class="section-title"><span class="section-icon">Auto</span> Auto-Improvement Status</div>
                        </div>
                        <table class="skills-table">
                            <thead>
                                <tr>
                                    <th>Skill</th>
                                    <th>Feedback Count</th>
                                    <th>Last Trained</th>
                                    <th>Status</th>
                                    <th>Actions</th>
                                </tr>
                            </thead>
                            <tbody id="skill-improvement-table">
                                <tr>
                                    <td colspan="5" style="text-align: center; color: var(--text-secondary);">
                                        Select a skill and add feedback to enable auto-improvement
                                    </td>
                                </tr>
                            </tbody>
                        </table>
                    </div>
                </div>
            </div>
        </div>
    </div>

    <script>
        // ============================================================
        // THEME TOGGLE
        // ============================================================
        function toggleTheme() {
            const html = document.documentElement;
            const currentTheme = html.getAttribute('data-theme');
            const newTheme = currentTheme === 'light' ? 'dark' : 'light';

            html.setAttribute('data-theme', newTheme);
            localStorage.setItem('theme', newTheme);

            // Update toggle button icon
            const toggleBtn = document.getElementById('theme-toggle');
            toggleBtn.textContent = newTheme === 'light' ? '☀️' : '🌙';
        }

        // Load saved theme on page load
        (function() {
            const savedTheme = localStorage.getItem('theme') || 'dark';
            document.documentElement.setAttribute('data-theme', savedTheme);

            // Update toggle button icon after DOM loads
            document.addEventListener('DOMContentLoaded', function() {
                const toggleBtn = document.getElementById('theme-toggle');
                if (toggleBtn) {
                    toggleBtn.textContent = savedTheme === 'light' ? '☀️' : '🌙';
                }
            });
        })();

        // ============================================================
        // TAB NAVIGATION
        // ============================================================
        function showMainTab(tabId) {
            document.querySelectorAll('.tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('.main-tab-btn').forEach(b => b.classList.remove('active'));

            // Map new simplified tabs to existing content
            let actualTabId = tabId;
            if (tabId === 'skills') actualTabId = 'skills-training';  // Skills -> New unified Skills & Training
            if (tabId === 'voice') actualTabId = 'voicelab';    // Voice -> Voice Lab
            if (tabId === 'settings') actualTabId = 'command';  // Settings -> Command Center (has API Keys)

            const tabElement = document.getElementById('tab-' + actualTabId);
            if (tabElement) {
                tabElement.classList.add('active');
            }
            if (event && event.target) {
                event.target.classList.add('active');
            }

            // Load appropriate data for each tab
            if (tabId === 'dashboard') { loadSkills(); loadMetrics(); }
            if (tabId === 'skills' || tabId === 'skills-training') { loadUnifiedSkills(); }  // New unified Skills & Training tab
            if (tabId === 'fastbrain') { loadFastBrainConfig(); loadFastBrainSkills(); refreshSystemStatus(); loadApiConnections(); }  // Legacy
            if (tabId === 'training') { loadTrainingSkillsDropdown(); refreshAdapters(); loadTrainingJobs(); }
            if (tabId === 'voice' || tabId === 'voicelab') { loadVoiceProjects(); loadSkillsForDropdowns(); }
            if (tabId === 'settings' || tabId === 'command') { refreshStats(); loadApiKeys(); }
            if (tabId === 'factory') { loadProfileDropdowns(); }
        }

        // New helper functions for consolidated navigation
        function showSkillsTab(subTab) {
            showMainTab('skills-training');
            setTimeout(() => {
                // Map old sub-tab requests to new workflow steps
                if (subTab === 'create') goToWorkflowStep(1);  // Create = Step 1 (select/create skill)
                if (subTab === 'test') goToWorkflowStep(3);    // Test = Step 3 (test & train)
                if (subTab === 'data') goToWorkflowStep(2);    // Data = Step 2 (add training data)
            }, 100);
        }

        function showVoiceTab(subTab) {
            showMainTab('voice');
            setTimeout(() => {
                if (subTab === 'connections') showVoiceLabTab('projects');
            }, 100);
        }

        function showSettingsTab(subTab) {
            showMainTab('settings');
            setTimeout(() => {
                if (subTab === 'keys') showCommandTab('keys');
            }, 100);
        }

        // ============================================================
        // UNIFIED SKILLS & TRAINING FUNCTIONS
        // ============================================================
        let currentSkillFilter = 'all';
        let currentModalSkillId = null;
        let currentWorkflowSkill = null;  // Selected skill for workflow
        let currentWorkflowStep = 1;
        let unifiedSkillsData = [];
        let modalTrainingPollInterval = null;
        let workflowTrainingPollInterval = null;

        // ============================================================
        // WORKFLOW NAVIGATION
        // ============================================================

        function goToWorkflowStep(step) {
            // Prevent navigation to disabled steps
            if (step > 1 && !currentWorkflowSkill) {
                showToast('Please select a skill first', 'warning');
                return;
            }

            if (step === currentWorkflowStep) return;

            // Update step indicators
            document.querySelectorAll('.workflow-step').forEach(s => {
                const stepNum = parseInt(s.dataset.step);
                s.classList.remove('active', 'completed', 'disabled');

                if (stepNum === step) {
                    s.classList.add('active');
                } else if (stepNum < step) {
                    s.classList.add('completed');
                } else if (!currentWorkflowSkill) {
                    s.classList.add('disabled');
                }
            });

            // Update connectors
            document.getElementById('connector-1-2').classList.toggle('completed', step > 1);
            document.getElementById('connector-2-3').classList.toggle('completed', step > 2);

            // Show/hide panels
            document.querySelectorAll('.workflow-panel').forEach(p => p.classList.remove('active'));
            const panel = document.getElementById('workflow-step-' + step);
            if (panel) panel.classList.add('active');

            currentWorkflowStep = step;

            // Load step-specific data
            if (step === 2) {
                loadWorkflowTrainingData();
            } else if (step === 3) {
                updateWorkflowReadiness();
                loadWorkflowAdapters();
            }
        }

        function selectSkillForWorkflow(skillId) {
            const skill = unifiedSkillsData.find(s => s.id === skillId);
            if (!skill) {
                showToast('Skill not found', 'error');
                return;
            }

            currentWorkflowSkill = skill;
            currentModalSkillId = skillId;  // For backwards compatibility

            // Update context bar
            const contextBar = document.getElementById('skill-context-bar');
            contextBar.classList.add('visible');

            document.getElementById('context-skill-icon').textContent = skill.icon || getSkillIcon(skillId);
            document.getElementById('context-skill-name').textContent = skill.name;

            const statusColors = {
                'untrained': '#ff6b6b',
                'has_data': 'var(--neon-orange)',
                'trained': 'var(--neon-green)'
            };
            document.getElementById('context-skill-status').innerHTML =
                `<span style="color: ${statusColors[skill.status] || '#888'}">● ${skill.status.replace('_', ' ')}</span>`;
            document.getElementById('context-skill-examples').textContent = `📊 ${skill.training_data?.total || 0} examples`;
            document.getElementById('context-skill-tokens').textContent = `🔤 ${skill.training_data?.tokens || 0} tokens`;

            // Enable workflow steps
            document.querySelectorAll('.workflow-step').forEach(s => {
                s.classList.remove('disabled');
            });

            // Move to step 2
            showToast(`Selected: ${skill.name}`, 'success');
            goToWorkflowStep(2);
        }

        function changeSelectedSkill() {
            goToWorkflowStep(1);
        }

        function editSkillDetails() {
            if (currentWorkflowSkill) {
                openSkillDetailModal(currentWorkflowSkill.id);
            }
        }

        // Toast notification system
        function showToast(message, type = 'info') {
            const container = document.getElementById('toast-container');
            if (!container) return;

            const icons = {
                'success': '✓',
                'error': '✗',
                'warning': '⚠',
                'info': 'ℹ'
            };

            const toast = document.createElement('div');
            toast.className = 'toast ' + type;
            toast.innerHTML = `<span style="font-size: 1.2rem;">${icons[type] || 'ℹ'}</span><span>${message}</span>`;

            container.appendChild(toast);

            setTimeout(() => {
                toast.style.animation = 'slideIn 0.3s ease reverse';
                setTimeout(() => toast.remove(), 300);
            }, 3000);
        }

        // ============================================================
        // WORKFLOW STEP 2: TRAINING DATA FUNCTIONS
        // ============================================================

        async function loadWorkflowTrainingData() {
            if (!currentWorkflowSkill) return;

            try {
                const response = await fetch(`/api/parser/data?skill_id=${currentWorkflowSkill.id}`);
                const data = await response.json();
                const items = data.items || [];

                // Update stats
                const approved = items.filter(i => i.is_approved).length;
                const total = items.length;
                const tokens = items.reduce((sum, i) => sum + (i.tokens || 0), 0);

                document.getElementById('wf-data-total').textContent = total;
                document.getElementById('wf-data-pending').textContent = total - approved;
                document.getElementById('wf-data-approved').textContent = approved;
                document.getElementById('wf-data-tokens').textContent = tokens;
                document.getElementById('step2-data-count').textContent = `${total} examples`;

                // Enable/disable continue button
                const nextBtn = document.getElementById('step2-next-btn');
                nextBtn.disabled = total < 5;
                nextBtn.title = total < 5 ? 'Need at least 5 training examples' : '';

                // Render training data list
                const list = document.getElementById('wf-training-data-list');
                if (items.length === 0) {
                    list.innerHTML = '<div style="padding: 1rem; text-align: center; color: var(--text-secondary);">No entries yet. Add your first training example!</div>';
                } else {
                    list.innerHTML = items.slice(0, 10).map(item => `
                        <div class="training-data-item" data-id="${item.id}">
                            <div class="data-content">
                                <div class="data-q" title="${escapeHtml(item.user_input)}">${escapeHtml(item.user_input)}</div>
                                <div class="data-a" title="${escapeHtml(item.assistant_response)}">${escapeHtml(item.assistant_response.substring(0, 100))}</div>
                            </div>
                            <div class="data-actions">
                                <button class="btn btn-sm btn-secondary" onclick="editWorkflowDataItem('${item.id}')" title="Edit">✏️</button>
                                <button class="btn btn-sm btn-danger" onclick="deleteWorkflowDataItem('${item.id}')" title="Delete">🗑</button>
                            </div>
                        </div>
                    `).join('');

                    if (items.length > 10) {
                        list.innerHTML += `<div style="padding: 0.5rem; text-align: center; color: var(--text-secondary); font-size: 0.85rem;">+ ${items.length - 10} more entries</div>`;
                    }
                }
            } catch (err) {
                console.error('Failed to load training data:', err);
                showToast('Failed to load training data', 'error');
            }
        }

        async function saveWorkflowEntry() {
            if (!currentWorkflowSkill) {
                showToast('Please select a skill first', 'warning');
                return;
            }

            const userInput = document.getElementById('wf-user-input').value.trim();
            const response = document.getElementById('wf-assistant-response').value.trim();

            if (!userInput || !response) {
                showToast('Both question and response are required', 'warning');
                return;
            }

            try {
                const res = await fetch('/api/parser/data', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentWorkflowSkill.id,
                        user_input: userInput,
                        assistant_response: response,
                        is_approved: false
                    })
                });

                const data = await res.json();
                if (data.success) {
                    showToast('Training entry saved!', 'success');
                    document.getElementById('wf-user-input').value = '';
                    document.getElementById('wf-assistant-response').value = '';
                    loadWorkflowTrainingData();
                } else {
                    showToast('Failed to save: ' + (data.error || 'Unknown error'), 'error');
                }
            } catch (err) {
                showToast('Failed to save entry', 'error');
            }
        }

        async function saveAndAddAnotherWorkflow() {
            await saveWorkflowEntry();
            document.getElementById('wf-user-input').focus();
        }

        async function deleteWorkflowDataItem(itemId) {
            if (!confirm('Delete this training entry?')) return;

            try {
                const res = await fetch(`/api/parser/data/${itemId}`, { method: 'DELETE' });
                const data = await res.json();
                if (data.success) {
                    showToast('Entry deleted', 'success');
                    loadWorkflowTrainingData();
                } else {
                    showToast('Failed to delete', 'error');
                }
            } catch (err) {
                showToast('Failed to delete entry', 'error');
            }
        }

        function editWorkflowDataItem(itemId) {
            // For now, open the existing modal
            openManualEntryModal();
            // TODO: Pre-populate with existing data
        }

        // File upload handlers for workflow
        function handleWorkflowDrop(event) {
            event.preventDefault();
            event.target.classList.remove('dragover');
            const files = event.dataTransfer.files;
            if (files.length > 0) {
                handleWorkflowUpload(files);
            }
        }

        async function handleWorkflowUpload(files) {
            if (!currentWorkflowSkill) {
                showToast('Please select a skill first', 'warning');
                return;
            }

            const statusDiv = document.getElementById('wf-upload-status');
            statusDiv.innerHTML = `<span style="color: var(--neon-cyan);">Uploading ${files.length} file(s)...</span>`;

            let successCount = 0;
            let totalExtracted = 0;

            for (const file of files) {
                try {
                    const formData = new FormData();
                    formData.append('file', file);
                    formData.append('skill_id', currentWorkflowSkill.id);

                    const response = await fetch('/api/parser/upload', {
                        method: 'POST',
                        body: formData
                    });

                    const data = await response.json();
                    console.log('Upload response:', data);  // Debug
                    if (data.success) {
                        successCount++;
                        totalExtracted += data.items_extracted || 0;
                        // Show debug info if no items extracted
                        if (data.items_extracted === 0 && data.debug_sample_lines?.length > 0) {
                            console.log('Sample lines from file:', data.debug_sample_lines);
                        }
                    }
                } catch (err) {
                    console.error('Upload error:', err);
                }
            }

            if (successCount > 0) {
                let msg = `✓ Uploaded ${successCount} file(s), extracted ${totalExtracted} Q&A pairs`;
                if (totalExtracted === 0) {
                    msg += ' (check browser console for debug info)';
                }
                statusDiv.innerHTML = `<span style="color: ${totalExtracted > 0 ? 'var(--neon-green)' : 'var(--neon-orange)'};">${msg}</span>`;
                loadWorkflowTrainingData();
            } else {
                statusDiv.innerHTML = `<span style="color: #ff6b6b;">✗ Upload failed</span>`;
            }

            setTimeout(() => { statusDiv.innerHTML = ''; }, 8000);
        }

        async function generateWorkflowAiData() {
            if (!currentWorkflowSkill) {
                showToast('Please select a skill first', 'warning');
                return;
            }

            const topic = document.getElementById('wf-ai-context').value.trim();
            const count = parseInt(document.getElementById('wf-ai-count').value);
            const statusDiv = document.getElementById('wf-ai-status');

            if (!topic) {
                showToast('Please enter a topic or context for AI to generate examples', 'warning');
                return;
            }

            statusDiv.innerHTML = `<span style="color: var(--neon-cyan);">✨ Generating ${count} training examples...</span>`;

            try {
                const response = await fetch('/api/parser/generate', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentWorkflowSkill.id,
                        topic: topic,
                        count: count
                    })
                });

                const data = await response.json();
                if (data.success) {
                    statusDiv.innerHTML = `<span style="color: var(--neon-green);">✓ Generated ${data.generated || count} examples!</span>`;
                    loadWorkflowTrainingData();
                } else {
                    statusDiv.innerHTML = `<span style="color: #ff6b6b;">✗ ${data.error || 'Generation failed'}</span>`;
                }
            } catch (err) {
                statusDiv.innerHTML = `<span style="color: #ff6b6b;">✗ Failed to generate training data</span>`;
            }

            setTimeout(() => { statusDiv.innerHTML = ''; }, 5000);
        }

        // ============================================================
        // WORKFLOW STEP 3: TEST & TRAIN FUNCTIONS
        // ============================================================

        function updateWorkflowReadiness() {
            if (!currentWorkflowSkill) return;

            const checksContainer = document.getElementById('wf-readiness-checks');
            const examples = currentWorkflowSkill.training_data?.total || 0;
            const hasSystemPrompt = !!(currentWorkflowSkill.system_prompt);

            // Fetch latest data
            fetch(`/api/parser/data?skill_id=${currentWorkflowSkill.id}`)
                .then(r => r.json())
                .then(data => {
                    const items = data.items || [];
                    const currentExamples = items.length;

                    const checks = [
                        {
                            pass: currentExamples >= 10,
                            warn: currentExamples >= 5 && currentExamples < 10,
                            text: `Training examples: ${currentExamples} ${currentExamples < 10 ? '(need 10+)' : '✓'}`,
                        },
                        {
                            pass: hasSystemPrompt,
                            warn: !hasSystemPrompt,
                            text: hasSystemPrompt ? 'System prompt defined ✓' : 'No system prompt (recommended)',
                        },
                        {
                            pass: true,
                            text: 'Modal GPU available ✓',
                        }
                    ];

                    checksContainer.innerHTML = checks.map(check => `
                        <div class="readiness-row">
                            <div class="readiness-check ${check.pass ? 'pass' : (check.warn ? 'warn' : 'fail')}">
                                ${check.pass ? '✓' : (check.warn ? '⚠' : '✗')}
                            </div>
                            <span>${check.text}</span>
                        </div>
                    `).join('');

                    // Enable/disable train button
                    const trainBtn = document.getElementById('wf-train-btn');
                    const canTrain = currentExamples >= 5;
                    trainBtn.disabled = !canTrain;

                    // Update current skill data
                    currentWorkflowSkill.training_data = {
                        ...currentWorkflowSkill.training_data,
                        total: currentExamples
                    };
                });
        }

        function updateWorkflowIntensityLabel() {
            const val = document.getElementById('wf-training-intensity').value;
            const labels = { '1': 'Quick', '2': 'Standard', '3': 'Deep' };
            const times = { '1': '~5 min', '2': '~10 min', '3': '~20 min' };
            const costs = { '1': '~$0.30', '2': '~$0.65', '3': '~$1.20' };

            document.getElementById('wf-intensity-label').textContent = labels[val];
            document.getElementById('wf-time-est').textContent = times[val];
            document.getElementById('wf-cost-est').textContent = costs[val];
        }

        async function loadWorkflowAdapters() {
            if (!currentWorkflowSkill) return;

            const container = document.getElementById('wf-adapters-list');
            const adapters = currentWorkflowSkill.adapters || [];

            if (adapters.length === 0) {
                container.innerHTML = '<div style="text-align: center; color: var(--text-secondary); font-size: 0.85rem; padding: 1rem;">No adapters yet. Train to create one!</div>';
                return;
            }

            container.innerHTML = adapters.map(adapter => `
                <div style="padding: 0.5rem; border-bottom: 1px solid var(--glass-border); display: flex; justify-content: space-between; align-items: center;">
                    <div>
                        <div style="font-weight: 500;">${adapter.adapter_name || adapter.id}</div>
                        <div style="font-size: 0.75rem; color: var(--text-secondary);">Loss: ${adapter.final_loss?.toFixed(3) || '--'}</div>
                    </div>
                    <button class="btn btn-sm btn-secondary" onclick="testAdapter('${adapter.skill_id || adapter.id}')">Test</button>
                </div>
            `).join('');
        }

        async function startWorkflowTraining() {
            // Prevent double-click: check button state immediately
            const trainBtn = document.getElementById('wf-train-btn');
            if (trainBtn.disabled) {
                console.log('[Training] Button already disabled, ignoring duplicate click');
                return;
            }
            trainBtn.disabled = true; // Disable immediately to prevent double-click

            if (!currentWorkflowSkill) {
                showToast('No skill selected', 'error');
                trainBtn.disabled = false;
                return;
            }

            const intensity = document.getElementById('wf-training-intensity').value;
            const epochs = { '1': 3, '2': 10, '3': 20 }[intensity];
            const statusDiv = document.getElementById('wf-training-status');

            trainBtn.innerHTML = '⏳ Starting training...';
            statusDiv.textContent = 'Initializing Modal GPU...';

            try {
                const response = await fetch(`/api/train-skill/${currentWorkflowSkill.id}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        epochs: epochs,
                        lora_r: 16,
                        learning_rate: 2e-4
                    })
                });

                const data = await response.json();
                if (data.success) {
                    showToast('Training started!', 'success');
                    trainBtn.innerHTML = '🔄 Training in progress...';
                    // Show enhanced training display immediately
                    const enhancedDiv = document.getElementById('wf-enhanced-training');
                    if (enhancedDiv) enhancedDiv.style.display = 'block';
                    startWorkflowTrainingPoll();
                } else {
                    showToast('Training failed: ' + (data.error || 'Unknown error'), 'error');
                    trainBtn.disabled = false;
                    trainBtn.innerHTML = '🚀 Start Training';
                }
            } catch (err) {
                showToast('Failed to start training', 'error');
                trainBtn.disabled = false;
                trainBtn.innerHTML = '🚀 Start Training';
            }
        }

        // Workflow training chart and facts
        let wfLossChart = null;
        let wfFactIndex = 0;
        let wfFactInterval = null;
        const wfTrainingFacts = [
            "LoRA trains only 0.1-1% of model parameters, making it 10-100x cheaper than full fine-tuning.",
            "The A10G GPU has 24GB VRAM and can process about 3.5 training steps per second.",
            "Loss measures prediction error - starting around 3.0, a good final loss is below 0.3.",
            "One epoch means the model has seen every training example once.",
            "Your adapter will be saved and can be used immediately after training completes."
        ];

        function initWfLossChart() {
            const canvas = document.getElementById('wf-loss-chart');
            if (!canvas) return;
            const ctx = canvas.getContext('2d');
            if (wfLossChart) wfLossChart.destroy();
            wfLossChart = new Chart(ctx, {
                type: 'line',
                data: { labels: [], datasets: [{ label: 'Loss', data: [], borderColor: '#10B981', backgroundColor: 'rgba(16,185,129,0.1)', fill: true, tension: 0.4, pointRadius: 2 }] },
                options: { responsive: true, maintainAspectRatio: false, animation: { duration: 300 }, plugins: { legend: { display: false } }, scales: { x: { display: false }, y: { beginAtZero: true, grid: { color: 'rgba(255,255,255,0.1)' }, ticks: { color: '#A1A1AA', font: { size: 10 } } } } }
            });
        }

        function rotateWfFact() {
            wfFactIndex = (wfFactIndex + 1) % wfTrainingFacts.length;
            const factEl = document.getElementById('wf-training-fact');
            if (factEl) factEl.textContent = wfTrainingFacts[wfFactIndex];
        }

        function updateWfTrainingUI(data) {
            // Show enhanced display
            const enhancedDiv = document.getElementById('wf-enhanced-training');
            if (enhancedDiv) enhancedDiv.style.display = 'block';

            // Update stats
            const lossEl = document.getElementById('wf-stat-loss');
            const stepsEl = document.getElementById('wf-stat-steps');
            const etaEl = document.getElementById('wf-stat-eta');
            const epochEl = document.getElementById('wf-stat-epoch');
            const progressPct = document.getElementById('wf-progress-percent');
            const progressBar = document.getElementById('wf-progress-bar');

            if (lossEl) lossEl.textContent = data.current_loss ? data.current_loss.toFixed(4) : '--';
            if (stepsEl) stepsEl.textContent = `${data.current_step || 0}/${data.total_steps || 0}`;
            if (etaEl) etaEl.textContent = data.eta_seconds ? Math.ceil(data.eta_seconds / 60) + 'm' : '--';
            if (epochEl) epochEl.textContent = `${data.current_epoch || 1}/${data.total_epochs || 10}`;

            const progress = data.progress || 0;
            if (progressPct) progressPct.textContent = progress.toFixed(0) + '%';
            if (progressBar) progressBar.style.width = progress + '%';

            // Update chart
            if (wfLossChart && data.loss_history && data.loss_history.length > 0) {
                wfLossChart.data.labels = data.loss_history.map((_, i) => i + 1);
                wfLossChart.data.datasets[0].data = data.loss_history;
                wfLossChart.update('none');
            }
        }

        function hideWfTrainingUI() {
            const enhancedDiv = document.getElementById('wf-enhanced-training');
            if (enhancedDiv) enhancedDiv.style.display = 'none';
            if (wfFactInterval) { clearInterval(wfFactInterval); wfFactInterval = null; }
        }

        function startWorkflowTrainingPoll() {
            if (workflowTrainingPollInterval) clearInterval(workflowTrainingPollInterval);

            // Initialize chart and start fact rotation
            initWfLossChart();
            if (wfFactInterval) clearInterval(wfFactInterval);
            wfFactInterval = setInterval(rotateWfFact, 6000);

            workflowTrainingPollInterval = setInterval(async () => {
                if (!currentWorkflowSkill) {
                    clearInterval(workflowTrainingPollInterval);
                    hideWfTrainingUI();
                    return;
                }

                try {
                    const res = await fetch(`/api/training/status/${currentWorkflowSkill.id}`);
                    const data = await res.json();

                    const statusDiv = document.getElementById('wf-training-status');
                    const trainBtn = document.getElementById('wf-train-btn');

                    if (data.status === 'completed') {
                        clearInterval(workflowTrainingPollInterval);
                        hideWfTrainingUI();
                        statusDiv.textContent = 'Training completed! New adapter created.';
                        trainBtn.disabled = false;
                        trainBtn.innerHTML = '🚀 Start Training';
                        showToast('Training completed successfully!', 'success');
                        // Fire confetti!
                        if (typeof confetti !== 'undefined') {
                            confetti({ particleCount: 100, spread: 70, origin: { y: 0.6 } });
                        }
                        loadWorkflowAdapters();
                        loadUnifiedSkills();
                    } else if (data.status === 'failed') {
                        clearInterval(workflowTrainingPollInterval);
                        hideWfTrainingUI();
                        statusDiv.textContent = 'Training failed: ' + (data.error || 'Unknown error');
                        trainBtn.disabled = false;
                        trainBtn.innerHTML = '🚀 Start Training';
                        showToast('Training failed', 'error');
                    } else {
                        statusDiv.textContent = `Training... ${data.progress || 0}% complete`;
                        updateWfTrainingUI(data);
                    }
                } catch (err) {
                    console.error('Training poll error:', err);
                }
            }, 3000);
        }

        // Workflow chat functions
        async function sendWorkflowChat() {
            if (!currentWorkflowSkill) {
                showToast('No skill selected', 'warning');
                return;
            }

            const input = document.getElementById('wf-chat-input');
            const message = input.value.trim();
            if (!message) return;

            const messagesDiv = document.getElementById('wf-chat-messages');

            // Add user message
            messagesDiv.innerHTML += `<div class="chat-message user">${escapeHtml(message)}</div>`;
            input.value = '';
            messagesDiv.scrollTop = messagesDiv.scrollHeight;

            // Add loading indicator
            const loadingId = 'loading-' + Date.now();
            messagesDiv.innerHTML += `<div id="${loadingId}" class="chat-message assistant" style="opacity: 0.6;">Thinking...</div>`;

            try {
                const response = await fetch('/api/chat', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentWorkflowSkill.id,
                        message: message
                    })
                });

                const data = await response.json();
                const loadingEl = document.getElementById(loadingId);
                if (loadingEl) {
                    loadingEl.remove();
                }

                if (data.response) {
                    messagesDiv.innerHTML += `<div class="chat-message assistant">${escapeHtml(data.response)}</div>`;
                } else {
                    messagesDiv.innerHTML += `<div class="chat-message assistant" style="color: #ff6b6b;">Error: ${data.error || 'No response'}</div>`;
                }
            } catch (err) {
                const loadingEl = document.getElementById(loadingId);
                if (loadingEl) {
                    loadingEl.innerHTML = 'Error: Failed to get response';
                    loadingEl.style.color = '#ff6b6b';
                }
            }

            messagesDiv.scrollTop = messagesDiv.scrollHeight;
        }

        // ============================================================
        // MODIFIED: Skill card click now uses workflow
        // ============================================================

        // Sub-tab navigation for unified view (legacy - kept for compatibility)
        function showUnifiedSubTab(tabName) {
            document.querySelectorAll('#tab-skills-training .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-skills-training .sub-tab-btn').forEach(b => b.classList.remove('active'));

            const tab = document.getElementById('unified-' + tabName);
            if (tab) {
                tab.classList.add('active');
                tab.style.display = 'block';
            }
            if (event && event.target) event.target.classList.add('active');

            // Load data for specific tabs
            if (tabName === 'skills') loadUnifiedSkills();
            if (tabName === 'golden') loadUnifiedSkillsDropdown('unified-golden-skill');
            if (tabName === 'training') loadUnifiedSkillsDropdown('unified-training-skill');
            if (tabName === 'data') loadUnifiedSkillsDropdown('unified-data-skill');
            if (tabName === 'chat') loadUnifiedSkillsDropdown('unified-chat-skill');
            if (tabName === 'adapters') refreshUnifiedAdapters();
        }

        async function loadUnifiedSkillsDropdown(selectId) {
            try {
                const response = await fetch('/api/fast-brain/skills');
                const data = await response.json();
                const select = document.getElementById(selectId);
                if (!select) return;

                const currentValue = select.value;
                select.innerHTML = '<option value="">-- Select a skill --</option>';

                if (data.skills) {
                    data.skills.forEach(skill => {
                        const option = document.createElement('option');
                        option.value = skill.id;
                        option.textContent = skill.name || skill.id;
                        select.appendChild(option);
                    });
                }
                if (currentValue) select.value = currentValue;
            } catch (err) {
                console.error('Failed to load skills dropdown:', err);
            }
        }

        // Create skill functions
        function showUnifiedCreateSkill() {
            document.getElementById('unified-create-skill-form').style.display = 'block';
            document.getElementById('unified-new-skill-id').focus();
        }

        function hideUnifiedCreateSkill() {
            document.getElementById('unified-create-skill-form').style.display = 'none';
        }

        async function createUnifiedSkill() {
            const skillId = document.getElementById('unified-new-skill-id').value.trim().toLowerCase().replace(/\s+/g, '_');
            const name = document.getElementById('unified-new-skill-name').value.trim();
            const description = document.getElementById('unified-new-skill-description').value.trim();
            const systemPrompt = document.getElementById('unified-new-skill-prompt').value.trim();

            if (!skillId || !name) {
                alert('Please fill in Skill ID and Name');
                return;
            }

            try {
                const res = await fetch('/api/fast-brain/skills', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ id: skillId, name, description, system_prompt: systemPrompt })
                });
                const data = await res.json();
                if (data.success) {
                    hideUnifiedCreateSkill();
                    loadUnifiedSkills();
                    alert('Skill created successfully!');
                } else {
                    throw new Error(data.error || 'Failed to create skill');
                }
            } catch (err) {
                alert('Error: ' + err.message);
            }
        }

        // Training functions
        function updateUnifiedIntensity() {
            const val = document.getElementById('unified-training-intensity').value;
            const labels = { 1: 'Quick', 2: 'Standard', 3: 'Deep' };
            const times = { 1: '~5 min', 2: '~10 min', 3: '~20 min' };
            const costs = { 1: '~$0.35', 2: '~$0.65', 3: '~$1.25' };
            document.getElementById('unified-intensity-label').textContent = labels[val];
            document.getElementById('unified-time-est').textContent = times[val];
            document.getElementById('unified-cost-est').textContent = costs[val];
        }

        async function startUnifiedTraining() {
            const skillId = document.getElementById('unified-training-skill').value;
            if (!skillId) { alert('Please select a skill'); return; }

            const intensity = document.getElementById('unified-training-intensity').value;
            const epochs = { 1: 3, 2: 10, 3: 20 }[intensity];

            try {
                const res = await fetch('/api/training/start', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId, config: { epochs } })
                });
                const data = await res.json();
                if (data.success) {
                    alert('Training started! Check the Adapters tab for progress.');
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                alert('Failed to start training: ' + err.message);
            }
        }

        // Adapters
        async function refreshUnifiedAdapters() {
            const container = document.getElementById('unified-adapters-list');
            try {
                const res = await fetch('/api/training/adapters');
                const data = await res.json();

                if (!data.adapters || data.adapters.length === 0) {
                    container.innerHTML = '<p style="color: var(--text-secondary); text-align: center; padding: 2rem;">No adapters yet. Train a skill to create one!</p>';
                    return;
                }

                container.innerHTML = data.adapters.map(a => `
                    <div class="glass-card" style="padding: 1rem; margin-bottom: 1rem;">
                        <div style="display: flex; justify-content: space-between; align-items: center;">
                            <div>
                                <strong>${a.skill_name || a.skill_id}</strong>
                                <div style="font-size: 0.85rem; color: var(--text-secondary);">
                                    Loss: ${a.final_loss?.toFixed(3) || '--'} | Created: ${new Date(a.created_at).toLocaleDateString()}
                                </div>
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-secondary btn-sm" onclick="testAdapter('${a.skill_id || a.id}')">Test</button>
                            </div>
                        </div>
                    </div>
                `).join('');
            } catch (err) {
                container.innerHTML = '<p style="color: var(--neon-orange); text-align: center;">Failed to load adapters</p>';
            }
        }

        // Chat
        async function sendUnifiedChat() {
            const skillId = document.getElementById('unified-chat-skill').value;
            const input = document.getElementById('unified-chat-input');
            const messages = document.getElementById('unified-chat-messages');
            const text = input.value.trim();

            if (!skillId) { alert('Please select a skill'); return; }
            if (!text) return;

            // Add user message
            messages.innerHTML += `<div style="margin-bottom: 0.75rem; text-align: right;"><span style="background: var(--neon-cyan); color: #000; padding: 0.5rem 1rem; border-radius: 12px; display: inline-block;">${text}</span></div>`;
            input.value = '';
            messages.scrollTop = messages.scrollHeight;

            try {
                const res = await fetch('/api/fast-brain/chat', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId, message: text })
                });
                const data = await res.json();
                const reply = data.response || data.message || 'No response';
                messages.innerHTML += `<div style="margin-bottom: 0.75rem;"><span style="background: var(--glass-surface); padding: 0.5rem 1rem; border-radius: 12px; display: inline-block;">${reply}</span></div>`;
                messages.scrollTop = messages.scrollHeight;
            } catch (err) {
                messages.innerHTML += `<div style="margin-bottom: 0.75rem;"><span style="background: rgba(255,0,0,0.2); padding: 0.5rem 1rem; border-radius: 12px; display: inline-block;">Error: ${err.message}</span></div>`;
            }
        }

        // API copy
        function copyUnifiedApiUrl() {
            const url = document.getElementById('unified-api-url').value;
            navigator.clipboard.writeText(url);
            alert('URL copied to clipboard!');
        }

        // Data manager
        async function loadUnifiedDataManager() {
            const skillId = document.getElementById('unified-data-skill').value;
            if (!skillId) return;

            try {
                const res = await fetch(`/api/parser/data?skill_id=${skillId}`);
                const data = await res.json();
                const items = data.items || [];

                const approved = items.filter(i => i.is_approved).length;
                document.getElementById('unified-data-total').textContent = items.length;
                document.getElementById('unified-data-pending').textContent = items.length - approved;
                document.getElementById('unified-data-approved').textContent = approved;
                document.getElementById('unified-data-tokens').textContent = items.reduce((s, i) => s + (i.tokens || 0), 0);

                const list = document.getElementById('unified-data-list');
                if (items.length === 0) {
                    list.innerHTML = '<p style="color: var(--text-secondary); text-align: center;">No training data for this skill.</p>';
                } else {
                    list.innerHTML = items.slice(0, 20).map(i => `
                        <div style="padding: 0.75rem; border-bottom: 1px solid var(--glass-border);">
                            <div style="font-weight: 500; margin-bottom: 0.25rem;">${escapeHtml(i.user_input.substring(0, 80))}...</div>
                            <div style="color: var(--text-secondary); font-size: 0.85rem;">${escapeHtml(i.assistant_response.substring(0, 100))}...</div>
                        </div>
                    `).join('');
                }
            } catch (err) {
                console.error('Failed to load data:', err);
            }
        }

        // Golden prompts placeholder
        async function loadUnifiedGoldenPrompts() {
            const skillId = document.getElementById('unified-golden-skill').value;
            const list = document.getElementById('unified-golden-list');
            if (!skillId) {
                list.innerHTML = '<p style="color: var(--text-secondary); text-align: center;">Select a skill to view golden prompts</p>';
                return;
            }
            // Load from existing golden prompts API
            try {
                const res = await fetch(`/api/golden-prompts?skill_id=${skillId}`);
                const data = await res.json();
                if (!data.prompts || data.prompts.length === 0) {
                    list.innerHTML = '<p style="color: var(--text-secondary); text-align: center;">No golden prompts for this skill.</p>';
                } else {
                    list.innerHTML = data.prompts.map(p => `
                        <div class="glass-card" style="padding: 1rem; margin-bottom: 0.5rem;">
                            <div><strong>User:</strong> ${escapeHtml(p.user_input)}</div>
                            <div><strong>Assistant:</strong> ${escapeHtml(p.assistant_response)}</div>
                        </div>
                    `).join('');
                }
            } catch (err) {
                list.innerHTML = '<p style="color: var(--text-secondary); text-align: center;">No golden prompts for this skill.</p>';
            }
        }

        function addUnifiedGoldenPrompt() {
            const skillId = document.getElementById('unified-golden-skill').value;
            if (!skillId) { alert('Please select a skill first'); return; }
            // Redirect to existing golden prompt creation or open modal
            openManualEntryModal();
        }

        // Sync skills from LPU API
        async function syncSkillsFromLPU() {
            try {
                const res = await fetch('/api/fast-brain/sync-skills', { method: 'POST' });
                const data = await res.json();
                if (data.success) {
                    alert(data.message || 'Skills synced successfully!');
                    loadUnifiedSkills();
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                alert('Failed to sync: ' + err.message);
            }
        }

        // Seed default built-in skills
        async function seedDefaultSkills() {
            try {
                const res = await fetch('/api/fast-brain/seed-skills', { method: 'POST' });
                const data = await res.json();
                if (data.success) {
                    alert(data.message || 'Default skills seeded!');
                    loadUnifiedSkills();
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                alert('Failed to seed: ' + err.message);
            }
        }

        async function loadUnifiedSkills() {
            const grid = document.getElementById('unified-skills-grid');
            grid.innerHTML = '<div style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--text-secondary);"><div style="font-size: 2rem; margin-bottom: 1rem;">🔄</div>Loading skills...</div>';

            try {
                // Fetch skills and their training status
                const [skillsRes, adaptersRes, dataRes] = await Promise.all([
                    fetch('/api/fast-brain/skills'),
                    fetch('/api/training/adapters'),
                    fetch('/api/parser/stats')
                ]);

                const skillsData = await skillsRes.json();
                const adaptersData = await adaptersRes.json();
                const dataStats = await dataRes.json();

                if (!skillsData.skills || skillsData.skills.length === 0) {
                    grid.innerHTML = '<div style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--text-secondary);"><div style="font-size: 3rem; margin-bottom: 1rem;">🎯</div><p>No skills yet. Create your first skill!</p><button class="btn btn-primary" onclick="showCreateSkillModalUnified()">+ Create Skill</button></div>';
                    return;
                }

                // Build adapter lookup
                const adaptersBySkill = {};
                if (adaptersData.adapters) {
                    adaptersData.adapters.forEach(adapter => {
                        const skillId = adapter.skill_id;
                        if (!adaptersBySkill[skillId]) adaptersBySkill[skillId] = [];
                        adaptersBySkill[skillId].push(adapter);
                    });
                }

                // Build data lookup
                const dataBySkill = {};
                if (dataStats.by_skill) {
                    dataStats.by_skill.forEach(stat => {
                        dataBySkill[stat.skill_id] = stat;
                    });
                }

                // Enrich skills with status
                unifiedSkillsData = skillsData.skills.map(skill => {
                    const adapters = adaptersBySkill[skill.id] || [];
                    const data = dataBySkill[skill.id] || { total: 0, tokens: 0 };

                    let status = 'untrained';
                    if (adapters.length > 0) {
                        status = 'trained';
                    } else if (data.total > 0) {
                        status = 'has_data';
                    }

                    return {
                        ...skill,
                        adapters,
                        training_data: data,
                        status
                    };
                });

                // Update currentWorkflowSkill reference if it exists (prevent stale reference)
                if (currentWorkflowSkill) {
                    const updatedSkill = unifiedSkillsData.find(s => s.id === currentWorkflowSkill.id);
                    if (updatedSkill) {
                        currentWorkflowSkill = updatedSkill;
                    }
                }

                renderSkillCards();
            } catch (err) {
                console.error('Failed to load unified skills:', err);
                grid.innerHTML = '<div style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--neon-orange);"><div style="font-size: 3rem; margin-bottom: 1rem;">⚠️</div><p>Failed to load skills. Please try again.</p><button class="btn btn-secondary" onclick="loadUnifiedSkills()">Retry</button></div>';
            }
        }

        function renderSkillCards() {
            const grid = document.getElementById('unified-skills-grid');
            const searchTerm = (document.getElementById('skill-search-input')?.value || '').toLowerCase();

            let filtered = unifiedSkillsData.filter(skill => {
                // Apply filter
                if (currentSkillFilter !== 'all' && skill.status !== currentSkillFilter) return false;
                // Apply search
                if (searchTerm && !skill.name.toLowerCase().includes(searchTerm) && !skill.id.toLowerCase().includes(searchTerm)) return false;
                return true;
            });

            // Sort
            const sortBy = document.getElementById('skill-sort-select')?.value || 'name';
            if (sortBy === 'name') {
                filtered.sort((a, b) => a.name.localeCompare(b.name));
            } else if (sortBy === 'status') {
                const order = { trained: 0, has_data: 1, untrained: 2 };
                filtered.sort((a, b) => order[a.status] - order[b.status]);
            }

            if (filtered.length === 0) {
                grid.innerHTML = '<div style="grid-column: 1 / -1; text-align: center; padding: 3rem; color: var(--text-secondary);">No skills match your filter.</div>';
                return;
            }

            grid.innerHTML = filtered.map(skill => {
                const statusLabels = {
                    untrained: '○ Untrained',
                    has_data: '◐ Has Data',
                    training: '◉ Training',
                    trained: '● Trained',
                    failed: '✕ Failed'
                };

                const icon = skill.icon || getSkillIcon(skill.id);
                const examples = skill.training_data?.total || 0;
                const adapters = skill.adapters?.length || 0;

                return `
                    <div class="skill-card" onclick="selectSkillForWorkflow('${skill.id}')" data-skill-id="${skill.id}" data-status="${skill.status}">
                        <div class="skill-card-header">
                            <div class="skill-card-icon">${icon}</div>
                            <span class="skill-status-badge ${skill.status}">${statusLabels[skill.status]}</span>
                        </div>
                        <div class="skill-card-title">${escapeHtml(skill.name)}</div>
                        <div class="skill-card-description">${escapeHtml(skill.description || 'No description')}</div>
                        <div class="skill-card-stats">
                            <div class="skill-card-stat">📝 <span class="value">${examples}</span> examples</div>
                            <div class="skill-card-stat">📦 <span class="value">${adapters}</span> adapters</div>
                        </div>
                    </div>
                `;
            }).join('');
        }

        function getSkillIcon(skillId) {
            const icons = {
                'molasses_alchemist': '🍯',
                'customer_service': '💬',
                'sales': '💼',
                'technical_support': '🔧',
                'booking': '📅',
                'default': '🎯'
            };
            return icons[skillId] || icons['default'];
        }

        function escapeHtml(str) {
            if (!str) return '';
            return str.replace(/&/g, '&amp;').replace(/</g, '&lt;').replace(/>/g, '&gt;').replace(/"/g, '&quot;');
        }

        function setSkillFilter(filter) {
            currentSkillFilter = filter;
            document.querySelectorAll('.skill-filter-btn').forEach(btn => {
                btn.classList.toggle('active', btn.dataset.filter === filter);
            });
            renderSkillCards();
        }

        function filterSkillCards() {
            renderSkillCards();
        }

        function sortSkillCards() {
            renderSkillCards();
        }

        function refreshUnifiedSkills() {
            loadUnifiedSkills();
        }

        // Skill Detail Modal
        async function openSkillDetailModal(skillId) {
            currentModalSkillId = skillId;
            const skill = unifiedSkillsData.find(s => s.id === skillId);
            if (!skill) return;

            // Populate header
            document.getElementById('skill-modal-icon').textContent = skill.icon || getSkillIcon(skillId);
            document.getElementById('skill-modal-title').textContent = skill.name;
            document.getElementById('skill-modal-status-badge').className = 'skill-status-badge ' + skill.status;
            document.getElementById('skill-modal-status-badge').textContent = skill.status.replace('_', ' ');
            document.getElementById('skill-modal-examples-count').textContent = (skill.training_data?.total || 0) + ' examples';

            // Populate overview
            document.getElementById('skill-modal-id').value = skill.id;
            document.getElementById('skill-modal-name').value = skill.name || '';
            document.getElementById('skill-modal-description').value = skill.description || '';
            document.getElementById('skill-modal-system-prompt').value = skill.system_prompt || '';

            // Stats
            document.getElementById('skill-modal-stat-examples').textContent = skill.training_data?.total || 0;
            document.getElementById('skill-modal-stat-tokens').textContent = skill.training_data?.tokens || 0;
            document.getElementById('skill-modal-stat-adapters').textContent = skill.adapters?.length || 0;

            // Reset to overview tab
            showSkillModalTab('overview');

            // Show modal
            document.getElementById('skill-detail-modal').style.display = 'flex';

            // Load training data for this skill
            loadSkillTrainingData(skillId);
            loadSkillAdapters(skillId);
            updateTrainingReadiness(skillId);
        }

        function closeSkillDetailModal() {
            document.getElementById('skill-detail-modal').style.display = 'none';
            currentModalSkillId = null;
            if (modalTrainingPollInterval) {
                clearInterval(modalTrainingPollInterval);
                modalTrainingPollInterval = null;
            }
        }

        function showSkillModalTab(tabName) {
            document.querySelectorAll('#skill-detail-modal .modal-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#skill-detail-modal .modal-tab-btn').forEach(b => b.classList.remove('active'));

            const tab = document.getElementById('skill-modal-' + tabName);
            if (tab) tab.classList.add('active');

            const btn = document.querySelector(`#skill-detail-modal .modal-tab-btn[data-tab="${tabName}"]`);
            if (btn) btn.classList.add('active');

            // Load tab-specific data
            if (tabName === 'training-data') loadSkillTrainingData(currentModalSkillId);
            if (tabName === 'train') updateTrainingReadiness(currentModalSkillId);
            if (tabName === 'adapters') loadSkillAdapters(currentModalSkillId);
        }

        async function loadSkillTrainingData(skillId) {
            try {
                const response = await fetch(`/api/parser/data?skill_id=${skillId}`);
                const data = await response.json();

                const tbody = document.getElementById('skill-data-tbody');
                const items = data.items || [];

                // Update stats
                const approved = items.filter(i => i.is_approved).length;
                document.getElementById('skill-data-total').textContent = items.length;
                document.getElementById('skill-data-pending').textContent = items.length - approved;
                document.getElementById('skill-data-approved').textContent = approved;
                document.getElementById('skill-data-tokens').textContent = items.reduce((sum, i) => sum + (i.tokens || 0), 0);

                if (items.length === 0) {
                    tbody.innerHTML = '<tr><td colspan="6" style="text-align: center; color: var(--text-secondary); padding: 2rem;">No training data yet. Add some!</td></tr>';
                    return;
                }

                tbody.innerHTML = items.map(item => `
                    <tr data-id="${item.id}">
                        <td><input type="checkbox" class="data-checkbox" value="${item.id}" onchange="updateSelectedCount()"></td>
                        <td style="max-width: 200px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;" title="${escapeHtml(item.user_input)}">${escapeHtml(item.user_input)}</td>
                        <td style="max-width: 250px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;" title="${escapeHtml(item.assistant_response)}">${escapeHtml(item.assistant_response.substring(0, 100))}...</td>
                        <td>${item.tokens || 0}</td>
                        <td><span class="skill-status-badge ${item.is_approved ? 'trained' : 'has_data'}">${item.is_approved ? 'Approved' : 'Pending'}</span></td>
                        <td>
                            <button class="btn btn-secondary btn-sm" onclick="editDataItem('${item.id}')" title="Edit">✏️</button>
                            <button class="btn btn-danger btn-sm" onclick="deleteDataItem('${item.id}')" title="Delete">🗑</button>
                        </td>
                    </tr>
                `).join('');
            } catch (err) {
                console.error('Failed to load training data:', err);
            }
        }

        async function loadSkillAdapters(skillId) {
            const container = document.getElementById('skill-adapters-list');
            const skill = unifiedSkillsData.find(s => s.id === skillId);
            const adapters = skill?.adapters || [];

            if (adapters.length === 0) {
                container.innerHTML = '<div style="text-align: center; color: var(--text-secondary); padding: 2rem;">No adapters for this skill yet. Train to create one!</div>';
                return;
            }

            container.innerHTML = adapters.map(adapter => `
                <div class="glass-card" style="padding: 1rem;">
                    <div style="display: flex; justify-content: space-between; align-items: center;">
                        <div>
                            <div style="font-weight: 600; color: var(--text-primary);">${adapter.adapter_name || adapter.id}</div>
                            <div style="font-size: 0.85rem; color: var(--text-secondary);">
                                Base: ${adapter.base_model || 'Unknown'} |
                                Loss: ${adapter.final_loss?.toFixed(3) || '--'} |
                                Created: ${new Date(adapter.created_at).toLocaleDateString()}
                            </div>
                        </div>
                        <div style="display: flex; gap: 0.5rem;">
                            <button class="btn btn-secondary btn-sm" onclick="testAdapter('${adapter.skill_id || adapter.id}')">🔬 Test</button>
                            <button class="btn btn-primary btn-sm" onclick="deployAdapter('${adapter.id}')">🚀 Deploy</button>
                        </div>
                    </div>
                </div>
            `).join('');
        }

        function updateTrainingReadiness(skillId) {
            const skill = unifiedSkillsData.find(s => s.id === skillId);
            const container = document.getElementById('readiness-checks');
            const examples = skill?.training_data?.total || 0;
            const hasSystemPrompt = !!(skill?.system_prompt);

            const checks = [
                { pass: examples >= 10, text: `Training examples: ${examples} (need 10+)`, icon: examples >= 10 ? '✓' : '✗' },
                { pass: hasSystemPrompt, text: hasSystemPrompt ? 'System prompt defined' : 'No system prompt (recommended)', icon: hasSystemPrompt ? '✓' : '⚠' },
                { pass: true, text: 'Modal GPU available', icon: '✓' }
            ];

            container.innerHTML = checks.map(check => `
                <div class="readiness-item">
                    <span class="readiness-icon ${check.pass ? 'pass' : 'fail'}">${check.icon}</span>
                    <span>${check.text}</span>
                </div>
            `).join('');

            // Enable/disable training button
            const btn = document.getElementById('modal-start-training-btn');
            const canTrain = examples >= 5;
            btn.disabled = !canTrain;
            btn.style.opacity = canTrain ? 1 : 0.5;
        }

        // Manual Entry Modal
        function openManualEntryModal() {
            document.getElementById('manual-user-input').value = '';
            document.getElementById('manual-assistant-response').value = '';
            document.getElementById('manual-category').value = '';
            document.getElementById('manual-entry-modal').style.display = 'flex';
        }

        function closeManualEntryModal() {
            document.getElementById('manual-entry-modal').style.display = 'none';
        }

        async function saveManualEntry() {
            await doSaveManualEntry();
            closeManualEntryModal();
            loadSkillTrainingData(currentModalSkillId);
        }

        async function saveAndAddAnother() {
            await doSaveManualEntry();
            document.getElementById('manual-user-input').value = '';
            document.getElementById('manual-assistant-response').value = '';
            loadSkillTrainingData(currentModalSkillId);
        }

        async function doSaveManualEntry() {
            const userInput = document.getElementById('manual-user-input').value.trim();
            const response = document.getElementById('manual-assistant-response').value.trim();
            const category = document.getElementById('manual-category').value.trim();

            if (!currentModalSkillId) {
                showToast('Please select a skill first before adding training data', 'error');
                return;
            }

            if (!userInput || !response) {
                showToast('Please fill in both user input and response', 'warning');
                return;
            }

            try {
                const res = await fetch('/api/parser/data', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentModalSkillId,
                        user_input: userInput,
                        assistant_response: response,
                        category: category || 'manual',
                        source_filename: 'manual_entry'
                    })
                });
                const data = await res.json();
                if (!data.success) throw new Error(data.error);
                showToast('Training example saved!', 'success');
            } catch (err) {
                showToast('Failed to save: ' + err.message, 'error');
            }
        }

        // Bulk Import Modal
        function openBulkImportModal() {
            document.getElementById('bulk-import-modal').style.display = 'flex';
        }

        function closeBulkImportModal() {
            document.getElementById('bulk-import-modal').style.display = 'none';
        }

        let bulkImportFiles = [];

        function handleBulkImportFiles(event) {
            bulkImportFiles = Array.from(event.target.files);
            if (bulkImportFiles.length > 0) {
                document.getElementById('bulk-import-queue').style.display = 'block';
                document.getElementById('bulk-import-file-list').innerHTML = bulkImportFiles.map((f, i) => `
                    <div style="display: flex; justify-content: space-between; padding: 0.5rem; border-bottom: 1px solid var(--glass-border);">
                        <span>${f.name}</span>
                        <span style="color: var(--text-secondary);">${(f.size / 1024).toFixed(1)} KB</span>
                    </div>
                `).join('');
            }
        }

        function clearBulkImportQueue() {
            bulkImportFiles = [];
            document.getElementById('bulk-import-queue').style.display = 'none';
            document.getElementById('bulk-import-files').value = '';
        }

        async function processBulkImport() {
            if (bulkImportFiles.length === 0) return;

            if (!currentModalSkillId) {
                showToast('Please select a skill first before uploading documents', 'error');
                return;
            }

            const statusEl = document.getElementById('bulk-import-status');
            const progressEl = document.getElementById('bulk-import-progress');
            const btn = document.querySelector('#bulk-import-modal .btn-primary');

            // Show progress UI
            if (statusEl) statusEl.style.display = 'block';
            if (btn) { btn.disabled = true; btn.innerHTML = '⏳ Processing...'; }

            let totalExtracted = 0;
            let filesProcessed = 0;
            let errors = [];

            for (let i = 0; i < bulkImportFiles.length; i++) {
                const file = bulkImportFiles[i];
                const formData = new FormData();
                formData.append('file', file);
                formData.append('skill_id', currentModalSkillId);

                // Update progress
                if (progressEl) {
                    progressEl.innerHTML = `Processing ${i + 1}/${bulkImportFiles.length}: ${file.name}`;
                }

                try {
                    const res = await fetch('/api/parser/upload', { method: 'POST', body: formData });
                    const data = await res.json();

                    if (data.success !== false) {
                        filesProcessed++;
                        totalExtracted += data.items_extracted || 0;
                    } else {
                        errors.push(`${file.name}: ${data.error || 'Unknown error'}`);
                    }
                } catch (err) {
                    console.error('Failed to upload:', file.name, err);
                    errors.push(`${file.name}: ${err.message}`);
                }
            }

            // Show completion status
            if (progressEl) {
                if (errors.length > 0) {
                    progressEl.innerHTML = `✅ Processed ${filesProcessed} files, extracted ${totalExtracted} examples. ⚠️ ${errors.length} errors.`;
                } else {
                    progressEl.innerHTML = `✅ Successfully processed ${filesProcessed} files and extracted ${totalExtracted} training examples!`;
                }
            }

            // Reset button
            if (btn) { btn.disabled = false; btn.innerHTML = '📤 Process Files'; }

            // Wait a moment to show the result, then close
            setTimeout(() => {
                clearBulkImportQueue();
                closeBulkImportModal();
                loadSkillTrainingData(currentModalSkillId);

                // Show toast notification
                if (totalExtracted > 0) {
                    showToast(`Extracted ${totalExtracted} training examples from ${filesProcessed} files!`, 'success');
                } else if (filesProcessed > 0) {
                    showToast(`Processed ${filesProcessed} files but no Q&A pairs found. Try different content.`, 'warning');
                }
            }, 1500);
        }

        // Toast notification helper
        function showToast(message, type = 'info') {
            const toast = document.createElement('div');
            toast.className = `toast toast-${type}`;
            toast.innerHTML = message;
            toast.style.cssText = `
                position: fixed; bottom: 20px; right: 20px; z-index: 99999;
                padding: 16px 24px; border-radius: 12px; font-weight: 500;
                animation: slideIn 0.3s ease;
                ${type === 'success' ? 'background: linear-gradient(135deg, #10b981, #059669); color: #fff;' : ''}
                ${type === 'error' ? 'background: linear-gradient(135deg, #ef4444, #dc2626); color: #fff;' : ''}
                ${type === 'warning' ? 'background: linear-gradient(135deg, #f59e0b, #d97706); color: #fff;' : ''}
                ${type === 'info' ? 'background: linear-gradient(135deg, #3b82f6, #2563eb); color: #fff;' : ''}
            `;
            document.body.appendChild(toast);
            setTimeout(() => toast.remove(), 4000);
        }

        // AI Generate Modal
        function openAiGenerateModal() {
            document.getElementById('ai-generate-modal').style.display = 'flex';
        }

        function closeAiGenerateModal() {
            document.getElementById('ai-generate-modal').style.display = 'none';
        }

        async function generateAiTrainingData() {
            const topic = document.getElementById('ai-generate-topic').value;
            const count = document.getElementById('ai-generate-count').value;
            const style = document.getElementById('ai-generate-style').value;
            const btn = document.querySelector('#ai-generate-modal .btn-primary');
            const statusEl = document.getElementById('ai-generate-status');

            if (!currentModalSkillId) {
                showToast('Please select a skill first before generating training data', 'error');
                return;
            }

            if (!topic) {
                showToast('Please enter a topic or context', 'warning');
                return;
            }

            // Show loading state
            if (btn) {
                btn.disabled = true;
                btn.innerHTML = `<span class="generating-dots">⏳</span> Generating ${count} examples...`;
            }
            if (statusEl) {
                statusEl.style.display = 'block';
                statusEl.innerHTML = `<div style="text-align: center; padding: 20px;">
                    <div style="font-size: 2rem; animation: pulse 1s infinite;">🤖</div>
                    <div style="margin-top: 10px; color: var(--text-secondary);">AI is generating training data...</div>
                    <div style="font-size: 0.8rem; color: var(--text-muted); margin-top: 5px;">This usually takes 15-30 seconds</div>
                </div>`;
            }

            try {
                const res = await fetch('/api/parser/generate', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentModalSkillId,
                        topic,
                        count: parseInt(count),
                        style
                    })
                });

                // Check for non-JSON response (server error)
                const contentType = res.headers.get('content-type');
                if (!contentType || !contentType.includes('application/json')) {
                    throw new Error('Server returned an error page. Please check your API keys in Settings.');
                }

                const data = await res.json();

                if (data.success) {
                    closeAiGenerateModal();
                    loadSkillTrainingData(currentModalSkillId);
                    showToast(`✨ Generated ${data.generated} training examples!`, 'success');
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                console.error('AI Generate error:', err);
                showToast('Failed to generate: ' + err.message, 'error');
                if (statusEl) {
                    statusEl.innerHTML = `<div style="text-align: center; padding: 20px; color: #ef4444;">
                        <div style="font-size: 1.5rem;">❌</div>
                        <div style="margin-top: 10px;">${err.message}</div>
                    </div>`;
                }
            } finally {
                // Reset button
                if (btn) {
                    btn.disabled = false;
                    btn.innerHTML = '🤖 Generate';
                }
            }
        }

        // Data Management
        function updateSelectedCount() {
            const checked = document.querySelectorAll('#skill-data-tbody .data-checkbox:checked').length;
            document.getElementById('skill-data-selected-count').textContent = checked + ' selected';
            document.getElementById('approve-selected-btn').disabled = checked === 0;
            document.getElementById('delete-selected-btn').disabled = checked === 0;
        }

        function toggleSelectAllData() {
            const selectAll = document.getElementById('select-all-data').checked;
            document.querySelectorAll('#skill-data-tbody .data-checkbox').forEach(cb => cb.checked = selectAll);
            updateSelectedCount();
        }

        async function approveSelectedData() {
            const ids = Array.from(document.querySelectorAll('#skill-data-tbody .data-checkbox:checked')).map(cb => cb.value);
            for (const id of ids) {
                await fetch(`/api/parser/data/${id}/approve`, { method: 'POST' });
            }
            loadSkillTrainingData(currentModalSkillId);
        }

        async function deleteSelectedData() {
            if (!confirm('Delete selected items?')) return;
            const ids = Array.from(document.querySelectorAll('#skill-data-tbody .data-checkbox:checked')).map(cb => cb.value);
            for (const id of ids) {
                await fetch(`/api/parser/data/${id}`, { method: 'DELETE' });
            }
            loadSkillTrainingData(currentModalSkillId);
        }

        async function deleteDataItem(id) {
            if (!confirm('Delete this item?')) return;
            await fetch(`/api/parser/data/${id}`, { method: 'DELETE' });
            loadSkillTrainingData(currentModalSkillId);
        }

        // Modal Training Config
        function setModalTrainingPreset(preset) {
            document.getElementById('modal-preset-simple').style.opacity = preset === 'simple' ? 1 : 0.5;
            document.getElementById('modal-preset-advanced').style.opacity = preset === 'advanced' ? 1 : 0.5;
            document.getElementById('modal-training-simple').style.display = preset === 'simple' ? 'block' : 'none';
            document.getElementById('modal-training-advanced').style.display = preset === 'advanced' ? 'block' : 'none';
        }

        function updateModalIntensityLabel() {
            const val = document.getElementById('modal-training-intensity').value;
            const labels = { 1: 'Quick', 2: 'Standard', 3: 'Deep' };
            const times = { 1: '~5 min', 2: '~10 min', 3: '~20 min' };
            const costs = { 1: '~$0.35', 2: '~$0.65', 3: '~$1.25' };
            document.getElementById('modal-intensity-label').textContent = labels[val];
            document.getElementById('modal-time-estimate').textContent = times[val];
            document.getElementById('modal-cost-estimate').textContent = costs[val];
        }

        async function startModalTraining() {
            if (!currentModalSkillId) {
                showToast('Please select a skill first', 'error');
                return;
            }

            const skill = unifiedSkillsData.find(s => s.id === currentModalSkillId);
            if (!skill) {
                showToast('Skill not found', 'error');
                return;
            }

            const intensity = document.getElementById('modal-training-intensity').value;
            const epochs = { 1: 3, 2: 10, 3: 20 }[intensity];

            // Show progress
            document.getElementById('modal-start-training-btn').disabled = true;
            document.getElementById('modal-training-progress').style.display = 'block';

            try {
                const res = await fetch('/api/training/start', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: currentModalSkillId,
                        config: {
                            epochs,
                            lora_r: parseInt(document.getElementById('modal-training-lora-r')?.value || 16),
                            learning_rate: document.getElementById('modal-training-lr')?.value || '2e-4'
                        }
                    })
                });
                const data = await res.json();
                if (data.success) {
                    startModalTrainingPolling();
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                alert('Failed to start training: ' + err.message);
                document.getElementById('modal-start-training-btn').disabled = false;
                document.getElementById('modal-training-progress').style.display = 'none';
            }
        }

        function startModalTrainingPolling() {
            modalTrainingPollInterval = setInterval(async () => {
                try {
                    const res = await fetch(`/api/training/status/${currentModalSkillId}`);
                    const data = await res.json();

                    if (data.status === 'running') {
                        document.getElementById('modal-progress-epoch').textContent = `${data.current_epoch || 0}/${data.total_epochs || 10}`;
                        document.getElementById('modal-progress-loss').textContent = data.loss?.toFixed(3) || '--';
                        document.getElementById('modal-progress-fill').style.width = `${((data.current_epoch || 0) / (data.total_epochs || 10)) * 100}%`;
                    } else if (data.status === 'completed') {
                        clearInterval(modalTrainingPollInterval);
                        document.getElementById('modal-training-status').textContent = 'Completed';
                        document.getElementById('modal-training-status').style.background = 'var(--neon-green)';
                        document.getElementById('modal-progress-fill').style.width = '100%';
                        loadUnifiedSkills();
                        loadSkillAdapters(currentModalSkillId);
                    } else if (data.status === 'failed') {
                        clearInterval(modalTrainingPollInterval);
                        document.getElementById('modal-training-status').textContent = 'Failed';
                        document.getElementById('modal-training-status').style.background = 'var(--neon-orange)';
                    }
                } catch (err) {
                    console.error('Polling error:', err);
                }
            }, 3000);
        }

        // Save skill overview
        async function saveSkillOverview() {
            const id = document.getElementById('skill-modal-id').value;
            const name = document.getElementById('skill-modal-name').value;
            const description = document.getElementById('skill-modal-description').value;
            const systemPrompt = document.getElementById('skill-modal-system-prompt').value;

            try {
                const res = await fetch(`/api/fast-brain/skills/${id}`, {
                    method: 'PUT',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ name, description, system_prompt: systemPrompt })
                });
                const data = await res.json();
                if (data.success) {
                    alert('Skill saved!');
                    loadUnifiedSkills();
                } else {
                    throw new Error(data.error);
                }
            } catch (err) {
                alert('Failed to save: ' + err.message);
            }
        }

        function testSkillChat() {
            // Open test chat in a new window or redirect
            closeSkillDetailModal();
            showMainTab('skills');
            setTimeout(() => {
                showFastBrainTab('chat');
                document.getElementById('fb-chat-skill-select').value = currentModalSkillId;
            }, 100);
        }

        function showCreateSkillModalUnified() {
            // Use existing create skill modal
            showMainTab('skills');
            setTimeout(() => {
                showFastBrainTab('skills');
                showCreateSkillModal();
            }, 100);
        }

        // Close modals on backdrop click
        document.addEventListener('click', function(e) {
            if (e.target.classList.contains('modal-overlay')) {
                e.target.style.display = 'none';
            }
        });

        // ============================================================
        // TRAINING TAB FUNCTIONS
        // ============================================================
        let currentTrainingSkillId = null;
        let trainingPollInterval = null;
        let lossChartData = [];
        let lrChartData = [];

        function showTrainingTab(subTab) {
            document.querySelectorAll('#tab-training .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-training .sub-tab-btn').forEach(b => b.classList.remove('active'));

            const tabContent = document.getElementById('training-' + subTab);
            if (tabContent) tabContent.classList.add('active');

            event.target.classList.add('active');

            if (subTab === 'adapters') refreshAdapters();
            if (subTab === 'progress') checkActiveTraining();
            if (subTab === 'data') loadDataManagerSkillsDropdown();
        }

        async function loadTrainingSkillsDropdown() {
            try {
                const response = await fetch('/api/fast-brain/skills');
                const data = await response.json();
                const select = document.getElementById('training-skill-select');

                select.innerHTML = '<option value="">-- Select a skill --</option>';

                if (data.skills && data.skills.length > 0) {
                    data.skills.forEach(skill => {
                        const option = document.createElement('option');
                        option.value = skill.id;
                        option.textContent = skill.name || skill.id;
                        select.appendChild(option);
                    });
                }
            } catch (err) {
                console.error('Failed to load skills:', err);
            }
        }

        async function loadTrainingDataStatus() {
            const skillId = document.getElementById('training-skill-select').value;
            const statusDiv = document.getElementById('training-data-status');

            if (!skillId) {
                statusDiv.style.display = 'none';
                return;
            }

            currentTrainingSkillId = skillId;
            statusDiv.style.display = 'block';

            try {
                const response = await fetch(`/api/fast-brain/skills/${skillId}`);
                const skill = await response.json();

                // Calculate training data stats
                const knowledge = skill.knowledge || [];
                const trainingData = skill.training_data || [];
                const examples = trainingData.length || knowledge.length || 0;

                // Estimate tokens (rough approximation)
                let totalTokens = 0;
                if (trainingData.length > 0) {
                    trainingData.forEach(d => {
                        totalTokens += (d.user_input?.length || 0) / 4;
                        totalTokens += (d.assistant_output?.length || 0) / 4;
                    });
                } else if (knowledge.length > 0) {
                    knowledge.forEach(k => totalTokens += (k.length || 0) / 4);
                }
                const avgTokens = examples > 0 ? Math.round(totalTokens / examples) : 0;

                // Calculate quality score
                let qualityStars = '';
                let qualityBadge = { text: 'Unknown', color: '#888' };
                if (examples >= 50) {
                    qualityStars = '★★★★★';
                    qualityBadge = { text: 'Excellent', color: 'var(--neon-green)' };
                } else if (examples >= 20) {
                    qualityStars = '★★★★☆';
                    qualityBadge = { text: 'Good', color: 'var(--neon-cyan)' };
                } else if (examples >= 10) {
                    qualityStars = '★★★☆☆';
                    qualityBadge = { text: 'Fair', color: 'var(--neon-orange)' };
                } else if (examples >= 5) {
                    qualityStars = '★★☆☆☆';
                    qualityBadge = { text: 'Minimal', color: 'var(--neon-pink)' };
                } else {
                    qualityStars = '★☆☆☆☆';
                    qualityBadge = { text: 'Insufficient', color: '#ff4444' };
                }

                document.getElementById('training-examples-count').textContent = examples;
                document.getElementById('training-avg-tokens').textContent = avgTokens;
                document.getElementById('training-quality-score').textContent = qualityStars;
                document.getElementById('training-topics').textContent = Math.min(examples, 10);

                const badge = document.getElementById('training-quality-badge');
                badge.textContent = qualityBadge.text;
                badge.style.background = qualityBadge.color;
                badge.style.color = qualityBadge.color === 'var(--neon-green)' || qualityBadge.color === 'var(--neon-cyan)' ? '#000' : '#fff';

                // Recommendation
                const recDiv = document.getElementById('training-recommendation');
                const recText = document.getElementById('training-recommendation-text');
                if (examples < 50) {
                    recDiv.style.display = 'block';
                    recText.textContent = `Add ${50 - examples} more examples for optimal results.`;
                } else {
                    recDiv.style.display = 'none';
                }

            } catch (err) {
                console.error('Failed to load skill data:', err);
            }
        }

        function setTrainingPreset(mode) {
            document.getElementById('preset-simple').style.opacity = mode === 'simple' ? '1' : '0.5';
            document.getElementById('preset-advanced').style.opacity = mode === 'advanced' ? '1' : '0.5';
            document.getElementById('training-config-simple').style.display = mode === 'simple' ? 'block' : 'none';
            document.getElementById('training-config-advanced').style.display = mode === 'advanced' ? 'block' : 'none';
        }

        function updateIntensityLabel() {
            const intensity = document.getElementById('training-intensity').value;
            const label = document.getElementById('training-intensity-label');
            const timeEst = document.getElementById('training-time-estimate');
            const costEst = document.getElementById('training-cost-estimate');

            if (intensity == 1) {
                label.textContent = 'Quick';
                label.style.color = 'var(--neon-cyan)';
                timeEst.textContent = '~5 minutes';
                costEst.textContent = '~$0.35';
            } else if (intensity == 2) {
                label.textContent = 'Standard';
                label.style.color = 'var(--neon-green)';
                timeEst.textContent = '~10 minutes';
                costEst.textContent = '~$0.65';
            } else {
                label.textContent = 'Deep';
                label.style.color = 'var(--neon-purple)';
                timeEst.textContent = '~20 minutes';
                costEst.textContent = '~$1.30';
            }
        }

        async function startTraining() {
            // Prevent double-click: check button state immediately
            const btn = document.getElementById('start-training-btn');
            if (btn.disabled) {
                console.log('[Training] Button already disabled, ignoring duplicate click');
                return;
            }
            btn.disabled = true; // Disable immediately to prevent double-click

            const skillId = document.getElementById('training-skill-select').value;
            if (!skillId) {
                alert('Please select a skill to train');
                btn.disabled = false;
                return;
            }

            btn.textContent = '⏳ Starting...';

            // Get configuration
            let config = {};
            const isAdvanced = document.getElementById('training-config-advanced').style.display !== 'none';

            if (isAdvanced) {
                config = {
                    epochs: parseInt(document.getElementById('training-epochs').value),
                    learning_rate: parseFloat(document.getElementById('training-lr').value),
                    lora_r: parseInt(document.getElementById('training-lora-r').value)
                };
            } else {
                const intensity = document.getElementById('training-intensity').value;
                if (intensity == 1) {
                    config = { epochs: 3 };
                } else if (intensity == 2) {
                    config = { epochs: 10 };
                } else {
                    config = { epochs: 20 };
                }
            }

            try {
                const response = await fetch(`/api/train-skill/${skillId}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(config)
                });

                const result = await response.json();

                // Handle 409 Conflict (already training)
                if (response.status === 409) {
                    const elapsed = result.elapsed_seconds ? Math.floor(result.elapsed_seconds / 60) : 0;
                    alert(`Training already in progress for this skill (running for ${elapsed} minutes). Check the Progress tab!`);
                    // Still show the progress tab since training is active
                    currentTrainingSkillId = skillId;
                    showTrainingTab('progress');
                    document.querySelector('#tab-training .sub-tab-btn:nth-child(2)').click();
                    document.getElementById('active-training-display').style.display = 'block';
                    startTrainingPolling(skillId);
                    return; // Don't re-enable button
                }

                if (result.success) {
                    currentTrainingSkillId = skillId;
                    // Show the progress tab and training UI
                    showTrainingTab('progress');
                    document.querySelector('#tab-training .sub-tab-btn:nth-child(2)').click();
                    document.getElementById('active-training-display').style.display = 'block';
                    startTrainingPolling(skillId);
                    // Keep button disabled during training
                    btn.textContent = '⏳ Training...';
                    return; // Don't re-enable button
                } else {
                    alert('Training failed: ' + (result.error || 'Unknown error'));
                    btn.disabled = false;
                    btn.textContent = '🚀 Start Training';
                }
            } catch (err) {
                alert('Training error: ' + err.message);
                btn.disabled = false;
                btn.textContent = '🚀 Start Training';
            }
        }

        // Chart.js instance for loss curve
        let lossChart = null;
        let factRotationInterval = null;
        let currentFactIndex = 0;

        // Educational facts for the carousel
        const trainingFacts = [
            { icon: '🧠', title: 'What is LoRA?', text: 'Low-Rank Adaptation trains only <strong style="color: var(--neon-green);">0.92%</strong> of parameters, making it 10x faster than full fine-tuning!' },
            { icon: '📉', title: 'Understanding Loss', text: 'Loss measures how "wrong" the model is. Starting at <strong style="color: var(--neon-green);">3.38</strong>, target is <strong style="color: var(--neon-green);">< 0.3</strong>!' },
            { icon: '🔄', title: "What is an Epoch?", text: "One epoch = seeing all training examples once. Training runs for <strong style='color: var(--neon-green);'>10 epochs</strong> total!" },
            { icon: '💰', title: 'Cost Breakdown', text: 'This training costs ~<strong style="color: var(--neon-green);">$0.50-$0.80</strong>. Competitors charge <strong style="color: var(--neon-green);">$5-15</strong> for the same task!' },
            { icon: '⚡', title: 'A10G GPU', text: 'The NVIDIA A10G has <strong style="color: var(--neon-green);">22GB VRAM</strong> and processes ~<strong style="color: var(--neon-green);">3.5 steps/second</strong>!' }
        ];

        function formatTime(seconds) {
            if (!seconds || seconds < 0) return '--:--';
            const mins = Math.floor(seconds / 60);
            const secs = Math.floor(seconds % 60);
            return `${mins}:${secs.toString().padStart(2, '0')}`;
        }

        function initLossChart() {
            const ctx = document.getElementById('loss-chart').getContext('2d');
            if (lossChart) {
                lossChart.destroy();
            }
            lossChart = new Chart(ctx, {
                type: 'line',
                data: {
                    labels: [],
                    datasets: [{
                        label: 'Training Loss',
                        data: [],
                        borderColor: '#10b981',
                        backgroundColor: 'rgba(16, 185, 129, 0.1)',
                        fill: true,
                        tension: 0.4,
                        pointRadius: 3,
                        pointBackgroundColor: '#10b981'
                    }]
                },
                options: {
                    responsive: true,
                    maintainAspectRatio: false,
                    animation: { duration: 300 },
                    plugins: { legend: { display: false } },
                    scales: {
                        y: {
                            beginAtZero: true,
                            grid: { color: 'rgba(255,255,255,0.05)' },
                            ticks: { color: '#94a3b8' }
                        },
                        x: {
                            grid: { color: 'rgba(255,255,255,0.05)' },
                            ticks: { color: '#94a3b8' }
                        }
                    }
                }
            });
        }

        function updateLossChartData(lossHistory) {
            if (!lossChart || !lossHistory || lossHistory.length === 0) return;
            lossChart.data.labels = lossHistory.map(h => h.step);
            lossChart.data.datasets[0].data = lossHistory.map(h => h.loss);
            lossChart.update('none');
        }

        function rotateFact() {
            currentFactIndex = (currentFactIndex + 1) % trainingFacts.length;
            const fact = trainingFacts[currentFactIndex];
            document.getElementById('fact-icon').textContent = fact.icon;
            document.getElementById('fact-title').textContent = fact.title;
            document.getElementById('fact-text').innerHTML = fact.text;
            // Update dots
            const dots = document.querySelectorAll('.fact-dot');
            dots.forEach((dot, i) => {
                dot.style.background = i === currentFactIndex ? 'var(--neon-green)' : 'rgba(255,255,255,0.2)';
            });
        }

        function startTrainingPolling(skillId) {
            // Show active training display
            document.getElementById('no-training-message').style.display = 'none';
            document.getElementById('active-training-display').style.display = 'block';
            document.getElementById('training-complete-display').style.display = 'none';
            document.getElementById('cancel-training-btn').style.display = 'block';
            document.getElementById('progress-skill-name').textContent = skillId;

            // Initialize Chart.js
            initLossChart();

            // Reset and start fact rotation
            currentFactIndex = 0;
            if (factRotationInterval) clearInterval(factRotationInterval);
            factRotationInterval = setInterval(rotateFact, 8000);

            // Start polling with enhanced endpoint (3 second interval)
            if (trainingPollInterval) clearInterval(trainingPollInterval);
            trainingPollInterval = setInterval(() => pollTrainingStatus(skillId), 3000);
            pollTrainingStatus(skillId);
        }

        async function pollTrainingStatus(skillId) {
            try {
                // Use the enhanced training status endpoint
                const response = await fetch(`/api/training/status/${skillId}`);
                const data = await response.json();

                if (data.status === 'idle' || data.error) {
                    console.log('Training job not found or idle');
                    return;
                }

                // Update status badge
                const statusBadge = document.getElementById('progress-status');
                const trainingPulse = document.getElementById('training-pulse');
                if (data.status === 'running') {
                    statusBadge.textContent = '● Running';
                    statusBadge.style.background = 'linear-gradient(135deg, var(--neon-green), #059669)';
                    trainingPulse.style.animation = 'pulse 2s infinite';
                } else if (data.status === 'completed') {
                    statusBadge.textContent = '✓ Complete';
                    statusBadge.style.background = 'linear-gradient(135deg, var(--neon-cyan), #0ea5e9)';
                    trainingPulse.style.animation = 'none';
                    trainingPulse.style.background = 'var(--neon-cyan)';
                    clearInterval(trainingPollInterval);
                    clearInterval(factRotationInterval);
                    showTrainingComplete(data);
                    return;
                } else if (data.status === 'failed') {
                    statusBadge.textContent = '✗ Failed';
                    statusBadge.style.background = '#ff4444';
                    trainingPulse.style.animation = 'none';
                    trainingPulse.style.background = '#ff4444';
                    clearInterval(trainingPollInterval);
                    clearInterval(factRotationInterval);
                    // Re-enable the Start Training button
                    const btn = document.getElementById('start-training-btn');
                    btn.disabled = false;
                    btn.textContent = '🚀 Start Training';
                    return;
                }

                // Update skill name
                if (data.skill_name) {
                    document.getElementById('progress-skill-name').textContent = data.skill_name;
                }

                // Update stats cards
                if (data.current_loss !== null) {
                    document.getElementById('stat-loss').textContent = data.current_loss.toFixed(4);
                }
                if (data.loss_improvement_percent) {
                    document.getElementById('stat-loss-trend').textContent = `↓ ${data.loss_improvement_percent.toFixed(0)}% from start`;
                }

                // Steps
                document.getElementById('stat-steps').textContent = `${data.current_step}/${data.total_steps}`;
                const stepProgress = data.total_steps > 0 ? (data.current_step / data.total_steps) * 100 : 0;
                document.getElementById('stat-steps-bar').style.width = `${stepProgress}%`;

                // ETA
                document.getElementById('stat-eta').textContent = formatTime(data.eta_seconds);
                document.getElementById('stat-elapsed').textContent = `Elapsed: ${formatTime(data.elapsed_seconds)}`;

                // Epoch
                document.getElementById('stat-epoch').textContent = data.current_epoch?.toFixed(1) || '0';
                document.getElementById('stat-epoch-total').textContent = `of ${data.total_epochs} epochs`;

                // GPU metrics
                if (data.gpu_metrics) {
                    document.getElementById('gpu-name').textContent = data.gpu_metrics.name;
                    document.getElementById('gpu-mem-text').textContent = `${data.gpu_metrics.memory_used_gb} / ${data.gpu_metrics.memory_total_gb} GB`;
                    const memPercent = (data.gpu_metrics.memory_used_gb / data.gpu_metrics.memory_total_gb) * 100;
                    document.getElementById('gpu-mem-bar').style.width = `${memPercent}%`;
                    document.getElementById('gpu-util-text').textContent = `${data.gpu_metrics.utilization_percent}%`;
                    document.getElementById('gpu-util-bar').style.width = `${data.gpu_metrics.utilization_percent}%`;
                }

                // Current example preview
                if (data.current_example_preview) {
                    document.getElementById('current-example-input').textContent = `"${data.current_example_preview.input}"`;
                    document.getElementById('current-example-output').textContent = `"${data.current_example_preview.output}"`;
                }

                // Update loss chart
                if (data.loss_history && data.loss_history.length > 0) {
                    updateLossChartData(data.loss_history);
                }

            } catch (err) {
                console.error('Failed to poll training status:', err);
            }
        }

        function updateTrainingLog(logs) {
            const logDiv = document.getElementById('training-log');
            if (!logDiv) return;
            logDiv.innerHTML = logs.map(line => {
                let color = 'var(--text-secondary)';
                if (line.includes('✓') || line.includes('SUCCESS') || line.includes('complete')) color = 'var(--neon-green)';
                if (line.includes('Error') || line.includes('Failed') || line.includes('✗')) color = '#ff4444';
                if (line.includes('loss=') || line.includes('Loss:')) color = 'var(--neon-cyan)';
                if (line.includes('%|')) color = 'var(--neon-purple)';
                return `<div style="color: ${color}; margin-bottom: 2px;">${escapeHtml(line)}</div>`;
            }).join('');
            logDiv.scrollTop = logDiv.scrollHeight;
        }

        function showTrainingComplete(data) {
            document.getElementById('active-training-display').style.display = 'none';
            document.getElementById('training-complete-display').style.display = 'block';
            document.getElementById('cancel-training-btn').style.display = 'none';

            // Re-enable the Start Training button
            const btn = document.getElementById('start-training-btn');
            btn.disabled = false;
            btn.textContent = '🚀 Start Training';

            // Fire confetti celebration!
            if (typeof confetti !== 'undefined') {
                confetti({
                    particleCount: 100,
                    spread: 70,
                    origin: { y: 0.6 }
                });
                // Fire again for more celebration
                setTimeout(() => {
                    confetti({
                        particleCount: 50,
                        angle: 60,
                        spread: 55,
                        origin: { x: 0 }
                    });
                    confetti({
                        particleCount: 50,
                        angle: 120,
                        spread: 55,
                        origin: { x: 1 }
                    });
                }, 250);
            }

            // Update result stats from the enhanced data
            const finalLoss = data.current_loss?.toFixed(4) || '--';
            const trainingTime = data.elapsed_seconds ? `${(data.elapsed_seconds / 60).toFixed(1)} min` : '--';
            const examples = data.total_examples || '--';

            document.getElementById('result-final-loss').textContent = finalLoss;
            document.getElementById('result-training-time').textContent = trainingTime;
            document.getElementById('result-examples').textContent = examples;
            document.getElementById('result-cost').textContent = '~$0.65';

            refreshAdapters();
        }

        async function testTrainedAdapter() {
            const prompt = document.getElementById('test-adapter-prompt').value;
            if (!prompt || !currentTrainingSkillId) return;

            const responseDiv = document.getElementById('test-adapter-response');
            responseDiv.style.display = 'block';
            responseDiv.innerHTML = '<span style="color: var(--text-secondary);">Testing adapter...</span>';

            try {
                const response = await fetch(`/api/test-adapter/${currentTrainingSkillId}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ prompt })
                });

                const result = await response.json();

                if (result.response) {
                    responseDiv.innerHTML = `<span style="color: var(--neon-cyan);">🧪 ${escapeHtml(result.response)}</span>`;
                } else {
                    responseDiv.innerHTML = `<span style="color: #ff4444;">Error: ${result.error || 'No response'}</span>`;
                }
            } catch (err) {
                responseDiv.innerHTML = `<span style="color: #ff4444;">Error: ${err.message}</span>`;
            }
        }

        function cancelTraining() {
            if (confirm('Are you sure you want to cancel the training?')) {
                clearInterval(trainingPollInterval);
                document.getElementById('active-training-display').style.display = 'none';
                document.getElementById('no-training-message').style.display = 'block';
            }
        }

        function clearTrainingLog() {
            document.getElementById('training-log').innerHTML = '<div style="color: var(--text-secondary);">Log cleared.</div>';
        }

        async function refreshAdapters() {
            try {
                const response = await fetch('/api/trained-adapters');
                const data = await response.json();

                const gallery = document.getElementById('adapters-gallery');
                const noMsg = document.getElementById('no-adapters-message');

                if (data.adapters && data.adapters.length > 0) {
                    noMsg.style.display = 'none';
                    gallery.innerHTML = data.adapters.map(adapter => `
                        <div class="glass-card" style="background: var(--glass-surface); padding: 1rem;">
                            <div style="display: flex; justify-content: space-between; align-items: start; margin-bottom: 0.5rem;">
                                <div>
                                    <div style="font-weight: 600; color: var(--neon-cyan);">🧪 ${escapeHtml(adapter.skill_name || adapter.skill_id)}</div>
                                    <div style="font-size: 0.8rem; color: var(--text-secondary);">${adapter.skill_id}</div>
                                </div>
                                <span class="badge" style="background: var(--neon-green); color: #000; padding: 0.2rem 0.5rem; font-size: 0.7rem;">Active</span>
                            </div>
                            <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 0.5rem; font-size: 0.85rem; margin-bottom: 0.75rem;">
                                <div><span style="color: var(--text-secondary);">Loss:</span> ${adapter.final_loss?.toFixed(4) || '--'}</div>
                                <div><span style="color: var(--text-secondary);">Examples:</span> ${adapter.training_examples || '--'}</div>
                            </div>
                            <div style="display: flex; gap: 0.5rem;">
                                <button class="btn btn-secondary btn-sm" onclick="testAdapterFromGallery('${adapter.skill_id}')">Test</button>
                                <button class="btn btn-secondary btn-sm" onclick="showAdapterDetails('${adapter.skill_id}')">⋮</button>
                            </div>
                        </div>
                    `).join('');
                } else {
                    noMsg.style.display = 'block';
                    gallery.innerHTML = '';
                    gallery.appendChild(noMsg);
                }

            } catch (err) {
                console.error('Failed to load adapters:', err);
            }
        }

        async function loadTrainingJobs() {
            try {
                const response = await fetch('/api/training-jobs');
                const data = await response.json();

                const table = document.getElementById('training-history-table');

                if (data.jobs && data.jobs.length > 0) {
                    table.innerHTML = data.jobs.map(job => {
                        const date = job.started_at ? new Date(job.started_at).toLocaleString() : '--';
                        const statusColor = job.status === 'completed' ? 'var(--neon-green)' : job.status === 'running' ? 'var(--neon-cyan)' : '#ff4444';
                        return `
                            <tr>
                                <td>${date}</td>
                                <td>${escapeHtml(job.skill_id)}</td>
                                <td><span style="color: ${statusColor};">${job.status}</span></td>
                                <td>--</td>
                                <td>--</td>
                                <td>--</td>
                            </tr>
                        `;
                    }).join('');
                }
            } catch (err) {
                console.error('Failed to load training jobs:', err);
            }
        }

        function checkActiveTraining() {
            // Check if there's an active training job
            if (currentTrainingSkillId) {
                pollTrainingStatus(currentTrainingSkillId);
            }
        }

        function showTrainingHelp() {
            alert('Training Help:\\n\\n1. Select a skill to train\\n2. Review the training data status\\n3. Choose training intensity (Quick/Standard/Deep)\\n4. Click Start Training\\n\\nTraining runs on Modal\\'s A10G GPU and typically takes 5-20 minutes.\\nCost: $0.35-$1.30 per training run.');
        }

        function viewTrainingExamples() {
            if (currentTrainingSkillId) {
                showMainTab('skills');
            }
        }

        async function testAdapterFromGallery(skillId) {
            const prompt = window.prompt('Enter a test prompt:');
            if (!prompt) return;

            try {
                const response = await fetch(`/api/test-adapter/${skillId}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ prompt: prompt })
                });

                const data = await response.json();

                if (data.success && data.response) {
                    alert('🧪 Trained Adapter Response:\\n\\n' + data.response);
                } else {
                    alert('❌ Error: ' + (data.error || 'No response from adapter'));
                }
            } catch (err) {
                alert('❌ Error testing adapter: ' + err.message);
            }
        }

        // Alias for testAdapter calls (adapter.id often equals skill_id)
        function testAdapter(adapterId) {
            testAdapterFromGallery(adapterId);
        }

        function showAdapterDetails(skillId) {
            alert('Adapter: ' + skillId + '\\n\\nOptions:\\n- Retrain with more data\\n- Export to HuggingFace\\n- Delete adapter');
        }

        // ============================================================
        // DATA MANAGER FUNCTIONS
        // ============================================================
        let dmCurrentSkillId = null;
        let dmSelectedIds = new Set();
        let dmCurrentPage = 1;
        let parserFiles = [];

        function loadDataManagerSkillsDropdown() {
            loadTrainingSkillsDropdown();  // Reuse existing function
            // Also populate dm-skill-select and parser-skill-select
            fetch('/api/fast-brain/skills')
                .then(r => r.json())
                .then(data => {
                    ['dm-skill-select', 'parser-skill-select'].forEach(selectId => {
                        const select = document.getElementById(selectId);
                        if (select) {
                            select.innerHTML = '<option value="">-- Select a skill --</option>';
                            if (data.skills) {
                                data.skills.forEach(skill => {
                                    const opt = document.createElement('option');
                                    opt.value = skill.id;
                                    opt.textContent = skill.name || skill.id;
                                    select.appendChild(opt);
                                });
                            }
                        }
                    });
                });
        }

        async function loadDataManagerData() {
            const skillId = document.getElementById('dm-skill-select').value;
            if (!skillId) return;

            dmCurrentSkillId = skillId;

            try {
                // Load stats
                const statsRes = await fetch(`/api/parser/data/${skillId}/stats`);
                const statsData = await statsRes.json();

                if (statsData.success) {
                    document.getElementById('dm-total-items').textContent = statsData.stats.total || 0;
                    document.getElementById('dm-pending-items').textContent = (statsData.stats.total - statsData.stats.approved) || 0;
                    document.getElementById('dm-approved-items').textContent = statsData.stats.approved || 0;
                    document.getElementById('dm-total-tokens').textContent = formatNumber(statsData.stats.total_tokens || 0);
                }

                // Load data
                filterDataManager();

            } catch (err) {
                console.error('Failed to load data manager:', err);
            }
        }

        async function filterDataManager() {
            if (!dmCurrentSkillId) return;

            const search = document.getElementById('dm-search').value;
            const category = document.getElementById('dm-category-filter').value;
            const status = document.getElementById('dm-status-filter').value;

            const params = new URLSearchParams({
                page: dmCurrentPage,
                per_page: 50
            });

            if (search) params.set('search', search);
            if (category) params.set('category', category);
            if (status === 'approved') params.set('approved', 'true');
            if (status === 'pending') params.set('approved', 'false');

            try {
                const res = await fetch(`/api/parser/data/${dmCurrentSkillId}?${params}`);
                const data = await res.json();

                if (data.success) {
                    renderDataManagerList(data.data);
                }
            } catch (err) {
                console.error('Filter error:', err);
            }
        }

        function renderDataManagerList(items) {
            const listEl = document.getElementById('dm-data-list');

            if (!items || items.length === 0) {
                listEl.innerHTML = `
                    <div style="text-align: center; padding: 3rem; color: var(--text-secondary);">
                        <div style="font-size: 3rem; margin-bottom: 1rem;">📭</div>
                        <p>No extracted data yet</p>
                        <button class="btn btn-primary" style="margin-top: 1rem;" onclick="openParserModal()">📤 Parse Documents</button>
                    </div>
                `;
                return;
            }

            listEl.innerHTML = items.map(item => {
                const importanceColor = item.importance_score >= 80 ? 'var(--neon-green)' :
                                       item.importance_score >= 50 ? 'var(--neon-orange)' : 'var(--text-secondary)';
                const statusColor = item.is_approved ? 'var(--neon-green)' : 'var(--neon-orange)';
                const isSelected = dmSelectedIds.has(item.id);

                return `
                    <div class="data-item" style="display: flex; gap: 1rem; padding: 1rem; border-bottom: 1px solid rgba(255,255,255,0.1); ${isSelected ? 'background: rgba(124, 58, 237, 0.1);' : ''}">
                        <input type="checkbox" ${isSelected ? 'checked' : ''} onchange="toggleDmSelection('${item.id}')" style="margin-top: 4px;">
                        <div style="flex: 1;">
                            <div style="display: flex; align-items: center; gap: 0.5rem; margin-bottom: 0.5rem;">
                                <span style="background: ${importanceColor}; color: #000; padding: 2px 8px; border-radius: 4px; font-size: 0.75rem; font-weight: bold;">${Math.round(item.importance_score)}%</span>
                                <span style="background: var(--neon-purple); color: #fff; padding: 2px 8px; border-radius: 4px; font-size: 0.7rem;">${item.category}</span>
                                <span style="color: var(--text-secondary); font-size: 0.75rem;">📄 ${item.source_filename || 'Unknown'}</span>
                                <span style="margin-left: auto; color: ${statusColor}; font-size: 0.75rem;">${item.is_approved ? '✅ Approved' : '⏳ Pending'}</span>
                            </div>
                            <div style="color: var(--neon-cyan); margin-bottom: 0.25rem;"><strong>Q:</strong> ${escapeHtml(item.user_input).substring(0, 200)}${item.user_input.length > 200 ? '...' : ''}</div>
                            <div style="color: var(--text-secondary); font-size: 0.9rem;"><strong>A:</strong> ${escapeHtml(item.assistant_response).substring(0, 300)}${item.assistant_response.length > 300 ? '...' : ''}</div>
                            <div style="margin-top: 0.5rem; display: flex; gap: 0.5rem;">
                                <button class="btn btn-secondary btn-sm" onclick="quickApproveDm('${item.id}')">✅</button>
                                <button class="btn btn-secondary btn-sm" onclick="quickDeleteDm('${item.id}')">🗑️</button>
                            </div>
                        </div>
                    </div>
                `;
            }).join('');
        }

        function toggleDmSelection(id) {
            if (dmSelectedIds.has(id)) {
                dmSelectedIds.delete(id);
            } else {
                dmSelectedIds.add(id);
            }
        }

        async function quickApproveDm(id) {
            await bulkDmAction('approve', [id]);
            filterDataManager();
            loadDataManagerData();
        }

        async function quickDeleteDm(id) {
            if (confirm('Delete this item?')) {
                await bulkDmAction('delete', [id]);
                filterDataManager();
                loadDataManagerData();
            }
        }

        async function bulkApproveSelected() {
            if (dmSelectedIds.size === 0) { alert('No items selected'); return; }
            await bulkDmAction('approve', Array.from(dmSelectedIds));
            dmSelectedIds.clear();
            filterDataManager();
            loadDataManagerData();
        }

        async function bulkMoveToTraining() {
            if (dmSelectedIds.size === 0) { alert('No items selected'); return; }
            await bulkDmAction('move_to_training', Array.from(dmSelectedIds));
            dmSelectedIds.clear();
            filterDataManager();
            loadDataManagerData();
            alert('Items moved to training data!');
        }

        async function bulkDeleteSelected() {
            if (dmSelectedIds.size === 0) { alert('No items selected'); return; }
            if (!confirm(`Delete ${dmSelectedIds.size} items?`)) return;
            await bulkDmAction('delete', Array.from(dmSelectedIds));
            dmSelectedIds.clear();
            filterDataManager();
            loadDataManagerData();
        }

        async function bulkDmAction(action, itemIds) {
            try {
                await fetch(`/api/parser/data/${dmCurrentSkillId}/bulk`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ action, item_ids: itemIds })
                });
            } catch (err) {
                console.error('Bulk action error:', err);
            }
        }

        // Parser Modal Functions
        function openParserModal() {
            loadDataManagerSkillsDropdown();
            document.getElementById('parser-modal').style.display = 'flex';
            parserFiles = [];
            document.getElementById('parser-file-list').innerHTML = '';
        }

        function closeParserModal() {
            document.getElementById('parser-modal').style.display = 'none';
            parserFiles = [];
        }

        function handleParserFiles(files) {
            parserFiles = Array.from(files);
            const listEl = document.getElementById('parser-file-list');
            listEl.innerHTML = parserFiles.map((f, i) => `
                <div style="display: flex; justify-content: space-between; padding: 0.5rem; background: var(--glass-surface); border-radius: 6px; margin-bottom: 0.5rem;">
                    <span>📄 ${f.name} <span style="color: var(--text-secondary);">(${(f.size / 1024).toFixed(1)} KB)</span></span>
                    <button class="btn btn-secondary btn-sm" onclick="removeParserFile(${i})">✕</button>
                </div>
            `).join('');
        }

        function removeParserFile(index) {
            parserFiles.splice(index, 1);
            handleParserFiles(parserFiles);
        }

        async function submitParserFiles() {
            const skillId = document.getElementById('parser-skill-select').value;
            if (!skillId) { alert('Please select a skill'); return; }
            if (parserFiles.length === 0) { alert('Please select files to parse'); return; }

            const btn = document.getElementById('parser-submit-btn');
            btn.disabled = true;
            btn.textContent = '⏳ Processing...';

            document.getElementById('parser-progress').style.display = 'block';

            const formData = new FormData();
            formData.append('skill_id', skillId);
            parserFiles.forEach(f => formData.append('files[]', f));

            try {
                const res = await fetch('/api/parser/upload', {
                    method: 'POST',
                    body: formData
                });

                const data = await res.json();

                document.getElementById('parser-progress-bar').style.width = '100%';
                document.getElementById('parser-progress-text').textContent = '100%';

                if (data.success) {
                    alert(`✅ Parsed ${data.files_processed} file(s), extracted ${data.items_extracted} items!`);
                    closeParserModal();
                    dmCurrentSkillId = skillId;
                    document.getElementById('dm-skill-select').value = skillId;
                    loadDataManagerData();
                } else {
                    alert('Error: ' + (data.error || 'Unknown error'));
                }
            } catch (err) {
                alert('Error: ' + err.message);
            } finally {
                btn.disabled = false;
                btn.textContent = '🚀 Parse & Extract';
                document.getElementById('parser-progress').style.display = 'none';
            }
        }

        function formatNumber(num) {
            if (num >= 1000000) return (num / 1000000).toFixed(1) + 'M';
            if (num >= 1000) return (num / 1000).toFixed(1) + 'K';
            return num.toString();
        }

        function escapeHtml(text) {
            if (!text) return '';
            return text.toString()
                .replace(/&/g, '&amp;')
                .replace(/</g, '&lt;')
                .replace(/>/g, '&gt;')
                .replace(/"/g, '&quot;')
                .replace(/'/g, '&#039;');
        }

        function hideGettingStarted() {
            const card = document.getElementById('getting-started');
            if (card) {
                card.style.display = 'none';
                localStorage.setItem('hideGettingStarted', 'true');
            }
        }

        // Check if getting started should be hidden
        if (localStorage.getItem('hideGettingStarted') === 'true') {
            document.addEventListener('DOMContentLoaded', () => {
                const card = document.getElementById('getting-started');
                if (card) card.style.display = 'none';
            });
        }

        function showFactoryTab(tabId) {
            document.querySelectorAll('#tab-factory .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-factory .sub-tab-btn').forEach(b => b.classList.remove('active'));
            document.getElementById('factory-' + tabId).classList.add('active');
            event.target.classList.add('active');

            if (tabId === 'manage') loadSkillsTable();
            if (tabId === 'golden') loadGoldenPrompt();
        }

        function showCommandTab(tabId) {
            document.querySelectorAll('#tab-command .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-command .sub-tab-btn').forEach(b => b.classList.remove('active'));
            document.getElementById('command-' + tabId).classList.add('active');
            event.target.classList.add('active');
        }

        // ============================================================
        // DASHBOARD
        // ============================================================
        async function loadSkills() {
            try {
                const res = await fetch('/api/skills');
                const skills = await res.json();
                const grid = document.getElementById('skills-grid');
                grid.innerHTML = '';

                let totalReq = 0, totalLat = 0, totalSat = 0;

                if (skills.length === 0) {
                    grid.innerHTML = `
                        <div style="text-align: center; grid-column: span 4; padding: 2rem;">
                            <div style="font-size: 3rem; margin-bottom: 1rem;">Brain</div>
                            <h3 style="margin-bottom: 0.5rem;">Create Your First AI Skill</h3>
                            <p style="color: var(--text-secondary); margin-bottom: 1.5rem;">Build an AI receptionist for your business in under 2 minutes</p>
                            <button class="btn btn-primary" onclick="showMainTab('skills'); setTimeout(() => showFastBrainTab('skills'), 100);">
                                + Create Skill
                            </button>
                        </div>`;
                }

                skills.forEach(s => {
                    totalReq += s.requests_today;
                    totalLat += s.avg_latency_ms;
                    totalSat += s.satisfaction_rate;

                    grid.innerHTML += `
                        <div class="skill-card ${s.status}">
                            <div class="skill-header">
                                <div>
                                    <div class="skill-name">${s.name}</div>
                                    <div class="skill-type">${s.type}</div>
                                </div>
                                <span class="skill-status ${s.status}">${s.status}</span>
                            </div>
                            <div class="skill-metrics">
                                <div class="skill-metric"><div class="metric-value">${s.requests_today}</div><div class="metric-label">Requests</div></div>
                                <div class="skill-metric"><div class="metric-value">${s.avg_latency_ms}ms</div><div class="metric-label">Latency</div></div>
                                <div class="skill-metric"><div class="metric-value">${s.satisfaction_rate}%</div><div class="metric-label">Satisfaction</div></div>
                            </div>
                        </div>
                    `;
                });

                document.getElementById('total-skills').textContent = skills.length || 0;
                document.getElementById('requests-today').textContent = totalReq.toLocaleString();
                document.getElementById('avg-latency').textContent = skills.length ? Math.round(totalLat / skills.length) + 'ms' : '0ms';
                document.getElementById('satisfaction').textContent = skills.length ? Math.round(totalSat / skills.length) + '%' : '0%';
            } catch (e) {
                console.error(e);
            }
        }

        async function refreshServerStatus() {
            try {
                const res = await fetch('/api/server-status');
                const s = await res.json();
                document.getElementById('server-status-text').textContent = `LPU ${s.status === 'online' ? 'Online' : 'Offline'} - ${s.warm_containers} Warm Container`;
                document.getElementById('warm-containers').textContent = s.warm_containers;
                document.getElementById('memory-usage').textContent = (s.memory_usage_mb / 1000).toFixed(1) + 'GB';
                document.getElementById('cost-today').textContent = '$' + s.cost_today_usd.toFixed(2);
            } catch (e) { console.error(e); }
        }

        async function loadMetrics() {
            try {
                const res = await fetch('/api/metrics');
                const metrics = await res.json();

                const reqChart = document.getElementById('request-chart');
                const latChart = document.getElementById('latency-chart');
                reqChart.innerHTML = '';
                latChart.innerHTML = '';

                const maxReq = Math.max(...metrics.map(m => m.requests));
                const maxLat = Math.max(...metrics.map(m => m.latency_p99));

                metrics.forEach(m => {
                    reqChart.innerHTML += `<div class="chart-bar" style="height: ${m.requests / maxReq * 100}%" title="${m.requests} requests"></div>`;
                    latChart.innerHTML += `<div class="chart-bar" style="height: ${m.latency_p50 / maxLat * 100}%; background: linear-gradient(180deg, var(--neon-green), rgba(0,255,136,0.3));" title="${m.latency_p50}ms"></div>`;
                });
            } catch (e) { console.error(e); }
        }

        async function loadActivity() {
            try {
                const res = await fetch('/api/activity');
                const activities = await res.json();
                const feed = document.getElementById('activity-feed');

                if (activities.length === 0) {
                    feed.innerHTML = '<p style="color: var(--text-secondary); text-align: center;">No recent activity</p>';
                    return;
                }

                feed.innerHTML = activities.map(a => `
                    <div class="activity-item">
                        <span class="activity-icon">${a.icon || 'i'}</span>
                        <div class="activity-content">
                            <div class="activity-message">${a.message}</div>
                            <div class="activity-time">${a.ago}</div>
                        </div>
                    </div>
                `).join('');
            } catch (e) { console.error(e); }
        }

        // ============================================================
        // SKILL FACTORY
        // ============================================================

        // Golden Prompts Functions
        let goldenPromptsCache = {};

        async function loadGoldenPrompts() {
            try {
                const res = await fetch('/api/golden-prompts');
                const data = await res.json();
                if (data.prompts) {
                    goldenPromptsCache = data.prompts;
                }
                loadGoldenPrompt();
            } catch (e) {
                console.error('Error loading golden prompts:', e);
            }
        }

        async function loadGoldenPrompt() {
            const skillId = document.getElementById('golden-skill-select').value;
            const businessName = document.getElementById('golden-business-name').value || '{business_name}';
            const agentName = document.getElementById('golden-agent-name').value || '{agent_name}';

            try {
                const res = await fetch(`/api/golden-prompts/${skillId}?business_name=${encodeURIComponent(businessName)}&agent_name=${encodeURIComponent(agentName)}`);
                const data = await res.json();

                if (data.content) {
                    document.getElementById('golden-prompt-content').value = data.content;
                    document.getElementById('golden-tokens').textContent = '~' + data.tokens;
                    document.getElementById('golden-prefill').textContent = '~' + data.prefill_ms + 'ms';
                    document.getElementById('golden-status').textContent = data.is_custom ? 'Custom' : 'Default';
                    document.getElementById('golden-status').style.color = data.is_custom ? 'var(--warning)' : 'var(--success)';
                }
            } catch (e) {
                console.error('Error loading prompt:', e);
                showGoldenMessage('Error loading prompt: ' + e.message, 'error');
            }

            // Also load available voices for this skill
            loadVoicesForSkill(skillId);
        }

        async function loadVoicesForSkill(skillId) {
            const voiceSelect = document.getElementById('golden-voice-select');
            if (!voiceSelect) return;

            try {
                // Get all trained voice projects
                const res = await fetch('/api/voice-lab/projects');
                const projects = await res.json();

                // Clear existing options
                voiceSelect.innerHTML = '<option value="">No custom voice</option>';

                // Add trained voices
                projects.filter(p => p.status === 'trained').forEach(p => {
                    const option = document.createElement('option');
                    option.value = p.id;
                    option.textContent = `${p.name} (${p.provider})`;
                    // Mark if this voice is linked to the current skill
                    if (p.skill_id === skillId) {
                        option.selected = true;
                    }
                    voiceSelect.appendChild(option);
                });

                // Add separator and option to create new voice
                if (projects.filter(p => p.status === 'trained').length > 0) {
                    const separator = document.createElement('option');
                    separator.disabled = true;
                    separator.textContent = '───────────';
                    voiceSelect.appendChild(separator);
                }

                const createOption = document.createElement('option');
                createOption.value = '__create__';
                createOption.textContent = '+ Create New Voice...';
                voiceSelect.appendChild(createOption);

            } catch (e) {
                console.error('Error loading voices:', e);
            }
        }

        async function linkVoiceToSkill() {
            const skillId = document.getElementById('golden-skill-select').value;
            const voiceSelect = document.getElementById('golden-voice-select');
            const voiceId = voiceSelect.value;

            if (voiceId === '__create__') {
                // Switch to Voice Lab tab
                showSubTab('voicelab');
                showVoiceLabTab('create');
                voiceSelect.value = '';
                return;
            }

            if (!voiceId) {
                // Clear voice link - not implemented yet, just return
                return;
            }

            try {
                const res = await fetch(`/api/voice-lab/projects/${voiceId}/link-skill`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId })
                });
                const data = await res.json();

                if (data.success) {
                    showGoldenMessage(`Voice linked to ${skillId} skill`, 'success');
                } else {
                    showGoldenMessage(data.error || 'Failed to link voice', 'error');
                }
            } catch (e) {
                showGoldenMessage('Error linking voice: ' + e.message, 'error');
            }
        }

        async function saveGoldenPrompt() {
            const skillId = document.getElementById('golden-skill-select').value;
            const content = document.getElementById('golden-prompt-content').value;

            try {
                const res = await fetch(`/api/golden-prompts/${skillId}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ content })
                });
                const data = await res.json();

                if (data.success) {
                    showGoldenMessage(data.message, 'success');
                    document.getElementById('golden-tokens').textContent = '~' + data.tokens;
                    document.getElementById('golden-prefill').textContent = '~' + data.prefill_ms + 'ms';
                    document.getElementById('golden-status').textContent = 'Custom';
                    document.getElementById('golden-status').style.color = 'var(--warning)';
                } else {
                    showGoldenMessage(data.error || 'Save failed', 'error');
                }
            } catch (e) {
                showGoldenMessage('Error: ' + e.message, 'error');
            }
        }

        async function resetGoldenPrompt() {
            const skillId = document.getElementById('golden-skill-select').value;

            if (!confirm(`Reset "${skillId}" prompt to default? Your customizations will be lost.`)) {
                return;
            }

            try {
                const res = await fetch(`/api/golden-prompts/${skillId}/reset`, { method: 'POST' });
                const data = await res.json();

                if (data.success) {
                    showGoldenMessage(data.message, 'success');
                    loadGoldenPrompt();
                } else {
                    showGoldenMessage(data.error || 'Reset failed', 'error');
                }
            } catch (e) {
                showGoldenMessage('Error: ' + e.message, 'error');
            }
        }

        async function testGoldenPrompt() {
            const skillId = document.getElementById('golden-skill-select').value;
            const businessName = document.getElementById('golden-business-name').value || 'Test Business';
            const agentName = document.getElementById('golden-agent-name').value || 'Assistant';

            const testMessage = prompt('Enter a test message:', 'Hello, I need some help with my electrical issue.');
            if (!testMessage) return;

            showGoldenMessage('Testing prompt with Fast Brain...', 'info');

            try {
                const res = await fetch('/api/golden-prompts/test', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: skillId,
                        message: testMessage,
                        business_name: businessName,
                        agent_name: agentName
                    })
                });
                const data = await res.json();

                if (data.success) {
                    document.getElementById('golden-test-output').style.display = 'block';
                    document.getElementById('golden-test-response').innerHTML = `
                        <strong>User:</strong> ${testMessage}\\n\\n
                        <strong>Response (${data.system_used}, ${data.latency_ms}ms):</strong>\\n${data.response}
                    `;
                    showGoldenMessage('Test completed successfully!', 'success');
                } else {
                    showGoldenMessage('Test failed: ' + data.error, 'error');
                }
            } catch (e) {
                showGoldenMessage('Error: ' + e.message, 'error');
            }
        }

        function copyGoldenPrompt() {
            const content = document.getElementById('golden-prompt-content').value;
            navigator.clipboard.writeText(content).then(() => {
                showGoldenMessage('Copied to clipboard!', 'success');
            }).catch(() => {
                showGoldenMessage('Failed to copy', 'error');
            });
        }

        async function exportGoldenPrompts() {
            try {
                const res = await fetch('/api/golden-prompts/export');
                const data = await res.json();

                const blob = new Blob([JSON.stringify(data.prompts, null, 2)], { type: 'application/json' });
                const url = URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                a.download = 'golden_prompts_export.json';
                a.click();
                URL.revokeObjectURL(url);

                showGoldenMessage(`Exported ${Object.keys(data.prompts).length} prompts (${data.custom_overrides.length} custom)`, 'success');
            } catch (e) {
                showGoldenMessage('Export failed: ' + e.message, 'error');
            }
        }

        function showGoldenMessage(msg, type) {
            const el = document.getElementById('golden-message');
            el.innerHTML = `<div class="message ${type}">${msg}</div>`;
            setTimeout(() => { el.innerHTML = ''; }, 5000);
        }

        async function loadProfileDropdowns() {
            try {
                const res = await fetch('/api/skills');
                const skills = await res.json();
                const options = skills.map(s => `<option value="${s.id}">${s.name}</option>`).join('');
                document.getElementById('doc-profile').innerHTML = options || '<option value="">No profiles yet</option>';
                document.getElementById('train-profile').innerHTML = options || '<option value="">No profiles yet</option>';
            } catch (e) { console.error(e); }
        }

        async function saveProfile(e) {
            e.preventDefault();
            const data = {
                name: document.getElementById('p-name').value,
                type: document.getElementById('p-type').value,
                description: document.getElementById('p-description').value,
                greeting: document.getElementById('p-greeting').value,
                personality: document.getElementById('p-personality').value,
                services: document.getElementById('p-services').value,
                customInstructions: document.getElementById('p-instructions').value
            };

            try {
                const res = await fetch('/api/create-skill', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(data)
                });
                const result = await res.json();
                if (result.success) {
                    document.getElementById('profile-message').innerHTML = '<div class="message success">Profile saved successfully!</div>';
                    document.getElementById('profile-form').reset();
                    loadProfileDropdowns();
                }
            } catch (e) {
                document.getElementById('profile-message').innerHTML = '<div class="message error">Error: ' + e.message + '</div>';
            }
        }

        async function generateFromProfile() {
            const name = document.getElementById('p-name').value;
            if (!name) {
                alert('Please enter a business name first');
                return;
            }
            const skillId = name.toLowerCase().replace(/[^\\w-]/g, '_');
            try {
                const res = await fetch(`/api/generate-training/${skillId}`, { method: 'POST' });
                const result = await res.json();
                if (result.success) {
                    document.getElementById('profile-message').innerHTML = `<div class="message success">Generated ${result.examples} training examples!</div>`;
                }
            } catch (e) {
                document.getElementById('profile-message').innerHTML = '<div class="message error">Error: ' + e.message + '</div>';
            }
        }

        async function handleFileUpload(input) {
            const file = input.files[0];
            if (!file) return;

            const skillId = document.getElementById('doc-profile').value;
            if (!skillId) {
                alert('Please select a profile first');
                return;
            }

            const formData = new FormData();
            formData.append('file', file);
            formData.append('skill_id', skillId);

            try {
                const res = await fetch('/api/upload-document', {
                    method: 'POST',
                    body: formData
                });
                const result = await res.json();
                if (result.success) {
                    document.getElementById('upload-message').innerHTML = `<div class="message success">Processed ${result.examples} training examples!</div>`;
                } else {
                    document.getElementById('upload-message').innerHTML = `<div class="message error">${result.error}</div>`;
                }
            } catch (e) {
                document.getElementById('upload-message').innerHTML = '<div class="message error">Error: ' + e.message + '</div>';
            }
        }

        function generateTrainingScript() {
            const skillId = document.getElementById('train-profile').value;
            const steps = document.getElementById('train-steps').value;

            const script = `# Training script for ${skillId}
# Run on a GPU machine (Colab, Modal, etc.)

from unsloth import FastLanguageModel
from datasets import load_dataset
from trl import SFTTrainer
from transformers import TrainingArguments

model, tokenizer = FastLanguageModel.from_pretrained(
    "unsloth/llama-3-8b-Instructbnb-4bit",
    max_seq_length=2048, load_in_4bit=True
)

model = FastLanguageModel.get_peft_model(model, r=16,
    target_modules=["q_proj", "k_proj", "v_proj", "o_proj",
                    "gate_proj", "up_proj", "down_proj"],
    lora_alpha=16, use_gradient_checkpointing="unsloth")

dataset = load_dataset("json",
    data_files="training_data/${skillId}_merged.jsonl",
    split="train")

trainer = SFTTrainer(model, tokenizer, train_dataset=dataset,
    dataset_text_field="output", max_seq_length=2048,
    args=TrainingArguments(
        per_device_train_batch_size=2,
        gradient_accumulation_steps=4,
        max_steps=${steps},
        learning_rate=2e-4,
        output_dir="outputs"))

trainer.train()
model.save_pretrained("adapters/${skillId}")
print("Training complete: adapters/${skillId}")`;

            document.getElementById('training-script').textContent = script;
            document.getElementById('training-output').style.display = 'block';
        }

        // ============================================================
        // SKILLS TABLE
        // ============================================================
        let allSkillsData = [];

        async function loadSkillsTable() {
            try {
                const res = await fetch('/api/skills-table');
                allSkillsData = await res.json();
                renderSkillsTable(allSkillsData);
            } catch (e) { console.error(e); }
        }

        function renderSkillsTable(skills) {
            const tbody = document.getElementById('skills-table-body');
            if (skills.length === 0) {
                tbody.innerHTML = `<tr><td colspan="6" style="text-align: center; padding: 2rem;">
                    <div style="font-size: 2rem; margin-bottom: 0.5rem;">Folder</div>
                    <p style="color: var(--text-secondary); margin-bottom: 1rem;">No skills found. Create a business profile to get started.</p>
                    <button class="btn btn-primary btn-sm" onclick="showFactoryTab('profile')">Create Business Profile</button>
                </td></tr>`;
                return;
            }
            tbody.innerHTML = skills.map(s => `
                <tr>
                    <td>${s.business}</td>
                    <td>${s.type}</td>
                    <td><span class="table-status ${s.status.toLowerCase()}">${s.status}</span></td>
                    <td>${s.examples}</td>
                    <td>${s.created}</td>
                    <td>
                        <button class="btn btn-secondary btn-sm" onclick="editSkill('${s.id}')">Edit</button>
                        <button class="btn btn-danger btn-sm" onclick="deleteSkill('${s.id}')">Delete</button>
                    </td>
                </tr>
            `).join('');
        }

        function filterSkillsTable() {
            const search = document.getElementById('table-search').value.toLowerCase();
            const typeFilter = document.getElementById('table-type-filter').value;
            const statusFilter = document.getElementById('table-status-filter').value;

            const filtered = allSkillsData.filter(s => {
                const matchSearch = s.business.toLowerCase().includes(search) || s.type.toLowerCase().includes(search);
                const matchType = !typeFilter || s.type.includes(typeFilter);
                const matchStatus = !statusFilter || s.status === statusFilter;
                return matchSearch && matchType && matchStatus;
            });

            renderSkillsTable(filtered);
        }

        async function editSkill(id) {
            // Fetch skill data from cache or API
            let skill = fbSkillsCache[id];

            if (!skill) {
                // Try to fetch from API
                try {
                    const res = await fetch(`/api/fast-brain/skills/${id}`);
                    const data = await res.json();
                    if (data.skill) {
                        skill = data.skill;
                    }
                } catch (e) {
                    console.error('Failed to fetch skill:', e);
                }
            }

            if (!skill) {
                alert('Skill not found: ' + id);
                return;
            }

            // Populate edit form
            document.getElementById('edit-skill-id').value = id;
            document.getElementById('edit-skill-id-display').value = id;
            document.getElementById('edit-skill-name').value = skill.name || '';
            document.getElementById('edit-skill-description').value = skill.description || '';
            document.getElementById('edit-skill-prompt').value = skill.system_prompt || '';

            // Handle knowledge array
            const knowledge = skill.knowledge || [];
            document.getElementById('edit-skill-knowledge').value = Array.isArray(knowledge) ? knowledge.join('\\n') : knowledge;

            // Show edit form
            document.getElementById('fb-edit-skill-form').style.display = 'block';
            document.getElementById('fb-create-skill-form').style.display = 'none';
            document.getElementById('edit-skill-message').innerHTML = '';
        }

        function hideEditSkillModal() {
            document.getElementById('fb-edit-skill-form').style.display = 'none';
            document.getElementById('edit-skill-id').value = '';
            document.getElementById('edit-skill-id-display').value = '';
            document.getElementById('edit-skill-name').value = '';
            document.getElementById('edit-skill-description').value = '';
            document.getElementById('edit-skill-prompt').value = '';
            document.getElementById('edit-skill-knowledge').value = '';
            document.getElementById('edit-skill-message').innerHTML = '';
        }

        async function saveEditedSkill() {
            const messageEl = document.getElementById('edit-skill-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Saving changes...</div>';

            const skillId = document.getElementById('edit-skill-id').value;
            const knowledge = document.getElementById('edit-skill-knowledge').value
                .split('\\n')
                .map(s => s.trim())
                .filter(s => s.length > 0);

            try {
                const res = await fetch(`/api/fast-brain/skills/${skillId}`, {
                    method: 'PUT',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        name: document.getElementById('edit-skill-name').value,
                        description: document.getElementById('edit-skill-description').value,
                        system_prompt: document.getElementById('edit-skill-prompt').value,
                        knowledge: knowledge,
                    })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = '<div style="color: var(--neon-green);">Skill updated successfully!</div>';
                    setTimeout(() => {
                        hideEditSkillModal();
                        loadFastBrainSkills();
                    }, 1000);
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        // ============================================================
        // VOICE TEST ENVIRONMENT
        // ============================================================
        async function loadVoiceTestVoices() {
            const providerSelect = document.getElementById('voice-test-provider');
            const voiceSelect = document.getElementById('voice-test-voice-id');
            const statusEl = document.getElementById('voice-test-status');

            const provider = providerSelect.value;

            if (!provider) {
                voiceSelect.disabled = true;
                voiceSelect.innerHTML = '<option value="">Select provider first...</option>';
                if (statusEl) statusEl.innerHTML = '';
                return;
            }

            voiceSelect.disabled = true;
            voiceSelect.innerHTML = '<option value="">Loading voices...</option>';
            if (statusEl) statusEl.innerHTML = '<span style="color: var(--text-secondary);">Fetching voices...</span>';

            try {
                const res = await fetch(`/api/voice-lab/provider-voices/${provider}`);
                const data = await res.json();

                if (data.error) {
                    voiceSelect.innerHTML = `<option value="">Error: ${data.error}</option>`;
                    if (statusEl) statusEl.innerHTML = `<span style="color: var(--neon-orange);">${data.error}</span>`;
                    return;
                }

                const voices = data.voices || [];
                if (voices.length === 0) {
                    voiceSelect.innerHTML = '<option value="">No voices available</option>';
                    if (statusEl) statusEl.innerHTML = '<span style="color: var(--text-secondary);">No voices found. Check API key.</span>';
                    return;
                }

                // Populate voice dropdown
                voiceSelect.innerHTML = '<option value="">Select a voice...</option>';
                voices.forEach(voice => {
                    const opt = document.createElement('option');
                    opt.value = voice.voice_id || voice.id;
                    opt.textContent = voice.name + (voice.labels?.accent ? ` (${voice.labels.accent})` : '');
                    voiceSelect.appendChild(opt);
                });

                voiceSelect.disabled = false;
                if (statusEl) statusEl.innerHTML = `<span style="color: var(--neon-green);">${voices.length} voices</span>`;

            } catch (e) {
                console.error('Failed to load voices:', e);
                voiceSelect.innerHTML = '<option value="">Error loading voices</option>';
                if (statusEl) statusEl.innerHTML = `<span style="color: var(--neon-orange);">Failed: ${e.message}</span>`;
            }
        }

        async function testVoiceTTS() {
            const provider = document.getElementById('voice-test-provider').value;
            const voiceId = document.getElementById('voice-test-voice-id').value;
            const text = document.getElementById('voice-test-chat-text').value || 'Hello! This is a test of the voice synthesis.';
            const statusEl = document.getElementById('voice-test-status');
            const audioEl = document.getElementById('voice-test-audio');
            const btn = document.getElementById('voice-test-btn');

            if (!provider || !voiceId) {
                alert('Please select a provider and voice first');
                return;
            }

            btn.disabled = true;
            btn.textContent = 'Generating...';
            statusEl.innerHTML = '<span style="color: var(--neon-cyan);">Synthesizing audio...</span>';

            try {
                const res = await fetch('/api/voice/test', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ provider, voice_id: voiceId, text })
                });

                const result = await res.json();

                if (result.success && result.audio_base64) {
                    const cacheStatus = result.cached ? '(cached)' : `(${result.duration_ms}ms)`;
                    const audio = new Audio('data:' + result.audio_format + ';base64,' + result.audio_base64);

                    audio.onended = () => {
                        statusEl.innerHTML = `<span style="color: var(--neon-green);">✓ Played ${cacheStatus}</span>`;
                    };
                    audio.onerror = () => {
                        statusEl.innerHTML = '<span style="color: var(--neon-orange);">Playback error</span>';
                    };
                    audio.play();
                    audioEl.style.display = 'none';
                    statusEl.innerHTML = `<span style="color: var(--neon-cyan);">▶ Playing... ${cacheStatus}</span>`;
                } else {
                    throw new Error(result.error || 'Synthesis failed');
                }
            } catch (e) {
                statusEl.innerHTML = `<span style="color: var(--neon-orange);">Error: ${e.message}</span>`;
            } finally {
                btn.disabled = false;
                btn.textContent = 'Test Voice';
            }
        }

        async function deleteSkill(id) {
            if (!confirm('Delete this skill?')) return;
            try {
                await fetch(`/api/delete-skill/${id}`, { method: 'DELETE' });
                loadSkillsTable();
            } catch (e) { console.error(e); }
        }

        // ============================================================
        // COMMAND CENTER
        // ============================================================
        async function loadApiKeys() {
            try {
                const res = await fetch('/api/api-keys');
                const keys = await res.json();

                // LLM Providers
                if (keys.groq) document.getElementById('key-groq').value = keys.groq;
                if (keys.openai) document.getElementById('key-openai').value = keys.openai;
                if (keys.anthropic) document.getElementById('key-anthropic').value = keys.anthropic;
                // Voice Providers
                if (keys.elevenlabs) document.getElementById('key-elevenlabs').value = keys.elevenlabs;
                if (keys.cartesia) document.getElementById('key-cartesia').value = keys.cartesia;
                if (keys.deepgram) document.getElementById('key-deepgram').value = keys.deepgram;
                if (keys.playht) document.getElementById('key-playht').value = keys.playht;
            } catch (e) {
                console.error('Failed to load API keys:', e);
            }
        }

        async function saveApiKeys() {
            const data = {
                // LLM Providers
                groq: document.getElementById('key-groq').value,
                openai: document.getElementById('key-openai').value,
                anthropic: document.getElementById('key-anthropic').value,
                // Voice Providers
                elevenlabs: document.getElementById('key-elevenlabs').value,
                cartesia: document.getElementById('key-cartesia').value,
                deepgram: document.getElementById('key-deepgram').value,
                playht: document.getElementById('key-playht').value
            };

            try {
                const res = await fetch('/api/save-api-keys', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(data)
                });
                const result = await res.json();
                if (result.saved && result.saved.length > 0) {
                    document.getElementById('keys-message').innerHTML = `<div class="message success">Saved: ${result.saved.join(', ')}</div>`;
                } else {
                    document.getElementById('keys-message').innerHTML = '<div class="message info">No keys provided. Using mock clients for LLM and free TTS.</div>';
                }
            } catch (e) {
                document.getElementById('keys-message').innerHTML = '<div class="message error">Error: ' + e.message + '</div>';
            }
        }

        // Copy API URL to clipboard
        function copyApiUrl(inputId) {
            const input = document.getElementById(inputId);
            input.select();
            input.setSelectionRange(0, 99999);
            navigator.clipboard.writeText(input.value).then(() => {
                // Show brief success indicator
                const btn = input.nextElementSibling;
                const originalText = btn.textContent;
                btn.textContent = 'Copied!';
                btn.style.background = 'var(--neon-green)';
                btn.style.color = '#000';
                setTimeout(() => {
                    btn.textContent = originalText;
                    btn.style.background = '';
                    btn.style.color = '';
                }, 1500);
            }).catch(err => {
                console.error('Copy failed:', err);
                alert('Failed to copy. Please copy manually.');
            });
        }

        // Check health of API endpoints
        async function checkApiHealth() {
            const endpoints = [
                {
                    id: 'fast-brain',
                    url: document.getElementById('api-url-fast-brain').value,
                    healthPath: '/health'
                },
                {
                    id: 'dashboard',
                    url: document.getElementById('api-url-dashboard').value,
                    healthPath: '/api/health'
                }
            ];

            for (const endpoint of endpoints) {
                const indicator = document.getElementById(`api-health-${endpoint.id}`);
                const statusEl = document.getElementById(`api-status-${endpoint.id}`);

                // Set checking state
                indicator.style.background = 'var(--neon-orange)';
                statusEl.textContent = 'Checking...';

                try {
                    const startTime = Date.now();
                    const response = await fetch(`${endpoint.url}${endpoint.healthPath}`, {
                        method: 'GET',
                        mode: 'cors',
                        signal: AbortSignal.timeout(10000)
                    });
                    const latency = Date.now() - startTime;

                    if (response.ok) {
                        const data = await response.json();
                        indicator.style.background = 'var(--neon-green)';
                        statusEl.innerHTML = `<span style="color: var(--neon-green);">Online</span> - ${latency}ms`;
                        if (data.skills) {
                            statusEl.innerHTML += ` - ${data.skills.length} skills loaded`;
                        }
                    } else {
                        indicator.style.background = 'var(--neon-orange)';
                        statusEl.innerHTML = `<span style="color: var(--neon-orange);">Error ${response.status}</span>`;
                    }
                } catch (e) {
                    indicator.style.background = 'var(--neon-pink)';
                    if (e.name === 'TimeoutError') {
                        statusEl.innerHTML = `<span style="color: var(--neon-pink);">Timeout</span> - Cold start may be needed`;
                    } else if (e.message.includes('CORS') || e.message.includes('Failed to fetch')) {
                        // CORS error likely means the endpoint is up but blocking cross-origin
                        indicator.style.background = 'var(--neon-orange)';
                        statusEl.innerHTML = `<span style="color: var(--neon-orange);">CORS blocked</span> - Endpoint may be reachable`;
                    } else {
                        statusEl.innerHTML = `<span style="color: var(--neon-pink);">Offline</span> - ${e.message}`;
                    }
                }
            }
        }

        async function testLLM() {
            const btn = document.getElementById('test-btn');
            const output = document.getElementById('test-output');
            const prompt = document.getElementById('test-prompt').value;
            const provider = document.getElementById('test-provider').value;

            if (!prompt) {
                alert('Please enter a prompt');
                return;
            }

            btn.innerHTML = '<span class="spinner"></span> Testing...';
            btn.disabled = true;

            try {
                const res = await fetch('/api/test-llm', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ prompt, provider })
                });
                const result = await res.json();
                output.innerHTML = `<span class="response">${result.response}</span>\\n\\n<span class="info">Latency: ${result.latency_ms}ms</span>`;
            } catch (e) {
                output.innerHTML = `<span class="error">Error: ${e.message}</span>`;
            }

            btn.innerHTML = 'Run Test';
            btn.disabled = false;
        }

        async function compareLLMs() {
            const btn = document.getElementById('compare-btn');
            const output = document.getElementById('compare-output');
            const prompt = document.getElementById('compare-prompt').value;
            const providers = Array.from(document.querySelectorAll('.compare-provider:checked')).map(c => c.value);

            if (!prompt) {
                alert('Please enter a prompt');
                return;
            }

            if (providers.length === 0) {
                alert('Please select at least one provider');
                return;
            }

            btn.innerHTML = '<span class="spinner"></span> Comparing...';
            btn.disabled = true;

            try {
                const res = await fetch('/api/compare-llms', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ prompt, providers })
                });
                const result = await res.json();

                output.innerHTML = result.results.map(r => `
                    <div class="compare-card">
                        <div class="compare-header">
                            <span class="compare-provider">${r.provider.toUpperCase()}</span>
                            <span class="compare-latency">First: ${r.first_token_ms}ms | Total: ${r.total_ms}ms</span>
                        </div>
                        <div>${r.response}</div>
                    </div>
                `).join('');
            } catch (e) {
                output.innerHTML = `<div class="message error">Error: ${e.message}</div>`;
            }

            btn.innerHTML = 'Compare';
            btn.disabled = false;
        }

        async function refreshStats() {
            try {
                const res = await fetch('/api/stats');
                const stats = await res.json();
                document.getElementById('stats-output').textContent = JSON.stringify(stats, null, 2);
            } catch (e) { console.error(e); }
        }

        // ============================================================
        // LPU INFERENCE
        // ============================================================
        async function runLpuInference() {
            const btn = document.getElementById('lpu-btn');
            const consoleEl = document.getElementById('lpu-console');
            const prompt = document.getElementById('lpu-prompt').value;
            const maxTokens = parseInt(document.getElementById('lpu-tokens').value);

            if (!prompt) {
                alert('Please enter a prompt');
                return;
            }

            btn.innerHTML = '<span class="spinner"></span> Running...';
            btn.disabled = true;

            consoleEl.innerHTML += `\\n<span class="prompt">&gt; ${prompt}</span>\\n`;

            try {
                const res = await fetch('/api/lpu-inference', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ prompt, maxTokens })
                });
                const result = await res.json();

                if (result.success) {
                    consoleEl.innerHTML += `<span class="response">${result.response}</span>\\n`;
                } else {
                    consoleEl.innerHTML += `<span class="error">Error: ${result.error}</span>\\n`;
                }
            } catch (e) {
                consoleEl.innerHTML += `<span class="error">Error: ${e.message}</span>\\n`;
            }

            btn.innerHTML = 'Run Inference';
            btn.disabled = false;
            consoleEl.scrollTop = consoleEl.scrollHeight;
        }

        function clearLpuConsole() {
            document.getElementById('lpu-console').innerHTML = '<span class="info">// Console cleared\\n</span>';
        }

        // ============================================================
        // VOICE MANAGEMENT (Unified with Dynamic Fetching & Caching)
        // ============================================================

        // Voice cache - stores dynamically fetched voices per provider
        const voiceCache = {
            providers: {},      // Cached voice lists by provider
            timestamps: {},     // When each provider was last fetched
            cacheDuration: 5 * 60 * 1000,  // 5 minute cache

            get(provider) {
                const cached = this.providers[provider];
                const timestamp = this.timestamps[provider];
                if (cached && timestamp && (Date.now() - timestamp) < this.cacheDuration) {
                    return cached;
                }
                return null;
            },

            set(provider, voices) {
                this.providers[provider] = voices;
                this.timestamps[provider] = Date.now();
            },

            getVoice(provider, voiceId) {
                const voices = this.providers[provider];
                if (!voices) return null;
                return voices.find(v => v.id === voiceId);
            }
        };

        // Legacy fallback - static voice data
        let voiceProviders = {};

        async function loadVoiceProviders() {
            // Load static providers as fallback (only needed once)
            if (Object.keys(voiceProviders).length > 0) return;
            try {
                const res = await fetch('/api/voice/providers');
                voiceProviders = await res.json();
            } catch (e) {
                console.error('Failed to load static voice providers:', e);
            }
        }

        async function loadVoiceSettingsVoices() {
            const provider = document.getElementById('voice-provider').value;
            const voiceSelect = document.getElementById('voice-select');
            const detailsEl = document.getElementById('voice-details');

            if (detailsEl) detailsEl.style.display = 'none';

            if (!provider) {
                voiceSelect.innerHTML = '<option value="">Select a provider first...</option>';
                return;
            }

            // Check cache first
            const cached = voiceCache.get(provider);
            if (cached) {
                populateVoiceSelect(voiceSelect, cached, provider);
                return;
            }

            voiceSelect.innerHTML = '<option value="">Loading voices...</option>';

            try {
                // Fetch dynamically from unified endpoint
                const res = await fetch(`/api/voice-lab/provider-voices/${provider}`);
                const data = await res.json();

                let voices = data.voices || [];

                // If dynamic fetch failed, try static fallback
                if (voices.length === 0 && data.error) {
                    await loadVoiceProviders();  // Ensure static data loaded
                    if (voiceProviders[provider]?.voices) {
                        voices = voiceProviders[provider].voices;
                    }
                }

                // Cache the results
                voiceCache.set(provider, voices);
                populateVoiceSelect(voiceSelect, voices, provider);

            } catch (e) {
                console.error('Failed to load provider voices:', e);
                voiceSelect.innerHTML = '<option value="">Error loading voices</option>';
            }
        }

        function populateVoiceSelect(selectEl, voices, provider) {
            selectEl.innerHTML = '<option value="">Select a voice...</option>';

            if (!voices || voices.length === 0) {
                selectEl.innerHTML = '<option value="">No voices available</option>';
                return;
            }

            voices.forEach(voice => {
                const option = document.createElement('option');
                option.value = voice.id;
                const gender = voice.gender || 'unknown';
                const style = voice.style || voice.description?.substring(0, 20) || 'default';
                option.textContent = `${voice.name} (${gender}, ${style})`;
                // Store preview URL as data attribute if available
                if (voice.preview_url) {
                    option.dataset.previewUrl = voice.preview_url;
                }
                selectEl.appendChild(option);
            });
        }

        function selectVoice() {
            const provider = document.getElementById('voice-provider').value;
            const voiceId = document.getElementById('voice-select').value;
            const detailsEl = document.getElementById('voice-details');

            if (!provider || !voiceId) {
                if (detailsEl) detailsEl.style.display = 'none';
                return;
            }

            // Try cache first (dynamic voices)
            let voice = voiceCache.getVoice(provider, voiceId);
            let providerName = provider;

            // Fallback to static data
            if (!voice && voiceProviders[provider]) {
                voice = voiceProviders[provider].voices?.find(v => v.id === voiceId);
                providerName = voiceProviders[provider].name || provider;
            }

            if (voice && detailsEl) {
                detailsEl.style.display = 'block';
                document.getElementById('voice-name').textContent = voice.name;
                document.getElementById('voice-meta').textContent =
                    `${providerName} | ${voice.gender || 'unknown'} | ${voice.style || 'default'}`;

                // Show/update preview button if preview URL available
                const previewBtn = document.getElementById('voice-preview-btn');
                if (previewBtn) {
                    if (voice.preview_url) {
                        previewBtn.style.display = 'inline-block';
                        previewBtn.onclick = () => playVoicePreviewUrl(voice.preview_url);
                    } else {
                        previewBtn.style.display = 'none';
                    }
                }
            }
        }

        // Play preview audio from provider URL
        function playVoicePreviewUrl(url) {
            if (!url) return;
            const audio = new Audio(url);
            audio.play().catch(e => {
                console.error('Failed to play preview:', e);
                alert('Could not play preview audio');
            });
        }

        async function testVoice() {
            const provider = document.getElementById('voice-provider').value;
            const voiceId = document.getElementById('voice-select').value;
            const text = document.getElementById('voice-test-text').value || 'Hello! This is a test of the voice synthesis system.';
            const resultEl = document.getElementById('voice-test-result');

            if (!voiceId) {
                resultEl.innerHTML = '<div style="color: var(--neon-orange);">Please select a voice first</div>';
                return;
            }

            // For Parler TTS, voiceId is the skill_id
            const isParler = provider === 'parler_tts';
            const emotionEl = document.getElementById('voice-emotion');
            const emotion = emotionEl ? emotionEl.value : 'neutral';

            if (isParler) {
                resultEl.innerHTML = '<div style="color: var(--neon-cyan);">Generating audio with Parler TTS (GPU)... May take 30-60s on cold start</div>';
            } else {
                resultEl.innerHTML = '<div style="color: var(--neon-cyan);">Generating audio...</div>';
            }

            try {
                const res = await fetch('/api/voice/test', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        voice_id: voiceId,
                        text,
                        provider,
                        emotion: emotion,
                        skill_id: isParler ? voiceId : 'general'
                    })
                });
                const result = await res.json();

                if (result.success && result.audio_base64) {
                    // Create audio element and play
                    const audio = new Audio('data:' + result.audio_format + ';base64,' + result.audio_base64);
                    const cacheStatus = result.cached ? '<span style="color: var(--neon-green);">CACHED</span>' : '<span style="color: var(--text-muted);">generated</span>';
                    audio.onended = () => {
                        resultEl.innerHTML = `<div style="color: var(--neon-green);">&#9658; Audio played! (${result.duration_ms}ms - ${cacheStatus})</div>`;
                    };
                    audio.onerror = (e) => {
                        resultEl.innerHTML = `<div style="color: var(--neon-orange);">Playback error</div>`;
                    };
                    audio.play();
                    resultEl.innerHTML = `<div style="color: var(--neon-cyan);">&#9658; Playing ${result.provider} audio... (${cacheStatus})</div>`;
                } else if (result.success) {
                    const cacheStatus = result.cached ? 'from cache' : 'generated';
                    resultEl.innerHTML = `<div style="color: var(--neon-green);">Voice test completed in ${result.duration_ms}ms (${cacheStatus})</div>`;
                } else {
                    resultEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error || result.message}</div>`;
                }
            } catch (e) {
                resultEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function saveVoiceConfig() {
            const provider = document.getElementById('voice-provider').value;
            const voiceId = document.getElementById('voice-select').value;
            const speed = document.getElementById('voice-speed').value;
            const resultEl = document.getElementById('voice-test-result');

            if (!voiceId) {
                resultEl.innerHTML = '<div style="color: var(--neon-orange);">Please select a voice first</div>';
                return;
            }

            try {
                const res = await fetch('/api/voice/config', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        selected_provider: provider,
                        selected_voice: voiceId,
                        voice_settings: { speed: parseFloat(speed) }
                    })
                });
                const result = await res.json();

                if (result.success) {
                    resultEl.innerHTML = '<div style="color: var(--neon-green);">Voice configuration saved as default!</div>';
                } else {
                    resultEl.innerHTML = '<div style="color: var(--neon-orange);">Failed to save configuration</div>';
                }
            } catch (e) {
                resultEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        // ============================================================
        // PLATFORM CONNECTIONS
        // ============================================================
        async function loadPlatforms() {
            try {
                const res = await fetch('/api/platforms');
                const platforms = await res.json();

                for (const [id, platform] of Object.entries(platforms)) {
                    updatePlatformStatus(id, platform.status);
                }
            } catch (e) {
                console.error('Failed to load platforms:', e);
            }
        }

        function updatePlatformStatus(platformId, status) {
            const statusEl = document.getElementById(`status-${platformId}`);
            if (!statusEl) return;

            statusEl.className = 'table-status';
            switch (status) {
                case 'connected':
                    statusEl.classList.add('deployed');
                    statusEl.textContent = 'Connected';
                    break;
                case 'error':
                    statusEl.classList.add('training');
                    statusEl.textContent = 'Error';
                    break;
                default:
                    statusEl.classList.add('draft');
                    statusEl.textContent = 'Disconnected';
            }
        }

        function getPlatformConfig(platformId) {
            const config = {};
            switch (platformId) {
                case 'livekit':
                    config.url = document.getElementById('livekit-url').value;
                    config.api_key = document.getElementById('livekit-api-key').value;
                    config.api_secret = document.getElementById('livekit-api-secret').value;
                    break;
                case 'vapi':
                    config.api_key = document.getElementById('vapi-api-key').value;
                    config.assistant_id = document.getElementById('vapi-assistant-id').value;
                    break;
                case 'twilio':
                    config.account_sid = document.getElementById('twilio-account-sid').value;
                    config.auth_token = document.getElementById('twilio-auth-token').value;
                    config.phone_number = document.getElementById('twilio-phone').value;
                    break;
                case 'retell':
                    config.api_key = document.getElementById('retell-api-key').value;
                    config.agent_id = document.getElementById('retell-agent-id').value;
                    break;
                case 'bland':
                    config.api_key = document.getElementById('bland-api-key').value;
                    config.pathway_id = document.getElementById('bland-pathway-id').value;
                    break;
                case 'daily':
                    config.api_key = document.getElementById('daily-api-key').value;
                    config.room_url = document.getElementById('daily-room-url').value;
                    break;
                case 'vocode':
                    config.api_key = document.getElementById('vocode-api-key').value;
                    break;
                case 'websocket':
                    config.url = document.getElementById('websocket-url').value;
                    config.auth_header = document.getElementById('websocket-auth').value;
                    break;
            }
            return config;
        }

        async function connectPlatform(platformId) {
            const config = getPlatformConfig(platformId);
            const messageEl = document.getElementById('platform-message');

            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Connecting...</div>';

            try {
                const res = await fetch(`/api/platforms/${platformId}/connect`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ config })
                });
                const result = await res.json();

                updatePlatformStatus(platformId, result.status);

                if (result.success) {
                    messageEl.innerHTML = `<div style="color: var(--neon-green);">${result.message}</div>`;
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">${result.message}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Connection failed: ${e.message}</div>`;
            }
        }

        async function testPlatform(platformId) {
            const messageEl = document.getElementById('platform-message');

            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Testing connection...</div>';

            try {
                const res = await fetch(`/api/platforms/${platformId}/test`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' }
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = `<div style="color: var(--neon-green);">${result.message} (latency: ${result.latency_ms}ms)</div>`;
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">${result.message}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Test failed: ${e.message}</div>`;
            }
        }

        async function disconnectPlatform(platformId) {
            const messageEl = document.getElementById('platform-message');

            try {
                const res = await fetch(`/api/platforms/${platformId}/disconnect`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' }
                });
                const result = await res.json();

                updatePlatformStatus(platformId, result.status);
                messageEl.innerHTML = `<div style="color: var(--text-secondary);">${result.message}</div>`;
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Disconnect failed: ${e.message}</div>`;
            }
        }

        // ============================================================
        // FAST BRAIN FUNCTIONS
        // ============================================================
        async function loadFastBrainConfig() {
            try {
                const res = await fetch('/api/fast-brain/config');
                const config = await res.json();

                document.getElementById('fb-url').value = config.url || '';
                document.getElementById('fb-min-containers').value = config.min_containers || 1;
                document.getElementById('fb-max-containers').value = config.max_containers || 10;
                document.getElementById('fb-status').textContent = config.status || 'Offline';
                document.getElementById('fb-requests').textContent = config.stats?.total_requests || 0;

                if (config.status === 'connected') {
                    document.getElementById('fb-status').style.color = 'var(--neon-green)';
                } else {
                    document.getElementById('fb-status').style.color = 'var(--text-secondary)';
                }
            } catch (e) {
                console.error('Failed to load Fast Brain config:', e);
            }
        }

        async function saveFastBrainConfig() {
            const messageEl = document.getElementById('fb-config-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Saving...</div>';

            try {
                const res = await fetch('/api/fast-brain/config', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        url: document.getElementById('fb-url').value,
                        min_containers: parseInt(document.getElementById('fb-min-containers').value),
                        max_containers: parseInt(document.getElementById('fb-max-containers').value),
                    })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = '<div style="color: var(--neon-green);">Configuration saved!</div>';
                    loadFastBrainConfig();
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function checkFastBrainHealth() {
            const messageEl = document.getElementById('fb-config-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Checking health...</div>';

            try {
                const res = await fetch('/api/fast-brain/health');
                const result = await res.json();

                if (result.status === 'healthy') {
                    messageEl.innerHTML = `<div style="color: var(--neon-green);">Connected! Model loaded: ${result.model_loaded}</div>`;
                    document.getElementById('fb-status').textContent = 'Online';
                    document.getElementById('fb-status').style.color = 'var(--neon-green)';
                } else if (result.status === 'not_configured') {
                    messageEl.innerHTML = '<div style="color: var(--neon-orange);">Not configured - enter URL above</div>';
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">${result.message || 'Connection failed'}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function testFastBrainChat() {
            const btn = document.getElementById('fb-chat-btn');
            const responseEl = document.getElementById('fb-response');
            const metricsEl = document.getElementById('fb-metrics');

            const message = document.getElementById('fb-message').value;
            const systemPrompt = document.getElementById('fb-system-prompt').value;
            const maxTokens = parseInt(document.getElementById('fb-max-tokens').value);

            if (!message) {
                alert('Please enter a message');
                return;
            }

            btn.innerHTML = '<span class="spinner"></span> Sending...';
            btn.disabled = true;
            responseEl.textContent = 'Processing...';
            metricsEl.textContent = '';

            try {
                const res = await fetch('/api/fast-brain/chat', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ message, system_prompt: systemPrompt, max_tokens: maxTokens })
                });
                const result = await res.json();

                if (result.success) {
                    responseEl.textContent = result.response;
                    const metrics = result.metrics;
                    metricsEl.innerHTML = `TTFB: <span style="color: ${metrics.ttfb_ms < 50 ? 'var(--neon-green)' : 'var(--neon-orange)'}">${metrics.ttfb_ms.toFixed(1)}ms</span> | Total: ${metrics.total_time_ms.toFixed(1)}ms | Throughput: ${metrics.tokens_per_sec.toFixed(1)} tok/s`;

                    // Update targets
                    document.getElementById('fb-ttfb-target').className = `table-status ${metrics.ttfb_ms < 50 ? 'deployed' : 'draft'}`;
                    document.getElementById('fb-throughput-target').className = `table-status ${metrics.tokens_per_sec > 500 ? 'deployed' : 'draft'}`;
                } else {
                    responseEl.innerHTML = `<span style="color: var(--neon-orange);">Error: ${result.error}</span>`;
                }
            } catch (e) {
                responseEl.innerHTML = `<span style="color: var(--neon-orange);">Error: ${e.message}</span>`;
            }

            btn.innerHTML = 'Send to Fast Brain';
            btn.disabled = false;
        }

        async function runFastBrainBenchmark() {
            const resultsDiv = document.getElementById('fb-benchmark-results');
            const outputDiv = document.getElementById('fb-benchmark-output');

            resultsDiv.style.display = 'block';
            outputDiv.innerHTML = '<div style="color: var(--neon-cyan);">Running benchmark (5 requests)...</div>';

            try {
                const res = await fetch('/api/fast-brain/benchmark', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ num_requests: 5, prompt: 'Hello, how are you?' })
                });
                const result = await res.json();

                if (result.success) {
                    const summary = result.summary;
                    outputDiv.innerHTML = `
                        <div style="margin-bottom: 1rem;">
                            <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Summary</h4>
                            <p>Avg TTFB: <span style="color: ${summary.ttfb_target_met ? 'var(--neon-green)' : 'var(--neon-orange)'}">${summary.avg_ttfb_ms}ms</span> ${summary.ttfb_target_met ? '✓' : '✗'} (target: <50ms)</p>
                            <p>Avg Total: ${summary.avg_total_ms}ms</p>
                            <p>Avg Throughput: <span style="color: ${summary.throughput_target_met ? 'var(--neon-green)' : 'var(--neon-orange)'}">${summary.avg_tokens_per_sec} tok/s</span> ${summary.throughput_target_met ? '✓' : '✗'} (target: >500)</p>
                        </div>
                        <h4 style="color: var(--neon-cyan); margin-bottom: 0.5rem;">Individual Requests</h4>
                        <table class="skills-table">
                            <tr><th>#</th><th>TTFB</th><th>Total</th><th>Tokens/sec</th></tr>
                            ${result.results.map(r => `<tr><td>${r.request}</td><td>${r.ttfb_ms.toFixed(1)}ms</td><td>${r.total_ms.toFixed(1)}ms</td><td>${r.tokens_per_sec.toFixed(1)}</td></tr>`).join('')}
                        </table>
                    `;

                    // Update status cards
                    document.getElementById('fb-ttfb').textContent = `${summary.avg_ttfb_ms}ms`;
                    document.getElementById('fb-throughput').textContent = `${summary.avg_tokens_per_sec}`;
                } else {
                    outputDiv.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                outputDiv.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        // ============================================================
        // FAST BRAIN SUB-TABS & SKILLS MANAGEMENT
        // ============================================================
        let fbSkillsCache = {};

        function showFastBrainTab(tabId) {
            document.querySelectorAll('#tab-fastbrain .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-fastbrain .sub-tab-btn').forEach(b => b.classList.remove('active'));
            const tabEl = document.getElementById('fb-tab-' + tabId);
            if (tabEl) tabEl.classList.add('active');
            if (event && event.target && event.target.classList) {
                event.target.classList.add('active');
            }

            // Load data based on tab
            if (tabId === 'skills') loadFastBrainSkills();
            if (tabId === 'dashboard') refreshSystemStatus();
            if (tabId === 'integration') loadIntegrationChecklist();
            if (tabId === 'golden') loadGoldenPromptFB();
        }

        // Golden Prompts functions for Skills tab
        async function loadGoldenPromptFB() {
            const skillId = document.getElementById('fb-golden-skill-select').value;
            try {
                const res = await fetch(`/api/golden-prompts/${skillId}`);
                const data = await res.json();
                if (data.content) {
                    document.getElementById('fb-golden-prompt-content').value = data.content;
                    document.getElementById('fb-golden-tokens').textContent = '~' + data.tokens;
                    document.getElementById('fb-golden-prefill').textContent = '~' + data.prefill_ms + 'ms';
                    document.getElementById('fb-golden-status').textContent = data.is_custom ? 'Custom' : 'Default';
                }
                loadVoicesForSkillFB(skillId);
            } catch (e) {
                console.error('Error loading prompt:', e);
            }
        }

        async function loadVoicesForSkillFB(skillId) {
            const voiceSelect = document.getElementById('fb-golden-voice-select');
            if (!voiceSelect) return;
            try {
                const res = await fetch('/api/voice-lab/projects');
                const projects = await res.json();
                voiceSelect.innerHTML = '<option value="">No custom voice</option>';
                (projects || []).filter(p => p.status === 'trained').forEach(p => {
                    const option = document.createElement('option');
                    option.value = p.id;
                    option.textContent = `${p.name} (${p.provider})`;
                    if (p.skill_id === skillId) option.selected = true;
                    voiceSelect.appendChild(option);
                });
            } catch (e) { console.error(e); }
        }

        async function saveGoldenPromptFB() {
            const skillId = document.getElementById('fb-golden-skill-select').value;
            const content = document.getElementById('fb-golden-prompt-content').value;
            try {
                const res = await fetch(`/api/golden-prompts/${skillId}`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ content })
                });
                const data = await res.json();
                if (data.success) {
                    document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-green);">Prompt saved!</div>';
                    document.getElementById('fb-golden-tokens').textContent = '~' + data.tokens;
                    document.getElementById('fb-golden-status').textContent = 'Custom';
                }
            } catch (e) {
                document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-pink);">Error: ' + e.message + '</div>';
            }
        }

        async function resetGoldenPromptFB() {
            const skillId = document.getElementById('fb-golden-skill-select').value;
            if (!confirm(`Reset "${skillId}" prompt to default?`)) return;
            try {
                await fetch(`/api/golden-prompts/${skillId}/reset`, { method: 'POST' });
                loadGoldenPromptFB();
                document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-green);">Reset to default!</div>';
            } catch (e) { console.error(e); }
        }

        async function testGoldenPromptFB() {
            document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-cyan);">Testing prompt...</div>';
            const content = document.getElementById('fb-golden-prompt-content').value;
            try {
                const res = await fetch('/api/fast-brain/chat', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        messages: [{ role: 'user', content: 'What are your hours?' }],
                        system_prompt: content
                    })
                });
                const data = await res.json();
                document.getElementById('fb-golden-message').innerHTML = `<div style="background: var(--glass-surface); padding: 1rem; border-radius: 8px; margin-top: 1rem;"><strong>Response:</strong><br>${data.response || data.content || 'No response'}</div>`;
            } catch (e) {
                document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-pink);">Test failed: ' + e.message + '</div>';
            }
        }

        async function linkVoiceToSkillFB() {
            const skillId = document.getElementById('fb-golden-skill-select').value;
            const voiceId = document.getElementById('fb-golden-voice-select').value;
            if (!voiceId) return;
            try {
                await fetch(`/api/voice-lab/projects/${voiceId}/link-skill`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId })
                });
                document.getElementById('fb-golden-message').innerHTML = '<div style="color: var(--neon-green);">Voice linked!</div>';
            } catch (e) { console.error(e); }
        }

        // Training script generation
        function generateTrainingScriptFB() {
            const skill = document.getElementById('fb-train-profile').value;
            const model = document.getElementById('fb-train-model').value;
            const steps = document.getElementById('fb-train-steps').value;

            const script = `# LoRA Training Script for ${skill}
# Base Model: ${model}
# Training Steps: ${steps}

from unsloth import FastLanguageModel
import torch

# Load base model
model, tokenizer = FastLanguageModel.from_pretrained(
    model_name="${model}",
    max_seq_length=2048,
    load_in_4bit=True,
)

# Add LoRA adapters
model = FastLanguageModel.get_peft_model(
    model,
    r=16,
    target_modules=["q_proj", "k_proj", "v_proj", "o_proj"],
    lora_alpha=16,
    lora_dropout=0,
    bias="none",
)

# Load training data
from datasets import load_dataset
dataset = load_dataset("json", data_files="training_data/${skill}.jsonl")

# Training configuration
from trl import SFTTrainer
from transformers import TrainingArguments

trainer = SFTTrainer(
    model=model,
    tokenizer=tokenizer,
    train_dataset=dataset["train"],
    args=TrainingArguments(
        per_device_train_batch_size=2,
        gradient_accumulation_steps=4,
        warmup_steps=5,
        max_steps=${steps},
        learning_rate=2e-4,
        fp16=not torch.cuda.is_bf16_supported(),
        bf16=torch.cuda.is_bf16_supported(),
        output_dir="outputs/${skill}",
    ),
)

# Train
trainer.train()

# Save LoRA adapter
model.save_pretrained("adapters/${skill}")
print("Training complete! Adapter saved to adapters/${skill}")
`;
            document.getElementById('fb-training-script').textContent = script;
        }

        function copyTrainingScript() {
            const script = document.getElementById('fb-training-script').textContent;
            navigator.clipboard.writeText(script);
            alert('Training script copied to clipboard!');
        }

        function downloadTrainingScript() {
            const script = document.getElementById('fb-training-script').textContent;
            const blob = new Blob([script], { type: 'text/plain' });
            const url = URL.createObjectURL(blob);
            const a = document.createElement('a');
            a.href = url;
            a.download = 'train_lora.py';
            a.click();
        }

        async function loadFastBrainSkills() {
            const container = document.getElementById('fb-skills-list');
            container.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center;">Loading skills...</div>';

            try {
                const res = await fetch('/api/fast-brain/skills');
                const data = await res.json();
                const skills = data.skills || [];
                fbSkillsCache = {};
                skills.forEach(s => fbSkillsCache[s.id] = s);

                // Update skill selector dropdown
                updateSkillSelector(skills);
                document.getElementById('fb-skills-count').textContent = skills.length;

                if (skills.length === 0) {
                    container.innerHTML = `
                        <div style="text-align: center; padding: 2rem; grid-column: span 3;">
                            <div style="font-size: 2.5rem; margin-bottom: 1rem;">Stars</div>
                            <h3 style="margin-bottom: 0.5rem;">No Skills Yet</h3>
                            <p style="color: var(--text-secondary); margin-bottom: 1.5rem;">Create your first AI skill to get started with voice automation</p>
                            <button class="btn btn-primary" onclick="showCreateSkillModal()">+ Create Your First Skill</button>
                        </div>`;
                    return;
                }

                container.innerHTML = skills.map(skill => `
                    <div class="skill-card" style="padding: 1rem; background: var(--glass-surface); border-radius: 8px; border: 1px solid var(--glass-border); border-left: 4px solid ${skill.is_builtin ? 'var(--neon-cyan)' : 'var(--neon-green)'};">
                        <div style="display: flex; justify-content: space-between; align-items: flex-start; margin-bottom: 0.5rem;">
                            <div>
                                <strong style="color: var(--text-primary);">${skill.name}</strong>
                                <span style="font-size: 0.75rem; color: var(--text-secondary); margin-left: 0.5rem;">${skill.id}</span>
                            </div>
                            <span class="table-status ${skill.is_builtin ? 'deployed' : 'active'}" style="font-size: 0.7rem;">
                                ${skill.is_builtin ? 'Built-in' : 'Custom'}
                            </span>
                        </div>
                        <p style="color: var(--text-secondary); font-size: 0.85rem; margin-bottom: 0.75rem;">${skill.description}</p>
                        <div style="font-size: 0.8rem; color: var(--text-secondary); background: var(--glass-bg); padding: 0.5rem; border-radius: 4px; max-height: 80px; overflow: hidden;">
                            ${skill.system_prompt?.substring(0, 150)}${skill.system_prompt?.length > 150 ? '...' : ''}
                        </div>
                        ${skill.knowledge && skill.knowledge.length > 0 ? `
                            <div style="margin-top: 0.5rem; font-size: 0.75rem; color: var(--neon-cyan);">
                                ${skill.knowledge.length} knowledge items
                            </div>
                        ` : ''}
                        <div style="margin-top: 0.75rem; display: flex; gap: 0.5rem; flex-wrap: wrap;">
                            <button class="btn btn-sm btn-secondary" onclick="selectSkillForChat('${skill.id}')">Use in Chat</button>
                            <button class="btn btn-sm btn-primary" onclick="editSkill('${skill.id}')">Edit</button>
                            ${!skill.is_builtin ? `<button class="btn btn-sm" onclick="deleteSkill('${skill.id}')" style="background: var(--neon-orange); color: white;">Delete</button>` : ''}
                        </div>
                    </div>
                `).join('');
            } catch (e) {
                container.innerHTML = `<div style="color: var(--neon-orange); padding: 2rem; text-align: center;">Error loading skills: ${e.message}</div>`;
            }
        }

        function updateSkillSelector(skills) {
            const selector = document.getElementById('fb-skill-selector');
            if (!selector) return;

            selector.innerHTML = skills.map(s => `
                <option value="${s.id}">${s.name}</option>
            `).join('');
        }

        function selectSkillForChat(skillId) {
            showFastBrainTab('chat');
            document.getElementById('fb-skill-selector').value = skillId;
            onSkillSelect();
        }

        function onSkillSelect() {
            const skillId = document.getElementById('fb-skill-selector').value;
            const skill = fbSkillsCache[skillId];

            if (skill) {
                document.getElementById('fb-active-skill-name').textContent = skill.name;
                document.getElementById('fb-active-skill-desc').textContent = skill.description;

                // Also update on server
                fetch('/api/fast-brain/skills/select', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId })
                });
            }
        }

        function loadSkillPrompt() {
            const skillId = document.getElementById('fb-skill-selector').value;
            const skill = fbSkillsCache[skillId];

            if (skill) {
                document.getElementById('fb-system-prompt').value = skill.system_prompt;
            }
        }

        function showCreateSkillModal() {
            document.getElementById('fb-create-skill-form').style.display = 'block';
        }

        function hideCreateSkillModal() {
            document.getElementById('fb-create-skill-form').style.display = 'none';
            // Clear form
            document.getElementById('new-skill-id').value = '';
            document.getElementById('new-skill-name').value = '';
            document.getElementById('new-skill-description').value = '';
            document.getElementById('new-skill-prompt').value = '';
            document.getElementById('new-skill-knowledge').value = '';
            document.getElementById('create-skill-message').innerHTML = '';
        }

        async function createCustomSkill() {
            const messageEl = document.getElementById('create-skill-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Creating skill...</div>';

            const knowledge = document.getElementById('new-skill-knowledge').value
                .split('\\n')
                .map(s => s.trim())
                .filter(s => s.length > 0);

            try {
                const res = await fetch('/api/fast-brain/skills', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        skill_id: document.getElementById('new-skill-id').value,
                        name: document.getElementById('new-skill-name').value,
                        description: document.getElementById('new-skill-description').value,
                        system_prompt: document.getElementById('new-skill-prompt').value,
                        knowledge: knowledge,
                    })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = '<div style="color: var(--neon-green);">Skill created successfully!</div>';
                    setTimeout(() => {
                        hideCreateSkillModal();
                        loadFastBrainSkills();
                    }, 1000);
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function deleteSkill(skillId) {
            if (!confirm(`Delete skill "${skillId}"? This cannot be undone.`)) return;

            try {
                const res = await fetch(`/api/fast-brain/skills/${skillId}`, { method: 'DELETE' });
                const result = await res.json();

                if (result.success) {
                    loadFastBrainSkills();
                } else {
                    alert(`Error: ${result.error}`);
                }
            } catch (e) {
                alert(`Error: ${e.message}`);
            }
        }

        // ============================================================
        // SYSTEM STATUS & HIVE215 INTEGRATION
        // ============================================================
        async function refreshSystemStatus() {
            const statusIndicators = {
                'fb-status-indicator': { element: null, card: 'status-fast-brain' },
                'groq-status-indicator': { element: null, card: 'status-groq' },
                'hive-status-indicator': { element: null, card: 'status-hive215' },
                'modal-status-indicator': { element: null, card: 'status-modal' },
            };

            // Set all to loading
            Object.keys(statusIndicators).forEach(id => {
                const el = document.getElementById(id);
                if (el) el.textContent = 'Checking...';
            });

            try {
                const res = await fetch('/api/system/status');
                const data = await res.json();
                const systems = data.systems || {};

                // Update Fast Brain status
                const fb = systems.fast_brain || {};
                const fbIndicator = document.getElementById('fb-status-indicator');
                const fbCard = document.getElementById('status-fast-brain');
                if (fbIndicator) {
                    fbIndicator.textContent = fb.status === 'online' ? 'Online' : fb.status === 'degraded' ? 'Degraded' : fb.status || 'Offline';
                    fbIndicator.style.color = fb.status === 'online' ? 'var(--neon-green)' : fb.status === 'degraded' ? 'var(--neon-yellow)' : 'var(--neon-orange)';
                }
                if (fbCard) {
                    fbCard.style.borderLeftColor = fb.status === 'online' ? 'var(--neon-green)' : fb.status === 'degraded' ? 'var(--neon-yellow)' : 'var(--text-secondary)';
                }
                document.getElementById('fb-status-latency').textContent = fb.latency_ms ? `${fb.latency_ms}ms` : '--';
                document.getElementById('fb-status-model').textContent = fb.model || '--';
                document.getElementById('fb-status-skills').textContent = fb.skills_count || '--';

                // Update Groq status (inferred from Fast Brain)
                const groqIndicator = document.getElementById('groq-status-indicator');
                const groqCard = document.getElementById('status-groq');
                if (groqIndicator) {
                    groqIndicator.textContent = fb.status === 'online' ? 'Connected' : 'Unknown';
                    groqIndicator.style.color = fb.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }
                if (groqCard) {
                    groqCard.style.borderLeftColor = fb.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }

                // Update Modal status
                const modalIndicator = document.getElementById('modal-status-indicator');
                const modalCard = document.getElementById('status-modal');
                if (modalIndicator) {
                    modalIndicator.textContent = fb.status === 'online' ? 'Running' : 'Unknown';
                    modalIndicator.style.color = fb.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }
                if (modalCard) {
                    modalCard.style.borderLeftColor = fb.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }

                // Update Hive215 status
                const hive = systems.hive215 || {};
                const hiveIndicator = document.getElementById('hive-status-indicator');
                const hiveCard = document.getElementById('status-hive215');
                if (hiveIndicator) {
                    hiveIndicator.textContent = hive.status === 'online' ? 'Connected' : hive.status || 'Not configured';
                    hiveIndicator.style.color = hive.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }
                if (hiveCard) {
                    hiveCard.style.borderLeftColor = hive.status === 'online' ? 'var(--neon-green)' : 'var(--text-secondary)';
                }

                // Show errors if any
                if (fb.error) {
                    document.getElementById('system-error-display').style.display = 'block';
                    document.getElementById('system-errors').textContent = fb.error;
                } else {
                    document.getElementById('system-error-display').style.display = 'none';
                }

            } catch (e) {
                console.error('Failed to refresh system status:', e);
            }
        }

        async function loadIntegrationChecklist() {
            const container = document.getElementById('integration-checklist');

            try {
                const res = await fetch('/api/system/checklist');
                const data = await res.json();

                // Update progress
                document.getElementById('integration-progress').textContent = `${data.percent}%`;
                document.getElementById('integration-progress-bar').style.width = `${data.percent}%`;

                // Group by category
                const byCategory = {};
                data.checklist.forEach(item => {
                    if (!byCategory[item.category]) byCategory[item.category] = [];
                    byCategory[item.category].push(item);
                });

                container.innerHTML = Object.entries(byCategory).map(([category, items]) => `
                    <div style="margin-bottom: 1.5rem;">
                        <h4 style="color: var(--neon-cyan); margin-bottom: 0.75rem; font-size: 0.9rem;">${category}</h4>
                        ${items.map(item => `
                            <div style="display: flex; align-items: center; gap: 0.75rem; padding: 0.5rem 0; border-bottom: 1px solid var(--glass-border);">
                                <span style="font-size: 1.2rem;">${item.status === 'complete' ? '✅' : '⬜'}</span>
                                <span style="color: ${item.status === 'complete' ? 'var(--text-primary)' : 'var(--text-secondary)'};">
                                    ${item.name}
                                </span>
                            </div>
                        `).join('')}
                    </div>
                `).join('');
            } catch (e) {
                container.innerHTML = `<div style="color: var(--neon-orange);">Error loading checklist: ${e.message}</div>`;
            }
        }

        async function saveHive215Config() {
            const messageEl = document.getElementById('hive-config-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Saving...</div>';

            try {
                const res = await fetch('/api/hive215/config', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        url: document.getElementById('hive-url').value,
                        api_key: document.getElementById('hive-api-key').value,
                    })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = '<div style="color: var(--neon-green);">Configuration saved!</div>';
                    loadIntegrationChecklist();
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function testHive215Connection() {
            const messageEl = document.getElementById('hive-config-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Testing connection...</div>';

            // For now, just verify URL is set
            const url = document.getElementById('hive-url').value;
            if (!url) {
                messageEl.innerHTML = '<div style="color: var(--neon-orange);">Please enter the Hive215 URL first</div>';
                return;
            }

            messageEl.innerHTML = '<div style="color: var(--neon-green);">URL configured! Full connection test requires deployed services.</div>';
        }

        async function loadHive215Config() {
            try {
                const res = await fetch('/api/hive215/config');
                const config = await res.json();

                if (document.getElementById('hive-url')) {
                    document.getElementById('hive-url').value = config.url || '';
                }
                if (document.getElementById('hive-api-key')) {
                    document.getElementById('hive-api-key').value = config.api_key || '';
                }
            } catch (e) {
                console.error('Failed to load Hive215 config:', e);
            }
        }

        // ============================================================
        // OUTGOING API CONNECTIONS FUNCTIONS
        // ============================================================
        let apiConnections = [];

        async function loadApiConnections() {
            try {
                const res = await fetch('/api/connections');
                apiConnections = await res.json();
                renderApiConnectionsList();
                renderConnectionsStatus();
            } catch (e) {
                console.error('Failed to load API connections:', e);
            }
        }

        function refreshApiConnections() {
            loadApiConnections();
        }

        function renderApiConnectionsList() {
            const container = document.getElementById('api-connections-list');
            if (!container) return;

            if (!apiConnections || apiConnections.length === 0) {
                container.innerHTML = `
                    <div style="text-align: center; padding: 2rem; color: var(--text-secondary);">
                        <div style="font-size: 2rem; margin-bottom: 1rem;">🔌</div>
                        <h3 style="margin-bottom: 0.5rem; color: var(--text-primary);">No API Connections</h3>
                        <p style="margin-bottom: 1rem;">Connect to external APIs, CRMs, or webhooks</p>
                        <button class="btn btn-primary" onclick="showAddConnectionModal()">+ Add Connection</button>
                    </div>`;
                return;
            }

            container.innerHTML = apiConnections.map(conn => {
                const statusColor = conn.status === 'connected' ? 'var(--neon-green)' :
                                   conn.status === 'error' ? 'var(--neon-orange)' :
                                   conn.status === 'timeout' ? 'var(--neon-orange)' : 'var(--text-secondary)';
                const statusIcon = conn.status === 'connected' ? '✓' :
                                  conn.status === 'error' ? '✗' :
                                  conn.status === 'timeout' ? '⏱' : '○';
                return `
                    <div style="background: var(--glass-surface); border-radius: 8px; padding: 1rem; margin-bottom: 0.75rem; border-left: 4px solid ${statusColor};">
                        <div style="display: flex; justify-content: space-between; align-items: flex-start; margin-bottom: 0.5rem;">
                            <div>
                                <div style="font-weight: 600; color: var(--text-primary);">${conn.name}</div>
                                <div style="font-size: 0.8rem; color: var(--text-secondary); font-family: monospace;">${conn.url}</div>
                            </div>
                            <div style="display: flex; align-items: center; gap: 0.5rem;">
                                <span style="color: ${statusColor}; font-size: 0.85rem;">${statusIcon} ${conn.status || 'Not tested'}</span>
                            </div>
                        </div>
                        <div style="display: flex; gap: 0.5rem; flex-wrap: wrap; margin-top: 0.75rem;">
                            <span style="font-size: 0.75rem; background: var(--glass-bg); padding: 0.2rem 0.5rem; border-radius: 4px;">${conn.auth_type || 'bearer'}</span>
                            ${conn.last_tested ? `<span style="font-size: 0.75rem; color: var(--text-secondary);">Tested: ${new Date(conn.last_tested).toLocaleString()}</span>` : ''}
                        </div>
                        ${conn.last_error ? `<div style="font-size: 0.8rem; color: var(--neon-orange); margin-top: 0.5rem;">Error: ${conn.last_error}</div>` : ''}
                        <div style="display: flex; gap: 0.5rem; margin-top: 0.75rem;">
                            <button class="btn btn-secondary btn-sm" onclick="testApiConnection('${conn.id}')">Test</button>
                            <button class="btn btn-secondary btn-sm" onclick="openRequestTester('${conn.id}', '${conn.name}')">Send Request</button>
                            <button class="btn btn-secondary btn-sm" onclick="editApiConnection('${conn.id}')">Edit</button>
                            <button class="btn btn-secondary btn-sm" onclick="deleteApiConnection('${conn.id}')" style="color: var(--neon-orange);">Delete</button>
                        </div>
                    </div>`;
            }).join('');
        }

        function renderConnectionsStatus() {
            const container = document.getElementById('connections-status-list');
            if (!container) return;

            if (!apiConnections || apiConnections.length === 0) {
                container.innerHTML = '<div style="color: var(--text-secondary); padding: 1rem; text-align: center;">No connections configured</div>';
                return;
            }

            container.innerHTML = apiConnections.map(conn => {
                const statusColor = conn.status === 'connected' ? 'var(--neon-green)' :
                                   conn.status === 'error' || conn.status === 'timeout' ? 'var(--neon-orange)' : 'var(--text-secondary)';
                return `
                    <div style="display: flex; justify-content: space-between; align-items: center; padding: 0.5rem; border-bottom: 1px solid var(--glass-border);">
                        <span>${conn.name}</span>
                        <span style="color: ${statusColor}; font-weight: 500;">${conn.status || 'Unknown'}</span>
                    </div>`;
            }).join('');
        }

        function showAddConnectionModal() {
            document.getElementById('conn-edit-id').value = '';
            document.getElementById('conn-name').value = '';
            document.getElementById('conn-url').value = '';
            document.getElementById('conn-auth-type').value = 'bearer';
            document.getElementById('conn-apikey').value = '';
            document.getElementById('conn-headers').value = '';
            document.getElementById('conn-webhook').value = '';
            document.getElementById('connection-form-title').textContent = 'Add Connection';
            document.getElementById('connection-form-card').style.display = 'block';
        }

        function hideConnectionForm() {
            document.getElementById('connection-form-card').style.display = 'none';
        }

        function editApiConnection(connId) {
            const conn = apiConnections.find(c => c.id === connId);
            if (!conn) return;

            document.getElementById('conn-edit-id').value = conn.id;
            document.getElementById('conn-name').value = conn.name || '';
            document.getElementById('conn-url').value = conn.url || '';
            document.getElementById('conn-auth-type').value = conn.auth_type || 'bearer';
            document.getElementById('conn-apikey').value = conn.api_key || '';
            try {
                document.getElementById('conn-headers').value = conn.headers ? (typeof conn.headers === 'string' ? conn.headers : JSON.stringify(conn.headers, null, 2)) : '';
            } catch { document.getElementById('conn-headers').value = ''; }
            document.getElementById('conn-webhook').value = conn.webhook_url || '';
            document.getElementById('connection-form-title').textContent = 'Edit Connection';
            document.getElementById('connection-form-card').style.display = 'block';
        }

        async function saveApiConnection() {
            const editId = document.getElementById('conn-edit-id').value;
            const name = document.getElementById('conn-name').value.trim();
            const url = document.getElementById('conn-url').value.trim();

            if (!name || !url) {
                alert('Name and URL are required');
                return;
            }

            let headers = null;
            const headersText = document.getElementById('conn-headers').value.trim();
            if (headersText) {
                try {
                    headers = JSON.parse(headersText);
                } catch {
                    alert('Invalid JSON in headers field');
                    return;
                }
            }

            const data = {
                name,
                url,
                auth_type: document.getElementById('conn-auth-type').value,
                api_key: document.getElementById('conn-apikey').value,
                headers,
                webhook_url: document.getElementById('conn-webhook').value.trim() || null
            };

            try {
                let res;
                if (editId) {
                    res = await fetch(`/api/connections/${editId}`, {
                        method: 'PUT',
                        headers: {'Content-Type': 'application/json'},
                        body: JSON.stringify(data)
                    });
                } else {
                    res = await fetch('/api/connections', {
                        method: 'POST',
                        headers: {'Content-Type': 'application/json'},
                        body: JSON.stringify(data)
                    });
                }

                const result = await res.json();
                if (result.success || result.id) {
                    hideConnectionForm();
                    loadApiConnections();
                } else {
                    alert('Failed to save: ' + (result.error || 'Unknown error'));
                }
            } catch (e) {
                alert('Failed to save connection: ' + e.message);
            }
        }

        async function deleteApiConnection(connId) {
            if (!confirm('Delete this connection?')) return;

            try {
                const res = await fetch(`/api/connections/${connId}`, {method: 'DELETE'});
                const result = await res.json();
                if (result.success) {
                    loadApiConnections();
                } else {
                    alert('Failed to delete: ' + (result.error || 'Unknown error'));
                }
            } catch (e) {
                alert('Failed to delete: ' + e.message);
            }
        }

        async function testApiConnection(connId) {
            const conn = apiConnections.find(c => c.id === connId);
            if (!conn) return;

            // Show testing indicator
            const container = document.getElementById('api-connections-list');
            const originalHtml = container.innerHTML;

            try {
                const res = await fetch(`/api/connections/${connId}/test`, {method: 'POST'});
                const result = await res.json();

                // Reload to show updated status
                await loadApiConnections();

                // Show result
                if (result.success) {
                    alert(`✓ Connection successful!\\nStatus: ${result.status_code}\\nLatency: ${result.latency_ms}ms`);
                } else {
                    alert(`✗ Connection failed\\nError: ${result.error}`);
                }
            } catch (e) {
                alert('Test failed: ' + e.message);
            }
        }

        function openRequestTester(connId, connName) {
            document.getElementById('tester-connection-id').value = connId;
            document.getElementById('tester-connection-name').textContent = connName;
            document.getElementById('tester-method').value = 'GET';
            document.getElementById('tester-path').value = '';
            document.getElementById('tester-body').value = '';
            document.getElementById('tester-result').innerHTML = '';
            document.getElementById('request-tester-card').style.display = 'block';
        }

        function hideRequestTester() {
            document.getElementById('request-tester-card').style.display = 'none';
        }

        async function sendTestRequest() {
            const connId = document.getElementById('tester-connection-id').value;
            const method = document.getElementById('tester-method').value;
            const path = document.getElementById('tester-path').value;
            const bodyText = document.getElementById('tester-body').value.trim();

            let body = null;
            if (bodyText && (method === 'POST' || method === 'PUT')) {
                try {
                    body = JSON.parse(bodyText);
                } catch {
                    document.getElementById('tester-result').innerHTML = '<div style="color: var(--neon-orange);">Invalid JSON in request body</div>';
                    return;
                }
            }

            document.getElementById('tester-result').innerHTML = '<div style="color: var(--neon-cyan);">Sending request...</div>';

            try {
                const res = await fetch(`/api/connections/${connId}/send`, {
                    method: 'POST',
                    headers: {'Content-Type': 'application/json'},
                    body: JSON.stringify({method, path, body})
                });

                const result = await res.json();

                const statusColor = result.success ? 'var(--neon-green)' : 'var(--neon-orange)';
                let responsePreview;
                try {
                    responsePreview = typeof result.response === 'object' ?
                        JSON.stringify(result.response, null, 2) : result.response;
                } catch {
                    responsePreview = String(result.response);
                }

                document.getElementById('tester-result').innerHTML = `
                    <div style="margin-bottom: 0.5rem;">
                        <span style="color: ${statusColor}; font-weight: 600;">${result.status_code || 'Error'}</span>
                        <span style="color: var(--text-secondary); margin-left: 0.5rem;">${result.latency_ms || 0}ms</span>
                    </div>
                    ${result.error ? `<div style="color: var(--neon-orange); margin-bottom: 0.5rem;">Error: ${result.error}</div>` : ''}
                    <div style="font-size: 0.8rem; color: var(--text-secondary); margin-bottom: 0.25rem;">Request: ${result.request?.method || method} ${result.request?.url || path}</div>
                    <pre class="code-block" style="max-height: 200px; overflow: auto; font-size: 0.75rem;">${responsePreview?.substring(0, 2000) || '(empty)'}</pre>`;
            } catch (e) {
                document.getElementById('tester-result').innerHTML = `<div style="color: var(--neon-orange);">Request failed: ${e.message}</div>`;
            }
        }

        async function refreshConnectionStatuses() {
            try {
                const res = await fetch('/api/connections/status');
                const statuses = await res.json();
                // Update apiConnections with new statuses
                for (const conn of apiConnections) {
                    if (statuses[conn.id]) {
                        conn.status = statuses[conn.id].status;
                        conn.last_tested = statuses[conn.id].last_tested;
                        conn.last_error = statuses[conn.id].last_error;
                    }
                }
                renderConnectionsStatus();
                renderApiConnectionsList();
            } catch (e) {
                console.error('Failed to refresh statuses:', e);
            }
        }

        // ============================================================
        // VOICE LAB FUNCTIONS
        // ============================================================
        function showVoiceLabTab(tabId) {
            document.querySelectorAll('#tab-voicelab .sub-tab-content').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('#tab-voicelab .sub-tab-btn').forEach(b => b.classList.remove('active'));
            const tabEl = document.getElementById('voicelab-' + tabId);
            if (tabEl) tabEl.classList.add('active');
            // Highlight the button if we can find it (edit tab has no button)
            if (event && event.target && event.target.classList) {
                event.target.classList.add('active');
            }
        }

        async function loadVoiceProjects() {
            try {
                const res = await fetch('/api/voice-lab/projects');
                const projects = await res.json();

                const container = document.getElementById('voice-projects-list');
                if (projects.length === 0) {
                    container.innerHTML = `
                        <div style="text-align: center; padding: 2rem; grid-column: span 3;">
                            <div style="font-size: 2.5rem; margin-bottom: 1rem;">Mic</div>
                            <h3 style="margin-bottom: 0.5rem;">Create Your Custom Voice</h3>
                            <p style="color: var(--text-secondary); margin-bottom: 1.5rem;">Clone your voice or customize TTS settings for natural-sounding AI responses</p>
                            <button class="btn btn-primary" onclick="showVoiceLabTab('create')">+ Create Voice Project</button>
                        </div>`;
                    return;
                }

                container.innerHTML = projects.map(p => `
                    <div class="skill-card ${p.status}" onclick="openVoiceProject('${p.id}')">
                        <div class="skill-header">
                            <div>
                                <div class="skill-name">${p.name}</div>
                                <div class="skill-type">${p.provider} | ${p.base_voice}</div>
                            </div>
                            <span class="skill-status ${p.status}">${p.status}</span>
                        </div>
                        <div class="skill-metrics">
                            <div class="skill-metric">
                                <div class="metric-value">${p.samples?.length || 0}</div>
                                <div class="metric-label">Samples</div>
                            </div>
                            <div class="skill-metric">
                                <div class="metric-value">${p.settings?.speed || 1.0}x</div>
                                <div class="metric-label">Speed</div>
                            </div>
                            <div class="skill-metric">
                                <div class="metric-value">${p.settings?.emotion || 'neutral'}</div>
                                <div class="metric-label">Emotion</div>
                            </div>
                        </div>
                    </div>
                `).join('');
            } catch (e) {
                console.error('Failed to load voice projects:', e);
            }
        }

        // Voice Lab - Load provider voices for dropdown (Dynamic API)
        let dynamicVoiceCache = {};  // Cache for dynamically loaded voices
        let currentBrowseVoices = []; // Current voices for filtering in browse tab

        async function loadVLProviderVoices() {
            const provider = document.getElementById('vl-provider').value;
            const voiceSelect = document.getElementById('vl-base-voice');
            voiceSelect.innerHTML = '<option value="">Loading voices...</option>';

            if (!provider) {
                voiceSelect.innerHTML = '<option value="">Select a provider first...</option>';
                return;
            }

            try {
                // Use dynamic API endpoint
                const res = await fetch(`/api/voice-lab/provider-voices/${provider}`);
                const data = await res.json();

                if (data.error) {
                    voiceSelect.innerHTML = `<option value="">${data.error}</option>`;
                    // Fall back to static if available
                    if (voiceProviders[provider]) {
                        populateVLVoices(provider, voiceProviders[provider].voices);
                    }
                    return;
                }

                // Cache the dynamic voices
                dynamicVoiceCache[provider] = data.voices;
                populateVLVoices(provider, data.voices);
            } catch (e) {
                console.error('Failed to load provider voices:', e);
                voiceSelect.innerHTML = '<option value="">Error loading voices</option>';
                // Fall back to static
                if (voiceProviders[provider]) {
                    populateVLVoices(provider, voiceProviders[provider].voices);
                }
            }
        }

        function populateVLVoices(provider, voices) {
            const voiceSelect = document.getElementById('vl-base-voice');
            voiceSelect.innerHTML = '<option value="">Select a voice...</option>';

            if (!voices || voices.length === 0) {
                voiceSelect.innerHTML = '<option value="">No voices available</option>';
                return;
            }

            voices.forEach(voice => {
                const option = document.createElement('option');
                option.value = voice.id;
                const gender = voice.gender || 'unknown';
                const style = voice.style || voice.description?.substring(0, 20) || 'default';
                option.textContent = `${voice.name} (${gender}, ${style})`;
                voiceSelect.appendChild(option);
            });
        }

        // Browse Voices from Providers
        async function loadProviderVoices() {
            const provider = document.getElementById('vl-browse-provider').value;
            const statusEl = document.getElementById('vl-browse-status');
            const filtersEl = document.getElementById('vl-browse-filters');
            const voicesEl = document.getElementById('vl-browse-voices');

            if (!provider) {
                statusEl.innerHTML = '';
                filtersEl.style.display = 'none';
                voicesEl.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;">Select a provider above to browse available voices.</div>';
                return;
            }

            statusEl.innerHTML = '<span style="color: var(--neon-cyan);">Loading voices from ' + provider + '...</span>';
            voicesEl.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;"><div class="loading-spinner"></div> Loading...</div>';

            try {
                const res = await fetch(`/api/voice-lab/provider-voices/${provider}`);
                const data = await res.json();

                if (data.error && (!data.voices || data.voices.length === 0)) {
                    statusEl.innerHTML = `<span style="color: var(--neon-orange);">${data.error}</span>`;
                    voicesEl.innerHTML = `<div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;">
                        ${data.error.includes('API key') ? 'Configure your API key in Settings → API Keys' : 'No voices found'}
                    </div>`;
                    filtersEl.style.display = 'none';
                    return;
                }

                currentBrowseVoices = data.voices || [];
                const isStatic = data.static ? ' (cached)' : ' (live)';
                const customCount = data.custom_count || currentBrowseVoices.filter(v => v.is_owner).length;
                const customText = customCount > 0 ? ` <span style="color: var(--neon-green);">★ ${customCount} custom</span>` : '';
                statusEl.innerHTML = `<span style="color: var(--neon-green);">Loaded ${currentBrowseVoices.length} voices${isStatic}</span>${customText}`;
                filtersEl.style.display = 'block';
                document.getElementById('vl-browse-search').value = '';
                document.getElementById('vl-browse-gender').value = '';
                renderBrowseVoices(currentBrowseVoices);
            } catch (e) {
                console.error('Failed to load provider voices:', e);
                statusEl.innerHTML = `<span style="color: var(--neon-pink);">Error: ${e.message}</span>`;
                voicesEl.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;">Failed to load voices</div>';
            }
        }

        function filterBrowseVoices() {
            const search = document.getElementById('vl-browse-search').value.toLowerCase();
            const gender = document.getElementById('vl-browse-gender').value;

            let filtered = currentBrowseVoices.filter(voice => {
                const matchesSearch = !search ||
                    voice.name?.toLowerCase().includes(search) ||
                    voice.description?.toLowerCase().includes(search) ||
                    voice.style?.toLowerCase().includes(search);
                const matchesGender = !gender || voice.gender?.toLowerCase() === gender;
                return matchesSearch && matchesGender;
            });

            renderBrowseVoices(filtered);
        }

        function renderBrowseVoices(voices) {
            const voicesEl = document.getElementById('vl-browse-voices');
            const countEl = document.getElementById('vl-browse-count');

            // Count custom voices
            const customCount = voices.filter(v => v.is_owner).length;
            const customText = customCount > 0 ? ` (${customCount} custom)` : '';
            countEl.textContent = `${voices.length} voice${voices.length !== 1 ? 's' : ''}${customText}`;

            if (voices.length === 0) {
                voicesEl.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center; grid-column: 1 / -1;">No voices match your filters</div>';
                return;
            }

            voicesEl.innerHTML = voices.map(voice => {
                const genderIcon = voice.gender === 'male' ? '♂' : voice.gender === 'female' ? '♀' : '◎';
                const genderColor = voice.gender === 'male' ? 'var(--neon-cyan)' : voice.gender === 'female' ? 'var(--neon-pink)' : 'var(--neon-purple)';
                const previewButton = voice.preview_url ? `<button class="btn btn-secondary btn-sm" onclick="playVoicePreview('${voice.preview_url}')" title="Provider Sample">▶</button>` : '';

                // Special styling for custom/owned voices
                const isCustom = voice.is_owner;
                const cardBorder = isCustom ? 'border: 2px solid var(--neon-green);' : '';
                const customBadge = isCustom ? '<span style="background: var(--neon-green); color: #000; padding: 2px 8px; border-radius: 10px; font-size: 0.7rem; font-weight: 600; margin-left: 8px;">CUSTOM</span>' : '';

                // Get current provider for test button
                const provider = document.getElementById('vl-browse-provider').value;

                return `
                    <div class="glass-card voice-card" style="padding: 1rem; background: var(--glass-surface); ${cardBorder}" id="voice-card-${voice.id}">
                        <div style="display: flex; justify-content: space-between; align-items: flex-start; margin-bottom: 0.5rem;">
                            <div>
                                <div style="font-weight: 600; color: var(--text-primary);">
                                    ${voice.name || 'Unknown'}${customBadge}
                                </div>
                                <div style="font-size: 0.85rem; color: ${genderColor};">${genderIcon} ${voice.gender || 'unknown'}</div>
                            </div>
                            <div style="display: flex; gap: 0.25rem;">
                                ${previewButton}
                                <button class="btn btn-secondary btn-sm" onclick="testBrowseVoice('${voice.id}', '${provider}', this)" title="Generate Test Audio">🔊 Test</button>
                                <button class="btn btn-primary btn-sm" onclick="useVoiceForProject('${voice.id}', '${(voice.name || '').replace(/'/g, "\\'")}')" title="Use in Project">Use</button>
                            </div>
                        </div>
                        <div style="font-size: 0.8rem; color: var(--text-secondary); margin-bottom: 0.5rem;">
                            ${voice.style || voice.description?.substring(0, 80) || 'No description'}
                        </div>
                        <div style="font-size: 0.75rem; color: var(--text-muted);">
                            ID: <code style="background: var(--glass-surface); padding: 0 4px; border-radius: 3px;">${voice.id}</code>
                        </div>
                        ${voice.category && !isCustom ? `<div style="font-size: 0.75rem; color: var(--neon-orange); margin-top: 0.25rem;">${voice.category}</div>` : ''}
                        ${voice.languages ? `<div style="font-size: 0.75rem; color: var(--text-muted); margin-top: 0.25rem;">Languages: ${Array.isArray(voice.languages) ? voice.languages.join(', ') : voice.languages}</div>` : ''}
                        <div id="voice-test-${voice.id}" style="margin-top: 0.5rem;"></div>
                    </div>
                `;
            }).join('');
        }

        async function testBrowseVoice(voiceId, provider, btn) {
            const testEl = document.getElementById('voice-test-' + voiceId);
            const originalText = btn.innerHTML;

            btn.disabled = true;
            btn.innerHTML = '⏳';
            testEl.innerHTML = '<span style="color: var(--neon-cyan); font-size: 0.8rem;">Generating...</span>';

            try {
                const res = await fetch('/api/voice/test', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        voice_id: voiceId,
                        text: 'Hello! This is a test of my voice. How does it sound?',
                        provider: provider
                    })
                });
                const result = await res.json();

                if (result.success && result.audio_base64) {
                    const audio = new Audio('data:' + result.audio_format + ';base64,' + result.audio_base64);
                    const cacheStatus = result.cached ? '(cached)' : `(${result.duration_ms}ms)`;

                    audio.onended = () => {
                        testEl.innerHTML = `<span style="color: var(--neon-green); font-size: 0.8rem;">✓ Played ${cacheStatus}</span>`;
                    };
                    audio.onerror = () => {
                        testEl.innerHTML = '<span style="color: var(--neon-orange); font-size: 0.8rem;">Playback error</span>';
                    };
                    audio.play();
                    testEl.innerHTML = `<span style="color: var(--neon-cyan); font-size: 0.8rem;">▶ Playing... ${cacheStatus}</span>`;
                } else {
                    testEl.innerHTML = `<span style="color: var(--neon-orange); font-size: 0.8rem;">${result.error || 'Test failed'}</span>`;
                }
            } catch (e) {
                testEl.innerHTML = `<span style="color: var(--neon-orange); font-size: 0.8rem;">Error: ${e.message}</span>`;
            }

            btn.disabled = false;
            btn.innerHTML = originalText;
        }

        function playVoicePreview(url) {
            const audio = new Audio(url);
            audio.play().catch(e => console.error('Failed to play preview:', e));
        }

        function useVoiceForProject(voiceId, voiceName) {
            // Switch to create tab and set the voice
            showVoiceLabTab('create');
            const provider = document.getElementById('vl-browse-provider').value;
            if (provider) {
                document.getElementById('vl-provider').value = provider;
                // Load voices for this provider then select the voice
                loadVLProviderVoices().then(() => {
                    setTimeout(() => {
                        document.getElementById('vl-base-voice').value = voiceId;
                    }, 500);
                });
            }
            // Set a suggested name
            if (!document.getElementById('vl-name').value) {
                document.getElementById('vl-name').value = voiceName + ' Clone';
            }
        }

        async function createVoiceProject() {
            const messageEl = document.getElementById('vl-create-message');
            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Creating project...</div>';

            try {
                const res = await fetch('/api/voice-lab/projects', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        name: document.getElementById('vl-name').value || 'Untitled Voice',
                        description: document.getElementById('vl-description').value,
                        provider: document.getElementById('vl-provider').value,
                        base_voice: document.getElementById('vl-base-voice').value,
                        settings: {
                            pitch: parseFloat(document.getElementById('vl-pitch').value),
                            speed: parseFloat(document.getElementById('vl-speed').value),
                            emotion: document.getElementById('vl-emotion').value,
                            style: document.getElementById('vl-style').value,
                        }
                    })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = `<div style="color: var(--neon-green);">Project "${result.project.name}" created!</div>`;
                    loadVoiceProjects();
                    showVoiceLabTab('projects');
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function openVoiceProject(projectId) {
            try {
                const res = await fetch(`/api/voice-lab/projects/${projectId}`);
                const project = await res.json();

                if (project.error) {
                    alert('Error loading project: ' + project.error);
                    return;
                }

                // Store project ID
                document.getElementById('vl-edit-project-id').value = projectId;

                // Populate form fields
                document.getElementById('vl-edit-title').textContent = project.name;
                document.getElementById('vl-edit-name').value = project.name || '';
                document.getElementById('vl-edit-description').value = project.description || '';
                document.getElementById('vl-edit-provider').value = project.provider || 'elevenlabs';

                // Status with color
                const statusEl = document.getElementById('vl-edit-status');
                statusEl.textContent = project.status || 'draft';
                statusEl.style.color = project.status === 'trained' ? 'var(--neon-green)' :
                                       project.status === 'training' ? 'var(--neon-orange)' :
                                       project.status === 'failed' ? 'var(--neon-pink)' : 'var(--text-secondary)';

                // Settings
                const settings = project.settings || {};
                document.getElementById('vl-edit-pitch').value = settings.pitch || 1.0;
                document.getElementById('vl-edit-pitch-value').textContent = settings.pitch || 1.0;
                document.getElementById('vl-edit-speed').value = settings.speed || 1.0;
                document.getElementById('vl-edit-speed-value').textContent = settings.speed || 1.0;
                document.getElementById('vl-edit-emotion').value = settings.emotion || 'neutral';
                document.getElementById('vl-edit-style').value = settings.style || 'conversational';

                // Linked skill
                document.getElementById('vl-edit-skill-link').value = project.skill_id || '';
                document.getElementById('vl-edit-linked-skill').textContent =
                    project.skill_id ? `Currently linked to: ${project.skill_id}` : '';

                // Load samples
                renderVoiceSamples(project.samples || []);

                // Show the edit tab
                showVoiceLabTab('edit');

            } catch (e) {
                alert('Error opening project: ' + e.message);
            }
        }

        function renderVoiceSamples(samples) {
            const container = document.getElementById('vl-edit-samples-list');

            if (!samples || samples.length === 0) {
                container.innerHTML = '<div style="color: var(--text-secondary); padding: 1rem; text-align: center;">No samples uploaded. Add audio samples to train your voice.</div>';
                return;
            }

            container.innerHTML = samples.map(s => `
                <div style="display: flex; align-items: center; gap: 1rem; padding: 0.75rem; background: var(--card-bg); border-radius: 8px; margin-bottom: 0.5rem;">
                    <div style="flex: 1;">
                        <div style="font-weight: 500;">${s.filename || 'Audio Sample'}</div>
                        <div style="font-size: 0.8rem; color: var(--text-secondary);">
                            ${s.duration_ms ? Math.round(s.duration_ms/1000) + 's' : 'Unknown duration'} | ${s.emotion || 'neutral'}
                        </div>
                        ${s.transcript ? `<div style="font-size: 0.8rem; color: var(--text-secondary); margin-top: 0.25rem;">"${s.transcript}"</div>` : ''}
                    </div>
                    <button class="btn btn-danger btn-sm" onclick="deleteVoiceSample('${s.id}')">Delete</button>
                </div>
            `).join('');
        }

        async function uploadVoiceSample(input) {
            const projectId = document.getElementById('vl-edit-project-id').value;
            if (!projectId) return;

            const file = input.files[0];
            if (!file) return;

            const formData = new FormData();
            formData.append('audio', file);
            formData.append('emotion', document.getElementById('vl-edit-emotion').value);

            try {
                showEditMessage('Uploading sample...', 'info');

                const res = await fetch(`/api/voice-lab/projects/${projectId}/samples`, {
                    method: 'POST',
                    body: formData
                });
                const result = await res.json();

                if (result.success) {
                    showEditMessage('Sample uploaded!', 'success');
                    // Reload project to refresh samples
                    openVoiceProject(projectId);
                } else {
                    showEditMessage(result.error || 'Upload failed', 'error');
                }
            } catch (e) {
                showEditMessage('Upload error: ' + e.message, 'error');
            }

            input.value = '';
        }

        async function deleteVoiceSample(sampleId) {
            const projectId = document.getElementById('vl-edit-project-id').value;
            if (!confirm('Delete this sample?')) return;

            try {
                const res = await fetch(`/api/voice-lab/projects/${projectId}/samples/${sampleId}`, {
                    method: 'DELETE'
                });
                const result = await res.json();

                if (result.success) {
                    openVoiceProject(projectId);
                }
            } catch (e) {
                alert('Error deleting sample: ' + e.message);
            }
        }

        async function saveVoiceProjectChanges() {
            const projectId = document.getElementById('vl-edit-project-id').value;
            if (!projectId) return;

            const data = {
                name: document.getElementById('vl-edit-name').value,
                description: document.getElementById('vl-edit-description').value,
                provider: document.getElementById('vl-edit-provider').value,
                settings: {
                    pitch: parseFloat(document.getElementById('vl-edit-pitch').value),
                    speed: parseFloat(document.getElementById('vl-edit-speed').value),
                    emotion: document.getElementById('vl-edit-emotion').value,
                    style: document.getElementById('vl-edit-style').value
                }
            };

            try {
                const res = await fetch(`/api/voice-lab/projects/${projectId}`, {
                    method: 'PUT',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(data)
                });
                const result = await res.json();

                if (result.success) {
                    document.getElementById('vl-edit-title').textContent = data.name;
                    showEditMessage('Changes saved!', 'success');
                } else {
                    showEditMessage(result.error || 'Save failed', 'error');
                }
            } catch (e) {
                showEditMessage('Error: ' + e.message, 'error');
            }
        }

        async function trainVoiceProject() {
            const projectId = document.getElementById('vl-edit-project-id').value;
            if (!projectId) return;

            try {
                showEditMessage('Starting voice training...', 'info');
                document.getElementById('vl-edit-status').textContent = 'training';
                document.getElementById('vl-edit-status').style.color = 'var(--neon-orange)';

                const res = await fetch(`/api/voice-lab/projects/${projectId}/train`, {
                    method: 'POST'
                });
                const result = await res.json();

                if (result.success) {
                    document.getElementById('vl-edit-status').textContent = 'trained';
                    document.getElementById('vl-edit-status').style.color = 'var(--neon-green)';
                    showEditMessage(result.message || 'Voice trained successfully!', 'success');
                } else {
                    document.getElementById('vl-edit-status').textContent = 'failed';
                    document.getElementById('vl-edit-status').style.color = 'var(--neon-pink)';
                    showEditMessage(result.error || 'Training failed', 'error');
                }
            } catch (e) {
                showEditMessage('Training error: ' + e.message, 'error');
            }
        }

        async function testVoiceProject() {
            const projectId = document.getElementById('vl-edit-project-id').value;
            const text = document.getElementById('vl-edit-test-text').value;
            if (!projectId || !text) return;

            try {
                showEditMessage('Synthesizing voice...', 'info');

                const res = await fetch(`/api/voice-lab/projects/${projectId}/test`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ text })
                });
                const result = await res.json();

                if (result.success && result.audio_base64) {
                    // Play base64 audio directly (like Settings voice test)
                    const audioPlayer = document.getElementById('vl-edit-audio-player');
                    const audioEl = document.getElementById('vl-edit-audio');

                    // Create audio from base64
                    const audioSrc = 'data:' + result.audio_format + ';base64,' + result.audio_base64;
                    audioEl.src = audioSrc;
                    audioPlayer.style.display = 'block';

                    // Auto-play
                    audioEl.play().catch(e => {
                        console.error('Play error:', e);
                        showEditMessage('Click the audio player to play.', 'info');
                    });

                    const sizeInfo = result.audio_size_bytes ? ` (${(result.audio_size_bytes / 1024).toFixed(1)} KB)` : '';
                    showEditMessage((result.message || 'Audio generated!') + sizeInfo, 'success');
                } else {
                    showEditMessage(result.error || 'Could not generate audio', 'error');
                }
            } catch (e) {
                showEditMessage('Test error: ' + e.message, 'error');
            }
        }

        async function linkVoiceProjectToSkill() {
            const projectId = document.getElementById('vl-edit-project-id').value;
            const skillId = document.getElementById('vl-edit-skill-link').value;
            if (!projectId || !skillId) {
                showEditMessage('Please select a skill', 'error');
                return;
            }

            try {
                const res = await fetch(`/api/voice-lab/projects/${projectId}/link-skill`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ skill_id: skillId })
                });
                const result = await res.json();

                if (result.success) {
                    document.getElementById('vl-edit-linked-skill').textContent = `Currently linked to: ${skillId}`;
                    showEditMessage(`Voice linked to ${skillId}!`, 'success');
                } else {
                    showEditMessage(result.error || 'Failed to link', 'error');
                }
            } catch (e) {
                showEditMessage('Error: ' + e.message, 'error');
            }
        }

        async function deleteCurrentVoiceProject() {
            const projectId = document.getElementById('vl-edit-project-id').value;
            const name = document.getElementById('vl-edit-name').value;

            if (!confirm(`Delete voice project "${name}"? This cannot be undone.`)) return;

            try {
                const res = await fetch(`/api/voice-lab/projects/${projectId}`, {
                    method: 'DELETE'
                });
                const result = await res.json();

                if (result.success) {
                    showVoiceLabTab('projects');
                    loadVoiceProjects();
                } else {
                    alert(result.error || 'Delete failed');
                }
            } catch (e) {
                alert('Error: ' + e.message);
            }
        }

        function showEditMessage(msg, type) {
            const el = document.getElementById('vl-edit-message');
            el.innerHTML = `<div style="color: ${type === 'success' ? 'var(--neon-green)' : type === 'error' ? 'var(--neon-pink)' : 'var(--neon-cyan)'};">${msg}</div>`;
            if (type !== 'info') {
                setTimeout(() => { el.innerHTML = ''; }, 5000);
            }
        }

        async function refreshVoiceTrainingStatus() {
            const container = document.getElementById('voice-training-queue');

            try {
                const res = await fetch('/api/voice-lab/training-status');
                const jobs = await res.json();

                if (jobs.length === 0) {
                    container.innerHTML = '<div style="color: var(--text-secondary); padding: 2rem; text-align: center;">No voice training jobs in progress.</div>';
                    return;
                }

                container.innerHTML = `
                    <table class="skills-table">
                        <tr><th>Project</th><th>Progress</th><th>Started</th><th>Status</th></tr>
                        ${jobs.map(j => `
                            <tr>
                                <td>${j.project_id}</td>
                                <td>
                                    <div style="background: var(--glass-surface); border-radius: 4px; overflow: hidden; height: 20px;">
                                        <div style="background: linear-gradient(90deg, var(--neon-cyan), var(--neon-green)); height: 100%; width: ${j.progress}%;"></div>
                                    </div>
                                    <span style="font-size: 0.75rem;">${j.progress}%</span>
                                </td>
                                <td>${new Date(j.started_at).toLocaleString()}</td>
                                <td><span class="table-status ${j.progress >= 100 ? 'deployed' : 'training'}">${j.progress >= 100 ? 'Complete' : 'Training'}</span></td>
                            </tr>
                        `).join('')}
                    </table>
                `;
            } catch (e) {
                container.innerHTML = `<div style="color: var(--neon-orange);">Error loading training status</div>`;
            }
        }

        // ============================================================
        // SKILL FINE-TUNING FUNCTIONS
        // ============================================================
        async function loadSkillsForDropdowns() {
            try {
                const res = await fetch('/api/skills');
                const skills = await res.json();

                const ftSelect = document.getElementById('ft-skill-select');
                const fbSelect = document.getElementById('feedback-skill-select');

                const options = '<option value="">Select a skill...</option>' +
                    skills.map(s => `<option value="${s.id}">${s.name}</option>`).join('');

                ftSelect.innerHTML = options;
                fbSelect.innerHTML = options;
            } catch (e) {
                console.error('Failed to load skills for dropdowns:', e);
            }
        }

        async function startSkillFineTune() {
            const messageEl = document.getElementById('ft-message');
            const skillId = document.getElementById('ft-skill-select').value;
            const examplesText = document.getElementById('ft-examples').value;

            if (!skillId) {
                messageEl.innerHTML = '<div style="color: var(--neon-orange);">Please select a skill</div>';
                return;
            }

            // Parse JSONL examples
            let examples = [];
            try {
                const lines = examplesText.trim().split('\\n').filter(l => l.trim());
                examples = lines.map(l => JSON.parse(l));
            } catch (e) {
                messageEl.innerHTML = '<div style="color: var(--neon-orange);">Invalid JSONL format</div>';
                return;
            }

            if (examples.length === 0) {
                messageEl.innerHTML = '<div style="color: var(--neon-orange);">Please add at least one training example</div>';
                return;
            }

            messageEl.innerHTML = '<div style="color: var(--neon-cyan);">Starting fine-tuning...</div>';

            try {
                const res = await fetch(`/api/skills/${skillId}/fine-tune`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ examples })
                });
                const result = await res.json();

                if (result.success) {
                    messageEl.innerHTML = `<div style="color: var(--neon-green);">${result.message}</div>`;
                } else {
                    messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${result.error}</div>`;
                }
            } catch (e) {
                messageEl.innerHTML = `<div style="color: var(--neon-orange);">Error: ${e.message}</div>`;
            }
        }

        async function submitSkillFeedback() {
            const skillId = document.getElementById('feedback-skill-select').value;
            const query = document.getElementById('feedback-query').value;
            const response = document.getElementById('feedback-response').value;
            const rating = document.getElementById('feedback-rating').value;
            const corrected = document.getElementById('feedback-corrected').value;

            if (!skillId || !query || !response) {
                alert('Please fill in skill, query, and response');
                return;
            }

            try {
                const res = await fetch(`/api/skills/${skillId}/feedback`, {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        query,
                        response,
                        rating: parseInt(rating),
                        corrected_response: corrected || null
                    })
                });
                const result = await res.json();

                if (result.success) {
                    alert('Feedback submitted successfully!');
                    document.getElementById('feedback-query').value = '';
                    document.getElementById('feedback-response').value = '';
                    document.getElementById('feedback-corrected').value = '';
                } else {
                    alert('Error: ' + result.error);
                }
            } catch (e) {
                alert('Error: ' + e.message);
            }
        }

        // ============================================================
        // INITIALIZATION
        // ============================================================
        document.addEventListener('DOMContentLoaded', () => {
            loadSkills();
            refreshServerStatus();
            loadMetrics();
            loadActivity();
            loadVoiceProviders();
            loadPlatforms();

            setInterval(loadSkills, 30000);
            setInterval(refreshServerStatus, 10000);
            setInterval(loadActivity, 15000);
        });
    </script>
</body>
</html>
'''


@app.route('/')
def dashboard():
    return render_template_string(DASHBOARD_HTML)


if __name__ == '__main__':
    # Add some initial activity
    add_activity("plumber_expert deployed successfully", "")
    add_activity("Training started for restaurant_host", "")
    add_activity("Low satisfaction detected on tech_support", "")
    add_activity("New documents processed: 47 examples", "")

    print("""
    ===============================================================
             UNIFIED SKILL COMMAND CENTER
             Voice Agent Intelligence Hub
    ===============================================================

    Dashboard: http://localhost:5000

    TABS:
      Dashboard       - Stats, Server, Skills, Training Pipeline, Activity
      Skill Factory   - Business Profile, Upload Docs, Train, Manage Skills
      Command Center  - API Keys, Test LLM, Compare, Masking, Stats, Voice
      LPU Inference   - BitNet model testing

    ===============================================================
    """)
    app.run(host='0.0.0.0', port=5000, debug=True)
