"""
Universal Parser - Parses 70+ file types
"""

import os
import io
import re
import json
import csv
import hashlib
import tempfile
import base64
import zipfile
import tarfile
import xml.etree.ElementTree as ET
from typing import Dict, Tuple, List
from pathlib import Path

# Document parsing
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
import html2text
from bs4 import BeautifulSoup
import email
from email import policy
import yaml
import toml

# ═══════════════════════════════════════════════════════════════════════════════
# SUPPORTED FILE TYPES
# ═══════════════════════════════════════════════════════════════════════════════

FILE_TYPES = {
    # Documents
    'pdf': {'category': 'document', 'parser': 'parse_pdf', 'name': 'PDF Document'},
    'docx': {'category': 'document', 'parser': 'parse_docx', 'name': 'Word Document'},
    'doc': {'category': 'document', 'parser': 'parse_doc', 'name': 'Word Document (Legacy)'},
    'rtf': {'category': 'document', 'parser': 'parse_rtf', 'name': 'Rich Text Format'},
    'odt': {'category': 'document', 'parser': 'parse_odt', 'name': 'OpenDocument Text'},
    'pages': {'category': 'document', 'parser': 'parse_pages', 'name': 'Apple Pages'},
    
    # Spreadsheets
    'xlsx': {'category': 'spreadsheet', 'parser': 'parse_xlsx', 'name': 'Excel Spreadsheet'},
    'xls': {'category': 'spreadsheet', 'parser': 'parse_xls', 'name': 'Excel (Legacy)'},
    'csv': {'category': 'spreadsheet', 'parser': 'parse_csv', 'name': 'CSV File'},
    'tsv': {'category': 'spreadsheet', 'parser': 'parse_tsv', 'name': 'TSV File'},
    'ods': {'category': 'spreadsheet', 'parser': 'parse_ods', 'name': 'OpenDocument Spreadsheet'},
    
    # Presentations
    'pptx': {'category': 'presentation', 'parser': 'parse_pptx', 'name': 'PowerPoint'},
    'ppt': {'category': 'presentation', 'parser': 'parse_ppt', 'name': 'PowerPoint (Legacy)'},
    'odp': {'category': 'presentation', 'parser': 'parse_odp', 'name': 'OpenDocument Presentation'},
    'key': {'category': 'presentation', 'parser': 'parse_key', 'name': 'Apple Keynote'},
    
    # Text & Code
    'txt': {'category': 'text', 'parser': 'parse_text', 'name': 'Plain Text'},
    'md': {'category': 'text', 'parser': 'parse_text', 'name': 'Markdown'},
    'rst': {'category': 'text', 'parser': 'parse_text', 'name': 'reStructuredText'},
    'log': {'category': 'text', 'parser': 'parse_text', 'name': 'Log File'},
    'py': {'category': 'code', 'parser': 'parse_code', 'name': 'Python'},
    'js': {'category': 'code', 'parser': 'parse_code', 'name': 'JavaScript'},
    'ts': {'category': 'code', 'parser': 'parse_code', 'name': 'TypeScript'},
    'java': {'category': 'code', 'parser': 'parse_code', 'name': 'Java'},
    'cpp': {'category': 'code', 'parser': 'parse_code', 'name': 'C++'},
    'c': {'category': 'code', 'parser': 'parse_code', 'name': 'C'},
    'cs': {'category': 'code', 'parser': 'parse_code', 'name': 'C#'},
    'go': {'category': 'code', 'parser': 'parse_code', 'name': 'Go'},
    'rs': {'category': 'code', 'parser': 'parse_code', 'name': 'Rust'},
    'rb': {'category': 'code', 'parser': 'parse_code', 'name': 'Ruby'},
    'php': {'category': 'code', 'parser': 'parse_code', 'name': 'PHP'},
    'sql': {'category': 'code', 'parser': 'parse_code', 'name': 'SQL'},
    'sh': {'category': 'code', 'parser': 'parse_code', 'name': 'Shell Script'},
    'bash': {'category': 'code', 'parser': 'parse_code', 'name': 'Bash Script'},
    
    # Data & Config
    'json': {'category': 'data', 'parser': 'parse_json', 'name': 'JSON'},
    'jsonl': {'category': 'data', 'parser': 'parse_jsonl', 'name': 'JSON Lines'},
    'xml': {'category': 'data', 'parser': 'parse_xml', 'name': 'XML'},
    'yaml': {'category': 'data', 'parser': 'parse_yaml', 'name': 'YAML'},
    'yml': {'category': 'data', 'parser': 'parse_yaml', 'name': 'YAML'},
    'toml': {'category': 'data', 'parser': 'parse_toml', 'name': 'TOML'},
    'ini': {'category': 'data', 'parser': 'parse_ini', 'name': 'INI Config'},
    'env': {'category': 'data', 'parser': 'parse_env', 'name': 'Environment File'},
    
    # Web
    'html': {'category': 'web', 'parser': 'parse_html', 'name': 'HTML'},
    'htm': {'category': 'web', 'parser': 'parse_html', 'name': 'HTML'},
    'xhtml': {'category': 'web', 'parser': 'parse_html', 'name': 'XHTML'},
    'css': {'category': 'web', 'parser': 'parse_code', 'name': 'CSS'},
    'scss': {'category': 'web', 'parser': 'parse_code', 'name': 'SCSS'},
    'less': {'category': 'web', 'parser': 'parse_code', 'name': 'LESS'},
    
    # Email
    'eml': {'category': 'email', 'parser': 'parse_email', 'name': 'Email Message'},
    'msg': {'category': 'email', 'parser': 'parse_msg', 'name': 'Outlook Message'},
    'mbox': {'category': 'email', 'parser': 'parse_mbox', 'name': 'Email Archive'},
    
    # eBooks
    'epub': {'category': 'ebook', 'parser': 'parse_epub', 'name': 'EPUB eBook'},
    'mobi': {'category': 'ebook', 'parser': 'parse_mobi', 'name': 'Kindle eBook'},
    
    # Images (OCR)
    'png': {'category': 'image', 'parser': 'parse_image', 'name': 'PNG Image'},
    'jpg': {'category': 'image', 'parser': 'parse_image', 'name': 'JPEG Image'},
    'jpeg': {'category': 'image', 'parser': 'parse_image', 'name': 'JPEG Image'},
    'gif': {'category': 'image', 'parser': 'parse_image', 'name': 'GIF Image'},
    'webp': {'category': 'image', 'parser': 'parse_image', 'name': 'WebP Image'},
    'bmp': {'category': 'image', 'parser': 'parse_image', 'name': 'Bitmap Image'},
    'tiff': {'category': 'image', 'parser': 'parse_image', 'name': 'TIFF Image'},
    'tif': {'category': 'image', 'parser': 'parse_image', 'name': 'TIFF Image'},
    'svg': {'category': 'image', 'parser': 'parse_svg', 'name': 'SVG Image'},
    'heic': {'category': 'image', 'parser': 'parse_image', 'name': 'HEIC Image'},
    
    # Audio
    'mp3': {'category': 'audio', 'parser': 'parse_audio', 'name': 'MP3 Audio'},
    'wav': {'category': 'audio', 'parser': 'parse_audio', 'name': 'WAV Audio'},
    'm4a': {'category': 'audio', 'parser': 'parse_audio', 'name': 'M4A Audio'},
    'ogg': {'category': 'audio', 'parser': 'parse_audio', 'name': 'OGG Audio'},
    'flac': {'category': 'audio', 'parser': 'parse_audio', 'name': 'FLAC Audio'},
    'aac': {'category': 'audio', 'parser': 'parse_audio', 'name': 'AAC Audio'},
    'wma': {'category': 'audio', 'parser': 'parse_audio', 'name': 'WMA Audio'},
    
    # Video
    'mp4': {'category': 'video', 'parser': 'parse_video', 'name': 'MP4 Video'},
    'mkv': {'category': 'video', 'parser': 'parse_video', 'name': 'MKV Video'},
    'avi': {'category': 'video', 'parser': 'parse_video', 'name': 'AVI Video'},
    'mov': {'category': 'video', 'parser': 'parse_video', 'name': 'MOV Video'},
    'wmv': {'category': 'video', 'parser': 'parse_video', 'name': 'WMV Video'},
    'webm': {'category': 'video', 'parser': 'parse_video', 'name': 'WebM Video'},
    'flv': {'category': 'video', 'parser': 'parse_video', 'name': 'FLV Video'},
    
    # Archives
    'zip': {'category': 'archive', 'parser': 'parse_archive', 'name': 'ZIP Archive'},
    'tar': {'category': 'archive', 'parser': 'parse_archive', 'name': 'TAR Archive'},
    'gz': {'category': 'archive', 'parser': 'parse_archive', 'name': 'GZIP Archive'},
    'tgz': {'category': 'archive', 'parser': 'parse_archive', 'name': 'TAR.GZ Archive'},
    'rar': {'category': 'archive', 'parser': 'parse_rar', 'name': 'RAR Archive'},
    '7z': {'category': 'archive', 'parser': 'parse_7z', 'name': '7-Zip Archive'},
    
    # Subtitles
    'srt': {'category': 'subtitle', 'parser': 'parse_srt', 'name': 'SRT Subtitles'},
    'vtt': {'category': 'subtitle', 'parser': 'parse_vtt', 'name': 'WebVTT Subtitles'},
    'ass': {'category': 'subtitle', 'parser': 'parse_ass', 'name': 'ASS Subtitles'},
    'ssa': {'category': 'subtitle', 'parser': 'parse_ass', 'name': 'SSA Subtitles'},
}

ALLOWED_EXTENSIONS = set(FILE_TYPES.keys())
TEMP_DIR = tempfile.gettempdir()


class UniversalParser:
    """Parses ANY supported file type and extracts text content."""
    
    def __init__(self):
        self.http_client = None
    
    def parse(self, file_path: str, filename: str) -> Dict:
        """Parse any file and return extracted text."""
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        
        if ext not in FILE_TYPES:
            return {'text': '', 'pages': 0, 'errors': [f'Unsupported: {ext}'], 'metadata': {}}
        
        file_info = FILE_TYPES[ext]
        parser_method = getattr(self, file_info['parser'], None)
        
        if not parser_method:
            return {'text': '', 'pages': 0, 'errors': [f'Parser not implemented: {ext}'], 'metadata': {}}
        
        try:
            result = parser_method(file_path, filename)
            result['source_type'] = file_info['category']
            result['file_type'] = ext
            return result
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [f'Parse error: {str(e)}'], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # DOCUMENT PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_pdf(self, file_path: str, filename: str) -> Dict:
        text_parts = []
        errors = []
        metadata = {}
        
        try:
            doc = fitz.open(file_path)
            metadata = dict(doc.metadata) if doc.metadata else {}
            pages = len(doc)
            
            for i, page in enumerate(doc):
                try:
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"[Page {i+1}]\n{text}")
                except Exception as e:
                    errors.append(f"Page {i+1}: {str(e)}")
            
            doc.close()
            return {'text': '\n\n'.join(text_parts), 'pages': pages, 'errors': errors, 'metadata': metadata}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_docx(self, file_path: str, filename: str) -> Dict:
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            
            return {'text': '\n\n'.join(text_parts), 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_doc(self, file_path: str, filename: str) -> Dict:
        return self.parse_text(file_path, filename)
    
    def parse_rtf(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            text = re.sub(r'\\[a-z]+\d* ?', '', content)
            text = re.sub(r'[{}]', '', text)
            return {'text': text.strip(), 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_odt(self, file_path: str, filename: str) -> Dict:
        try:
            with zipfile.ZipFile(file_path) as z:
                with z.open('content.xml') as f:
                    content = f.read()
            soup = BeautifulSoup(content, 'xml')
            text = soup.get_text(separator='\n')
            return {'text': text.strip(), 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_pages(self, file_path: str, filename: str) -> Dict:
        return self.parse_odt(file_path, filename)
    
    # ─────────────────────────────────────────────────────────────────────────
    # SPREADSHEET PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_xlsx(self, file_path: str, filename: str) -> Dict:
        try:
            wb = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip(' |'):
                        text_parts.append(row_text)
            
            wb.close()
            return {'text': '\n'.join(text_parts), 'pages': len(wb.sheetnames), 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_xls(self, file_path: str, filename: str) -> Dict:
        return self.parse_xlsx(file_path, filename)
    
    def parse_csv(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.DictReader(f)
                rows = [' | '.join(f"{k}: {v}" for k, v in row.items() if v) for row in reader]
            return {'text': '\n'.join(rows), 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_tsv(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.DictReader(f, delimiter='\t')
                rows = [' | '.join(f"{k}: {v}" for k, v in row.items() if v) for row in reader]
            return {'text': '\n'.join(rows), 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_ods(self, file_path: str, filename: str) -> Dict:
        return self.parse_odt(file_path, filename)
    
    # ─────────────────────────────────────────────────────────────────────────
    # PRESENTATION PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_pptx(self, file_path: str, filename: str) -> Dict:
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {i+1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join(cell.text for cell in row.cells)
                            if row_text.strip():
                                slide_text.append(row_text)
                
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            
            return {'text': '\n\n'.join(text_parts), 'pages': len(prs.slides), 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_ppt(self, file_path: str, filename: str) -> Dict:
        return self.parse_pptx(file_path, filename)
    
    def parse_odp(self, file_path: str, filename: str) -> Dict:
        return self.parse_odt(file_path, filename)
    
    def parse_key(self, file_path: str, filename: str) -> Dict:
        return self.parse_odt(file_path, filename)
    
    # ─────────────────────────────────────────────────────────────────────────
    # TEXT & CODE PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_text(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            return {'text': text, 'pages': 1, 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_code(self, file_path: str, filename: str) -> Dict:
        result = self.parse_text(file_path, filename)
        ext = filename.rsplit('.', 1)[-1].lower()
        result['metadata']['language'] = ext
        return result
    
    # ─────────────────────────────────────────────────────────────────────────
    # DATA FORMAT PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_json(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {'text': json.dumps(data, indent=2), 'pages': 1, 'errors': [], 'metadata': {'format': 'json'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_jsonl(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                items = [json.loads(line) for line in f if line.strip()]
            return {'text': json.dumps(items, indent=2), 'pages': 1, 'errors': [], 'metadata': {'format': 'jsonl', 'count': len(items)}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_xml(self, file_path: str, filename: str) -> Dict:
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            text = ET.tostring(root, encoding='unicode', method='text')
            return {'text': text.strip(), 'pages': 1, 'errors': [], 'metadata': {'format': 'xml'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_yaml(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            return {'text': yaml.dump(data, default_flow_style=False), 'pages': 1, 'errors': [], 'metadata': {'format': 'yaml'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_toml(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = toml.load(f)
            return {'text': toml.dumps(data), 'pages': 1, 'errors': [], 'metadata': {'format': 'toml'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_ini(self, file_path: str, filename: str) -> Dict:
        return self.parse_text(file_path, filename)
    
    def parse_env(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                lines = []
                for line in f:
                    if '=' in line and not line.strip().startswith('#'):
                        key = line.split('=')[0]
                        lines.append(f"{key}=<redacted>")
                    else:
                        lines.append(line.strip())
            return {'text': '\n'.join(lines), 'pages': 1, 'errors': [], 'metadata': {'format': 'env'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # WEB PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_html(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            text = h.handle(html_content)
            
            soup = BeautifulSoup(html_content, 'html.parser')
            title = soup.title.string if soup.title else ''
            
            return {'text': text.strip(), 'pages': 1, 'errors': [], 'metadata': {'title': title}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # EMAIL PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_email(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'rb') as f:
                msg = email.message_from_binary_file(f, policy=policy.default)
            
            subject = msg.get('subject', '')
            from_addr = msg.get('from', '')
            to_addr = msg.get('to', '')
            date = msg.get('date', '')
            
            body = ''
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        body = part.get_content()
                        break
            else:
                body = msg.get_content()
            
            text = f"Subject: {subject}\nFrom: {from_addr}\nTo: {to_addr}\nDate: {date}\n\n{body}"
            return {'text': text, 'pages': 1, 'errors': [], 'metadata': {'subject': subject, 'from': from_addr}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_msg(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 0, 'errors': ['MSG requires extract-msg library'], 'metadata': {}}
    
    def parse_mbox(self, file_path: str, filename: str) -> Dict:
        try:
            import mailbox
            mbox = mailbox.mbox(file_path)
            texts = []
            
            for msg in mbox:
                subject = msg.get('subject', 'No Subject')
                body = ''
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == 'text/plain':
                            body = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                            break
                else:
                    body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                texts.append(f"[Email: {subject}]\n{body}")
            
            mbox.close()
            return {'text': '\n\n---\n\n'.join(texts), 'pages': len(texts), 'errors': [], 'metadata': {'email_count': len(texts)}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # EBOOK PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_epub(self, file_path: str, filename: str) -> Dict:
        try:
            with zipfile.ZipFile(file_path) as z:
                text_parts = []
                for name in z.namelist():
                    if name.endswith(('.html', '.xhtml', '.htm')):
                        with z.open(name) as f:
                            soup = BeautifulSoup(f.read(), 'html.parser')
                            text = soup.get_text(separator='\n')
                            if text.strip():
                                text_parts.append(text.strip())
            return {'text': '\n\n'.join(text_parts), 'pages': len(text_parts), 'errors': [], 'metadata': {}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_mobi(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 0, 'errors': ['MOBI requires mobi library'], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # IMAGE/AUDIO/VIDEO PARSERS (flags for async processing)
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_image(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 1, 'errors': [], 'metadata': {'needs_ocr': True, 'file_path': file_path}}
    
    def parse_svg(self, file_path: str, filename: str) -> Dict:
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            texts = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
            return {'text': '\n'.join(texts), 'pages': 1, 'errors': [], 'metadata': {'format': 'svg'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_audio(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 1, 'errors': [], 'metadata': {'needs_transcription': True, 'file_path': file_path}}
    
    def parse_video(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 1, 'errors': [], 'metadata': {'needs_transcription': True, 'file_path': file_path}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # ARCHIVE PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_archive(self, file_path: str, filename: str) -> Dict:
        extracted_texts = []
        errors = []
        
        try:
            if filename.endswith('.zip'):
                with zipfile.ZipFile(file_path) as z:
                    for name in z.namelist():
                        if name.endswith('/'):
                            continue
                        ext = name.rsplit('.', 1)[-1].lower() if '.' in name else ''
                        if ext in ALLOWED_EXTENSIONS and ext not in ['zip', 'tar', 'gz', 'rar', '7z']:
                            try:
                                with z.open(name) as f:
                                    temp_path = os.path.join(TEMP_DIR, f"temp_{hashlib.md5(name.encode()).hexdigest()[:8]}")
                                    with open(temp_path, 'wb') as tf:
                                        tf.write(f.read())
                                    result = self.parse(temp_path, name)
                                    if result['text']:
                                        extracted_texts.append(f"[{name}]\n{result['text']}")
                                    os.remove(temp_path)
                            except Exception as e:
                                errors.append(f"{name}: {str(e)}")
            
            elif filename.endswith(('.tar', '.tar.gz', '.tgz')):
                mode = 'r:gz' if filename.endswith(('.tar.gz', '.tgz')) else 'r'
                with tarfile.open(file_path, mode) as t:
                    for member in t.getmembers():
                        if not member.isfile():
                            continue
                        ext = member.name.rsplit('.', 1)[-1].lower() if '.' in member.name else ''
                        if ext in ALLOWED_EXTENSIONS:
                            try:
                                f = t.extractfile(member)
                                if f:
                                    temp_path = os.path.join(TEMP_DIR, f"temp_{hashlib.md5(member.name.encode()).hexdigest()[:8]}")
                                    with open(temp_path, 'wb') as tf:
                                        tf.write(f.read())
                                    result = self.parse(temp_path, member.name)
                                    if result['text']:
                                        extracted_texts.append(f"[{member.name}]\n{result['text']}")
                                    os.remove(temp_path)
                            except Exception as e:
                                errors.append(f"{member.name}: {str(e)}")
            
            return {'text': '\n\n---\n\n'.join(extracted_texts), 'pages': len(extracted_texts), 'errors': errors, 'metadata': {'archive': True}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_rar(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 0, 'errors': ['RAR requires rarfile library'], 'metadata': {}}
    
    def parse_7z(self, file_path: str, filename: str) -> Dict:
        return {'text': '', 'pages': 0, 'errors': ['7z requires py7zr library'], 'metadata': {}}
    
    # ─────────────────────────────────────────────────────────────────────────
    # SUBTITLE PARSERS
    # ─────────────────────────────────────────────────────────────────────────
    
    def parse_srt(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            lines = [line.strip() for line in content.split('\n') 
                     if line.strip() and not line.strip().isdigit() and '-->' not in line]
            return {'text': ' '.join(lines), 'pages': 1, 'errors': [], 'metadata': {'format': 'srt'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
    
    def parse_vtt(self, file_path: str, filename: str) -> Dict:
        return self.parse_srt(file_path, filename)
    
    def parse_ass(self, file_path: str, filename: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            lines = []
            for line in content.split('\n'):
                if line.startswith('Dialogue:'):
                    parts = line.split(',', 9)
                    if len(parts) >= 10:
                        text = re.sub(r'\{[^}]+\}', '', parts[9])
                        if text.strip():
                            lines.append(text.strip())
            return {'text': ' '.join(lines), 'pages': 1, 'errors': [], 'metadata': {'format': 'ass'}}
        except Exception as e:
            return {'text': '', 'pages': 0, 'errors': [str(e)], 'metadata': {}}
